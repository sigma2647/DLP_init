import os
import yaml
import magic
import time
import logging
import threading
import uuid
import shutil
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any, Iterator
from concurrent.futures import ThreadPoolExecutor, as_completed
from markitdown import MarkItDown
from dataclasses import dataclass, field
import re
from pdfminer.high_level import extract_text as pdf_extract_text
import pandas as pd
import openpyxl
import csv
import json
from datetime import datetime
import argparse
import sys
import win32com.client
import pythoncom
import psutil
import patoolib
from pptx import Presentation

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('sensitive_detector.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class ProcessingResult:
    """数据类，用于存储文件处理结果"""
    file_path: str
    mime_type: str
    content: Dict[str, Any]
    sensitive_words: List[Tuple[str, List[int]]] = field(default_factory=list)
    error: Optional[str] = None
    processing_time: float = 0.0

class FileTypeDetector:
    """文件类型检测器，使用 python-magic 库"""
    
    MIME_TYPES = {
        ".txt": "text/plain", ".csv": "text/csv", ".xml": "text/xml", ".html": "text/html",
        ".htm": "text/html", ".json": "application/json", ".yaml": "application/yaml",
        ".yml": "application/yaml", ".md": "text/markdown", ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".zip": "application/zip", ".rar": "application/x-rar", ".7z": "application/x-7z-compressed",
        ".tar": "application/x-tar", ".gz": "application/gzip", ".bz2": "application/x-bzip2",
        ".pdf": "application/pdf", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
        ".gif": "image/gif", ".bmp": "image/bmp", ".svg": "image/svg+xml", ".webp": "image/webp",
        ".mp3": "audio/mpeg", ".mp4": "video/mp4", ".avi": "video/x-msvideo", ".mov": "video/quicktime",
        ".wav": "audio/wav", ".bin": "application/octet-stream", ".exe": "application/x-msdownload",
        ".dll": "application/x-msdownload",
    }
    MIME_TO_EXT = {mime: ext for ext, mime in MIME_TYPES.items()}
    SKIP_EXTENSIONS = {'.dwg', '.mp3', '.wav', '.mp4', '.avi', '.mkv', '.flv', '.mov'}

    def __init__(self):
        """初始化文件类型检测器"""
        self.mime = magic.Magic(mime=True)
    
    def get_all_files(self, directory: str) -> List[str]:
        """递归获取目录中的所有文件，忽略 ~ 开头或内部流文件及跳过文件"""
        files = []
        for f in Path(directory).rglob('*'):
            if f.is_file() and not f.name.startswith('~$') and not self._is_internal_stream(f.name):
                ext = f.suffix.lower()
                if ext not in self.SKIP_EXTENSIONS:
                    files.append(str(f))
        return files
    
    def _is_internal_stream(self, filename: str) -> bool:
        """判断文件是否为 Office 内部流文件（如 [5]SummaryInformation）"""
        return bool(re.match(r'^\[\d+\].+', filename))
    
    def get_file_info(self, file_path: str) -> Dict:
        """获取文件信息，包括 MIME 类型和文件头"""
        try:
            mime_type = self.detect_file_type(file_path)
            with open(file_path, "rb") as f:
                file_header = f.read(16).hex().upper()
            file_extension = Path(file_path).suffix.lower()
            return {
                "file_path": file_path,
                "mime_type": mime_type,
                "file_extension": file_extension,
                "file_header": file_header,
                "size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"获取文件信息失败: {e} - {file_path}")
            return {"file_path": file_path, "mime_type": "unknown", "error": str(e)}
    
    def detect_file_type(self, file_path: str) -> str:
        """检测文件 MIME 类型，优化文件类型检测"""
        try:
            normalized_path = os.path.normpath(file_path.encode('utf-8', errors='replace').decode('utf-8'))
            mime_type = self.mime.from_file(normalized_path)
            
            ext = Path(file_path).suffix.lower()
            
            # First check by extension for known Office files
            if ext == '.docx':
                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif ext == '.xlsx':
                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif ext == '.pptx':
                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif ext == '.doc':
                return 'application/msword'
            elif ext == '.xls':
                return 'application/vnd.ms-excel'
            elif ext == '.ppt':
                return 'application/vnd.ms-powerpoint'
            
            # For other files, check the file header
            with open(file_path, 'rb') as f:
                header = f.read(8)
                
                # Check for Office Open XML formats (they're ZIP-based)
                if header.startswith(b'PK\x03\x04'):
                    try:
                        import zipfile
                        with zipfile.ZipFile(file_path) as zf:
                            file_list = zf.namelist()
                            if 'word/document.xml' in file_list:
                                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                            elif 'xl/workbook.xml' in file_list:
                                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                            elif 'ppt/presentation.xml' in file_list:
                                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                    except:
                        pass
                    if ext in self.MIME_TYPES:
                        return self.MIME_TYPES[ext]
                    return 'application/zip'
                    
                elif header.startswith(b'\xD0\xCF\x11\xE0'):
                    if ext == '.doc':
                        return 'application/msword'
                    elif ext == '.xls':
                        return 'application/vnd.ms-excel'
                    elif ext == '.ppt':
                        return 'application/vnd.ms-powerpoint'
                    return 'application/msword' if mime_type == 'application/octet-stream' else mime_type
            
            if mime_type not in ('application/octet-stream', 'text/plain'):
                return mime_type
            
            return self.MIME_TYPES.get(ext, 'application/octet-stream')
            
        except Exception as e:
            logger.error(f"文件类型检测失败 {file_path}: {e}")
            return self.MIME_TYPES.get(Path(file_path).suffix.lower(), "application/octet-stream")

class ContentExtractor:
    """优化后的文件内容提取器，提高容错性，确保错误跳过而不停止执行"""
    
    MIME_TYPE = {
        'TEXT': 'text/plain',
        'CSV': 'text/csv',
        'PDF': 'application/pdf',
        'MARKDOWN': 'text/markdown',
        'DOCX': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'XLSX': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'XLS': 'application/vnd.ms-excel',
        'DOC': 'application/msword',
        'PPT': 'application/vnd.ms-powerpoint',
        'PPTX': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'ZIP': 'application/zip',
        'RAR': 'application/x-rar',
        'SEVENZ': 'application/x-7z-compressed'
    }
    
    def __init__(self, detector, is_windows: bool = True):
        self.detector = detector
        self.md = MarkItDown()
        self.is_windows = is_windows
        self.word_app = None
        self.word_lock = threading.Lock()
        
        # 动态构建提取器字典，确保即使某些库不可用也能处理其他文件类型
        self.extractors = {}
        
        # 基本文本格式的提取器，几乎总是可用
        self._register_basic_extractors()
        
        # 尝试注册文档格式提取器，出错则记录日志
        self._register_document_extractors()
        
        # 根据系统类型注册Windows特有的提取器
        if self.is_windows:
            self._register_windows_extractors()
    
    def _register_basic_extractors(self):
        """注册基本文本格式提取器"""
        self.extractors.update({
            self.MIME_TYPE['TEXT']: self._extract_text_content,
            self.MIME_TYPE['CSV']: self._extract_csv_content,
            self.MIME_TYPE['MARKDOWN']: self._extract_markdown_content,
        })
        
        # 尝试添加PDF提取器
        try:
            import pdfminer
            self.extractors[self.MIME_TYPE['PDF']] = self._extract_pdf_content
        except ImportError:
            logger.warning("PDF提取器未注册: pdfminer库不可用")
    
    def _register_document_extractors(self):
        """注册文档格式提取器"""
        # 尝试添加Office Open XML格式提取器
        try:
            self.extractors.update({
                self.MIME_TYPE['DOCX']: self._extract_docx_content,
                self.MIME_TYPE['XLSX']: self._extract_xlsx_content,
                self.MIME_TYPE['PPTX']: self._extract_pptx_content,
            })
        except Exception as e:
            logger.warning(f"部分Office Open XML提取器注册失败: {str(e)}")
        
        # 尝试添加压缩文件提取器
        try:
            import patoolib
            self.extractors.update({
                self.MIME_TYPE['ZIP']: self._extract_archive_content,
                self.MIME_TYPE['RAR']: self._extract_archive_content,
                self.MIME_TYPE['SEVENZ']: self._extract_archive_content,
            })
        except ImportError:
            logger.warning("压缩文件提取器未注册: patoolib库不可用")
    
    def _register_windows_extractors(self):
        """注册Windows特有的提取器"""
        try:
            import win32com.client
            self.extractors.update({
                self.MIME_TYPE['DOC']: self._extract_doc_content,
                self.MIME_TYPE['PPT']: self._extract_ppt_content,
                self.MIME_TYPE['XLS']: self._extract_xls_content,
            })
        except ImportError:
            logger.warning("Windows特有提取器未注册: win32com库不可用")
    
    def _init_word_app(self):
        """初始化 Word 应用程序实例，增加错误处理和超时机制"""
        if not self.is_windows:
            return False
            
        with self.word_lock:
            if self.word_app is None:
                try:
                    # 设置超时机制
                    timer = threading.Timer(10.0, self._force_close_office_processes)
                    timer.start()
                    
                    pythoncom.CoInitialize()
                    self.word_app = win32com.client.Dispatch("Word.Application")
                    self.word_app.Visible = False
                    
                    # 取消超时
                    timer.cancel()
                    
                    logger.info("成功初始化 Word 应用程序")
                    return True
                except Exception as e:
                    logger.error(f"初始化 Word 应用程序失败: {str(e)}")
                    self.word_app = None
                    
                    # 确保清理COM资源
                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass
                        
                    return False
                finally:
                    try:
                        timer.cancel()
                    except:
                        pass
            return True
    
    def _cleanup_word_app(self):
        """清理 Word 应用程序实例，增强可靠性"""
        with self.word_lock:
            if self.word_app is not None:
                try:
                    # 设置超时，防止Quit卡住
                    timer = threading.Timer(5.0, lambda: self._force_close_office_processes('WINWORD.EXE'))
                    timer.start()
                    
                    self.word_app.Quit(SaveChanges=False)
                    self.word_app = None
                    
                    # 取消超时
                    timer.cancel()
                except:
                    pass
                finally:
                    try:
                        timer.cancel()
                    except:
                        pass
                    
                    # 确保进程被终止
                    self._force_close_office_processes('WINWORD.EXE')
                    self.word_app = None
                    
                    # 清理COM资源
                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass
                    
                    logger.info("成功清理 Word 应用程序")
    
    def _create_empty_result(self, file_type: str) -> Dict[str, Any]:
        """创建空结果对象"""
        return {
            'content': '',
            'metadata': {'file_type': file_type},
            'error': None,
            'is_empty': True
        }
    
    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果对象"""
        return {
            'content': '',
            'metadata': {'file_type': file_type},
            'error': error_msg
        }
    
    def _get_file_header(self, file_path: str, size: int = 16) -> str:
        """获取文件头部字节的十六进制表示，增加容错处理"""
        try:
            with open(file_path, 'rb') as f:
                return f.read(size).hex().upper()
        except Exception as e:
            logger.error(f"无法读取文件头: {file_path} - {e}")
            return ""
    
    def _precheck_file_type(self, file_path: str, detected_mime: str) -> Tuple[bool, str]:
        """预检查文件类型和格式，增加容错处理"""
        try:
            ext = Path(file_path).suffix.lower()
            expected_mime = self.detector.MIME_TYPES.get(ext, detected_mime)
            
            # 计算文件大小，过滤空文件
            try:
                file_size = os.path.getsize(file_path)
                if file_size == 0:
                    return False, "文件为空"
                    
                # 超大文件可能会导致处理问题，标记但不阻止继续处理
                if file_size > 50 * 1024 * 1024:  # 50MB
                    logger.warning(f"文件过大 ({file_size / 1024 / 1024:.1f} MB): {file_path}")
            except Exception as e:
                logger.warning(f"无法获取文件大小: {file_path} - {e}")
            
            # 预读文件头部以验证文件类型
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(8)
                    
                    if expected_mime == self.MIME_TYPE['PPT'] and not header.startswith(b'\xD0\xCF\x11\xE0'):
                        return False, "文件头不符合 .ppt 格式"
                    elif expected_mime == self.MIME_TYPE['PPTX'] and not header.startswith(b'PK\x03\x04'):
                        return False, "文件头不符合 .pptx 格式"
                    elif expected_mime == self.MIME_TYPE['XLSX'] and not header.startswith(b'PK\x03\x04'):
                        return False, "文件头不符合 .xlsx 格式"
                    elif expected_mime == self.MIME_TYPE['DOC'] and not header.startswith(b'\xD0\xCF\x11\xE0'):
                        return False, "文件头不符合 .doc 格式"
            except Exception as e:
                logger.warning(f"文件头检查失败: {file_path} - {e}")
                # 尽管预检查失败，仍然继续处理
                return True, f"文件头检查失败，但继续处理: {str(e)}"
                
            return True, ""
        except Exception as e:
            logger.warning(f"预检查异常: {file_path} - {e}")
            # 出错不阻止继续处理
            return True, f"预检查异常，但继续处理: {str(e)}"

    def _is_valid_xlsx(self, file_path: str) -> Tuple[bool, str]:
        """检查 .xlsx 文件是否有效，使用更容错的方法"""
        try:
            with open(file_path, 'rb') as f:
                header = f.read(4)
                if not header.startswith(b'PK\x03\x04'):
                    return False, "文件头不符合 .xlsx 格式"
            
            # 尝试使用openpyxl验证
            try:
                wb = openpyxl.load_workbook(file_path, read_only=True)
                has_data = False
                
                # 限制检查行数，避免大文件消耗过多资源
                for i, row in enumerate(wb.active.iter_rows(max_rows=10, values_only=True)):
                    if any(cell is not None for cell in row):
                        has_data = True
                        break
                    if i >= 10:  # 额外限制，确保不会无限循环
                        break
                
                wb.close()
                
                if not has_data:
                    return False, "文件为空或无有效数据"
                
                return True, ""
            except openpyxl.utils.exceptions.InvalidFileException as e:
                # 尝试使用xlrd作为备选方案
                try:
                    import xlrd
                    wb = xlrd.open_workbook(file_path, on_demand=True)
                    # 如果能打开，说明文件有效
                    return True, "使用xlrd验证成功"
                except:
                    return False, f"文件格式无效或损坏: {str(e)}"
            except Exception as e:
                return False, f"文件检查失败: {str(e)}"
        except Exception as e:
            logger.warning(f"XLSX检查异常: {file_path} - {e}")
            # 出错也返回有效，让后续提取器决定是否能处理
            return True, "文件检查异常，但继续处理"
    
    def _extract_text_content(self, file_path: str) -> Dict[str, Any]:
        """提取文本文件内容，增强容错性"""
        result = self._create_empty_result('text')
        
        for encoding in ['utf-8', 'gbk', 'latin1', 'cp1252']:
            try:
                with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                    content = f.read(1024000)  # 限制最大读取1MB
                
                if content:
                    result = {
                        'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                        'metadata': {'file_type': 'text', 'encoding': encoding},
                        'error': None
                    }
                    return result
            except Exception as e:
                continue
        
        # 如果所有编码都失败，使用二进制模式读取并转换为十六进制
        try:
            with open(file_path, 'rb') as f:
                binary_content = f.read(4096)  # 仅读取前4KB
                hex_content = ' '.join('{:02X}'.format(b) for b in binary_content)
                
            return {
                'content': f"[无法以文本形式读取，显示前4KB十六进制内容]\n{hex_content}",
                'metadata': {'file_type': 'binary', 'original_mime': 'text/plain'},
                'error': "无法以常见编码读取文本内容"
            }
        except Exception as e:
            return self._create_error_result('text', f"读取文本文件失败: {str(e)}")
    
    def _extract_csv_content(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容，增强容错性和处理异常编码"""
        try:
            # 尝试多种编码和解析方式
            encodings = ['utf-8', 'latin1', 'gbk', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    # 尝试不同分隔符
                    for sep in [',', ';', '\t']:
                        try:
                            df = pd.read_csv(file_path, encoding=encoding, sep=sep, engine='python')
                            if len(df.columns) > 1:  # 确认确实解析出了多列
                                break
                        except:
                            continue
                    
                    if df is not None and len(df.columns) > 1:
                        break
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    logger.debug(f"CSV读取尝试失败 (编码: {encoding}): {str(e)}")
                    continue
            
            # 如果所有尝试都失败，使用最宽松的设置
            if df is None or len(df.columns) <= 1:
                df = pd.read_csv(file_path, encoding='utf-8', errors='ignore', engine='python', sep=None)
            
            # 限制大小
            if len(df) > 10000:
                df = df.head(10000)
                notice = "\n\n[注意: 文件过大，仅显示前10000行]"
            else:
                notice = ""
            
            # 确保输出美观
            pd.set_option('display.max_columns', 50)
            pd.set_option('display.width', 1000)
            
            return {
                'content': df.to_string(index=False) + notice,
                'metadata': {'file_type': 'csv', 'rows': len(df), 'columns': len(df.columns)},
                'error': None
            }
        except Exception as e:
            logger.warning(f"CSV处理失败: {file_path} - {e}")
            
            # 作为文本文件处理
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(10240)  # 读取前10KB作为样本
                
                return {
                    'content': content + '\n\n[注意: CSV解析失败，以纯文本显示]',
                    'metadata': {'file_type': 'csv', 'fallback': 'text'},
                    'error': f"CSV解析失败: {str(e)}"
                }
            except Exception as e2:
                return self._create_error_result('csv', f"CSV处理完全失败: {str(e)}; 文本读取失败: {str(e2)}")
    
    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文件内容，增强容错性和性能"""
        try:
            # 检查PDF文件有效性
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(5)
                    if not header.startswith(b'%PDF-'):
                        return self._create_error_result('pdf', "不是有效的PDF文件")
            except Exception as e:
                logger.warning(f"PDF头检查失败: {file_path} - {e}")
                # 继续尝试处理
            
            # 使用pdfminer提取文本，包含超时控制
            def extract_with_timeout():
                return pdf_extract_text(file_path)
            
            # 创建线程执行提取
            result = [None, None]  # [文本, 错误]
            def extract_worker():
                try:
                    result[0] = extract_with_timeout()
                except Exception as e:
                    result[1] = str(e)
            
            thread = threading.Thread(target=extract_worker)
            thread.daemon = True
            thread.start()
            thread.join(timeout=60)  # 60秒超时
            
            if thread.is_alive():
                # 提取超时
                return self._create_error_result('pdf', "PDF提取超时(>60秒)")
            
            if result[1]:
                # 提取出错
                return self._create_error_result('pdf', f"PDF处理失败: {result[1]}")
            
            text = result[0] or ""
            
            # 限制内容长度
            if len(text) > 100000:
                text = text[:100000] + "\n\n[注意: 文件过大，内容已截断]"
            
            # 如果提取的文本为空，尝试metadata报告
            if not text.strip():
                try:
                    # 尝试使用PyPDF2作为备选
                    from PyPDF2 import PdfReader
                    reader = PdfReader(file_path)
                    info = reader.metadata
                    num_pages = len(reader.pages)
                    
                    return {
                        'content': f"[PDF内容为空或无法提取文本]\n页数: {num_pages}\n元数据: {info}",
                        'metadata': {'file_type': 'pdf', 'pages': num_pages},
                        'error': "未提取到文本内容"
                    }
                except:
                    return {
                        'content': "",
                        'metadata': {'file_type': 'pdf'},
                        'error': "PDF为空或不包含可提取文本"
                    }
            
            return {
                'content': text,
                'metadata': {'file_type': 'pdf'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('pdf', f"PDF处理失败: {str(e)}")
    
    def _extract_markdown_content(self, file_path: str) -> Dict[str, Any]:
        """提取Markdown文件内容，增强容错性"""
        try:
            for encoding in ['utf-8', 'gbk', 'latin1']:
                try:
                    with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                        content = f.read(1024000)  # 限制最大读取1MB
                    
                    return {
                        'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                        'metadata': {'file_type': 'markdown', 'encoding': encoding},
                        'error': None
                    }
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    logger.debug(f"Markdown读取失败 (编码: {encoding}): {str(e)}")
            
            # 所有编码都失败，使用最宽松的设置
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(1024000)
            
            return {
                'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                'metadata': {'file_type': 'markdown', 'encoding': 'utf-8 (替换错误)'},
                'error': "文件编码检测失败，使用替换模式"
            }
        except Exception as e:
            return self._create_error_result('markdown', f"Markdown文件处理失败: {str(e)}")
    
    def _extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOCX文件内容，优先使用markitdown，增强容错性"""
        try:
            # 首先尝试使用MarkItDown
            try:
                result = self.md.convert(file_path)
                return {
                    'content': result.text_content,
                    'metadata': {'file_type': 'docx', 'extractor': 'markitdown'},
                    'error': None
                }
            except Exception as e:
                logger.warning(f"MarkItDown处理DOCX失败: {file_path} - {str(e)}")
                
                # 尝试使用python-docx作为备选
                try:
                    from docx import Document
                    doc = Document(file_path)
                    
                    paragraphs = []
                    for para in doc.paragraphs:
                        if para.text.strip():
                            paragraphs.append(para.text)
                    
                    # 处理表格内容
                    tables = []
                    for table in doc.tables:
                        table_text = []
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text)
                            table_text.append(' | '.join(row_text))
                        tables.append('\n'.join(table_text))
                    
                    content = '\n\n'.join(paragraphs)
                    if tables:
                        content += '\n\n表格内容:\n' + '\n\n'.join(tables)
                    
                    return {
                        'content': content,
                        'metadata': {'file_type': 'docx', 'extractor': 'python-docx'},
                        'error': f"MarkItDown失败，使用python-docx: {str(e)}"
                    }
                except Exception as e2:
                    if self.is_windows:
                        # 尝试使用Win32COM作为最后手段
                        try:
                            return self._extract_doc_content(file_path)
                        except Exception as e3:
                            return self._create_error_result('docx', 
                                f"所有提取方法失败: MarkItDown: {str(e)}; python-docx: {str(e2)}; COM: {str(e3)}")
                    return self._create_error_result('docx', 
                        f"所有提取方法失败: MarkItDown: {str(e)}; python-docx: {str(e2)}")
        except Exception as e:
            return self._create_error_result('docx', f"DOCX处理失败: {str(e)}")
    
    def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLSX文件内容，增强容错性和多种尝试方法"""
        try:
            # 首先尝试使用MarkItDown
            try:
                result = self.md.convert(file_path)
                return {
                    'content': result.text_content,
                    'metadata': {'file_type': 'xlsx', 'extractor': 'markitdown'},
                    'error': None
                }
            except Exception as e:
                logger.warning(f"MarkItDown处理XLSX失败: {file_path} - {str(e)}")
                
                # 尝试使用pandas作为备选
                try:
                    # 先检查有哪些sheet
                    sheet_names = pd.ExcelFile(file_path).sheet_names
                    
                    content = []
                    for sheet_name in sheet_names:
                        try:
                            # 只读取前5000行，避免内存问题
                            df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=5000)
                            
                            if len(df) > 0:
                                if len(df) > 1000:
                                    df = df.head(1000)
                                    sheet_content = f"Sheet: {sheet_name} (显示前1000行，共{len(df)}行)\n{df.to_string(index=False)}"
                                else:
                                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                                content.append(sheet_content)
                        except Exception as sheet_e:
                            content.append(f"Sheet: {sheet_name} - 读取错误: {str(sheet_e)}")
                    
                    if content:
                        return {
                            'content': '\n\n'.join(content),
                            'metadata': {'file_type': 'xlsx', 'extractor': 'pandas', 'sheets': len(sheet_names)},
                            'error': f"MarkItDown失败，使用pandas: {str(e)}"
                        }
                except Exception as e2:
                    # 尝试使用openpyxl作为最后手段
                    try:
                        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                        content = []
                        
                        for sheet_name in wb.sheetnames:
                            ws = wb[sheet_name]
                            sheet_content = [f"Sheet: {sheet_name}"]
                            
                            row_count = 0
                            for i, row in enumerate(ws.iter_rows(values_only=True)):
                                if i >= 1000:  # 限制行数
                                    sheet_content.append(f"... (超过1000行，已截断)")
                                    break
                                
                                row_count += 1
                                if any(cell is not None for cell in row):  # 跳过空行
                                    sheet_content.append(' | '.join(str(cell) if cell is not None else '' for cell in row))
                            
                            if row_count > 0:
                                content.append('\n'.join(sheet_content))
                        
                        return {
                            'content': '\n\n'.join(content),
                            'metadata': {'file_type': 'xlsx', 'extractor': 'openpyxl', 'sheets': len(wb.sheetnames)},
                            'error': f"MarkItDown和pandas失败，使用openpyxl: {str(e)}; {str(e2)}"
                        }
                    except Exception as e3:
                        return self._create_error_result('xlsx', 
                            f"所有提取方法失败: MarkItDown: {str(e)}; pandas: {str(e2)}; openpyxl: {str(e3)}")
        except Exception as e:
            return self._create_error_result('xlsx', f"XLSX处理失败: {str(e)}")
    
    def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPTX文件内容，优先使用markitdown，增强容错性"""
        try:
            # 首先尝试使用MarkItDown
            try:
                result = self.md.convert(file_path)
                return {
                    'content': result.text_content,
                    'metadata': {'file_type': 'pptx', 'extractor': 'markitdown'},
                    'error': None
                }
            except Exception as e:
                logger.warning(f"MarkItDown处理PPTX失败: {file_path} - {str(e)}")
                
                # 尝试使用python-pptx作为备选
                try:
                    prs = Presentation(file_path)
                    text_content = []
                    
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        slide_text.append(f"=== 幻灯片 {i+1} ===")
                        
                        # 处理标题
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        
                        # 检查是否有备注
                        if hasattr(slide, 'notes_slide') and slide.notes_slide:
                            for shape in slide.notes_slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    slide_text.append(f"备注: {shape.text.strip()}")
                        
                        if len(slide_text) > 1:  # 确保至少有内容
                            text_content.append('\n'.join(slide_text))
                    
                    if text_content:
                        content = "\n\n".join(text_content)
                        return {
                            'content': content,
                            'metadata': {'file_type': 'pptx', 'extractor': 'python-pptx', 'slides': len(prs.slides)},
                            'error': f"MarkItDown失败，使用python-pptx: {str(e)}"
                        }
                    else:
                        return {
                            'content': f"[PPTX文件未提取到文本内容]\n幻灯片数: {len(prs.slides)}",
                            'metadata': {'file_type': 'pptx', 'extractor': 'python-pptx', 'slides': len(prs.slides)},
                            'error': "未提取到文本内容"
                        }
                except Exception as e2:
                    if self.is_windows:
                        # 尝试使用Win32COM作为最后手段
                        try:
                            return self._extract_ppt_content(file_path)
                        except Exception as e3:
                            return self._create_error_result('pptx', 
                                f"所有提取方法失败: MarkItDown: {str(e)}; python-pptx: {str(e2)}; COM: {str(e3)}")
                    return self._create_error_result('pptx', 
                        f"所有提取方法失败: MarkItDown: {str(e)}; python-pptx: {str(e2)}")
        except Exception as e:
            return self._create_error_result('pptx', f"PPTX处理失败: {str(e)}")
    
    def _extract_xls_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLS文件内容，使用传统方法，增强容错性"""
        try:
            # 尝试使用xlrd读取
            try:
                sheets = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
                content = []
                
                for sheet_name, df in sheets.items():
                    if len(df) > 5000:
                        df = df.head(5000)
                        sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                    else:
                        sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                    content.append(sheet_content)
                
                return {
                    'content': '\n\n'.join(content),
                    'metadata': {'file_type': 'xls', 'extractor': 'xlrd', 'sheets': len(sheets)},
                    'error': None
                }
            except Exception as e:
                logger.warning(f"xlrd处理XLS失败: {file_path} - {str(e)}")
                
                # 尝试使用Win32COM作为备选
                if self.is_windows:
                    try:
                        pythoncom.CoInitialize()
                        excel_app = win32com.client.Dispatch("Excel.Application")
                        excel_app.Visible = False
                        excel_app.DisplayAlerts = False
                        
                        # 设置超时
                        timer = threading.Timer(30.0, lambda: self._force_close_office_processes('EXCEL.EXE'))
                        timer.start()
                        
                        workbook = excel_app.Workbooks.Open(os.path.abspath(file_path), ReadOnly=True)
                        content = []
                        
                        for i in range(1, workbook.Sheets.Count + 1):
                            sheet = workbook.Sheets(i)
                            sheet_content = [f"Sheet: {sheet.Name}"]
                            
                            # 获取已使用区域范围
                            used_range = sheet.UsedRange
                            if used_range:
                                # 获取内容为表格形式
                                table_data = []
                                for row in range(1, min(1001, used_range.Rows.Count + 1)):
                                    row_data = []
                                    for col in range(1, used_range.Columns.Count + 1):
                                        try:
                                            cell_value = sheet.Cells(row, col).Value
                                            row_data.append(str(cell_value) if cell_value is not None else '')
                                        except:
                                            row_data.append('')
                                    table_data.append(' | '.join(row_data))
                                
                                sheet_content.extend(table_data)
                                if used_range.Rows.Count > 1000:
                                    sheet_content.append(f"... (超过1000行，已截断)")
                            
                            content.append('\n'.join(sheet_content))
                        
                        workbook.Close(False)
                        excel_app.Quit()
                        
                        # 取消超时
                        timer.cancel()
                        
                        return {
                            'content': '\n\n'.join(content),
                            'metadata': {'file_type': 'xls', 'extractor': 'win32com', 'sheets': workbook.Sheets.Count},
                            'error': f"xlrd失败，使用COM: {str(e)}"
                        }
                    except Exception as e2:
                        return self._create_error_result('xls', f"所有提取方法失败: xlrd: {str(e)}; COM: {str(e2)}")
                    finally:
                        try:
                            timer.cancel()
                        except:
                            pass
                        
                        try:
                            if 'workbook' in locals():
                                workbook.Close(False)
                            if 'excel_app' in locals():
                                excel_app.Quit()
                        except:
                            pass
                        
                        self._force_close_office_processes('EXCEL.EXE')
                        pythoncom.CoUninitialize()
                
                return self._create_error_result('xls', f"XLS处理失败: {str(e)}")
        except Exception as e:
            return self._create_error_result('xls', f"XLS处理失败: {str(e)}")
    
    def _extract_archive_content(self, file_path: str) -> Dict[str, Any]:
        """提取压缩文件内容，使用安全的临时目录和路径验证，增强容错性"""
        # 生成唯一的临时目录，确保不冲突
        temp_dir = os.path.join(
            os.path.dirname(file_path), 
            f"temp_extract_{uuid.uuid4().hex}"
        )
        
        try:
            os.makedirs(temp_dir, exist_ok=True)
            
            # 使用子进程和超时控制提取压缩文件
            def extract_archive_with_timeout():
                try:
                    patoolib.extract_archive(file_path, outdir=temp_dir, verbosity=-1)
                    return True, None
                except Exception as e:
                    return False, str(e)
            
            # 创建线程执行提取
            result = [False, None]  # [成功标志, 错误]
            def extract_worker():
                result[0], result[1] = extract_archive_with_timeout()
            
            thread = threading.Thread(target=extract_worker)
            thread.daemon = True
            thread.start()
            thread.join(timeout=60)  # 60秒超时
            
            if thread.is_alive():
                # 提取超时
                return self._create_error_result(
                    Path(file_path).suffix[1:],
                    "压缩文件提取超时(>60秒)"
                )
            
            if not result[0]:
                # 提取出错，尝试读取部分内容作为HEX显示
                try:
                    with open(file_path, 'rb') as f:
                        header = f.read(1024)
                        hex_content = ' '.join('{:02X}'.format(b) for b in header)
                    
                    return {
                        'content': f"[压缩文件提取失败]\n错误: {result[1]}\n文件头(HEX): {hex_content[:200]}...",
                        'metadata': {'file_type': Path(file_path).suffix[1:]},
                        'error': f"压缩文件提取失败: {result[1]}"
                    }
                except:
                    return self._create_error_result(
                        Path(file_path).suffix[1:],
                        f"压缩文件处理失败: {result[1]} - 跳过"
                    )
            
            content = []
            max_files = 100
            max_size = 10 * 1024 * 1024  # 10MB
            total_size = 0
            files_processed = 0
            
            for root, _, files in os.walk(temp_dir):
                if files_processed >= max_files:
                    content.append(f"\n[警告: 超过最大处理文件数 {max_files}, 仅展示部分内容]")
                    break
                
                # 按名称排序，确保输出一致性
                sorted_files = sorted(files)
                for file_name in sorted_files:
                    if file_name.startswith('._') or file_name.startswith('__MACOSX'):
                        continue
                    
                    if files_processed >= max_files:
                        break
                    
                    # 安全路径检查
                    file_path_full = os.path.join(root, file_name)
                    normalized_path = os.path.normpath(file_path_full)
                    if not normalized_path.startswith(os.path.normpath(temp_dir)):
                        content.append(f"[警告: 跳过不安全路径: {file_name}]")
                        continue
                    
                    try:
                        # 文件大小检查
                        file_size = os.path.getsize(file_path_full)
                        if total_size + file_size > max_size:
                            content.append("\n[警告: 达到总大小限制，跳过剩余文件]")
                            break
                        
                        # 相对路径
                        rel_path = os.path.relpath(file_path_full, temp_dir)
                        
                        # 文本文件检查
                        if file_size < 1024 * 1024:  # 1MB
                            mime = magic.Magic(mime=True).from_file(file_path_full)
                            
                            # 尝试读取文本内容
                            if mime.startswith('text/') or mime in ['application/json', 'application/xml']:
                                try:
                                    with open(file_path_full, 'r', encoding='utf-8', errors='ignore') as f:
                                        text = f.read(10240)
                                        if text.strip():
                                            content.append(f"文件: {rel_path}\n{text[:1024]}...")
                                            files_processed += 1
                                            total_size += file_size
                                            continue
                                except:
                                    pass
                            
                            # 非文本或读取失败，仅显示文件信息
                            content.append(f"文件: {rel_path} (大小: {file_size} 字节, 类型: {mime})")
                            files_processed += 1
                    except OSError as oe:
                        content.append(f"[警告: 无法读取文件: {file_name} - {str(oe)}]")
                    except Exception as e:
                        content.append(f"[警告: 处理文件出错: {file_name} - {str(e)}]")
            
            if not content:
                content.append("[压缩文件解压成功，但未找到有效内容]")
            
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': Path(file_path).suffix[1:], 'files_processed': files_processed},
                'error': None
            }
        except patoolib.util.PatoolError as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过（请确保安装 7-Zip 或 unrar）"
            )
        except Exception as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过"
            )
        finally:
            try:
                # 强制清理临时目录
                if os.path.exists(temp_dir):
                    def remove_readonly(func, path, _):
                        try:
                            os.chmod(path, stat.S_IWRITE)
                            func(path)
                        except:
                            pass
                    
                    shutil.rmtree(temp_dir, onerror=remove_readonly)
            except Exception as e:
                logger.warning(f"清理临时目录失败: {temp_dir} - {str(e)}")
    
    def _force_close_office_processes(self, process_name: str = 'WINWORD.EXE'):
        """强制关闭指定的Office进程，加强可靠性"""
        try:
            logger.debug(f"尝试强制关闭进程: {process_name}")
            closed_count = 0
            
            for attempt in range(3):
                for proc in psutil.process_iter(['pid', 'name']):
                    try:
                        if proc.info['name'].upper() == process_name:
                            # 首先尝试温和终止
                            proc.terminate()
                            try:
                                proc.wait(timeout=0.5)  # 等待进程终止
                            except psutil.TimeoutExpired:
                                # 然后尝试强制终止
                                proc.kill()
                            
                            closed_count += 1
                            logger.info(f"强制关闭进程: {proc.info['name']} (PID: {proc.info['pid']})")
                    except (psutil.NoSuchProcess, psutil.AccessDenied, Exception) as e:
                        logger.debug(f"关闭进程时出错: {e}")
                        continue
                
                # 检查是否还有相关进程运行
                if not any(p.info['name'].upper() == process_name for p in psutil.process_iter(['pid', 'name'])):
                    break
                
                time.sleep(0.5)
            
            logger.debug(f"共关闭 {closed_count} 个 {process_name} 进程")
        except Exception as e:
            logger.warning(f"关闭进程失败: {e}")

    def __del__(self):
        """在对象销毁时清理资源"""
        try:
            self._cleanup_word_app()
        except Exception as e:
            # 对象销毁时不抛出异常
            logger.debug(f"清理资源时出错: {e}")
    
    def _extract_doc_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOC文件内容，使用复用的 Word 实例，增强容错性和可靠性"""
        logger.info(f'处理DOC文件: {file_path}')
        
        max_retries = 3
        last_error = None
        
        for attempt in range(max_retries):
            # 确保每次尝试都是独立的COM环境
            pythoncom.CoInitialize()
            doc = None
            timer = None
            
            try:
                if not self._init_word_app():
                    pythoncom.CoUninitialize()  # 清理COM资源
                    return self._create_error_result('doc', "无法初始化 Word 应用程序")
                
                # 设置超时，防止打开文件卡住
                timer = threading.Timer(20.0, lambda: self._force_close_office_processes('WINWORD.EXE'))
                timer.start()
                
                with self.word_lock:
                    # 通过绝对路径打开，避免相对路径问题
                    abs_path = os.path.abspath(file_path)
                    logger.debug(f"尝试打开文档: {abs_path}")
                    
                    try:
                        doc = self.word_app.Documents.Open(
                            abs_path,
                            ReadOnly=True,
                            PasswordDocument="",
                            Visible=False,
                            NoEncodingDialog=True
                        )
                    except Exception as open_error:
                        logger.warning(f"打开DOC文件失败: {str(open_error)}")
                        raise
                    
                    # 限制提取文本的大小
                    try:
                        if doc.Words.Count > 100000:
                            content = doc.Range(0, 100000).Text + "\n\n[注意: 文档过大，内容已截断]"
                        else:
                            content = doc.Range().Text
                    except:
                        # 备选方法：按段落提取
                        paragraphs = []
                        try:
                            for i in range(1, min(1000, doc.Paragraphs.Count + 1)):
                                paragraphs.append(doc.Paragraphs(i).Range.Text)
                            content = '\n'.join(paragraphs)
                            if doc.Paragraphs.Count > 1000:
                                content += "\n\n[注意: 文档过大，仅显示前1000个段落]"
                        except:
                            # 最后手段：尝试读取一些文本
                            try:
                                content = doc.Content.Text[:50000]
                            except:
                                content = "[无法提取文档内容]"
                    
                    # 安全关闭文档
                    try:
                        doc.Close(False)
                        doc = None
                    except:
                        pass
                    
                    # 取消超时
                    if timer:
                        timer.cancel()
                        timer = None
                    
                    result = {
                        'content': content,
                        'metadata': {'file_type': 'doc', 'extractor': 'win32com', 'attempt': attempt + 1},
                        'error': None
                    }
                    
                    pythoncom.CoUninitialize()  # 清理COM资源
                    return result
                    
            except Exception as e:
                last_error = str(e)
                logger.warning(f"第 {attempt + 1} 次尝试处理DOC失败: {file_path} - {last_error}")
                
                # 确保文档关闭
                if doc:
                    try:
                        doc.Close(False)
                    except:
                        pass
                
                # 取消超时
                if timer:
                    try:
                        timer.cancel()
                    except:
                        pass
                
                # 清理Word实例并重置COM
                self._cleanup_word_app()
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
                
                if attempt < max_retries - 1:
                    time.sleep(1)  # 休息一下再尝试
                
        # 所有尝试都失败，返回错误结果
        return self._create_error_result('doc', f"DOC处理失败 ({max_retries}次尝试): {last_error}")
    
    def _extract_ppt_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPT文件内容，使用传统方法，增强容错性"""
        logger.info(f'处理PPT文件: {file_path}')
        
        # 初始化COM
        pythoncom.CoInitialize()
        ppt_app = None
        presentation = None
        timer = None
        
        try:
            # 设置超时，防止打开文件卡住
            timer = threading.Timer(30.0, lambda: self._force_close_office_processes('POWERPNT.EXE'))
            timer.start()
            
            # 创建PowerPoint实例
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            
            # 通过绝对路径打开，避免相对路径问题
            abs_path = os.path.abspath(file_path)
            presentation = ppt_app.Presentations.Open(
                abs_path, 
                WithWindow=False,
                ReadOnly=True
            )
            
            # 取消超时
            timer.cancel()
            timer = None
            
            # 提取内容
            text_content = []
            slide_count = 0
            
            try:
                slide_count = presentation.Slides.Count
                for slide_idx in range(1, min(101, slide_count + 1)):
                    try:
                        slide = presentation.Slides(slide_idx)
                        slide_text = [f"=== 幻灯片 {slide_idx} ==="]
                        
                        # 提取形状中的文本
                        shape_count = 0
                        try:
                            shape_count = slide.Shapes.Count
                            for shape_idx in range(1, shape_count + 1):
                                try:
                                    shape = slide.Shapes(shape_idx)
                                    if shape.HasTextFrame:
                                        try:
                                            text = shape.TextFrame.TextRange.Text
                                            if text.strip():
                                                slide_text.append(text.strip())
                                        except:
                                            continue
                                except:
                                    continue
                        except:
                            slide_text.append("[无法读取幻灯片形状]")
                        
                        # 提取备注
                        try:
                            if hasattr(slide, 'NotesPage'):
                                notes_shapes = slide.NotesPage.Shapes
                                for shape_idx in range(1, notes_shapes.Count + 1):
                                    try:
                                        shape = notes_shapes(shape_idx)
                                        if shape.HasTextFrame:
                                            notes_text = shape.TextFrame.TextRange.Text
                                            if notes_text.strip() and "Click to edit Master text styles" not in notes_text:
                                                slide_text.append(f"备注: {notes_text.strip()}")
                                    except:
                                        continue
                        except:
                            pass
                        
                        if len(slide_text) > 1:  # 确保至少有内容
                            text_content.append('\n'.join(slide_text))
                    except Exception as slide_e:
                        logger.debug(f"处理幻灯片 {slide_idx} 出错: {str(slide_e)}")
                        text_content.append(f"=== 幻灯片 {slide_idx} ===\n[处理出错: {str(slide_e)}]")
                
                if slide_count > 100:
                    text_content.append(f"\n[注意: 幻灯片过多，仅显示前100张，共{slide_count}张]")
            except Exception as e:
                logger.warning(f"提取PPT内容时出错: {str(e)}")
                text_content.append(f"[提取PPT内容时出错: {str(e)}]")
            
            # 如果没有内容，提供基本信息
            if not text_content:
                try:
                    text_content.append(f"[PPT文件未提取到文本内容]\n标题: {presentation.Name}\n幻灯片数: {slide_count}")
                except:
                    text_content.append("[PPT文件未提取到文本内容]")
            
            content = "\n\n".join(text_content)
            
            return {
                'content': content,
                'metadata': {'file_type': 'ppt', 'extractor': 'win32com', 'slides': slide_count},
                'error': None
            }
        except Exception as e:
            logger.warning(f"Win32COM处理PPT失败: {file_path} - {str(e)}")
            return self._create_error_result('ppt', f"PPT处理失败: {str(e)}")
        finally:
            # 取消超时计时器
            if timer:
                try:
                    timer.cancel()
                except:
                    pass
            
            # 确保PPT关闭
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            
            # 确保应用程序退出
            if ppt_app:
                try:
                    ppt_app.Quit()
                except:
                    pass
            
            # 强制关闭可能残留的进程
            self._force_close_office_processes('POWERPNT.EXE')
            
            # 清理COM资源
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    
    def extract_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """提取文件内容，根据MIME类型选择合适的提取器，增强容错性"""
        # 检查文件是否存在
        if not Path(file_path).exists():
            return {'error': f'文件不存在: {file_path}', 'content': '', 'metadata': {}}
        
        # 检查文件是否为空
        try:
            if Path(file_path).stat().st_size == 0:
                return self._create_empty_result(mime_type)
        except Exception as e:
            return {'error': f'文件状态检查失败: {str(e)}', 'content': '', 'metadata': {}}
        
        # 处理特殊情况：加密文件
        if mime_type == "application/encrypted":
            return self._create_error_result(mime_type, "加密文件，跳过处理")
        
        # 预检查文件类型
        is_valid, error_msg = self._precheck_file_type(file_path, mime_type)
        if not is_valid:
            file_header = self._get_file_header(file_path)
            # 虽然预检查失败，但记录警告后仍然尝试提取
            logger.warning(f"文件预检查失败: {file_path} - {error_msg}")
        
        # 获取文件扩展名
        ext = Path(file_path).suffix.lower()
        
        # 首先尝试使用基于MIME类型的专用提取器
        if mime_type in self.extractors:
            try:
                result = self.extractors[mime_type](file_path)
                # 检查提取是否成功
                if result.get('content'):
                    return result
                
                # 如果提取失败但没有明确错误，记录日志并继续尝试其他方法
                if not result.get('error'):
                    logger.warning(f"使用 {mime_type} 提取器未能提取内容: {file_path}")
                
            except Exception as e:
                logger.warning(f"主提取器出错: {mime_type} - {file_path} - {str(e)}")
                # 不返回错误，继续尝试其他方法
        
        # 尝试根据扩展名选择提取器
        if ext in ('.xlsx', '.xlsm', '.xltx', '.xltm') and self.MIME_TYPE['XLSX'] in self.extractors:
            try:
                return self.extractors[self.MIME_TYPE['XLSX']](file_path)
            except Exception as e:
                logger.warning(f"XLSX提取器出错: {file_path} - {str(e)}")
                
        elif ext in ('.docx', '.docm') and self.MIME_TYPE['DOCX'] in self.extractors:
            try:
                return self.extractors[self.MIME_TYPE['DOCX']](file_path)
            except Exception as e:
                logger.warning(f"DOCX提取器出错: {file_path} - {str(e)}")
                
        elif ext in ('.pptx', '.pptm') and self.MIME_TYPE['PPTX'] in self.extractors:
            try:
                return self.extractors[self.MIME_TYPE['PPTX']](file_path)
            except Exception as e:
                logger.warning(f"PPTX提取器出错: {file_path} - {str(e)}")
                
        elif self.is_windows:
            if ext in ('.doc', '.dot') and self.MIME_TYPE['DOC'] in self.extractors:
                try:
                    return self.extractors[self.MIME_TYPE['DOC']](file_path)
                except Exception as e:
                    logger.warning(f"DOC提取器出错: {file_path} - {str(e)}")
                    
            elif ext in ('.ppt', '.pot') and self.MIME_TYPE['PPT'] in self.extractors:
                try:
                    return self.extractors[self.MIME_TYPE['PPT']](file_path)
                except Exception as e:
                    logger.warning(f"PPT提取器出错: {file_path} - {str(e)}")
                    
            elif ext in ('.xls', '.xlt') and self.MIME_TYPE['XLS'] in self.extractors:
                try:
                    return self.extractors[self.MIME_TYPE['XLS']](file_path)
                except Exception as e:
                    logger.warning(f"XLS提取器出错: {file_path} - {str(e)}")
        
        # 作为最后手段，尝试使用MarkItDown
        try:
            result = self.md.convert(file_path)
            return {
                'content': result.text_content,
                'metadata': {'file_type': mime_type, 'converter': 'markitdown'},
                'error': None
            }
        except Exception as e:
            logger.warning(f"MarkItDown转换失败: {file_path} - {str(e)}")
            
            # 如果所有方法都失败，尝试以文本方式读取
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(10240)  # 读取前10KB
                
                if content.strip():
                    return {
                        'content': content[:10000] + '\n\n[注意: 以文本模式读取，可能不准确]',
                        'metadata': {'file_type': mime_type, 'fallback': 'text'},
                        'error': "所有特定格式提取器失败，使用文本模式"
                    }
            except:
                pass
            
            # 最后尝试读取文件头部的二进制表示
            try:
                with open(file_path, 'rb') as f:
                    binary = f.read(1024)
                    hex_content = ' '.join('{:02X}'.format(b) for b in binary)
                
                return {
                    'content': f"[无法提取文件内容，显示文件头部十六进制]\n{hex_content}",
                    'metadata': {'file_type': mime_type, 'fallback': 'binary'},
                    'error': f"无法提取内容: {str(e)}"
                }
            except Exception as final_e:
                return self._create_error_result(mime_type, f"所有提取方法失败: {str(e)}; 读取二进制失败: {str(final_e)}")

class SensitiveChecker:
    """敏感内容检查器，使用正则表达式替代 Aho-Corasick"""
    
    def __init__(self, config_path: str = "sensitive_config.yaml"):
        """初始化敏感词配置"""
        self.config = self._load_config(config_path)
        self.all_keywords = self.config.get('security_marks', []) + \
                           [kw for cat in self.config.get('sensitive_patterns', {}).values() 
                            for kw in cat.get('keywords', [])]
        escaped_keywords = [re.escape(kw) for kw in self.all_keywords]
        self.keyword_pattern = re.compile('|'.join(escaped_keywords))
    
    def _load_config(self, config_path: str) -> Dict:
        """加载敏感词配置文件"""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                return yaml.safe_load(f)
        except Exception as e:
            logger.error(f"加载敏感词配置失败: {e}")
            return {}
    
    def check_content(self, text: str) -> List[Tuple[str, List[int]]]:
        """检查文本中的敏感词"""
        keyword_matches = {}
        for match in self.keyword_pattern.finditer(text or ''):
            keyword = match.group()
            if keyword not in keyword_matches:
                keyword_matches[keyword] = []
            keyword_matches[keyword].append(match.start())
        keyword_results = list(keyword_matches.items())
        
        structured_results = []
        for pattern, weight in self.config.get('structured_patterns', {}).items():
            matches = list(re.finditer(pattern, text or ''))
            if matches:
                positions = [m.start() for m in matches]
                structured_results.append((pattern, positions))
        
        number_results = []
        for pattern in self.config.get('number_patterns', []):
            matches = list(re.finditer(pattern, text or ''))
            if matches:
                positions = [m.start() for m in matches]
                number_results.append((pattern, positions))
        
        return keyword_results + structured_results + number_results

class ResultExporter:
    """处理结果导出器"""
    
    def export_to_json(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 JSON 文件"""
        export_data = [
            {
                'file_path': r.file_path,
                'mime_type': r.mime_type,
                'content_preview': r.content.get('content', '')[:200],
                'sensitive_words': [{'word': w, 'positions': p} for w, p in r.sensitive_words],
                'error': r.error,
                'processing_time': r.processing_time
            } for r in results
        ]
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)
    
    def export_to_excel(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 Excel 文件"""
        data = [
            {
                '文件路径': r.file_path,
                '文件类型': r.mime_type,
                '敏感词统计': '; '.join([f"{w}({len(p)}次)" for w, p in r.sensitive_words]),
                '处理时间(秒)': round(r.processing_time, 3),
                '错误信息': r.error or ''
            } for r in results
        ]
        df = pd.DataFrame(data)
        df.to_excel(output_path, index=False, engine='openpyxl')

class ResultMonitor:
    """结果监控器，用于实时输出处理进度和结果"""
    
    def __init__(self, output_csv: str = "processing_results.csv"):
        self.output_csv = output_csv
        self._init_csv()
    
    def _init_csv(self):
        """初始化 CSV 文件"""
        with open(self.output_csv, 'w', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([
                '处理时间', '文件路径', '文件类型', '文件大小(bytes)', '是否为空',
                '敏感词数量', '敏感词列表', '处理时长(秒)', '状态', '错误信息'
            ])
    
    def record_result(self, result: ProcessingResult):
        """记录单个处理结果"""
        try:
            file_size = Path(result.file_path).stat().st_size
        except:
            file_size = 0
        
        sensitive_words_count = len(result.sensitive_words)
        sensitive_words_list = '; '.join([f"{w}({len(p)}处)" for w, p in result.sensitive_words])
        
        status = '失败' if result.error else ('空文件' if result.content.get('is_empty') else 
                                              '已跳过' if result.content.get('skipped') else '成功')
        
        with open(self.output_csv, 'a', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([
                datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                result.file_path,
                result.mime_type,
                file_size,
                'Yes' if result.content.get('is_empty') else 'No',
                sensitive_words_count,
                sensitive_words_list,
                round(result.processing_time, 3),
                status,
                result.error or ''
            ])
        
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 处理文件: {result.file_path}")
        print(f"  类型: {result.mime_type}")
        print(f"  大小: {file_size} bytes")
        print(f"  状态: {status}")
        if sensitive_words_count > 0:
            print(f"  发现敏感词: {sensitive_words_list}")
        if result.error:
            print(f"  错误: {result.error}")
        print("-" * 80)

class FileProcessor:
    """优化后的文件处理器主类"""
    
    def __init__(self, config_path: str = "sensitive_config.yaml", 
                 monitor_output: str = "processing_results.csv",
                 chunk_size: int = 1000,
                 max_workers: Optional[int] = None,
                 is_windows: bool = True):
        self.detector = FileTypeDetector()
        self.extractor = ContentExtractor(detector=self.detector, is_windows=is_windows)
        self.checker = SensitiveChecker(config_path)
        self.exporter = ResultExporter()
        self.monitor = ResultMonitor(monitor_output)
        self.chunk_size = chunk_size
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self._mime_cache = {}
        self._file_size_cache = {}
    
    def _scan_directory(self, directory: str) -> Iterator[List[str]]:
        """使用生成器分批扫描目录文件"""
        current_chunk = []
        for entry in os.scandir(directory):
            try:
                if entry.is_file(follow_symlinks=False) and not entry.name.startswith('~$') and not self.detector._is_internal_stream(entry.name):
                    ext = Path(entry.path).suffix.lower()
                    if ext not in self.detector.SKIP_EXTENSIONS:
                        current_chunk.append(entry.path)
                        if len(current_chunk) >= self.chunk_size:
                            yield current_chunk
                            current_chunk = []
                elif entry.is_dir(follow_symlinks=False):
                    for sub_chunk in self._scan_directory(entry.path):
                        yield sub_chunk
            except (PermissionError, OSError) as e:
                logger.warning(f"访问文件/目录出错 {entry.path}: {e}")
                continue
        if current_chunk:
            yield current_chunk
    
    def _preload_file_info(self, file_paths: List[str]):
        """预加载文件信息以减少 I/O 操作"""
        for file_path in file_paths:
            try:
                stat = os.stat(file_path)
                self._file_size_cache[file_path] = stat.st_size
                mime_type = self.detector.detect_file_type(file_path)
                self._mime_cache[file_path] = mime_type
            except Exception as e:
                logger.warning(f"预加载文件信息失败 {file_path}: {e}")
    
    def _process_file_batch(self, file_paths: List[str]) -> List[ProcessingResult]:
        """处理一批文件"""
        results = []
        self._preload_file_info(file_paths)
        
        office_files = [fp for fp in file_paths if Path(fp).suffix.lower() in ('.doc', '.ppt')]
        other_files = [fp for fp in file_paths if fp not in office_files]
        
        for file_path in office_files:
            result = self.process_file(file_path)
            results.append(result)
            self.monitor.record_result(result)
        
        current_workers = min(len(other_files), self.max_workers)
        with ThreadPoolExecutor(max_workers=current_workers) as executor:
            future_to_file = {executor.submit(self.process_file, fp): fp for fp in other_files}
            for future in as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    result = future.result()
                    results.append(result)
                    self.monitor.record_result(result)
                except Exception as e:
                    logger.error(f"处理文件异常 {file_path}: {e}")
        
        if len(self._mime_cache) > 10000:
            self._mime_cache.clear()
            self._file_size_cache.clear()
        
        return results
    
    def process_directory(self, directory: str) -> List[ProcessingResult]:
        """处理目录下的所有文件"""
        results = []
        total_files = len(self.detector.get_all_files(directory))
        completed = 0
        
        print(f"\n开始处理目录: {directory}")
        print(f"共发现 {total_files} 个文件（已忽略 ~$ 开头、内部流文件及 .dwg/音频视频文件）")
        print("=" * 80)
        
        try:
            for file_chunk in self._scan_directory(directory):
                chunk_results = self._process_file_batch(file_chunk)
                results.extend(chunk_results)
                completed += len(file_chunk)
                print(f"\n已完成: {completed}/{total_files} ({completed/total_files*100:.1f}%)")
        except KeyboardInterrupt:
            print("\n用户中断处理...")
        except Exception as e:
            logger.error(f"处理目录时发生错误: {e}")
        finally:
            if results:
                try:
                    self.exporter.export_to_json(results, "interrupted_results.json")
                    self.exporter.export_to_excel(results, "interrupted_results.xlsx")
                except Exception as e:
                    logger.error(f"导出中断结果失败: {e}")
        
        print(f"\n处理完成！共处理 {len(results)}/{total_files} 个文件")
        print(f"详细结果已保存至: {self.monitor.output_csv}")
        print("=" * 80)
        return results
    
    def process_file(self, file_path: str) -> ProcessingResult:
        start_time = time.time()
        try:
            mime_type = self._mime_cache.get(file_path) or self.detector.detect_file_type(file_path)
            logger.info(f"Processing file: {file_path} (type: {mime_type})")
            
            file_size = self._file_size_cache.get(file_path, 0)
            if file_size == 0:
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={'is_empty': True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time
                )
            
            content = self.extractor.extract_content(file_path, mime_type)
            if content.get('error'):
                logger.warning(f"提取内容失败: {file_path} - MIME: {mime_type} - 错误: {content.get('error')}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content=content,
                    sensitive_words=[],
                    error=content.get('error'),
                    processing_time=time.time() - start_time
                )
            
            sensitive_words = self.checker.check_content(content.get('content', ''))
            if sensitive_words:
                logger.info(f"Found {len(sensitive_words)} sensitive words in {file_path}")
            
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type,
                content=content,
                sensitive_words=sensitive_words,
                error=None,
                processing_time=time.time() - start_time
            )
        except Exception as e:
            error_msg = f"处理文件失败: {str(e)}"
            logger.error(f"{error_msg} - {file_path}")
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type if 'mime_type' in locals() else "unknown",
                content={'error': error_msg},
                sensitive_words=[],
                error=error_msg,
                processing_time=time.time() - start_time
            )

def main():
    """主函数"""
    parser = argparse.ArgumentParser(description='文件敏感内容检测工具')
    parser.add_argument('path', help='要处理的文件或目录路径')
    parser.add_argument('--config', default='sensitive_config.yaml', help='敏感词配置文件路径')
    parser.add_argument('--output', default='results', help='输出结果文件名(不含扩展名)')
    parser.add_argument('--chunk-size', type=int, default=1000, help='每批处理的文件数量')
    parser.add_argument('--workers', type=int, default=None, help='最大工作线程数')
    parser.add_argument('--no-windows', action='store_true', help='指定非Windows平台，禁用win32com相关功能')
    
    args = parser.parse_args()
    
    try:
        processor = FileProcessor(
            config_path=args.config,
            monitor_output=f"{args.output}_processing.csv",
            chunk_size=args.chunk_size,
            max_workers=args.workers,
            is_windows=not args.no_windows
        )
        
        results = []
        if Path(args.path).is_file():
            results = [processor.process_file(args.path)]
        else:
            results = processor.process_directory(args.path)
        
        processor.exporter.export_to_json(results, f"{args.output}.json")
        processor.exporter.export_to_excel(results, f"{args.output}.xlsx")
        
        logger.info(f"处理完成，共处理 {len(results)} 个文件")
    except KeyboardInterrupt:
        print("\n用户中断程序执行")
        sys.exit(1)
    except Exception as e:
        logger.error(f"程序执行出错: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
