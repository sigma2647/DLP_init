import os
import yaml
import magic
import time
import logging
import threading
import uuid
import shutil
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any, Iterator
from concurrent.futures import ThreadPoolExecutor, as_completed
from markitdown import MarkItDown
from dataclasses import dataclass, field
import re
from pdfminer.high_level import extract_text as pdf_extract_text
import pandas as pd
import openpyxl
import csv
import json
from datetime import datetime
import argparse
import sys
import win32com.client
import pythoncom
import psutil
import patoolib
from pptx import Presentation

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('sensitive_detector.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class ProcessingResult:
    """数据类，用于存储文件处理结果"""
    file_path: str
    mime_type: str
    content: Dict[str, Any]
    sensitive_words: List[Tuple[str, List[int]]] = field(default_factory=list)
    error: Optional[str] = None
    processing_time: float = 0.0

class FileTypeDetector:
    """文件类型检测器，使用 python-magic 库"""
    
    MIME_TYPES = {
        ".txt": "text/plain", ".csv": "text/csv", ".xml": "text/xml", ".html": "text/html",
        ".htm": "text/html", ".json": "application/json", ".yaml": "application/yaml",
        ".yml": "application/yaml", ".md": "text/markdown", ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".zip": "application/zip", ".rar": "application/x-rar", ".7z": "application/x-7z-compressed",
        ".tar": "application/x-tar", ".gz": "application/gzip", ".bz2": "application/x-bzip2",
        ".pdf": "application/pdf", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
        ".gif": "image/gif", ".bmp": "image/bmp", ".svg": "image/svg+xml", ".webp": "image/webp",
        ".mp3": "audio/mpeg", ".mp4": "video/mp4", ".avi": "video/x-msvideo", ".mov": "video/quicktime",
        ".wav": "audio/wav", ".bin": "application/octet-stream", ".exe": "application/x-msdownload",
        ".dll": "application/x-msdownload",
    }
    MIME_TO_EXT = {mime: ext for ext, mime in MIME_TYPES.items()}
    SKIP_EXTENSIONS = {'.dwg', '.mp3', '.wav', '.mp4', '.avi', '.mkv', '.flv', '.mov'}

    def __init__(self):
        """初始化文件类型检测器"""
        self.mime = magic.Magic(mime=True)
    
    def get_all_files(self, directory: str) -> List[str]:
        """递归获取目录中的所有文件，忽略 ~ 开头或内部流文件及跳过文件"""
        files = []
        for f in Path(directory).rglob('*'):
            if f.is_file() and not f.name.startswith('~$') and not self._is_internal_stream(f.name):
                ext = f.suffix.lower()
                if ext not in self.SKIP_EXTENSIONS:
                    files.append(str(f))
        return files
    
    def _is_internal_stream(self, filename: str) -> bool:
        """判断文件是否为 Office 内部流文件（如 [5]SummaryInformation）"""
        return bool(re.match(r'^\[\d+\].+', filename))
    
    def get_file_info(self, file_path: str) -> Dict:
        """获取文件信息，包括 MIME 类型和文件头"""
        try:
            mime_type = self.detect_file_type(file_path)
            with open(file_path, "rb") as f:
                file_header = f.read(16).hex().upper()
            file_extension = Path(file_path).suffix.lower()
            return {
                "file_path": file_path,
                "mime_type": mime_type,
                "file_extension": file_extension,
                "file_header": file_header,
                "size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"获取文件信息失败: {e} - {file_path}")
            return {"file_path": file_path, "mime_type": "unknown", "error": str(e)}
    
    def detect_file_type(self, file_path: str) -> str:
        """检测文件 MIME 类型，优化文件类型检测"""
        try:
            normalized_path = os.path.normpath(file_path.encode('utf-8', errors='replace').decode('utf-8'))
            mime_type = self.mime.from_file(normalized_path)
            
            ext = Path(file_path).suffix.lower()
            
            # First check by extension for known Office files
            if ext == '.docx':
                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif ext == '.xlsx':
                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif ext == '.pptx':
                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif ext == '.doc':
                return 'application/msword'
            elif ext == '.xls':
                return 'application/vnd.ms-excel'
            elif ext == '.ppt':
                return 'application/vnd.ms-powerpoint'
            
            # For other files, check the file header
            with open(file_path, 'rb') as f:
                header = f.read(8)
                
                # Check for Office Open XML formats (they're ZIP-based)
                if header.startswith(b'PK\x03\x04'):
                    try:
                        import zipfile
                        with zipfile.ZipFile(file_path) as zf:
                            file_list = zf.namelist()
                            if 'word/document.xml' in file_list:
                                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                            elif 'xl/workbook.xml' in file_list:
                                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                            elif 'ppt/presentation.xml' in file_list:
                                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                    except:
                        pass
                    if ext in self.MIME_TYPES:
                        return self.MIME_TYPES[ext]
                    return 'application/zip'
                    
                elif header.startswith(b'\xD0\xCF\x11\xE0'):
                    if ext == '.doc':
                        return 'application/msword'
                    elif ext == '.xls':
                        return 'application/vnd.ms-excel'
                    elif ext == '.ppt':
                        return 'application/vnd.ms-powerpoint'
                    return 'application/msword' if mime_type == 'application/octet-stream' else mime_type
            
            if mime_type not in ('application/octet-stream', 'text/plain'):
                return mime_type
            
            return self.MIME_TYPES.get(ext, 'application/octet-stream')
            
        except Exception as e:
            logger.error(f"文件类型检测失败 {file_path}: {e}")
            return self.MIME_TYPES.get(Path(file_path).suffix.lower(), "application/octet-stream")

class ContentExtractor:
    """优化后的文件内容提取器"""
    
    MIME_TYPE = {
        'TEXT': 'text/plain',
        'CSV': 'text/csv',
        'PDF': 'application/pdf',
        'MARKDOWN': 'text/markdown',
        'DOCX': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'XLSX': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'XLS': 'application/vnd.ms-excel',
        'DOC': 'application/msword',
        'PPT': 'application/vnd.ms-powerpoint',
        'PPTX': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'ZIP': 'application/zip',
        'RAR': 'application/x-rar',
        'SEVENZ': 'application/x-7z-compressed'
    }
    
    def __init__(self, detector, is_windows: bool = True):
        self.detector = detector
        try:
            self.md = MarkItDown()
        except Exception as e:
            self.md = None
            logger.error(f"初始化 MarkItDown 失败: {str(e)}")
        self.is_windows = is_windows
        self.word_app = None
        self.word_lock = threading.Lock()
        
        self.extractors = {
            self.MIME_TYPE['TEXT']: self._extract_text_content,
            self.MIME_TYPE['CSV']: self._extract_csv_content,
            self.MIME_TYPE['PDF']: self._extract_pdf_content,
            self.MIME_TYPE['MARKDOWN']: self._extract_markdown_content,
            self.MIME_TYPE['ZIP']: self._extract_archive_content,
            self.MIME_TYPE['RAR']: self._extract_archive_content,
            self.MIME_TYPE['SEVENZ']: self._extract_archive_content,
            self.MIME_TYPE['DOCX']: self._extract_docx_content,
            self.MIME_TYPE['XLSX']: self._extract_xlsx_content,
            self.MIME_TYPE['PPTX']: self._extract_pptx_content,
        }
        
        if self.is_windows:
            self.extractors.update({
                self.MIME_TYPE['DOC']: self._extract_doc_content,
                self.MIME_TYPE['PPT']: self._extract_ppt_content,
                self.MIME_TYPE['XLS']: self._extract_xls_content,
            })
    
    def _init_word_app(self):
        """初始化 Word 应用程序实例"""
        if not self.is_windows:
            return False
        with self.word_lock:
            if self.word_app is None:
                try:
                    pythoncom.CoInitialize()
                    self.word_app = win32com.client.Dispatch("Word.Application")
                    self.word_app.Visible = False
                    logger.info("成功初始化 Word 应用程序")
                    return True
                except Exception as e:
                    logger.error(f"初始化 Word 应用程序失败: {str(e)}")
                    self.word_app = None
                    return False
            return True
    
    def _cleanup_word_app(self):
        """清理 Word 应用程序实例"""
        with self.word_lock:
            if self.word_app is not None:
                try:
                    self.word_app.Quit()
                except:
                    pass
                self._force_close_office_processes('WINWORD.EXE')
                self.word_app = None
                pythoncom.CoUninitialize()
                logger.info("成功清理 Word 应用程序")
    
    def _create_empty_result(self, file_type: str) -> Dict[str, Any]:
        """创建空结果对象"""
        return {
            'content': '',
            'metadata': {'file_type': file_type},
            'error': None,
            'is_empty': True
        }
    
    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果对象"""
        return {
            'content': '',
            'metadata': {'file_type': file_type},
            'error': error_msg
        }
    
    def _get_file_header(self, file_path: str, size: int = 16) -> str:
        """获取文件头部字节的十六进制表示"""
        try:
            with open(file_path, 'rb') as f:
                return f.read(size).hex().upper()
        except Exception as e:
            logger.error(f"无法读取文件头: {file_path} - {e}")
            return ""
    
    def _precheck_file_type(self, file_path: str, detected_mime: str) -> Tuple[bool, str]:
        """预检查文件类型和格式"""
        try:
            ext = Path(file_path).suffix.lower()
            expected_mime = self.detector.MIME_TYPES.get(ext, detected_mime)
            
            if detected_mime != expected_mime:
                logger.warning(f"检测到的 MIME 类型 {detected_mime} 与基于扩展名推断的 {expected_mime} 不匹配，文件: {file_path}")
            
            with open(file_path, 'rb') as f:
                header = f.read(8)
                
                if expected_mime == self.MIME_TYPE['PPT'] and not header.startswith(b'\xD0\xCF\x11\xE0'):
                    logger.warning(f"文件头不符合 .ppt 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE['PPTX'] and not header.startswith(b'PK\x03\x04'):
                    logger.warning(f"文件头不符合 .pptx 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE['XLSX'] and not header.startswith(b'PK\x03\x04'):
                    logger.warning(f"文件头不符合 .xls 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE['DOC'] and not header.startswith(b'\xD0\xCF\x11\xE0'):
                    logger.warning(f"文件头不符合 .doc 格式，文件: {file_path}")
                
            return True, ""
        except Exception as e:
            return False, f"预检查失败: {str(e)}"

    def _is_valid_xlsx(self, file_path: str) -> Tuple[bool, str]:
        """检查 .xlsx 文件是否有效"""
        try:
            with open(file_path, 'rb') as f:
                header = f.read(4)
                if not header.startswith(b'PK\x03\x04'):
                    return False, "文件头不符合 .xlsx 格式"
            
            wb = openpyxl.load_workbook(file_path, read_only=True)
            ws = wb.active
            has_data = False
            
            for row in ws.iter_rows(max_rows=10, values_only=True):
                if any(cell is not None for cell in row):
                    has_data = True
                    break
            
            wb.close()
            
            if not has_data:
                return False, "文件为空或无有效数据"
            
            return True, ""
        except openpyxl.utils.exceptions.InvalidFileException as e:
            return False, f"文件格式无效或损坏: {str(e)}"
        except Exception as e:
            return False, f"文件检查失败: {str(e)}"
    
    def _extract_text_content(self, file_path: str) -> Dict[str, Any]:
        """提取文本文件内容"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                'metadata': {'file_type': 'text'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('text', f"读取文本文件失败: {str(e)}")
    
    def _extract_csv_content(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容"""
        try:
            encodings = ['utf-8', 'latin1', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, engine='python')
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                df = pd.read_csv(file_path, encoding='utf-8', errors='ignore')
            
            if len(df) > 10000:
                df = df.head(10000)
                notice = "\n\n[注意: 文件过大，仅显示前10000行]"
            else:
                notice = ""
                
            return {
                'content': df.to_string(index=False) + notice,
                'metadata': {'file_type': 'csv'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('csv', f"CSV处理失败: {str(e)}")
    
    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文件内容"""
        try:
            text = pdf_extract_text(file_path)
            if len(text) > 100000:
                text = text[:100000] + "\n\n[注意: 文件过大，内容已截断]"
            return {
                'content': text,
                'metadata': {'file_type': 'pdf'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('pdf', f"PDF处理失败: {str(e)}")
    
    def _extract_markdown_content(self, file_path: str) -> Dict[str, Any]:
        """提取Markdown文件内容"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                'metadata': {'file_type': 'markdown'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('markdown', f"Markdown文件处理失败: {str(e)}")
    
    def _extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOCX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result('docx', "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                'content': result.text_content,
                'metadata': {'file_type': 'docx', 'extractor': 'markitdown'},
                'error': None
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理DOCX失败: {file_path} - {str(e)}")
            return self._create_error_result('docx', f"DOCX处理失败: {str(e)}")
    
    def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLSX文件内容"""
        try:
            sheets = pd.read_excel(file_path, sheet_name=None)
            content = []
            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': 'excel'},
                'error': None
            }
        except ImportError as e:
            return self._create_error_result('excel', f"缺少必要库: {str(e)}")
        except ValueError as e:
            return self._create_error_result('excel', f"文件格式错误: {str(e)}")
        except Exception as e:
            return self._create_error_result('excel', f"Excel处理失败: {str(e)}，类型: {type(e).__name__}")
    
    def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPTX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result('pptx', "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                'content': result.text_content,
                'metadata': {'file_type': 'pptx', 'extractor': 'markitdown'},
                'error': None
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理PPTX失败: {file_path} - {str(e)}")
            if self.is_windows:
                try:
                    prs = Presentation(file_path)
                    text_content = []
                    
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        
                        if slide_text:
                            text_content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
                    
                    content = "\n\n".join(text_content)
                    return {
                        'content': content,
                        'metadata': {'file_type': 'pptx', 'extractor': 'python-pptx'},
                        'error': None
                    }
                except Exception as e2:
                    return self._create_error_result('pptx', f"所有提取方法失败: MarkItDown: {str(e)}; python-pptx: {str(e2)}")
            return self._create_error_result('pptx', f"PPTX处理失败: {str(e)}")
    
    def _extract_xls_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLS文件内容"""
        return self._extract_xlsx_content(file_path)  # 统一处理 Excel 文件
    
    def _extract_archive_content(self, file_path: str) -> Dict[str, Any]:
        """提取压缩文件内容，使用安全的临时目录和路径验证"""
        temp_dir = os.path.join(
            os.path.dirname(file_path), 
            f"temp_extract_{uuid.uuid4().hex}"
        )
        
        try:
            os.makedirs(temp_dir, exist_ok=True)
            patoolib.extract_archive(file_path, outdir=temp_dir)
            
            content = []
            max_files = 100
            max_size = 10 * 1024 * 1024  # 10MB
            total_size = 0
            files_processed = 0
            
            for root, _, files in os.walk(temp_dir):
                if files_processed >= max_files:
                    content.append(f"\n[警告: 超过最大处理文件数 {max_files}, 仅展示部分内容]")
                    break
                
                for file_name in files:
                    if file_name.startswith('._') or file_name.startswith('__MACOSX'):
                        continue
                    
                    if files_processed >= max_files:
                        break
                    
                    file_path_full = os.path.join(root, file_name)
                    normalized_path = os.path.normpath(file_path_full)
                    if not normalized_path.startswith(temp_dir):
                        content.append(f"[警告: 跳过不安全路径: {file_name}]")
                        continue
                    
                    try:
                        file_size = os.path.getsize(file_path_full)
                        if total_size + file_size > max_size:
                            content.append("\n[警告: 达到总大小限制，跳过剩余文件]")
                            break
                            
                        if file_size < 1024 * 1024:  # 1MB
                            try:
                                with open(file_path_full, 'r', encoding='utf-8', errors='ignore') as f:
                                    text = f.read(10240)
                                    if text.strip():
                                        content.append(f"File: {file_name}\n{text[:1024]}...")
                                        files_processed += 1
                                        total_size += file_size
                            except (UnicodeDecodeError, IOError):
                                pass
                    except OSError:
                        content.append(f"[警告: 无法读取文件: {file_name}]")
            
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': Path(file_path).suffix[1:]},
                'error': None
            }
        except patoolib.util.PatoolError as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过（请确保安装 7-Zip 或 unrar）"
            )
        except Exception as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过"
            )
        finally:
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"清理临时目录失败: {temp_dir} - {str(e)}")
    
    def _extract_doc_content(self, file_path: str) -> Dict[str, Any]:
            """提取 DOC 文件内容，使用类级别的 Word 实例"""
            logger.info(f'处理 DOC 文件: {file_path}')
            max_retries = 3
            abs_path = os.path.abspath(file_path)  # 缓存绝对路径

            for attempt in range(max_retries):
                try:
                    doc = self.word_app.Documents.Open(
                        abs_path,
                        ReadOnly=True,
                        PasswordDocument="",
                        Visible=False,
                        NoEncodingDialog=True
                    )
                    content = doc.Range().Text
                    doc.Close(False)
                    return {
                        'content': content,
                        'metadata': {'file_type': 'doc', 'extractor': 'win32com'},
                        'error': None
                    }
                except Exception as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"第 {attempt + 1} 次尝试处理 {file_path} 失败: {str(e)}, 重试中...")
                        time.sleep(1)
                        continue
                    logger.error(f"Win32COM 处理失败: {str(e)} - 文件: {file_path}")
                    return self._create_error_result('doc', f"Win32COM 处理失败: {str(e)}")

    def _force_close_office_processes(self, process_name: str = 'WINWORD.EXE') -> None:
        """优化后的强制关闭指定 Office 进程方法，仅在必要时执行"""
        try:
            process_found = False
            for proc in psutil.process_iter(['pid', 'name']):
                try:
                    if proc.info['name'].upper() == process_name:
                        process_found = True
                        proc.terminate()
                        proc.wait(timeout=0.5)
                        if proc.is_running():
                            proc.kill()
                        logger.info(f"强制关闭进程: {proc.info['name']} (PID: {proc.info['pid']})")
                except (psutil.NoSuchProcess, psutil.AccessDenied) as e:
                    logger.debug(f"关闭进程时出错: {e}")
                    continue
            if not process_found:
                logger.debug(f"未找到 {process_name} 进程，无需关闭")
        except Exception as e:
            logger.warning(f"关闭进程失败: {e}")

    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果字典"""
        return {
            'content': None,
            'metadata': {'file_type': file_type, 'extractor': 'win32com'},
            'error': error_msg
        }

    def _cleanup_word_app(self) -> None:
        """清理 Word 应用程序"""
        if hasattr(self, 'word_app') and self.word_app:
            try:
                self.word_app.Quit()
            except Exception as e:
                logger.debug(f"关闭 Word 应用程序时出错: {e}")
                self._force_close_office_processes('WINWORD.EXE')

    def __del__(self):
        """对象销毁时清理资源"""
        try:
            self._cleanup_word_app()
            pythoncom.CoUninitialize()
        except Exception as e:
            logger.debug(f"清理资源时出错: {e}")
    
    def _extract_ppt_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPT文件内容，使用传统方法"""
        logger.info(f'处理PPT文件: {file_path}')
        
        try:
            pythoncom.CoInitialize()
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            
            presentation = ppt_app.Presentations.Open(
                os.path.abspath(file_path), 
                WithWindow=False
            )
            
            text_content = []
            for slide in presentation.Slides:
                slide_text = []
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        try:
                            text = shape.TextFrame.TextRange.Text
                            if text.strip():
                                slide_text.append(text.strip())
                        except:
                            continue
                
                if slide_text:
                    text_content.append(f"Slide {slide.SlideIndex}:\n" + "\n".join(slide_text))
            
            presentation.Close()
            content = "\n\n".join(text_content)
            
            return {
                'content': content,
                'metadata': {'file_type': 'ppt', 'extractor': 'win32com'},
                'error': None
            }
        except Exception as e:
            logger.warning(f"Win32COM处理PPT失败: {file_path} - {str(e)}")
            return self._create_error_result('ppt', f"PPT处理失败: {str(e)}")
        finally:
            if 'presentation' in locals():
                try:
                    presentation.Close()
                except:
                    pass
            if 'ppt_app' in locals():
                try:
                    ppt_app.Quit()
                except:
                    pass
            self._force_close_office_processes('POWERPNT.EXE')
            pythoncom.CoUninitialize()
    
    def extract_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        print(f"Extracting content for: {file_path}, MIME: {mime_type}")  # 打印提取开始
        if not Path(file_path).exists():
            print(f"File does not exist: {file_path}")
            return {'error': f'文件不存在: {file_path}', 'content': '', 'metadata': {}}
        
        if Path(file_path).stat().st_size == 0:
            print(f"File is empty: {file_path}")
            return self._create_empty_result(mime_type)
        
        if mime_type == "application/encrypted":
            print(f"Encrypted file, skipping: {file_path}")
            return self._create_error_result(mime_type, "加密文件，跳过处理")
        
        is_valid, error_msg = self._precheck_file_type(file_path, mime_type)
        if not is_valid:
            file_header = self._get_file_header(file_path)
            print(f"Precheck failed: {file_path}, Error: {error_msg}")
            return {
                'content': '',
                'metadata': {'file_type': mime_type, 'file_header': file_header},
                'error': f"预检查失败: {error_msg}"
            }
        
        ext = Path(file_path).suffix.lower()
        if mime_type in (self.MIME_TYPE['XLSX'], self.MIME_TYPE['DOCX'], self.MIME_TYPE['PPTX']) or \
           ext in ('.xlsx', '.xlsm', '.xltx', '.xltm', '.docx', '.pptx'):
            print(f"Using specific extractor for: {file_path}")
            return self.extractors[mime_type](file_path)
        elif self.is_windows and (mime_type in (self.MIME_TYPE['DOC'], self.MIME_TYPE['PPT'], self.MIME_TYPE['XLS']) or \
           ext in ('.doc', '.ppt', '.xls', '.xlt')):
            print(f"Using Windows-specific extractor for: {file_path}")
            return self.extractors[mime_type](file_path)
        
        if mime_type in self.extractors:
            print(f"Using registered extractor for: {file_path}")
            return self.extractors[mime_type](file_path)
        
        # 尝试使用 MarkItDown 转换，捕获异常
        if self.md is not None:
            try:
                print(f"Attempting MarkItDown conversion for: {file_path}")
                result = self.md.convert(file_path)
                print(f"MarkItDown conversion succeeded for: {file_path}")
                return {
                    'content': result.text_content,
                    'metadata': {'file_type': mime_type, 'converter': 'markitdown'},
                    'error': None
                }
            except Exception as e:
                error_msg = f"MarkItDown转换失败: {str(e)} - 跳过"
                print(f"MarkItDown failed: {file_path}, Error: {error_msg}")
                return self._create_error_result(mime_type, error_msg)
        else:
            error_msg = "MarkItDown 不可用，跳过处理"
            print(f"MarkItDown unavailable: {file_path}, Error: {error_msg}")
            return self._create_error_result(mime_type, error_msg)

class SensitiveChecker:
    """敏感内容检查器，使用正则表达式替代 Aho-Corasick"""
    
    def __init__(self, config_path: str = "sensitive_config.yaml"):
        """初始化敏感词配置"""
        self.config = self._load_config(config_path)
        self.all_keywords = self.config.get('security_marks', []) + \
                           [kw for cat in self.config.get('sensitive_patterns', {}).values() 
                            for kw in cat.get('keywords', [])]
        escaped_keywords = [re.escape(kw) for kw in self.all_keywords]
        self.keyword_pattern = re.compile('|'.join(escaped_keywords))
    
    def _load_config(self, config_path: str) -> Dict:
        """加载敏感词配置文件"""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                return yaml.safe_load(f)
        except Exception as e:
            logger.error(f"加载敏感词配置失败: {e}")
            return {}
    
    def check_content(self, text: str) -> List[Tuple[str, List[int]]]:
        """检查文本中的敏感词"""
        keyword_matches = {}
        for match in self.keyword_pattern.finditer(text or ''):
            keyword = match.group()
            if keyword not in keyword_matches:
                keyword_matches[keyword] = []
            keyword_matches[keyword].append(match.start())
        keyword_results = list(keyword_matches.items())
        
        structured_results = []
        for pattern, weight in self.config.get('structured_patterns', {}).items():
            matches = list(re.finditer(pattern, text or ''))
            if matches:
                positions = [m.start() for m in matches]
                structured_results.append((pattern, positions))
        
        number_results = []
        for pattern in self.config.get('number_patterns', []):
            matches = list(re.finditer(pattern, text or ''))
            if matches:
                positions = [m.start() for m in matches]
                number_results.append((pattern, positions))
        
        return keyword_results + structured_results + number_results

class ResultExporter:
    """处理结果导出器"""
    
    def export_to_json(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 JSON 文件"""
        export_data = [
            {
                'file_path': r.file_path,
                'mime_type': r.mime_type,
                'content_preview': r.content.get('content', '')[:200],
                'sensitive_words': [{'word': w, 'positions': p} for w, p in r.sensitive_words],
                'error': r.error,
                'processing_time': r.processing_time
            } for r in results
        ]
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)
    
    def export_to_excel(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 Excel 文件"""
        data = [
            {
                '文件路径': r.file_path,
                '文件类型': r.mime_type,
                '敏感词统计': '; '.join([f"{w}({len(p)}次)" for w, p in r.sensitive_words]),
                '处理时间(秒)': round(r.processing_time, 3),
                '错误信息': r.error or ''
            } for r in results
        ]
        df = pd.DataFrame(data)
        df.to_excel(output_path, index=False, engine='openpyxl')

class ResultMonitor:
    """结果监控器，用于实时输出处理进度和结果"""
    
    def __init__(self, output_csv: str = "processing_results.csv"):
        self.output_csv = output_csv
        self._init_csv()
    
    def _init_csv(self):
        """初始化 CSV 文件"""
        with open(self.output_csv, 'w', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([
                '处理时间', '文件路径', '文件类型', '文件大小(bytes)', '是否为空',
                '敏感词数量', '敏感词列表', '处理时长(秒)', '状态', '错误信息'
            ])
    
    def record_result(self, result: ProcessingResult):
        """记录单个处理结果"""
        try:
            file_size = Path(result.file_path).stat().st_size
        except:
            file_size = 0
        
        sensitive_words_count = len(result.sensitive_words)
        sensitive_words_list = '; '.join([f"{w}({len(p)}处)" for w, p in result.sensitive_words])
        
        status = '失败' if result.error else ('空文件' if result.content.get('is_empty') else 
                                              '已跳过' if result.content.get('skipped') else '成功')
        
        with open(self.output_csv, 'a', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([
                datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                result.file_path,
                result.mime_type,
                file_size,
                'Yes' if result.content.get('is_empty') else 'No',
                sensitive_words_count,
                sensitive_words_list,
                round(result.processing_time, 3),
                status,
                result.error or ''
            ])
        
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 处理文件: {result.file_path}")
        print(f"  类型: {result.mime_type}")
        print(f"  大小: {file_size} bytes")
        print(f"  状态: {status}")
        if sensitive_words_count > 0:
            print(f"  发现敏感词: {sensitive_words_list}")
        if result.error:
            print(f"  错误: {result.error}")
        print("-" * 80)

class FileProcessor:
    """优化后的文件处理器主类"""
    
    def __init__(self, config_path: str = "sensitive_config.yaml", 
                 monitor_output: str = "processing_results.csv",
                 chunk_size: int = 1000,
                 max_workers: Optional[int] = None,
                 is_windows: bool = True):
        self.detector = FileTypeDetector()
        self.extractor = ContentExtractor(detector=self.detector, is_windows=is_windows)
        self.checker = SensitiveChecker(config_path)
        self.exporter = ResultExporter()
        self.monitor = ResultMonitor(monitor_output)
        self.chunk_size = chunk_size
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self._mime_cache = {}
        self._file_size_cache = {}
    
    def _scan_directory(self, directory: str) -> Iterator[List[str]]:
        """使用生成器分批扫描目录文件"""
        current_chunk = []
        for entry in os.scandir(directory):
            try:
                if entry.is_file(follow_symlinks=False) and not entry.name.startswith('~$') and not self.detector._is_internal_stream(entry.name):
                    ext = Path(entry.path).suffix.lower()
                    if ext not in self.detector.SKIP_EXTENSIONS:
                        current_chunk.append(entry.path)
                        if len(current_chunk) >= self.chunk_size:
                            yield current_chunk
                            current_chunk = []
                elif entry.is_dir(follow_symlinks=False):
                    for sub_chunk in self._scan_directory(entry.path):
                        yield sub_chunk
            except (PermissionError, OSError) as e:
                logger.warning(f"访问文件/目录出错 {entry.path}: {e}")
                continue
        if current_chunk:
            yield current_chunk
    
    def _preload_file_info(self, file_paths: List[str]):
        """预加载文件信息以减少 I/O 操作"""
        for file_path in file_paths:
            try:
                stat = os.stat(file_path)
                self._file_size_cache[file_path] = stat.st_size
                mime_type = self.detector.detect_file_type(file_path)
                self._mime_cache[file_path] = mime_type
            except Exception as e:
                logger.warning(f"预加载文件信息失败 {file_path}: {e}")
    
    def _process_file_batch(self, file_paths: List[str]) -> List[ProcessingResult]:
        """处理一批文件"""
        results = []
        self._preload_file_info(file_paths)
        
        office_files = [fp for fp in file_paths if Path(fp).suffix.lower() in ('.doc', '.ppt')]
        other_files = [fp for fp in file_paths if fp not in office_files]
        
        for file_path in office_files:
            result = self.process_file(file_path)
            results.append(result)
            self.monitor.record_result(result)
        
        current_workers = min(len(other_files), self.max_workers)
        with ThreadPoolExecutor(max_workers=current_workers) as executor:
            future_to_file = {executor.submit(self.process_file, fp): fp for fp in other_files}
            for future in as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    result = future.result()
                    results.append(result)
                    self.monitor.record_result(result)
                except Exception as e:
                    logger.error(f"处理文件异常 {file_path}: {e}")
        
        if len(self._mime_cache) > 10000:
            self._mime_cache.clear()
            self._file_size_cache.clear()
        
        return results
    
    def process_directory(self, directory: str) -> List[ProcessingResult]:
        """处理目录下的所有文件"""
        results = []
        total_files = len(self.detector.get_all_files(directory))
        completed = 0
        
        print(f"\n开始处理目录: {directory}")
        print(f"共发现 {total_files} 个文件（已忽略 ~$ 开头、内部流文件及 .dwg/音频视频文件）")
        print("=" * 80)
        
        try:
            for file_chunk in self._scan_directory(directory):
                chunk_results = self._process_file_batch(file_chunk)
                results.extend(chunk_results)
                completed += len(file_chunk)
                print(f"\n已完成: {completed}/{total_files} ({completed/total_files*100:.1f}%)")
        except KeyboardInterrupt:
            print("\n用户中断处理...")
        except Exception as e:
            logger.error(f"处理目录时发生错误: {e}")
        finally:
            if results:
                try:
                    self.exporter.export_to_json(results, "interrupted_results.json")
                    self.exporter.export_to_excel(results, "interrupted_results.xlsx")
                except Exception as e:
                    logger.error(f"导出中断结果失败: {e}")
        
        print(f"\n处理完成！共处理 {len(results)}/{total_files} 个文件")
        print(f"详细结果已保存至: {self.monitor.output_csv}")
        print("=" * 80)
        return results
    

    def process_file(self, file_path: str) -> ProcessingResult:
        start_time = time.time()
        try:
            mime_type = self._mime_cache.get(file_path) or self.detector.detect_file_type(file_path)
            ext = Path(file_path).suffix.lower()
            print(f"Processing file: {file_path}, Extension: {ext}, MIME: {mime_type}")  # 打印文件信息
            
            file_size = self._file_size_cache.get(file_path, 0)
            
            # 定义常规文件的条件
            regular_file_conditions = (
                mime_type == 'image/vnd.dwg' or           # .dwg 文件（即使无后缀）
                mime_type == 'application/octet-stream' or # 无法识别的文件
                ext in self.detector.SKIP_EXTENSIONS       # 已定义的跳过扩展名
            )
            
            if regular_file_conditions:
                logger.info(f"识别为常规文件，跳过内容提取: {file_path} (type: {mime_type})")
                print(f"Skipping as regular file: {file_path}")  # 确认跳过
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={'content': '', 'metadata': {'file_type': 'regular'}, 'skipped': True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time
                )
            
            if file_size == 0:
                logger.info(f"空文件，跳过: {file_path}")
                print(f"Skipping empty file: {file_path}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={'is_empty': True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time
                )
            
            content = self.extractor.extract_content(file_path, mime_type)
            if content.get('error'):
                logger.warning(f"提取内容失败: {file_path} - MIME: {mime_type} - 错误: {content.get('error')}")
                print(f"Content extraction failed: {file_path}, Error: {content.get('error')}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content=content,
                    sensitive_words=[],
                    error=content.get('error'),
                    processing_time=time.time() - start_time
                )
            
            sensitive_result = self.checker.check_content(content.get('content', ''))
            sensitive_words = sensitive_result
            print(f"Processed successfully: {file_path}, Sensitive words: {len(sensitive_words)} found")
            
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type,
                content=content,
                sensitive_words=sensitive_words,
                error=None,
                processing_time=time.time() - start_time
            )
        except Exception as e:
            error_msg = f"处理文件失败: {str(e)}"
            logger.error(f"{error_msg} - {file_path}")
            print(f"Error occurred: {file_path}, Message: {error_msg}")  # 打印错误
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type if 'mime_type' in locals() else "unknown",
                content={'content': '', 'metadata': {'file_type': 'regular'}, 'skipped': True},
                sensitive_words=[],
                error=error_msg,
                processing_time=time.time() - start_time
            )

def main():
    """主函数"""
    parser = argparse.ArgumentParser(description='文件敏感内容检测工具')
    parser.add_argument('path', help='要处理的文件或目录路径')
    parser.add_argument('--config', default='sensitive_config.yaml', help='敏感词配置文件路径')
    parser.add_argument('--output', default='results', help='输出结果文件名(不含扩展名)')
    parser.add_argument('--chunk-size', type=int, default=1000, help='每批处理的文件数量')
    parser.add_argument('--workers', type=int, default=None, help='最大工作线程数')
    parser.add_argument('--no-windows', action='store_true', help='指定非Windows平台，禁用win32com相关功能')
    
    args = parser.parse_args()
    
    try:
        processor = FileProcessor(
            config_path=args.config,
            monitor_output=f"{args.output}_processing.csv",
            chunk_size=args.chunk_size,
            max_workers=args.workers,
            is_windows=not args.no_windows
        )
        
        results = []
        if Path(args.path).is_file():
            results = [processor.process_file(args.path)]
        else:
            results = processor.process_directory(args.path)
        
        processor.exporter.export_to_json(results, f"{args.output}.json")
        processor.exporter.export_to_excel(results, f"{args.output}.xlsx")
        
        logger.info(f"处理完成，共处理 {len(results)} 个文件")
    except KeyboardInterrupt:
        print("\n用户中断程序执行")
        sys.exit(1)
    except Exception as e:
        logger.error(f"程序执行出错: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
