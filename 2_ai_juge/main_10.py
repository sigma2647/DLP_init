import os
import yaml
import magic
import time
import xlrd
import logging
import threading
import uuid
import shutil
import binascii
import re
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any, Iterator
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pdfminer.high_level import extract_text as pdf_extract_text
import pandas as pd
from markitdown import MarkItDown
import openpyxl
import csv
import json
from datetime import datetime
import zipfile
import argparse
import sys
import win32com.client
import pythoncom
import psutil
from pptx import Presentation
from transformers import BertForSequenceClassification, BertTokenizer
import torch

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler("sensitive_detector.log", encoding="utf-8"),
        logging.StreamHandler(),
    ],
)
logger = logging.getLogger(__name__)


@dataclass
class ProcessingResult:
    """数据类，用于存储文件处理结果"""

    file_path: str
    mime_type: str
    content: Dict[str, Any]
    sensitive_words: List[Tuple[str, List[int]]] = field(default_factory=list)
    error: Optional[str] = None
    processing_time: float = 0.0


class FileSignatureDetector:
    """
    增强型文件签名检测器，使用文件头部签名、扩展名和内容特征识别文件类型
    """

    # FILE_SIGNATURES, MIME_TYPES, and SKIP_EXTENSIONS are kept as is

    def __init__(self):
        """初始化文件签名检测器，使用python-magic"""
        try:
            self.mime = magic.Magic(mime=True)
            logger.info("初始化magic库用于MIME类型检测")
        except ImportError:
            logger.warning("未安装python-magic库，部分文件检测功能将受限")
            self.mime = None
        except Exception as e:
            logger.error(f"初始化magic库出错: {e}")
            self.mime = None

        # 反向映射
        self.MIME_TO_EXT = {mime: ext for ext, mime in self.MIME_TYPES.items()}
        
        # 编译正则表达式，提高性能
        self.ole_pattern = re.compile(r'(WordDocument|Workbook|PowerPoint Document|Microsoft (?:Word|Excel|PowerPoint))', re.IGNORECASE)
        self.excel_cell_pattern = re.compile(r'Cell[^a-zA-Z0-9]', re.IGNORECASE)
        self.excel_sheet_pattern = re.compile(r'Sheet\d', re.IGNORECASE)
        self.internal_stream_pattern = re.compile(r'^\[\d+\].+|^~\$|^__SRP_\d+$|^__substg1\.0_|^__attach_version1\.0_|^__properties_version1\.0$|^MsoDataStore|Thumbs\.db$|\.tmp$')

    def _is_internal_stream(self, file_name: str) -> bool:
        """
        检查文件是否为内部流文件（如OLE文件中的那些）
        
        Args:
            file_name: 要检查的文件名
            
        Returns:
            如果文件是内部流，则为True，否则为False
        """
        return bool(self.internal_stream_pattern.search(file_name))

    def read_file_header(self, file_path: str, bytes_to_read: int = 32) -> str:
        """
        读取文件的前N个字节并返回十六进制字符串

        Args:
            file_path: 文件路径
            bytes_to_read: 要读取的字节数

        Returns:
            文件头的大写十六进制字符串
        """
        try:
            with open(file_path, "rb") as f:
                header = f.read(bytes_to_read)
            return binascii.hexlify(header).decode("ascii").upper()
        except Exception as e:
            logger.error(f"读取文件头出错: {file_path} - {e}")
            return ""

    def get_magic_mime_type(self, file_path: str) -> str:
        """
        使用python-magic检测MIME类型

        Args:
            file_path: 文件路径

        Returns:
            MIME类型字符串，失败时返回空字符串
        """
        if not self.mime:
            return ""

        try:
            # 处理非ASCII文件名或含空格的路径
            normalized_path = os.path.normpath(os.path.abspath(file_path))
            return self.mime.from_file(normalized_path)
        except Exception as e:
            logger.error(f"Magic MIME类型检测失败: {file_path} - {e}")
            return ""

    def inspect_zip_content(self, file_path: str) -> Optional[str]:
        """
        检查ZIP文件内容以确定是否为OOXML文档

        Args:
            file_path: ZIP文件路径

        Returns:
            如果确定，返回MIME类型，否则返回None
        """
        try:
            with zipfile.ZipFile(file_path) as zip_file:
                file_list = set(zip_file.namelist())  # 使用集合提高查找效率

                # 检查DOCX
                if "word/document.xml" in file_list:
                    return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                # 检查XLSX
                elif "xl/workbook.xml" in file_list:
                    return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                # 检查PPTX
                elif "ppt/presentation.xml" in file_list:
                    return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                # 检查ODT (OpenDocument Text)
                elif "mimetype" in file_list and "content.xml" in file_list:
                    try:
                        mimetype = zip_file.read("mimetype").decode('utf-8').strip()
                        if mimetype == "application/vnd.oasis.opendocument.text":
                            return "application/vnd.oasis.opendocument.text"
                        elif mimetype == "application/vnd.oasis.opendocument.spreadsheet":
                            return "application/vnd.oasis.opendocument.spreadsheet"
                        elif mimetype == "application/vnd.oasis.opendocument.presentation":
                            return "application/vnd.oasis.opendocument.presentation"
                    except:
                        pass

                # 如果没有找到特定标记，它是通用ZIP
                return "application/zip"
        except zipfile.BadZipFile:
            logger.warning(f"无效的ZIP文件: {file_path}")
            return None
        except Exception as e:
            logger.warning(f"ZIP内容检查失败: {file_path} - {e}")
            return None

    def inspect_ole_content(self, file_path: str) -> Optional[str]:
        """
        检查OLE文件内容以确定是否为特定Office文档类型

        Args:
            file_path: OLE文件路径

        Returns:
            如果确定，返回MIME类型，否则返回None
        """
        try:
            # 首先，通过扩展名快速判断
            ext = Path(file_path).suffix.lower()
            if ext == ".doc":
                return "application/msword"
            elif ext == ".xls":
                return "application/vnd.ms-excel"
            elif ext == ".ppt":
                return "application/vnd.ms-powerpoint"

            # 如果扩展名不能帮助，尝试深度分析
            # 读取较大的块以搜索标记
            with open(file_path, "rb") as f:
                content = f.read(16384)  # 读取16KB

            content_str = content.decode("latin1", errors="ignore")

            # 使用预编译正则表达式检查特定标记
            if self.ole_pattern.search(content_str):
                # 检查Word标记
                if "WordDocument" in content_str or "Microsoft Word" in content_str:
                    return "application/msword"
                
                # 检查Excel标记
                if "Workbook" in content_str or "Microsoft Excel" in content_str:
                    return "application/vnd.ms-excel"
                
                # 检查单元格和工作表模式（Excel特有）
                if self.excel_cell_pattern.search(content_str) and self.excel_sheet_pattern.search(content_str):
                    return "application/vnd.ms-excel"
                
                # 检查PowerPoint标记
                if "PowerPoint Document" in content_str or "Microsoft PowerPoint" in content_str:
                    return "application/vnd.ms-powerpoint"

            # 检查二进制模式，查找Excel特有的标记
            if b"Standard Jet DB" in content or b"Excel.Sheet" in content:
                return "application/vnd.ms-excel"

            # 如果没有找到特定标记，它是通用OLE文件
            return "application/x-ole-storage"
        except Exception as e:
            logger.warning(f"OLE内容检查失败: {file_path} - {e}")
            return None

    def detect_file_type(self, file_path: str) -> str:
        """
        综合检测文件类型，返回MIME类型

        Args:
            file_path: 文件路径

        Returns:
            检测到的MIME类型
        """
        try:
            # 基本文件信息
            file_size = os.path.getsize(file_path)
            file_extension = Path(file_path).suffix.lower()
            expected_mime = self.MIME_TYPES.get(
                file_extension, "application/octet-stream"
            )

            # 跳过非常小或空的文件
            if file_size == 0:
                return "application/octet-stream"

            # 缓存机制 - 使用文件路径和修改时间作为键
            file_mtime = os.path.getmtime(file_path)
            cache_key = f"{file_path}:{file_mtime}"
            
            # 可以在这里添加缓存查找逻辑
            # if cache_key in self.mime_cache:
            #    return self.mime_cache[cache_key]

            # 读取文件头部（前32字节）
            header_hex = self.read_file_header(file_path, 32)

            # 获取magic MIME类型（如果可用）
            magic_mime = self.get_magic_mime_type(file_path) if self.mime else ""

            # 通过文件头部签名检测
            detected_mime = "application/octet-stream"

            # 检查PDF签名
            if header_hex.startswith("25504446"):
                return "application/pdf"
            
            # 检查图像格式
            elif header_hex.startswith("FFD8FF"):
                return "image/jpeg"
            elif header_hex.startswith("89504E47"):
                return "image/png"
            elif header_hex.startswith("47494638"):
                return "image/gif"
            elif header_hex.startswith("424D"):
                return "image/bmp"
            
            # 检查音频/视频
            elif header_hex.startswith("494433") or header_hex.startswith("FFFB"):
                return "audio/mpeg"
            elif header_hex.startswith("52494646") and "57415645" in header_hex:
                return "audio/wav"
            elif header_hex.startswith("00000020667479704D3441"):
                return "video/mp4"
            
            # 检查ZIP签名
            elif header_hex.startswith("504B0304"):  # ZIP签名
                zip_mime = self.inspect_zip_content(file_path)
                if zip_mime:
                    return zip_mime

            # 检查OLE签名
            elif header_hex.startswith("D0CF11E0A1B11AE1"):  # OLE签名
                ole_mime = self.inspect_ole_content(file_path)
                if ole_mime:
                    return ole_mime

            # 检查压缩文件
            elif header_hex.startswith("526172"):  # RAR
                return "application/x-rar"
            elif header_hex.startswith("377ABCAF"):  # 7-Zip
                return "application/x-7z-compressed"
            elif header_hex.startswith("1F8B08"):  # GZIP
                return "application/gzip"

            # 如果通过签名无法识别，尝试使用magic库
            if magic_mime and magic_mime != "application/octet-stream":
                return magic_mime

            # 检查文本文件 - 尝试读取并解码前4KB
            if file_size < 1024 * 1024 * 5:  # 小于5MB的文件
                try:
                    with open(file_path, "rb") as f:
                        content = f.read(4096)
                    
                    # 尝试以UTF-8解码
                    try:
                        content.decode('utf-8')
                        
                        # 检查具体的文本格式
                        text_content = content.decode('utf-8')
                        
                        # 检查JSON
                        if file_extension == '.json' or (text_content.strip().startswith('{') and text_content.strip().endswith('}')):
                            return "application/json"
                        
                        # 检查XML
                        if file_extension == '.xml' or (text_content.strip().startswith('<') and text_content.strip().endswith('>')):
                            return "text/xml"
                        
                        # 检查CSV
                        if file_extension == '.csv' or ',' in text_content and '\n' in text_content:
                            comma_count = text_content.count(',')
                            newline_count = text_content.count('\n')
                            if comma_count > 0 and newline_count > 0 and comma_count / newline_count > 1:
                                return "text/csv"
                        
                        # 检查Markdown
                        if file_extension == '.md' or '##' in text_content or '```' in text_content:
                            return "text/markdown"
                        
                        # 检查YAML
                        if file_extension in ['.yml', '.yaml'] or ':' in text_content and '\n' in text_content:
                            return "application/yaml"
                        
                        # 默认为纯文本
                        return "text/plain"
                    except UnicodeDecodeError:
                        pass
                except:
                    pass

            # 最后尝试使用扩展名
            if file_extension and expected_mime != "application/octet-stream":
                return expected_mime

            # 如果一切都失败了，返回二进制流类型
            return "application/octet-stream"

        except Exception as e:
            logger.error(f"文件类型检测失败: {file_path} - {e}")
            return "application/octet-stream"

    def get_all_files(self, directory: str) -> List[str]:
        """
        递归获取目录中的所有文件，过滤掉某些类型

        Args:
            directory: 要扫描的目录路径

        Returns:
            文件路径列表
        """
        files = []
        try:
            # 更高效的文件遍历
            for root, dirs, file_names in os.walk(directory):
                # 跳过隐藏目录
                dirs[:] = [d for d in dirs if not d.startswith('.')]
                
                for file_name in file_names:
                    # 跳过隐藏文件和Office临时文件
                    if file_name.startswith(".") or file_name.startswith("~$"):
                        continue

                    # 跳过Office内部流文件
                    if self._is_internal_stream(file_name):
                        continue

                    file_path = os.path.join(root, file_name)
                    
                    # 跳过已知的二进制格式
                    ext = Path(file_path).suffix.lower()
                    if ext in self.SKIP_EXTENSIONS:
                        continue
                    
                    files.append(file_path)
        except Exception as e:
            logger.error(f"遍历目录失败: {directory} - {e}")
        
        return files

    def check_file_signature(self, file_path: str) -> Dict:
        """
        检查文件签名并返回详细信息

        Args:
            file_path: 文件路径

        Returns:
            包含详细文件签名信息的字典
        """
        try:
            mime_type = self.detect_file_type(file_path)

            # 获取文件头
            header_hex = self.read_file_header(file_path, 32)

            # 获取更多元数据
            file_size = os.path.getsize(file_path)
            creation_time = os.path.getctime(file_path)
            modification_time = os.path.getmtime(file_path)
            
            return {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_size": file_size,
                "extension": Path(file_path).suffix.lower(),
                "mime_type": mime_type,
                "header_hex": header_hex,
                "created_at": datetime.fromtimestamp(creation_time).strftime("%Y-%m-%d %H:%M:%S"),
                "modified_at": datetime.fromtimestamp(modification_time).strftime("%Y-%m-%d %H:%M:%S"),
            }
        except Exception as e:
            logger.error(f"检查文件签名失败: {file_path} - {e}")
            return {
                "file_path": file_path,
                "error": str(e)
            }

class ContentExtractor:
    """优化后的文件内容提取器"""

    MIME_TYPE = {
        "TEXT": "text/plain",
        """
        plantext csv 
        code like html json
        office documents like doc docx ppt pptx xls xlsx
        image & video
        zip file
        other
        - font
        """
    }

    def __init__(self, detector, is_windows: bool = True):
        self.detector = detector
        try:
            self.md = MarkItDown()
        except Exception as e:
            self.md = None
            logger.error(f"初始化 MarkItDown 失败: {str(e)}")
        self.is_windows = is_windows
        self.word_app = None
        self.word_lock = threading.Lock()

        self.extractors = {
            ["TEXT"]:      self._extract_text_content,
            ["CSV"]:       self._extract_csv_content,
            ["PDF"]:       self._extract_pdf_content,
            ["MARKDOWN"]:  self._extract_markdown_content,
            ["ZIP"]:       self._extract_archive_content,
            ["RAR"]:       self._extract_archive_content,
            ["SEVENZ"]:    self._extract_archive_content,
            ["DOCX"]:      self._extract_docx_content,
            ["XLSX"]:      self._extract_xlsx_content,
            ["PPTX"]:      self._extract_pptx_content,
        }

        if self.is_windows:
            # 尝试初始化Word应用
            if self._init_word_app():
                self.extractors.update(
                    {
                        self.MIME_TYPE["DOC"]: self._extract_doc_content,
                    }
                )
            else:
                logger.warning("Word应用程序初始化失败，.doc文件将使用备用方法处理")

            # PowerPoint不需要持久应用实例
            self.extractors.update(
                {
                    self.MIME_TYPE["PPT"]: self._extract_ppt_content,
                    self.MIME_TYPE["XLS"]: self._extract_xls_content,
                }
            )

    def _init_word_app(self):
        """初始化Word应用程序实例"""
        if not self.is_windows:
            return False

        with self.word_lock:
            if self.word_app is None:
                try:
                    # 强制关闭任何现有Word进程
                    self._force_close_office_processes("WINWORD.EXE")

                    # 初始化当前线程的COM
                    pythoncom.CoInitialize()

                    # 创建新的Word应用实例
                    self.word_app = win32com.client.Dispatch("Word.Application")
                    self.word_app.Visible = False
                    self.word_app.DisplayAlerts = False

                    logger.info("成功初始化Word应用程序")
                    return True
                except Exception as e:
                    logger.error(f"初始化Word应用程序失败: {str(e)}")
                    self.word_app = None

                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass

                    return False
            return True

    def _cleanup_word_app(self):
        """清理 Word 应用程序实例"""
        with self.word_lock:
            if self.word_app is not None:
                try:
                    self.word_app.Quit()
                except:
                    pass
                self._force_close_office_processes("WINWORD.EXE")
                self.word_app = None
                pythoncom.CoUninitialize()
                logger.info("成功清理 Word 应用程序")

    def _create_empty_result(self, file_type: str) -> Dict[str, Any]:
        """创建空结果对象"""
        return {
            "content": "",
            "metadata": {"file_type": file_type},
            "error": None,
            "is_empty": True,
        }

    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果对象"""
        return {"content": "", "metadata": {"file_type": file_type}, "error": error_msg}

    def _get_file_header(self, file_path: str, size: int = 16) -> str:
        """获取文件头部字节的十六进制表示"""
        try:
            with open(file_path, "rb") as f:
                return f.read(size).hex().upper()
        except Exception as e:
            logger.error(f"无法读取文件头: {file_path} - {e}")
            return ""

    def _precheck_file_type(
        self, file_path: str, detected_mime: str
    ) -> Tuple[bool, str]:
        """预检查文件类型和格式"""
        try:
            ext = Path(file_path).suffix.lower()
            expected_mime = self.detector.MIME_TYPES.get(ext, detected_mime)

            if detected_mime != expected_mime:
                logger.warning(
                    f"检测到的 MIME 类型 {detected_mime} 与基于扩展名推断的 {expected_mime} 不匹配，文件: {file_path}"
                )

            with open(file_path, "rb") as f:
                header = f.read(8)

                if expected_mime == self.MIME_TYPE["PPT"] and not header.startswith(
                    b"\xd0\xcf\x11\xe0"
                ):
                    logger.warning(f"文件头不符合 .ppt 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE["PPTX"] and not header.startswith(
                    b"PK\x03\x04"
                ):
                    logger.warning(f"文件头不符合 .pptx 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE["XLSX"] and not header.startswith(
                    b"PK\x03\x04"
                ):
                    logger.warning(f"文件头不符合 .xls 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE["DOC"] and not header.startswith(
                    b"\xd0\xcf\x11\xe0"
                ):
                    logger.warning(f"文件头不符合 .doc 格式，文件: {file_path}")

            return True, ""
        except Exception as e:
            return False, f"预检查失败: {str(e)}"

    def _is_valid_xlsx(self, file_path: str) -> Tuple[bool, str]:
        """检查 .xlsx 文件是否有效"""
        try:
            with open(file_path, "rb") as f:
                header = f.read(4)
                if not header.startswith(b"PK\x03\x04"):
                    return False, "文件头不符合 .xlsx 格式"

            wb = openpyxl.load_workbook(file_path, read_only=True)
            ws = wb.active
            has_data = False

            for row in ws.iter_rows(max_rows=10, values_only=True):
                if any(cell is not None for cell in row):
                    has_data = True
                    break

            wb.close()

            if not has_data:
                return False, "文件为空或无有效数据"

            return True, ""
        except openpyxl.utils.exceptions.InvalidFileException as e:
            return False, f"文件格式无效或损坏: {str(e)}"
        except Exception as e:
            return False, f"文件检查失败: {str(e)}"

    def _extract_text_content(self, file_path: str) -> Dict[str, Any]:
        """提取文本文件内容"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                "content": content[:1000000]
                + (
                    "\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""
                ),
                "metadata": {"file_type": "text"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result("text", f"读取文本文件失败: {str(e)}")

    def _extract_csv_content(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容"""
        try:
            encodings = ["utf-8", "latin1", "cp1252"]
            df = None

            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, engine="python")
                    break
                except UnicodeDecodeError:
                    continue

            if df is None:
                df = pd.read_csv(file_path, encoding="utf-8", errors="ignore")

            if len(df) > 10000:
                df = df.head(10000)
                notice = "\n\n[注意: 文件过大，仅显示前10000行]"
            else:
                notice = ""

            return {
                "content": df.to_string(index=False) + notice,
                "metadata": {"file_type": "csv"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result("csv", f"CSV处理失败: {str(e)}")

    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文件内容"""
        try:
            text = pdf_extract_text(file_path)
            if len(text) > 100000:
                text = text[:100000] + "\n\n[注意: 文件过大，内容已截断]"
            return {"content": text, "metadata": {"file_type": "pdf"}, "error": None}
        except Exception as e:
            return self._create_error_result("pdf", f"PDF处理失败: {str(e)}")

    def _extract_markdown_content(self, file_path: str) -> Dict[str, Any]:
        """提取Markdown文件内容"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                "content": content[:1000000]
                + (
                    "\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""
                ),
                "metadata": {"file_type": "markdown"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(
                "markdown", f"Markdown文件处理失败: {str(e)}"
            )

    def _extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOCX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result("docx", "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                "content": result.text_content,
                "metadata": {"file_type": "docx", "extractor": "markitdown"},
                "error": None,
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理DOCX失败: {file_path} - {str(e)}")
            return self._create_error_result("docx", f"DOCX处理失败: {str(e)}")

    def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLSX文件内容"""
        try:
            sheets = pd.read_excel(file_path, sheet_name=None)
            content = []
            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)
            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": "excel"},
                "error": None,
            }
        except ImportError as e:
            return self._create_error_result("excel", f"缺少必要库: {str(e)}")
        except ValueError as e:
            return self._create_error_result("excel", f"文件格式错误: {str(e)}")
        except Exception as e:
            return self._create_error_result(
                "excel", f"Excel处理失败: {str(e)}，类型: {type(e).__name__}"
            )

    def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPTX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result("pptx", "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                "content": result.text_content,
                "metadata": {"file_type": "pptx", "extractor": "markitdown"},
                "error": None,
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理PPTX失败: {file_path} - {str(e)}")
            if self.is_windows:
                try:
                    prs = Presentation(file_path)
                    text_content = []

                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())

                        if slide_text:
                            text_content.append(
                                f"Slide {i + 1}:\n" + "\n".join(slide_text)
                            )

                    content = "\n\n".join(text_content)
                    return {
                        "content": content,
                        "metadata": {"file_type": "pptx", "extractor": "python-pptx"},
                        "error": None,
                    }
                except Exception as e2:
                    return self._create_error_result(
                        "pptx",
                        f"所有提取方法失败: MarkItDown: {str(e)}; python-pptx: {str(e2)}",
                    )
            return self._create_error_result("pptx", f"PPTX处理失败: {str(e)}")

    def _extract_xls_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLS文件内容，增强错误处理和多级回退策略"""
        logger.info(f"处理XLS文件: {file_path}")

        # 策略1: 尝试使用pandas通用方法
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="xlrd")
            content = []

            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)

            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": "excel", "extractor": "pandas-xlrd"},
                "error": None,
            }
        except xlrd.biffh.XLRDError as xlrd_error:
            # 特殊处理XLRD的"Can't find workbook in OLE2 compound document"错误
            logger.warning(f"XLRD无法解析XLS文件: {file_path} - {str(xlrd_error)}")
            # 继续尝试其他方法
        except Exception as pandas_error:
            logger.warning(f"Pandas处理XLS失败: {file_path} - {str(pandas_error)}")
            # 继续尝试其他方法

        # 策略2: 如果在Windows上，尝试使用COM接口
        if self.is_windows:
            try:
                pythoncom.CoInitialize()
                excel_app = win32com.client.Dispatch("Excel.Application")
                excel_app.Visible = False
                excel_app.DisplayAlerts = False

                workbook = excel_app.Workbooks.Open(
                    os.path.abspath(file_path), ReadOnly=True, UpdateLinks=False
                )

                content = []
                for sheet_index in range(1, workbook.Sheets.Count + 1):
                    try:
                        worksheet = workbook.Sheets(sheet_index)
                        sheet_name = worksheet.Name

                        # 获取已使用范围
                        used_range = worksheet.UsedRange
                        if used_range.Rows.Count > 0 and used_range.Columns.Count > 0:
                            # 创建二维数组存储数据
                            data = []
                            for row in range(1, min(5001, used_range.Rows.Count + 1)):
                                row_data = []
                                for col in range(1, used_range.Columns.Count + 1):
                                    try:
                                        cell_value = worksheet.Cells(row, col).Value
                                        row_data.append(
                                            str(cell_value)
                                            if cell_value is not None
                                            else ""
                                        )
                                    except:
                                        row_data.append("")
                                data.append(row_data)

                            # 格式化为类似表格的文本
                            if data:
                                max_lengths = [
                                    max(
                                        len(str(row[i])) for row in data if i < len(row)
                                    )
                                    for i in range(len(data[0]))
                                ]
                                formatted_rows = []

                                for row in data:
                                    formatted_row = [
                                        str(val).ljust(max_lengths[i])
                                        for i, val in enumerate(row)
                                        if i < len(max_lengths)
                                    ]
                                    formatted_rows.append(" | ".join(formatted_row))

                                sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(
                                    formatted_rows
                                )
                                if used_range.Rows.Count > 5000:
                                    sheet_text += "\n[注意: 表格过大，仅显示前5000行]"

                                content.append(sheet_text)
                    except Exception as sheet_error:
                        logger.warning(
                            f"处理工作表 {sheet_index} 时出错: {str(sheet_error)}"
                        )
                        continue

                workbook.Close(SaveChanges=False)
                excel_app.Quit()

                return {
                    "content": "\n\n".join(content),
                    "metadata": {"file_type": "excel", "extractor": "win32com"},
                    "error": None,
                }
            except Exception as com_error:
                logger.warning(
                    f"使用COM接口处理XLS失败: {file_path} - {str(com_error)}"
                )
            finally:
                try:
                    if "workbook" in locals() and workbook is not None:
                        workbook.Close(SaveChanges=False)
                except:
                    pass

                try:
                    if "excel_app" in locals() and excel_app is not None:
                        excel_app.Quit()
                except:
                    pass

                self._force_close_office_processes("EXCEL.EXE")
                pythoncom.CoUninitialize()

        # 策略3: 尝试openpyxl (虽然通常用于xlsx，但有时也能处理某些xls)
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

            content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                data = []
                row_count = 0

                for row in sheet.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else "" for cell in row])
                    row_count += 1
                    if row_count >= 5000:
                        break

                if data:
                    max_lengths = [
                        max(len(str(row[i])) for row in data if i < len(row))
                        for i in range(len(data[0]))
                    ]
                    formatted_rows = []

                    for row in data:
                        formatted_row = [
                            str(val).ljust(max_lengths[i])
                            for i, val in enumerate(row)
                            if i < len(max_lengths)
                        ]
                        formatted_rows.append(" | ".join(formatted_row))

                    sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(formatted_rows)
                    if row_count >= 5000:
                        sheet_text += "\n[注意: 表格过大，仅显示前5000行]"

                    content.append(sheet_text)

            workbook.close()

            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": "excel", "extractor": "openpyxl"},
                "error": None,
            }
        except Exception as openpyxl_error:
            logger.warning(
                f"使用openpyxl处理XLS失败: {file_path} - {str(openpyxl_error)}"
            )

        # 最后策略: 尝试提取二进制文本
        try:
            binary_result = self._extract_binary_text_content(
                file_path, "application/vnd.ms-excel"
            )
            if not binary_result.get("error"):
                return binary_result
        except Exception as e:
            pass

        # 所有方法失败，返回错误
        return self._create_error_result(
            "excel", "无法提取XLS文件内容: 所有已知方法均已尝试并失败"
        )

    def _extract_archive_content(self, file_path: str) -> Dict[str, Any]:
        """提取压缩文件内容，使用安全的临时目录和路径验证"""
        temp_dir = os.path.join(
            os.path.dirname(file_path), f"temp_extract_{uuid.uuid4().hex}"
        )

        try:
            os.makedirs(temp_dir, exist_ok=True)
            patoolib.extract_archive(file_path, outdir=temp_dir)

            content = []
            max_files = 100
            max_size = 10 * 1024 * 1024  # 10MB
            total_size = 0
            files_processed = 0

            for root, _, files in os.walk(temp_dir):
                if files_processed >= max_files:
                    content.append(
                        f"\n[警告: 超过最大处理文件数 {max_files}, 仅展示部分内容]"
                    )
                    break

                for file_name in files:
                    if file_name.startswith("._") or file_name.startswith("__MACOSX"):
                        continue

                    if files_processed >= max_files:
                        break

                    file_path_full = os.path.join(root, file_name)
                    normalized_path = os.path.normpath(file_path_full)
                    if not normalized_path.startswith(temp_dir):
                        content.append(f"[警告: 跳过不安全路径: {file_name}]")
                        continue

                    try:
                        file_size = os.path.getsize(file_path_full)
                        if total_size + file_size > max_size:
                            content.append("\n[警告: 达到总大小限制，跳过剩余文件]")
                            break

                        if file_size < 1024 * 1024:  # 1MB
                            try:
                                with open(
                                    file_path_full,
                                    "r",
                                    encoding="utf-8",
                                    errors="ignore",
                                ) as f:
                                    text = f.read(10240)
                                    if text.strip():
                                        content.append(
                                            f"File: {file_name}\n{text[:1024]}..."
                                        )
                                        files_processed += 1
                                        total_size += file_size
                            except (UnicodeDecodeError, IOError):
                                pass
                    except OSError:
                        content.append(f"[警告: 无法读取文件: {file_name}]")

            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": Path(file_path).suffix[1:]},
                "error": None,
            }
        except patoolib.util.PatoolError as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过（请确保安装 7-Zip 或 unrar）",
            )
        except Exception as e:
            return self._create_error_result(
                Path(file_path).suffix[1:], f"压缩文件处理失败: {str(e)} - 跳过"
            )
        finally:
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"清理临时目录失败: {temp_dir} - {str(e)}")

    def _extract_doc_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOC文件内容，使用专用的Word实例，改进COM对象管理"""
        logger.info(f"处理DOC文件: {file_path}")
        max_retries = 3
        timeout = 10  # 默认超时时间(秒)

        return {
            "content": "",
            "metadata": {"file_type": "doc", "extractor": "win32com"},
            "error": "所有处理尝试均失败",
        }

    def _extract_ppt_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPT文件内容，处理Visible属性错误和多级回退策略"""
        logger.info(f"处理PPT文件: {file_path}")
        max_retries = 3
        timeout = 10  # 默认超时时间(秒)

        for attempt in range(max_retries):
            try:
                # 设置超时标志和结果变量
                timeout_occurred = [False]
                result = [None]
                exception = [None]
                processing_completed = [False]

                # 主处理函数
                def process_ppt():
                    try:
                        # 初始化COM
                        pythoncom.CoInitialize()
                        ppt_app = None
                        presentation = None

                        try:
                            # 创建PowerPoint应用实例
                            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                            # 不设置Visible属性，避免某些环境下的错误

                            # 打开演示文稿
        self._force_close_office_processes("POWERPNT.EXE")
        return {
            "content": "",
            "metadata": {"file_type": "ppt", "extractor": "win32com"},
            "error": "所有处理尝试均失败",
        }



    def _extract_ppt_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPT文件内容，处理Visible属性错误和多级回退策略"""
        logger.info(f"处理PPT文件: {file_path}")
        max_retries = 3
        timeout = 10  # 默认超时时间(秒)

        for attempt in range(max_retries):
            try:
                # 设置超时标志和结果变量
                timeout_occurred = [False]
                result = [None]
                exception = [None]
                processing_completed = [False]

                # 主处理函数
                def process_ppt():
                    try:
                        # 初始化COM
                        pythoncom.CoInitialize()
                        ppt_app = None
                        presentation = None

                        try:
                            # 创建PowerPoint应用实例
                            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                            # 不设置Visible属性，避免某些环境下的错误

                            # 打开演示文稿
                            abs_path = os.path.abspath(file_path)
                            presentation = ppt_app.Presentations.Open(
                                abs_path, WithWindow=False, ReadOnly=True
                            )

                            # 提取文本内容
                            text_content = []
                            slide_count = presentation.Slides.Count

                            for i in range(1, slide_count + 1):
                                try:
                                    slide = presentation.Slides.Item(i)
                                    slide_text = []

                                    for shape_idx in range(1, slide.Shapes.Count + 1):
                                        try:
                                            shape = slide.Shapes.Item(shape_idx)
                                            if shape.HasTextFrame:
                                                text_frame = shape.TextFrame
                                                if (
                                                    hasattr(text_frame, "TextRange")
                                                    and text_frame.TextRange.Text.strip()
                                                ):
                                                    slide_text.append(
                                                        text_frame.TextRange.Text.strip()
                                                    )
                                        except:
                                            continue

                                    if slide_text:
                                        text_content.append(
                                            f"Slide {i}:\n" + "\n".join(slide_text)
                                        )
                                except:
                                    continue

                            content = "\n\n".join(text_content)

                            # 设置结果
                            result[0] = {
                                "content": content,
                                "metadata": {
                                    "file_type": "ppt",
                                    "extractor": "win32com",
                                    "slide_count": slide_count,
                                },
                                "error": None
                                if content.strip()
                                else "未提取到任何文本内容",
                            }

                        except Exception as e:
                            # 捕获处理异常
                            exception[0] = e
                        finally:
                            # 确保资源释放
                            try:
                                if presentation:
                                    presentation.Close()
                            except:
                                pass
                            try:
                                if ppt_app:
                                    ppt_app.Quit()
                            except:
                                pass
                            try:
                                pythoncom.CoUninitialize()
                            except:
                                pass

                            # 标记处理完成
                            processing_completed[0] = True
                    except Exception as e:
                        # 捕获总体异常
                        exception[0] = e
                        processing_completed[0] = True

                # 超时处理函数
                def handle_timeout():
                    if not processing_completed[0]:
                        timeout_occurred[0] = True
                        logger.warning(f"处理文件超时 ({timeout}秒): {file_path}")
                        self._force_close_office_processes("POWERPNT.EXE")

                # 创建并启动处理线程
                process_thread = threading.Thread(target=process_ppt)
                process_thread.daemon = True
                process_thread.start()

                # 创建并启动超时定时器
                timer = threading.Timer(timeout, handle_timeout)
                timer.daemon = True
                timer.start()

                # 等待处理完成或超时
                process_thread.join(timeout + 2)  # 给超时处理留出额外时间
                timer.cancel()  # 取消定时器

                # 检查处理结果
                if timeout_occurred[0]:
                    logger.error(f"处理PPT文件超时: {file_path}")
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": f"处理超时 ({timeout}秒)",
                    }
                elif exception[0]:
                    # 处理过程中发生异常
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 失败: {str(exception[0])}, 重试中..."
                        )
                        time.sleep(2)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": f"Win32COM 处理失败: {str(exception[0])}",
                    }
                elif result[0]:
                    # 处理成功
                    return result[0]
                else:
                    # 未知错误
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 出现未知错误, 重试中..."
                        )
                        time.sleep(2)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": "处理过程遇到未知错误",
                    }

            except Exception as e:
                logger.error(f"处理PPT文件异常: {str(e)}")
                if attempt == max_retries - 1:
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": f"处理失败: {str(e)}",
                    }

        # 清理可能残留的进程
        self._force_close_office_processes("POWERPNT.EXE")
        return {
            "content": "",
            "metadata": {"file_type": "ppt", "extractor": "win32com"},
            "error": "所有处理尝试均失败",
        }

    def _force_close_office_processes(self, process_name: str) -> None:
        """强制关闭指定Office进程"""
        try:
            for proc in psutil.process_iter(["pid", "name"]):
                try:
                    if proc.info["name"].upper() == process_name:
                        logger.info(
                            f"关闭Office进程: {proc.info['name']} (PID: {proc.info['pid']})"
                        )

                        # 尝试优雅终止
                        proc.terminate()

                        # 等待进程终止
                        try:
                            proc.wait(timeout=3)
                        except psutil.TimeoutExpired:
                            # 强制终止
                            proc.kill()
                except:
                    continue
        except Exception as e:
            logger.debug(f"关闭Office进程失败: {e}")

    def _extract_binary_text_content(
        self, file_path: str, mime_type: str
    ) -> Dict[str, Any]:
        """从二进制文件中提取可能的文本内容，作为最后的回退方案"""
        try:
            with open(file_path, "rb") as f:
                binary_data = f.read()

            # 尝试不同编码
            potential_text = ""
            encodings = ["utf-8", "latin1", "cp1252", "gb18030", "big5"]

            for encoding in encodings:
                try:
                    decoded = binary_data.decode(encoding, errors="ignore")
                    # 检查解码是否产生有意义的内容
                    if len(decoded.strip()) > len(potential_text.strip()):
                        potential_text = decoded
                except:
                    continue

            # 只提取可打印ASCII字符和常见空白
            printable_text = ""
            for char in potential_text:
                if char.isprintable() or char in " \t\n\r":
                    printable_text += char

            # 移除连续空白
            cleaned_text = re.sub(r"\s+", " ", printable_text).strip()

            # 移除二进制垃圾（随机字符序列）
            final_text = re.sub(
                r"[^\w\s.,;:!?(){}\[\]\'\"<>@#$%^&*+=\-_\\|/]", "", cleaned_text
            )

            # 移除非常短的片段
            words = [w for w in final_text.split(" ") if len(w) > 1]
            result_text = " ".join(words)

            if len(result_text) < 50:  # 如果提取的文本不够多
                return self._create_error_result(
                    mime_type, f"二进制内容提取失败: 未能提取到足够的文本内容"
                )

            return {
                "content": result_text,
                "metadata": {"file_type": mime_type, "extractor": "binary"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(mime_type, f"二进制内容提取失败: {str(e)}")

    def extract_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """提取文件内容，包含多级失败回退策略，增强容错性"""

        # 将所有操作包装在一个大的try-except块中，确保不会崩溃
        try:
            if not Path(file_path).exists():
                return {
                    "error": f"文件不存在: {file_path}",
                    "content": "",
                    "metadata": {},
                }

            if Path(file_path).stat().st_size == 0:
                return self._create_empty_result(mime_type)

            # 检查文件扩展名
            ext = Path(file_path).suffix.lower()

            # 跳过不支持的文件类型
            skip_mime_types = [
                "image/jpeg",
                "image/png",
                "image/gif",
                "image/bmp",
                "audio/mpeg",
                "audio/wav",
                "video/mp4",
                "video/x-msvideo",
                "video/quicktime",
                "font/ttf",
                "font/otf",
                "font/woff",
                "font/woff2",
                "application/x-executable",
            ]

            # 对于dwg文件特别处理
            if mime_type == "image/vnd.dwg":
                return self._create_error_result(mime_type, "DWG图纸文件, 跳过处理")

            # 对于明确是二进制格式且不是文档/压缩文件的特殊处理
            if mime_type in skip_mime_types:
                return self._create_error_result(
                    mime_type, f"不支持的MIME类型: {mime_type}, 跳过处理"
                )

            # 首先尝试使用主要提取器
            if mime_type in self.extractors:
                try:
                    logger.info(f"使用主要提取器 '{mime_type}' 处理: {file_path}")
        except Exception as e:
            # 总体异常处理，确保函数永远不会崩溃
            logger.error(f"提取内容时发生严重错误: {str(e)} - 文件: {file_path}")
            return self._create_error_result(
                mime_type, f"提取过程发生严重错误: {str(e)}"
            )

    def __del__(self):
        """对象销毁时清理资源"""
        try:
            self._cleanup_word_app()
            pythoncom.CoUninitialize()
        except Exception as e:
            logger.debug(f"清理资源时出错: {e}")


class SensitiveChecker:
    """敏感内容检查器，根据提供的参数使用BERT模型或YAML配置"""

    def __init__(self, config_path="sensitive_config.yaml", model_path="best_model.pth", use_model=False):
        """
        初始化敏感内容检查器
        
        Args:
            config_path: YAML配置文件路径
            model_path: BERT模型文件路径
            use_model: 是否优先使用模型，True表示使用BERT模型，False表示使用YAML配置
        """
        self.config_path = config_path
        self.model_path = model_path
        self.use_model = use_model
        self.mode = None
        self.bert_model = None
        self.bert_tokenizer = None
        self.config = None
        self.all_keywords = []
        self.keyword_pattern = None
        self.max_length = 128
        self.device = 'cuda' if torch.cuda.is_available() else 'cpu'
        
        if use_model:
            # 使用BERT模型
            if os.path.exists(model_path):
                try:
                    # 加载分词器和模型
                    model_name = 'bert-base-multilingual-cased'
                    self.bert_tokenizer = BertTokenizer.from_pretrained(model_name)
                    self.bert_model = BertForSequenceClassification.from_pretrained(model_name, num_labels=2)
                    self.bert_model.load_state_dict(torch.load(model_path, map_location=self.device))
                    self.bert_model.to(self.device)
                    self.bert_model.eval()
                    
                    self.mode = 'bert'
                    logger.info(f"使用BERT模型进行敏感内容检测: {model_path}")
                except Exception as e:
                    logger.error(f"加载BERT模型失败: {str(e)}")
                    self.mode = 'empty'
            else:
                logger.error(f"指定的模型文件不存在: {model_path}")
                self.mode = 'empty'
        else:
            # 使用YAML配置
            if os.path.exists(config_path):
                try:
                    self._load_yaml_config()
                    self.mode = 'yaml'
                    logger.info(f"使用YAML配置进行敏感内容检测: {config_path}")
                except Exception as e:
                    logger.error(f"加载YAML配置失败: {str(e)}")
                    self.mode = 'empty'
            else:
                logger.error(f"指定的配置文件不存在: {config_path}")
                self.mode = 'empty'
        
        if self.mode == 'empty':
            logger.warning("检测器初始化失败，将使用空检测器（不会检测任何敏感内容）")
    
    def _load_yaml_config(self):
        """加载YAML配置文件"""
        with open(self.config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)
            
        self.all_keywords = self.config.get("security_marks", []) + [
            kw
            for cat in self.config.get("sensitive_patterns", {}).values()
            for kw in cat.get("keywords", [])
        ]
        escaped_keywords = [re.escape(kw) for kw in self.all_keywords]
        self.keyword_pattern = re.compile("|".join(escaped_keywords)) if escaped_keywords else None

    def preprocess_text(self, text):
        """预处理输入文本"""
        if not text or not isinstance(text, str):
            return None
            
        encoding = self.bert_tokenizer.encode_plus(
            text,
            add_special_tokens=True,
            max_length=self.max_length,
            truncation=True,
            padding='max_length',
            return_attention_mask=True,
            return_tensors='pt'
        )
        return encoding

    def predict(self, text):
        """对输入文本进行分类预测"""
        if self.bert_model is None or self.bert_tokenizer is None:
            logger.error("BERT模型未成功加载，无法进行预测")
            return 0  # 默认为非敏感
            
        # 预处理文本
        encoding = self.preprocess_text(text)
        if encoding is None:
            return 0
            
        input_ids = encoding['input_ids'].to(self.device)
        attention_mask = encoding['attention_mask'].to(self.device)

        # 进行预测
        with torch.no_grad():
            outputs = self.bert_model(input_ids=input_ids, attention_mask=attention_mask)
            logits = outputs.logits
            preds = torch.argmax(logits, dim=1).cpu().numpy()

        return preds[0]  # 返回预测结果: 0表示非敏感，1表示敏感

    def _check_content_yaml(self, text):
        """使用YAML配置的正则表达式检查敏感内容"""
        if not self.keyword_pattern or not text:
            return []
            
        keyword_matches = {}
        for match in self.keyword_pattern.finditer(text or ""):
            keyword = match.group()
            if keyword not in keyword_matches:
                keyword_matches[keyword] = []
            keyword_matches[keyword].append(match.start())
        keyword_results = list(keyword_matches.items())

        structured_results = []
        for pattern, weight in self.config.get("structured_patterns", {}).items():
            matches = list(re.finditer(pattern, text or ""))
            if matches:
                positions = [m.start() for m in matches]
                structured_results.append((pattern, positions))

        number_results = []
        for pattern in self.config.get("number_patterns", []):
            matches = list(re.finditer(pattern, text or ""))
            if matches:
                positions = [m.start() for m in matches]
                number_results.append((pattern, positions))

        return keyword_results + structured_results + number_results

    def check_content(self, text):
        """检查文本中的敏感内容，根据初始化时的模式选择检测方法"""
        if not text or not isinstance(text, str) or len(text.strip()) == 0:
            return []
            
        try:
            if self.mode == 'bert':
                # 使用BERT模型检测
                prediction = self.predict(text[:self.max_length])
                
                if prediction == 1:  # 如果预测为敏感
                    return [("BERT模型检测为敏感内容", [0])]
                else:
                    return []  # 非敏感返回空列表
            elif self.mode == 'yaml':
                # 使用YAML配置的正则表达式检测
                return self._check_content_yaml(text)
            else:
                # 空检测器
                return []
        except Exception as e:
            logger.error(f"敏感内容检查失败: {str(e)}")
            return []

class ResultExporter:
    """增强版处理结果导出器，添加文本内容到Excel"""

    def export_to_json(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 JSON 文件"""
        export_data = [
            {
                "file_path": r.file_path,
                "mime_type": r.mime_type,
                "content_preview": self._get_content_preview(r, 200),
                "sensitive_words": [
                    {"word": w, "positions": p} for w, p in r.sensitive_words
                ],
                "error": r.error,
                "processing_time": r.processing_time,
            }
            for r in results
        ]
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)

    def export_to_excel(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 Excel 文件，包含文本内容"""
        try:
            # 创建主工作表数据
            summary_data = [
                {
                    "文件路径": r.file_path,
                    "文件类型": r.mime_type,
                    "敏感词统计": "; ".join(
                        [f"{w}({len(p)}次)" for w, p in r.sensitive_words]
                    ),
                    "处理时间(秒)": round(r.processing_time, 3),
                    "错误信息": r.error or "",
                    "内容预览": self._get_content_preview(r, 500),
                }
                for r in results
            ]

            # 创建内容工作表数据
            content_data = [
                {
                    "文件路径": r.file_path,
                    "文件类型": r.mime_type,
                    "内容": self._get_full_content(r),
                }
                for r in results
            ]

            # 创建Excel工作簿并写入两个工作表
            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                # 摘要工作表
                pd.DataFrame(summary_data).to_excel(
                    writer, sheet_name="处理结果", index=False
                )
                # 内容工作表
                pd.DataFrame(content_data).to_excel(
                    writer, sheet_name="文本内容", index=False
                )

            logger.info(f"结果已导出到Excel（含文本内容）: {output_path}")
        except Exception as e:
            logger.error(f"导出到Excel失败: {e}")
            # 尝试导出简化版本，不包含内容
            try:
                simplified_path = output_path.replace(".xlsx", "_简化版.xlsx")
                # 创建简化版本
                simple_data = [
                    {
                        "文件路径": r.file_path,
                        "文件类型": r.mime_type,
                        "敏感词统计": "; ".join(
                            [f"{w}({len(p)}次)" for w, p in r.sensitive_words]
                        ),
                        "处理时间(秒)": round(r.processing_time, 3),
                        "错误信息": r.error or "",
                    }
                    for r in results
                ]
                pd.DataFrame(simple_data).to_excel(
                    simplified_path, index=False, engine="openpyxl"
                )
                logger.info(f"已导出简化版Excel: {simplified_path}")
            except Exception as e2:
                logger.error(f"导出简化版Excel也失败: {e2}")

    def _get_content_preview(
        self, result: ProcessingResult, max_length: int = 200
    ) -> str:
        """从处理结果中获取内容预览"""
        try:
            if isinstance(result.content, dict):
                content = result.content.get("content", "")
                if isinstance(content, str):
                    return content[:max_length] + (
                        "..." if len(content) > max_length else ""
                    )
            return ""
        except Exception:
            return ""

    def _get_full_content(self, result: ProcessingResult) -> str:
        """从处理结果中获取完整内容"""
        try:
            if isinstance(result.content, dict):
                content = result.content.get("content", "")
                if isinstance(content, str):
                    return content
            return ""
        except Exception:
            return ""


class ResultMonitor:
    """结果监控器，用于实时输出处理进度和结果"""

    def __init__(self, output_csv: str = "processing_results.csv"):
        self.output_csv = output_csv
        self._init_csv()

    def _init_csv(self):
        """初始化 CSV 文件"""
        with open(self.output_csv, "w", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow(["文件名", "识别类型"])

    def record_result(self, result: ProcessingResult):
        """记录单个处理结果"""
        file_name = os.path.basename(result.file_path)
        sensitive_words_count = len(result.sensitive_words)
        recognition_type = "敏感" if sensitive_words_count > 0 else "常规"

        with open(self.output_csv, "a", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow([file_name, recognition_type])

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 处理文件: {result.file_path}")
        print(f"  类型: {result.mime_type}")
        print(f"  识别类型: {recognition_type}")
        if sensitive_words_count > 0:
            print(f"  发现敏感词: {'; '.join([f'{w}({len(p)}次)' for w, p in result.sensitive_words])}")
        if result.error:
            print(f"  错误: {result.error}")
        print("-" * 80)


class FileProcessor:
    """优化后的文件处理器主类"""

    def __init__(
        self,
        config_path: str = "sensitive_config.yaml",
        model_path: str = "best_model.pth",
        use_model: bool = False,
        monitor_output: str = "processing_results.csv",
        chunk_size: int = 1000,
        max_workers: Optional[int] = None,
        is_windows: bool = True,
    ):
        self.detector = FileTypeDetector()
        self.extractor = ContentExtractor(detector=self.detector, is_windows=is_windows)
        self.checker = SensitiveChecker(config_path=config_path, model_path=model_path, use_model=use_model)
        self.exporter = ResultExporter()
        self.monitor = ResultMonitor(monitor_output)
        self.chunk_size = chunk_size
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self._mime_cache = {}
        self._file_size_cache = {}

    def _scan_directory(self, directory: str) -> Iterator[List[str]]:
        """使用生成器分批扫描目录文件，确保正确关闭scandir迭代器"""
        current_chunk = []

        try:
            with os.scandir(directory) as dir_iter:
                for entry in dir_iter:
                    try:
                        if (
                            entry.is_file(follow_symlinks=False)
                            and not entry.name.startswith("~$")
                            and not self.detector._is_internal_stream(entry.name)
                        ):
                            ext = Path(entry.path).suffix.lower()
                            if ext not in self.detector.SKIP_EXTENSIONS:
                                current_chunk.append(entry.path)
                                if len(current_chunk) >= self.chunk_size:
                                    yield current_chunk
                                    current_chunk = []
                        elif entry.is_dir(follow_symlinks=False):
                            for sub_chunk in self._scan_directory(entry.path):
                                yield sub_chunk
                    except (PermissionError, OSError) as e:
                        logger.warning(f"访问文件/目录出错 {entry.path}: {e}")
                        continue
        except (PermissionError, OSError) as e:
            logger.warning(f"访问目录出错 {directory}: {e}")

        if current_chunk:
            yield current_chunk

    def _preload_file_info(self, file_paths: List[str]):
        """预加载文件信息以减少 I/O 操作"""
        for file_path in file_paths:
            try:
                stat = os.stat(file_path)
                self._file_size_cache[file_path] = stat.st_size
                mime_type = self.detector.detect_file_type(file_path)
                self._mime_cache[file_path] = mime_type
            except Exception as e:
                logger.warning(f"预加载文件信息失败 {file_path}: {e}")

    def _process_file_batch(self, file_paths: List[str]) -> List[ProcessingResult]:
        """处理一批文件，Office文件串行处理，其他文件并行处理"""
        results = []
        self._preload_file_info(file_paths)

        # 按类型对文件进行分类
        office_files = []
        other_files = []

        for fp in file_paths:
            mime_type = self._mime_cache.get(fp)
            ext = Path(fp).suffix.lower()

            # 检查是否为需要串行处理的Office文件
            if mime_type in (
                self.extractor.MIME_TYPE["DOC"],
                self.extractor.MIME_TYPE["PPT"],
            ) or ext in (".doc", ".ppt", ".xls"):
                office_files.append(fp)
            else:
                other_files.append(fp)


        return results

    def process_directory(self, directory: str) -> List[ProcessingResult]:
        """处理目录下的所有文件"""
        results = []
        total_files = len(self.detector.get_all_files(directory))
        completed = 0

        print(f"\n开始处理目录: {directory}")
        print(
            f"共发现 {total_files} 个文件"
        )
        print("=" * 80)

        try:
            for file_chunk in self._scan_directory(directory):
                chunk_results = self._process_file_batch(file_chunk)
                results.extend(chunk_results)
                completed += len(file_chunk)
                print(
                    f"\n已完成: {completed}/{total_files} ({completed / total_files * 100:.1f}%)"
                )
        except KeyboardInterrupt:
            print("\n用户中断处理...")
        except Exception as e:
            logger.error(f"处理目录时发生错误: {e}")
        finally:
            if results:
                try:
                    self.exporter.export_to_json(results, "interrupted_results.json")
                    self.exporter.export_to_excel(results, "interrupted_results.xlsx")
                except Exception as e:
                    logger.error(f"导出中断结果失败: {e}")

        print(f"\n处理完成！共处理 {len(results)}/{total_files} 个文件")
        print(f"详细结果已保存至: {self.monitor.output_csv}")
        print("=" * 80)
        return results

    def process_file(self, file_path: str) -> ProcessingResult:
        """处理单个文件，添加超时功能和容错性增强"""
        start_time = time.time()
        try:
            mime_type = self._mime_cache.get(
                file_path
            ) or self.detector.detect_file_type(file_path)
            ext = Path(file_path).suffix.lower()

            # 检查文件头部以辅助判断文件类型
            file_header = ""
            try:
                with open(file_path, "rb") as f:
                    file_header = f.read(16).hex().upper()
            except Exception:
                pass

            # 扩展不支持处理的MIME类型列表
            skip_mime_types = [
                "image/vnd.dwg",  # DWG图纸文件
                "application/octet-stream",  # 二进制文件
                "application/x-msdownload",  # 可执行文件
                "application/font-sfnt",  # 字体文件
                "font/ttf",  # TTF字体
                "font/otf",  # OTF字体
                "font/woff",  # WOFF字体
                "font/woff2",  # WOFF2字体
                "text/css",  # CSS文件
                "application/encrypted",
                "text/javascript",  # JS文件
                "image/jpeg",
                "image/png",
                "image/gif",
                "image/bmp",
                "image/webp",
                "image/svg+xml",  # 图片文件
                "audio/mpeg",
                "audio/wav",  # 音频文件
                "video/mp4",
                "video/x-msvideo",
                "video/quicktime",  # 视频文件
            ]

            # 扩展不支持的MIME类型前缀
            skip_mime_prefixes = ["font/", "image/", "audio/", "video/"]

            # 定义常规文件的条件 - 不再主要依赖文件扩展名
            regular_file_conditions = (
                mime_type in skip_mime_types
                or any(mime_type.startswith(prefix) for prefix in skip_mime_prefixes)
                or (ext and ext in self.detector.SKIP_EXTENSIONS)
            )

            if regular_file_conditions:
                logger.info(
                    f"识别为无需处理的文件类型，跳过内容提取: {file_path} (type: {mime_type}, header: {file_header})"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={
                        "content": "",
                        "metadata": {
                            "file_type": "regular",
                            "file_header": file_header,
                        },
                        "skipped": True,
                    },
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time,
                )

            # 检查文件是否为空
            file_size = self._file_size_cache.get(file_path, 0) or os.path.getsize(
                file_path
            )
            if file_size == 0:
                logger.info(f"空文件，跳过: {file_path}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={"content": "", "is_empty": True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time,
                )

            # 检查文件大小，过大的文件可能会导致处理缓慢
            size_mb = file_size / (1024 * 1024)
            if size_mb > 50:  # 大于50MB的文件
                logger.warning(
                    f"文件过大 ({size_mb:.2f} MB)，可能影响处理性能: {file_path}"
                )

            # =====================================================================
            # 特殊处理Office文档，避免线程锁序列化问题 - 使用超时功能
            # =====================================================================
            doc_timeout = 10  # 设置DOC文件处理超时时间为10秒
            ppt_timeout = 10  # 设置PPT文件处理超时时间为10秒

            # 特殊处理Office文档，避免线程锁序列化问题 - 使用超时功能
            # 因为没有扩展名，改为基于MIME类型判断
            if mime_type == "application/msword":
                # 使用带超时功能的函数处理DOC文件
                content = self.extractor._extract_doc_content(file_path)
            elif mime_type == "application/vnd.ms-powerpoint":
                # 使用带超时功能的函数处理PPT文件
                content = self.extractor._extract_ppt_content(file_path)
            elif mime_type == "application/vnd.ms-excel":
                # 针对Excel文件使用专用处理函数
                try:
                    content = self.extractor._extract_xls_content(file_path)
                except Exception as excel_error:
                    logger.error(f"处理Excel文件失败: {str(excel_error)}")
                    content = {
                        "content": "",
                        "metadata": {"file_type": "excel"},
                        "error": f"Excel处理失败: {str(excel_error)}"
                    }
            else:
                # 使用容错增强版的提取内容方法处理其他文件
                try:
                    content = self.extractor.extract_content(file_path, mime_type)
                except Exception as extract_error:
                    # 即使extract_content方法出现未捕获的异常，也不会导致程序崩溃
                    logger.error(
                        f"提取内容时发生严重错误: {str(extract_error)} - 文件: {file_path}"
                    )
                    content = {
                        "content": "",
                        "metadata": {"file_type": mime_type},
                        "error": f"提取内容时发生严重错误: {str(extract_error)}",
                    }
            # =====================================================================

            # 检查提取是否失败
            if content.get("error") and "UnsupportedFormatException" in content.get(
                "error"
            ):
                # 如果是不支持的格式，标记为跳过
                logger.info(
                    f"不支持的文件格式，标记为跳过: {file_path} (type: {mime_type})"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={
                        "content": "",
                        "metadata": {"file_type": "unsupported"},
                        "skipped": True,
                    },
                    sensitive_words=[],
                    error=content.get("error"),
                    processing_time=time.time() - start_time,
                )
            elif content.get("error"):
                logger.warning(
                    f"提取内容失败: {file_path} - MIME: {mime_type} - 错误: {content.get('error')}"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content=content,
                    sensitive_words=[],
                    error=content.get("error"),
                    processing_time=time.time() - start_time,
                )

            # 敏感内容检查 - 包装在try-except中，确保不会因为敏感词检查失败而中断
            try:
                sensitive_result = self.checker.check_content(
                    content.get("content", "")
                )
                sensitive_words = sensitive_result
            except Exception as check_error:
                logger.error(f"敏感内容检查失败: {file_path} - {str(check_error)}")
                sensitive_words = []

            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type,
                content=content,
                sensitive_words=sensitive_words,
                error=None,
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            error_msg = f"处理文件失败: {str(e)}"
            logger.error(f"{error_msg} - {file_path}")
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type if "mime_type" in locals() else "unknown",
                content={
                    "content": "",
                    "metadata": {"file_type": "regular"},
                    "skipped": True,
                },
                sensitive_words=[],
                error=error_msg,
                processing_time=time.time() - start_time,
            )


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="文件敏感内容检测工具")
    parser.add_argument("path", help="要处理的文件或目录路径")
    parser.add_argument(
        "--config", default="sensitive_config.yaml", help="敏感词配置文件路径"
    )
    parser.add_argument(
        "--model", default="best_model.pth", help="BERT分类模型路径"
    )
    parser.add_argument(
        "--use-model", action="store_true", 
        help="使用BERT模型进行敏感内容检测（默认使用YAML配置）"
    )
    parser.add_argument(
        "--output", default="results", help="输出结果文件名(不含扩展名)"
    )
    parser.add_argument(
        "--chunk-size", type=int, default=1000, help="每批处理的文件数量"
    )
    parser.add_argument("--workers", type=int, default=None, help="最大工作线程数")
    parser.add_argument(
        "--no-windows",
        action="store_true",
        help="指定非Windows平台，禁用win32com相关功能",
    )

    args = parser.parse_args()

    try:
        processor = FileProcessor(
            config_path=args.config,
            model_path=args.model,
            use_model=args.use_model,
            monitor_output=f"{args.output}_processing.csv",
            chunk_size=args.chunk_size,
            max_workers=args.workers,
            is_windows=not args.no_windows,
        )

        results = []
        if Path(args.path).is_file():
            results = [processor.process_file(args.path)]
        else:
            results = processor.process_directory(args.path)

        processor.exporter.export_to_json(results, f"{args.output}.json")
        processor.exporter.export_to_excel(results, f"{args.output}.xlsx")

        logger.info(f"处理完成，共处理 {len(results)} 个文件")
    except KeyboardInterrupt:
        print("\n用户中断程序执行")
        sys.exit(1)
    except Exception as e:
        logger.error(f"程序执行出错: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
