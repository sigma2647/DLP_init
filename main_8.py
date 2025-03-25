import os
import yaml
import magic
import time
import xlrd
import logging
import threading
import uuid
import shutil
import binascii
import re
import struct
import tempfile
import zipfile
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any, Iterator
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pdfminer.high_level import extract_text as pdf_extract_text
import pandas as pd
from markitdown import MarkItDown
import openpyxl
import csv
import json
from datetime import datetime
import argparse
import sys
import win32com.client
import pythoncom
import psutil
from pptx import Presentation
from collections import OrderedDict, Counter

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler("sensitive_detector.log", encoding="utf-8"),
        logging.StreamHandler(),
    ],
)
logger = logging.getLogger(__name__)


@dataclass
class ProcessingResult:
    """数据类，用于存储文件处理结果"""

    file_path: str
    mime_type: str
    content: Dict[str, Any]
    sensitive_words: List[Tuple[str, List[int]]] = field(default_factory=list)
    error: Optional[str] = None
    processing_time: float = 0.0


class EnhancedFileSignatureDetector:
    """
    增强型文件签名检测器，使用文件头部签名识别文件类型，
    整合了foo2_3.py中的高级文件签名检测能力
    """

    def __init__(self, verbose=False):
        """初始化文件签名探测器"""
        self.verbose = verbose
        
        # 文件签名字典（十六进制 -> 文件类型及MIME类型）
        self.signatures = OrderedDict([
            # 办公文档格式
            (b'504B0304', {
                'description': 'ZIP格式 (可能是Office文档/APK/JAR等)',
                'mime_type': 'application/zip'
            }),
            (b'D0CF11E0A1B11AE1', {
                'description': 'Microsoft复合文档二进制格式 (DOC/XLS/PPT等)',
                'mime_type': 'application/x-ole-storage'
            }),
            
            # PDF文档
            (b'25504446', {
                'description': 'PDF文档',
                'mime_type': 'application/pdf'
            }),
            
            # 图像文件格式
            (b'FFD8FF', {
                'description': 'JPEG图像',
                'mime_type': 'image/jpeg'
            }),
            (b'89504E470D0A1A0A', {
                'description': 'PNG图像',
                'mime_type': 'image/png'
            }),
            (b'47494638', {
                'description': 'GIF图像',
                'mime_type': 'image/gif'
            }),
            
            # 压缩和存档格式
            (b'526172211A0700', {
                'description': 'RAR压缩包 (v1.5+)',
                'mime_type': 'application/x-rar'
            }),
            (b'526172211A070100', {
                'description': 'RAR压缩包 (v5.0+)',
                'mime_type': 'application/x-rar'
            }),
            (b'377ABCAF271C', {
                'description': 'RAR压缩包 (旧版本)',
                'mime_type': 'application/x-rar'
            }),
            (b'1F8B08', {
                'description': 'GZIP压缩包',
                'mime_type': 'application/gzip'
            }),
            (b'425A68', {
                'description': 'BZIP2压缩包',
                'mime_type': 'application/x-bzip2'
            }),
            (b'FD377A585A00', {
                'description': 'XZ压缩包',
                'mime_type': 'application/x-xz'
            }),
        ])
        
        # 常见文件扩展名到MIME类型的映射（用于回退）
        self.MIME_TYPES = {
            ".txt": "text/plain",
            ".csv": "text/csv",
            ".xml": "text/xml",
            ".html": "text/html",
            ".htm": "text/html",
            ".json": "application/json",
            ".yaml": "application/yaml",
            ".yml": "application/yaml",
            ".md": "text/markdown",
            ".doc": "application/msword",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".xls": "application/vnd.ms-excel",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".ppt": "application/vnd.ms-powerpoint",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".zip": "application/zip",
            ".rar": "application/x-rar",
            ".7z": "application/x-7z-compressed",
            ".tar": "application/x-tar",
            ".gz": "application/gzip",
            ".bz2": "application/x-bzip2",
            ".pdf": "application/pdf",
        }
        
        # 要跳过处理的扩展名列表
        self.SKIP_EXTENSIONS = {
            ".dwg",
            ".mp3",
            ".wav",
            ".mp4",
            ".avi",
            ".mkv",
            ".flv",
            ".mov",
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".tif",
            ".tiff",
            ".webp",
            ".exe",
            ".dll",
            ".so",
            ".class",
            ".pyc",
            ".pyd",
            ".wasm",
            ".ttf",
            ".otf",
            ".woff",
            ".woff2",
            ".eot",
        }
        
        # 反向映射
        self.MIME_TO_EXT = {mime: ext for ext, mime in self.MIME_TYPES.items()}

        try:
            self.mime = magic.Magic(mime=True)
            logger.info("初始化magic库用于MIME类型检测")
        except ImportError:
            logger.warning("未安装python-magic库，将使用更基础的文件签名检测")
            self.mime = None
        except Exception as e:
            logger.error(f"初始化magic库出错: {e}")
            self.mime = None

    def read_file_header(self, file_path: str, bytes_to_read: int = 32) -> str:
        """
        读取文件的前N个字节并返回十六进制字符串

        Args:
            file_path: 文件路径
            bytes_to_read: 要读取的字节数

        Returns:
            文件头的大写十六进制字符串
        """
        try:
            with open(file_path, "rb") as f:
                header = f.read(bytes_to_read)
            return binascii.hexlify(header).decode("ascii").upper()
        except Exception as e:
            logger.error(f"读取文件头出错: {file_path} - {e}")
            return ""

    def get_magic_mime_type(self, file_path: str) -> str:
        """
        使用python-magic检测MIME类型

        Args:
            file_path: 文件路径

        Returns:
            MIME类型字符串，失败时返回空字符串
        """
        if not self.mime:
            return ""

        try:
            # 处理非ASCII文件名或含空格的路径
            normalized_path = os.path.normpath(os.path.abspath(file_path))
            return self.mime.from_file(normalized_path)
        except Exception as e:
            logger.error(f"Magic MIME类型检测失败: {file_path} - {e}")
            return ""

    def _check_office_signatures(self, file_path):
        """专门检查Office文档签名，包括处理签名损坏的情况"""
        try:
            # 首先检查文件大小，Office文档通常不会太小
            file_size = os.path.getsize(file_path)
            if file_size < 4096:  # 小于4KB可能不是Office文档
                return None, None
                
            # 读取文件头部
            header_hex = self.read_file_header(file_path, 16)
            if not header_hex:
                return None, None
                
            # 首先检查标准签名
            if header_hex.startswith("D0CF11E0A1B11AE1"):
                # 这是Office二进制格式，进一步判断具体类型
                result, mime_type = self._analyze_ole_file(file_path)
                return result, mime_type
                
            if header_hex.startswith("504B0304"):
                # 这可能是Office Open XML格式，进一步判断
                result, mime_type = self._analyze_ooxml_file(file_path)
                return result, mime_type
                
            # 检查签名损坏的情况
            # 有时文件头部可能被损坏，但仍然保留了部分Office格式的特征
            
            # 创建临时文件副本进行修复尝试
            temp_dir = tempfile.mkdtemp()
            try:
                temp_file = os.path.join(temp_dir, "temp_file")
                shutil.copy2(file_path, temp_file)
                
                # 尝试修复OLE结构标记
                if self._try_repair_ole(temp_file):
                    result, mime_type = self._analyze_ole_file(temp_file)
                    if result:
                        return f"{result} (已修复签名)", mime_type
                
                # 尝试修复ZIP/OOXML结构
                if self._try_repair_ooxml(temp_file):
                    result, mime_type = self._analyze_ooxml_file(temp_file)
                    if result:
                        return f"{result} (已修复签名)", mime_type
                
            finally:
                # 清理临时文件
                shutil.rmtree(temp_dir, ignore_errors=True)
                
            # 如果没有明确匹配，尝试基于内容特征判断
            result, mime_type = self._guess_office_type_by_content(file_path)
            return result, mime_type
            
        except Exception as e:
            if self.verbose:
                logger.error(f"分析Office文档签名时出错: {e}")
            return None, None

    def _try_repair_ole(self, file_path):
        """尝试修复OLE格式的签名"""
        try:
            with open(file_path, 'r+b') as f:
                # OLE头部标记是 D0CF11E0A1B11AE1
                ole_signature = bytearray([0xD0, 0xCF, 0x11, 0xE0, 0xA1, 0xB1, 0x1A, 0xE1])
                f.seek(0)
                f.write(ole_signature)
                return True
        except:
            return False
    
    def _try_repair_ooxml(self, file_path):
        """尝试修复OOXML(ZIP)格式的签名"""
        try:
            with open(file_path, 'r+b') as f:
                # ZIP/OOXML头部标记是 504B0304
                zip_signature = bytearray([0x50, 0x4B, 0x03, 0x04])
                f.seek(0)
                f.write(zip_signature)
                return True
        except:
            return False
            
    def _analyze_ole_file(self, file_path):
        """深度分析OLE文件结构来确定具体的Office文件类型"""
        try:
            # 读取更多文件内容以便深度分析
            with open(file_path, 'rb') as f:
                content = f.read(32768)  # 读取前32KB内容进行分析
                
            # 特定二进制标记检测
            # 检查是否存在Word特有的二进制指纹
            word_markers = [
                b'\x57\x00\x6F\x00\x72\x00\x64\x00\x44\x00\x6F\x00\x63\x00',  # "WordDoc" in UTF-16
                b'\x57\x00\x6F\x00\x72\x00\x64\x00\x2E\x00\x44\x00\x6F\x00\x63\x00',  # "Word.Doc" in UTF-16
                b'\x4D\x53\x57\x6F\x72\x64\x44\x6F\x63',  # "MSWordDoc"
                b'\x42\x6A\x44\x77\x70'  # "BjDwp" 特有标记
            ]
            
            # 检查是否存在Excel特有的二进制指纹
            excel_markers = [
                b'\x57\x00\x6F\x00\x72\x00\x6B\x00\x62\x00\x6F\x00\x6F\x00\x6B\x00',  # "Workbook" in UTF-16
                b'\x45\x00\x78\x00\x63\x00\x65\x00\x6C\x00',  # "Excel" in UTF-16
                b'\x42\x00\x49\x00\x46\x00\x46',  # "BIFF" in UTF-16
                b'\x53\x68\x65\x65\x74',  # "Sheet"
                b'\xE5\x7A\x31\x1C'  # Excel 特有二进制标记
            ]
            
            # 检查是否存在PowerPoint特有的二进制指纹
            ppt_markers = [
                b'\x50\x00\x6F\x00\x77\x00\x65\x00\x72\x00\x50\x00\x6F\x00\x69\x00\x6E\x00\x74\x00',  # "PowerPoint" in UTF-16
                b'\x50\x52\x45\x53\x45\x4E\x54\x41\x54\x49\x4F\x4E',  # "PRESENTATION"
                b'\x50\x50\x30\x31',  # PPT 特有标记
                b'\x53\x6C\x69\x64\x65'  # "Slide"
            ]
            
            # 计算匹配的标记数量
            word_matches = sum(1 for marker in word_markers if marker in content)
            excel_matches = sum(1 for marker in excel_markers if marker in content)
            ppt_matches = sum(1 for marker in ppt_markers if marker in content)
            
            # 分析文件结构中的目录名称
            # Office文档通常有特定的流和存储名称
            ole_dirs = []
            pos = 0
            while True:
                pos = content.find(b'\x03\x00', pos)
                if pos == -1 or pos + 64 > len(content):
                    break
                    
                # 尝试提取目录名称 (通常是UTF-16编码)
                dir_name_bytes = content[pos+2:pos+64]
                try:
                    dir_name = dir_name_bytes.decode('utf-16-le').strip('\x00')
                    if dir_name and len(dir_name) > 1:
                        ole_dirs.append(dir_name)
                except:
                    pass
                pos += 1
            
            # 检查特定的目录名称
            word_dirs = ['WordDocument', 'Document', '1Table', 'Data', 'ObjectPool']
            excel_dirs = ['Workbook', 'Book', 'SummaryInformation', 'Worksheet']
            ppt_dirs = ['PowerPoint Document', 'Pictures', 'Current User', 'SummaryInformation']
            
            word_dir_matches = sum(1 for d in word_dirs if d in ole_dirs)
            excel_dir_matches = sum(1 for d in excel_dirs if d in ole_dirs)
            ppt_dir_matches = sum(1 for d in ppt_dirs if d in ole_dirs)
            
            # 扫描特定的十六进制模式
            # Word .doc 文件中的特定16进制模式
            if b'\xEC\xA5\xC1\x00' in content or b'\x42\x6F\x6F\x6B\x6D\x61\x72\x6B' in content:
                word_matches += 2
                
            # Excel .xls 文件中的特定16进制模式
            if b'\x09\x08\x10\x00\x00\x06\x05\x00' in content or b'\xFD\xFF\xFF\xFF\x10' in content:
                excel_matches += 2
                
            # PowerPoint .ppt 文件中的特定16进制模式
            if b'\xA0\x46\x1D\xF0' in content or b'\x00\x6E\x28\x00' in content:
                ppt_matches += 2
                
            # 分析文档头部的魔数特征
            # 读取文件的特定偏移位置
            try:
                with open(file_path, 'rb') as f:
                    f.seek(512)  # 跳过OLE头
                    sector_data = f.read(128)
                    
                    # Excel特有的扇区标记
                    if sector_data.startswith(b'\x09\x08') or b'\xFD\xFF\x00\x00' in sector_data:
                        excel_matches += 3
                        
                    # Word特有的扇区标记
                    if b'\xEC\xA5\xC1\x00' in sector_data or sector_data.startswith(b'\xDC\x00'):
                        word_matches += 3
                        
                    # PowerPoint特有的扇区标记
                    if b'\x00\x6E\x28\x00' in sector_data or b'\x0F\x00\x00\x00\xAF' in sector_data:
                        ppt_matches += 3
            except:
                pass
                
            # 统计字节频率分布，各种文档类型有不同的分布特征
            byte_freq = {}
            for b in content[:4096]:  # 只分析前4KB
                if b not in byte_freq:
                    byte_freq[b] = 0
                byte_freq[b] += 1
                
            # 某些字节频率是特定文档类型的指示器
            if byte_freq.get(0x57, 0) > 30 and byte_freq.get(0x6F, 0) > 30:  # 'W', 'o'
                word_matches += 1
            if byte_freq.get(0x45, 0) > 30 and byte_freq.get(0x78, 0) > 30:  # 'E', 'x'
                excel_matches += 1
            if byte_freq.get(0x50, 0) > 30 and byte_freq.get(0x53, 0) > 20:  # 'P', 'S'
                ppt_matches += 1
                
            # 综合分析结果，确定最可能的文件类型
            if self.verbose:
                logger.info(f"Word 匹配得分: {word_matches}")
                logger.info(f"Excel 匹配得分: {excel_matches}")
                logger.info(f"PowerPoint 匹配得分: {ppt_matches}")
                
            # 根据匹配分数判断文件类型和MIME类型
            scores = {
                "Microsoft Word文档 (.doc)": (word_matches, "application/msword"),
                "Microsoft Excel工作簿 (.xls)": (excel_matches, "application/vnd.ms-excel"),
                "Microsoft PowerPoint演示文稿 (.ppt)": (ppt_matches, "application/vnd.ms-powerpoint")
            }
            
            # 找出得分最高的类型
            max_score_type = max(scores.items(), key=lambda x: x[1][0])
            
            # 只有当得分超过阈值时才确定为该类型
            if max_score_type[1][0] >= 2:
                return max_score_type[0], max_score_type[1][1]
            elif max_score_type[1][0] > 0:
                return f"可能是{max_score_type[0]} (置信度低)", max_score_type[1][1]
            else:
                # 额外检查文件扩展名
                ext = os.path.splitext(file_path)[1].lower()
                if ext == '.doc':
                    return "Microsoft Word文档 (.doc) (基于文件扩展名)", "application/msword"
                elif ext == '.xls':
                    return "Microsoft Excel工作簿 (.xls) (基于文件扩展名)", "application/vnd.ms-excel"
                elif ext == '.ppt':
                    return "Microsoft PowerPoint演示文稿 (.ppt) (基于文件扩展名)", "application/vnd.ms-powerpoint"
                    
                return "Microsoft Office二进制文档 (无法确定具体类型)", "application/x-ole-storage"
                    
        except Exception as e:
            if self.verbose:
                logger.error(f"分析OLE文件时出错: {e}")
            return "Microsoft Office二进制文档 (分析出错)", "application/x-ole-storage"
    
    def _analyze_ooxml_file(self, file_path):
        """分析OOXML (Office Open XML)文件来确定具体的Office文件类型"""
        try:
            # 检查是否是有效的ZIP文件
            if not zipfile.is_zipfile(file_path):
                return None, None
                
            with zipfile.ZipFile(file_path) as zf:
                # 获取文件列表
                file_list = zf.namelist()
                
                # 检查[Content_Types].xml
                if '[Content_Types].xml' in file_list:
                    try:
                        content_types = zf.read('[Content_Types].xml').decode('utf-8', errors='ignore')
                        
                        # 检查Word文档
                        if 'application/vnd.openxmlformats-officedocument.wordprocessingml' in content_types:
                            return "Microsoft Word文档 (.docx)", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                            
                        # 检查Excel工作簿
                        if 'application/vnd.openxmlformats-officedocument.spreadsheetml' in content_types:
                            return "Microsoft Excel工作簿 (.xlsx)", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                            
                        # 检查PowerPoint演示文稿
                        if 'application/vnd.openxmlformats-officedocument.presentationml' in content_types:
                            return "Microsoft PowerPoint演示文稿 (.pptx)", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    except:
                        pass
                
                # 通过目录结构判断
                # Word文档特征
                if any('word/document.xml' in f for f in file_list) or any('/word/' in f for f in file_list):
                    return "Microsoft Word文档 (.docx)", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    
                # Excel工作簿特征
                if any('xl/workbook.xml' in f for f in file_list) or any('/xl/' in f for f in file_list):
                    return "Microsoft Excel工作簿 (.xlsx)", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    
                # PowerPoint演示文稿特征
                if any('ppt/presentation.xml' in f for f in file_list) or any('/ppt/' in f for f in file_list):
                    return "Microsoft PowerPoint演示文稿 (.pptx)", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                
            # 如果以上都没匹配，但确认是ZIP文件，可能是其他OOXML格式
            return "Microsoft Office Open XML文档 (未知类型)", "application/zip"
            
        except Exception as e:
            if self.verbose:
                logger.error(f"分析OOXML文件时出错: {e}")
            return "Microsoft Office Open XML文档 (可能已损坏)", "application/zip"
    

    def _guess_office_type_by_content(self, file_path):
        """通过文件内容特征猜测Office文档类型"""
        try:
            # 读取较多内容进行深度分析
            with open(file_path, 'rb') as f:
                content = f.read(32768)  # 读取前32KB
                
            # 首先检查是否包含Office文档特征
            office_markers = [b'Microsoft', b'Office', b'Document', b'Workbook', 
                              b'PowerPoint', b'Excel', b'Word', b'Worksheet']
                              
            office_markers_found = [marker for marker in office_markers if marker in content]
            if not office_markers_found:
                return None, None  # 不太可能是Office文档
                
            # 统计各种Office类型特征出现次数
            marker_counts = {
                'doc': 0,
                'xls': 0,
                'ppt': 0
            }
            
            # Word文档特征
            doc_markers = [b'Word', b'Document', b'DocFile', b'MSWordDoc', b'.docx', b'.doc',
                           b'w:document', b'word/document.xml', b'wordDocument']
            for marker in doc_markers:
                marker_counts['doc'] += content.count(marker)
                
            # Excel工作簿特征
            xls_markers = [b'Excel', b'Workbook', b'Worksheet', b'spreadsheet', b'.xlsx', b'.xls',
                           b'xl/workbook.xml', b'xl:workbook', b'sst', b'Sheet', b'x:sheet']
            for marker in xls_markers:
                marker_counts['xls'] += content.count(marker)
                
            # PowerPoint演示文稿特征
            ppt_markers = [b'PowerPoint', b'Presentation', b'Slide', b'.pptx', b'.ppt',
                           b'ppt/presentation.xml', b'p:presentation', b'slideshow']
            for marker in ppt_markers:
                marker_counts['ppt'] += content.count(marker)
                
            # 确定最可能的类型
            max_type = max(marker_counts.items(), key=lambda x: x[1])
            if max_type[1] > 0:
                if max_type[0] == 'doc':
                    # 再次判断是.doc还是.docx
                    if b'word/document.xml' in content or b'Content_Types' in content:
                        return "可能是Microsoft Word文档 (.docx) [基于内容特征判断]", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    else:
                        return "可能是Microsoft Word文档 (.doc) [基于内容特征判断]", "application/msword"
                        
                elif max_type[0] == 'xls':
                    # 再次判断是.xls还是.xlsx
                    if b'xl/workbook.xml' in content or b'Content_Types' in content:
                        return "可能是Microsoft Excel工作簿 (.xlsx) [基于内容特征判断]", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    else:
                        return "可能是Microsoft Excel工作簿 (.xls) [基于内容特征判断]", "application/vnd.ms-excel"
                        
                elif max_type[0] == 'ppt':
                    # 再次判断是.ppt还是.pptx
                    if b'ppt/presentation.xml' in content or b'Content_Types' in content:
                        return "可能是Microsoft PowerPoint演示文稿 (.pptx) [基于内容特征判断]", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    else:
                        return "可能是Microsoft PowerPoint演示文稿 (.ppt) [基于内容特征判断]", "application/vnd.ms-powerpoint"
            
            return "可能是Office文档 (具体类型未知) [基于内容特征判断]", "application/x-ole-storage"
            
        except Exception as e:
            if self.verbose:
                logger.error(f"通过内容猜测Office类型时出错: {e}")
            return None, None

    def inspect_zip_content(self, file_path: str) -> Tuple[str, str]:
        """
        检查ZIP文件内容以确定是否为OOXML文档

        Args:
            file_path: ZIP文件路径

        Returns:
            (文件类型描述, MIME类型)，无法确定时返回(None, None)
        """
        try:
            file_type, mime_type = self._analyze_ooxml_file(file_path)
            if file_type and mime_type:
                return file_type, mime_type
                
            # 如果不是Office格式，则作为普通ZIP返回
            return "ZIP压缩包", "application/zip"
        except Exception as e:
            logger.warning(f"ZIP内容检查失败: {file_path} - {e}")
            return None, None

    def inspect_ole_content(self, file_path: str) -> Tuple[str, str]:
        """
        检查OLE文件内容以确定是否为特定Office文档类型

        Args:
            file_path: OLE文件路径

        Returns:
            (文件类型描述, MIME类型)，无法确定时返回(None, None)
        """
        try:
            file_type, mime_type = self._analyze_ole_file(file_path)
            if file_type and mime_type:
                return file_type, mime_type
            
            # 如果没有找到特定标记，它是通用OLE文件
            return "OLE复合文档", "application/x-ole-storage"
        except Exception as e:
            logger.warning(f"OLE内容检查失败: {file_path} - {e}")
            return None, None

    def detect_file_type(self, file_path: str) -> str:
        """
        综合检测文件类型，返回MIME类型
        
        Args:
            file_path: 文件路径

        Returns:
            MIME类型字符串
        """
        try:
            # 基本文件信息
            file_size = os.path.getsize(file_path)
            file_extension = Path(file_path).suffix.lower()
            
            # 跳过非常小或空的文件
            if file_size == 0:
                return "application/octet-stream"

            # 读取文件头部（前32字节）
            header_hex = self.read_file_header(file_path, 32)
            if not header_hex:
                return "application/octet-stream"

            # 获取magic MIME类型（如果可用）
            magic_mime = self.get_magic_mime_type(file_path) if self.mime else ""
            
            # 通过文件头部签名检测 - 首先尝试特殊处理Office文档
            office_type, office_mime = self._check_office_signatures(file_path)
            if office_type and office_mime:
                logger.info(f"Office签名检测: {file_path} - {office_type} ({office_mime})")
                # 返回MIME类型 - 不再返回元组
                return office_mime

            # 检查ZIP签名
            if header_hex.startswith("504B0304"):  # ZIP签名
                zip_type, zip_mime = self.inspect_zip_content(file_path)
                if zip_type and zip_mime:
                    # 返回MIME类型 - 不再返回元组
                    return zip_mime

            # 检查OLE签名
            elif header_hex.startswith("D0CF11E0A1B11AE1"):  # OLE签名
                ole_type, ole_mime = self.inspect_ole_content(file_path)
                if ole_type and ole_mime:
                    # 返回MIME类型 - 不再返回元组
                    return ole_mime

            # 检查PDF签名
            elif header_hex.startswith("25504446"):
                return "application/pdf"
                
            # 检查RAR签名
            elif header_hex.startswith("526172211A07"):
                return "application/x-rar"
                
            # 检查7Z签名
            elif header_hex.startswith("377ABCAF"):
                return "application/x-7z-compressed"
                
            # 检查GZIP签名
            elif header_hex.startswith("1F8B08"):
                return "application/gzip"
                
            # 检查JPEG签名
            elif header_hex.startswith("FFD8FF"):
                return "image/jpeg"
                
            # 检查PNG签名
            elif header_hex.startswith("89504E470D0A1A0A"):
                return "image/png"
                
            # 检查GIF签名
            elif header_hex.startswith("47494638"):
                return "image/gif"

            # 如果通过签名无法识别，尝试使用magic库
            if magic_mime and magic_mime != "application/octet-stream":
                return magic_mime

            # 尝试根据文件内容猜测类型
            content_type, content_mime = self._guess_office_type_by_content(file_path)
            if content_type and content_mime:
                return content_mime

            # 最后尝试使用扩展名
            if file_extension in self.MIME_TYPES:
                return self.MIME_TYPES[file_extension]

            # 如果一切都失败了，返回二进制流类型
            return "application/octet-stream"

        except Exception as e:
            logger.error(f"文件类型检测失败: {file_path} - {e}")
            return "application/octet-stream"

    def get_all_files(self, directory: str) -> List[str]:
        """
        递归获取目录中的所有文件，过滤掉某些类型

        Args:
            directory: 要扫描的目录路径

        Returns:
            文件路径列表
        """
        files = []
        for f in Path(directory).rglob("*"):
            if f.is_file():
                # 跳过隐藏文件和Office临时文件
                if f.name.startswith(".") or f.name.startswith("~$"):
                    continue

                # 跳过Office内部流文件
                if bool(re.match(r"^\[\d+\].+", f.name)):
                    continue

                # 跳过已知的二进制格式
                if f.suffix.lower() in self.SKIP_EXTENSIONS:
                    continue

                files.append(str(f))
        return files

    def check_file_signature(self, file_path: str) -> Dict:
        """
        检查文件签名并返回详细信息

        Args:
            file_path: 文件路径

        Returns:
            包含详细文件签名信息的字典
        """
        file_type, mime_type = self.detect_file_type(file_path)
        file_size = os.path.getsize(file_path)
        header_hex = self.read_file_header(file_path, 32)
        
        return {
            "file_path": file_path,
            "file_size": file_size,
            "extension": Path(file_path).suffix.lower(),
            "file_type": file_type,
            "mime_type": mime_type,
            "header_hex": header_hex,
        }


class ContentExtractor:
    """优化后的文件内容提取器，基于文件签名而非扩展名"""

    MIME_TYPE = {
        "TEXT": "text/plain",
        "CSV": "text/csv",
        "PDF": "application/pdf",
        "MARKDOWN": "text/markdown",
        "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "XLSX": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "XLS": "application/vnd.ms-excel",
        "DOC": "application/msword",
        "PPT": "application/vnd.ms-powerpoint",
        "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ZIP": "application/zip",
        "RAR": "application/x-rar",
        "SEVENZ": "application/x-7z-compressed",
    }

    def __init__(self, detector, is_windows: bool = True):
        self.detector = detector
        try:
            self.md = MarkItDown()
        except Exception as e:
            self.md = None
            logger.error(f"初始化 MarkItDown 失败: {str(e)}")
        self.is_windows = is_windows
        self.word_app = None
        self.word_lock = threading.Lock()

        self.extractors = {
            self.MIME_TYPE["TEXT"]: self._extract_text_content,
            self.MIME_TYPE["CSV"]: self._extract_csv_content,
            self.MIME_TYPE["PDF"]: self._extract_pdf_content,
            self.MIME_TYPE["MARKDOWN"]: self._extract_markdown_content,
            self.MIME_TYPE["ZIP"]: self._extract_archive_content,
            self.MIME_TYPE["RAR"]: self._extract_archive_content,
            self.MIME_TYPE["SEVENZ"]: self._extract_archive_content,
            self.MIME_TYPE["DOCX"]: self._extract_docx_content,
            self.MIME_TYPE["XLSX"]: self._extract_xlsx_content,
            self.MIME_TYPE["PPTX"]: self._extract_pptx_content,
        }

        if self.is_windows:
            # 尝试初始化Word应用
            if self._init_word_app():
                self.extractors.update(
                    {
                        self.MIME_TYPE["DOC"]: self._extract_doc_content,
                    }
                )
            else:
                logger.warning("Word应用程序初始化失败，.doc文件将使用备用方法处理")

            # PowerPoint不需要持久应用实例
            self.extractors.update(
                {
                    self.MIME_TYPE["PPT"]: self._extract_ppt_content,
                    self.MIME_TYPE["XLS"]: self._extract_xls_content,
                }
            )

    def _init_word_app(self):
        """初始化Word应用程序实例"""
        if not self.is_windows:
            return False

        with self.word_lock:
            if self.word_app is None:
                try:
                    # 强制关闭任何现有Word进程
                    self._force_close_office_processes("WINWORD.EXE")

                    # 初始化当前线程的COM
                    pythoncom.CoInitialize()

                    # 创建新的Word应用实例
                    self.word_app = win32com.client.Dispatch("Word.Application")
                    self.word_app.Visible = False
                    self.word_app.DisplayAlerts = False

                    logger.info("成功初始化Word应用程序")
                    return True
                except Exception as e:
                    logger.error(f"初始化Word应用程序失败: {str(e)}")
                    self.word_app = None

                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass

                    return False
            return True

    def _cleanup_word_app(self):
        """清理 Word 应用程序实例"""
        with self.word_lock:
            if self.word_app is not None:
                try:
                    self.word_app.Quit()
                except:
                    pass
                self._force_close_office_processes("WINWORD.EXE")
                self.word_app = None
                pythoncom.CoUninitialize()
                logger.info("成功清理 Word 应用程序")

    def _create_empty_result(self, file_type: str) -> Dict[str, Any]:
        """创建空结果对象"""
        return {
            "content": "",
            "metadata": {"file_type": file_type},
            "error": None,
            "is_empty": True,
        }

    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果对象"""
        return {"content": "", "metadata": {"file_type": file_type}, "error": error_msg}

    def _extract_text_content(self, file_path: str) -> Dict[str, Any]:
        """提取文本文件内容"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                "content": content[:1000000]
                + (
                    "\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""
                ),
                "metadata": {"file_type": "text"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result("text", f"读取文本文件失败: {str(e)}")

    def _extract_csv_content(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容"""
        try:
            encodings = ["utf-8", "latin1", "cp1252", "gbk"]
            df = None

            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, engine="python")
                    break
                except UnicodeDecodeError:
                    continue

            if df is None:
                df = pd.read_csv(file_path, encoding="utf-8", errors="ignore")

            if len(df) > 10000:
                df = df.head(10000)
                notice = "\n\n[注意: 文件过大，仅显示前10000行]"
            else:
                notice = ""

            return {
                "content": df.to_string(index=False) + notice,
                "metadata": {"file_type": "csv"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result("csv", f"CSV处理失败: {str(e)}")

    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文件内容"""
        try:
            text = pdf_extract_text(file_path)
            if len(text) > 100000:
                text = text[:100000] + "\n\n[注意: 文件过大，内容已截断]"
            return {"content": text, "metadata": {"file_type": "pdf"}, "error": None}
        except Exception as e:
            return self._create_error_result("pdf", f"PDF处理失败: {str(e)}")

    def _extract_markdown_content(self, file_path: str) -> Dict[str, Any]:
        """提取Markdown文件内容"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                "content": content[:1000000]
                + (
                    "\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""
                ),
                "metadata": {"file_type": "markdown"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(
                "markdown", f"Markdown文件处理失败: {str(e)}"
            )

    def _extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOCX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result("docx", "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                "content": result.text_content,
                "metadata": {"file_type": "docx", "extractor": "markitdown"},
                "error": None,
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理DOCX失败: {file_path} - {str(e)}")
            return self._create_error_result("docx", f"DOCX处理失败: {str(e)}")

    def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLSX文件内容"""
        try:
            sheets = pd.read_excel(file_path, sheet_name=None)
            content = []
            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)
            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": "excel"},
                "error": None,
            }
        except ImportError as e:
            return self._create_error_result("excel", f"缺少必要库: {str(e)}")
        except ValueError as e:
            return self._create_error_result("excel", f"文件格式错误: {str(e)}")
        except Exception as e:
            return self._create_error_result(
                "excel", f"Excel处理失败: {str(e)}，类型: {type(e).__name__}"
            )

    def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPTX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result("pptx", "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                "content": result.text_content,
                "metadata": {"file_type": "pptx", "extractor": "markitdown"},
                "error": None,
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理PPTX失败: {file_path} - {str(e)}")
            if self.is_windows:
                try:
                    prs = Presentation(file_path)
                    text_content = []

                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())

                        if slide_text:
                            text_content.append(
                                f"Slide {i + 1}:\n" + "\n".join(slide_text)
                            )

                    content = "\n\n".join(text_content)
                    return {
                        "content": content,
                        "metadata": {"file_type": "pptx", "extractor": "python-pptx"},
                        "error": None,
                    }
                except Exception as e2:
                    return self._create_error_result(
                        "pptx",
                        f"所有提取方法失败: MarkItDown: {str(e)}; python-pptx: {str(e2)}",
                    )
            return self._create_error_result("pptx", f"PPTX处理失败: {str(e)}")


    def _extract_xls_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLS文件内容，增强错误处理和多级回退策略"""
        logger.info(f"处理XLS文件: {file_path}")

        # 策略1: 尝试使用pandas通用方法
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="xlrd")
            content = []

            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)

            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": "excel", "extractor": "pandas-xlrd"},
                "error": None,
            }
        except xlrd.biffh.XLRDError as xlrd_error:
            # 特殊处理XLRD的"Can't find workbook in OLE2 compound document"错误
            logger.warning(f"XLRD无法解析XLS文件: {file_path} - {str(xlrd_error)}")
            # 继续尝试其他方法
        except Exception as pandas_error:
            logger.warning(f"Pandas处理XLS失败: {file_path} - {str(pandas_error)}")
            # 继续尝试其他方法

        # 策略2: 如果在Windows上，尝试使用COM接口
        if self.is_windows:
            try:
                pythoncom.CoInitialize()
                excel_app = win32com.client.Dispatch("Excel.Application")
                excel_app.Visible = False
                excel_app.DisplayAlerts = False

                workbook = excel_app.Workbooks.Open(
                    os.path.abspath(file_path), ReadOnly=True, UpdateLinks=False
                )

                content = []
                for sheet_index in range(1, workbook.Sheets.Count + 1):
                    try:
                        worksheet = workbook.Sheets(sheet_index)
                        sheet_name = worksheet.Name

                        # 获取已使用范围
                        used_range = worksheet.UsedRange
                        if used_range.Rows.Count > 0 and used_range.Columns.Count > 0:
                            # 创建二维数组存储数据
                            data = []
                            for row in range(1, min(5001, used_range.Rows.Count + 1)):
                                row_data = []
                                for col in range(1, used_range.Columns.Count + 1):
                                    try:
                                        cell_value = worksheet.Cells(row, col).Value
                                        row_data.append(
                                            str(cell_value)
                                            if cell_value is not None
                                            else ""
                                        )
                                    except:
                                        row_data.append("")
                                data.append(row_data)

                            # 格式化为类似表格的文本
                            if data:
                                max_lengths = [
                                    max(
                                        len(str(row[i])) for row in data if i < len(row)
                                    )
                                    for i in range(len(data[0]))
                                ]
                                formatted_rows = []

                                for row in data:
                                    formatted_row = [
                                        str(val).ljust(max_lengths[i])
                                        for i, val in enumerate(row)
                                        if i < len(max_lengths)
                                    ]
                                    formatted_rows.append(" | ".join(formatted_row))

                                sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(
                                    formatted_rows
                                )
                                if used_range.Rows.Count > 5000:
                                    sheet_text += "\n[注意: 表格过大，仅显示前5000行]"

                                content.append(sheet_text)
                    except Exception as sheet_error:
                        logger.warning(
                            f"处理工作表 {sheet_index} 时出错: {str(sheet_error)}"
                        )
                        continue

                workbook.Close(SaveChanges=False)
                excel_app.Quit()

                return {
                    "content": "\n\n".join(content),
                    "metadata": {"file_type": "excel", "extractor": "win32com"},
                    "error": None,
                }
            except Exception as com_error:
                logger.warning(
                    f"使用COM接口处理XLS失败: {file_path} - {str(com_error)}"
                )
            finally:
                try:
                    if "workbook" in locals() and workbook is not None:
                        workbook.Close(SaveChanges=False)
                except:
                    pass

                try:
                    if "excel_app" in locals() and excel_app is not None:
                        excel_app.Quit()
                except:
                    pass

                self._force_close_office_processes("EXCEL.EXE")
                pythoncom.CoUninitialize()

        # 策略3: 尝试openpyxl (虽然通常用于xlsx，但有时也能处理某些xls)
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

            content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                data = []
                row_count = 0

                for row in sheet.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else "" for cell in row])
                    row_count += 1
                    if row_count >= 5000:
                        break

                if data:
                    max_lengths = [
                        max(len(str(row[i])) for row in data if i < len(row))
                        for i in range(len(data[0]))
                    ]
                    formatted_rows = []

                    for row in data:
                        formatted_row = [
                            str(val).ljust(max_lengths[i])
                            for i, val in enumerate(row)
                            if i < len(max_lengths)
                        ]
                        formatted_rows.append(" | ".join(formatted_row))

                    sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(formatted_rows)
                    if row_count >= 5000:
                        sheet_text += "\n[注意: 表格过大，仅显示前5000行]"

                    content.append(sheet_text)

            workbook.close()

            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": "excel", "extractor": "openpyxl"},
                "error": None,
            }
        except Exception as openpyxl_error:
            logger.warning(
                f"使用openpyxl处理XLS失败: {file_path} - {str(openpyxl_error)}"
            )

        # 最后策略: 尝试提取二进制文本
        try:
            binary_result = self._extract_binary_text_content(
                file_path, "application/vnd.ms-excel"
            )
            if not binary_result.get("error"):
                return binary_result
        except Exception as e:
            pass

        # 所有方法失败，返回错误
        return self._create_error_result(
            "excel", "无法提取XLS文件内容: 所有已知方法均已尝试并失败"
        )

    def _extract_archive_content(self, file_path: str) -> Dict[str, Any]:
        """提取压缩文件内容，使用安全的临时目录和路径验证"""
        temp_dir = os.path.join(
            os.path.dirname(file_path), f"temp_extract_{uuid.uuid4().hex}"
        )

        try:
            os.makedirs(temp_dir, exist_ok=True)
            
            # 使用Python内置模块处理zip文件
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(file_path) as zip_file:
                    # 限制解压文件数量，避免zip炸弹攻击
                    file_count = 0
                    max_files = 100
                    
                    # 解压文件内容
                    for name in zip_file.namelist():
                        if file_count >= max_files:
                            break
                        # 验证文件路径，防止路径穿越攻击
                        if ".." in name or name.startswith("/"):
                            continue
                        try:
                            zip_file.extract(name, temp_dir)
                            file_count += 1
                        except Exception as e:
                            logger.warning(f"解压文件失败: {name} - {e}")
            else:
                # 对于非zip格式，尝试使用patoolib库
                try:
                    import patoolib
                    patoolib.extract_archive(file_path, outdir=temp_dir)
                except ImportError:
                    return self._create_error_result(
                        Path(file_path).suffix[1:], 
                        "无法处理此压缩格式，缺少patoolib库"
                    )
                except Exception as e:
                    return self._create_error_result(
                        Path(file_path).suffix[1:],
                        f"解压失败: {str(e)}"
                    )

            content = []
            max_files = 100
            max_size = 10 * 1024 * 1024  # 10MB
            total_size = 0
            files_processed = 0

            for root, _, files in os.walk(temp_dir):
                if files_processed >= max_files:
                    content.append(
                        f"\n[警告: 超过最大处理文件数 {max_files}, 仅展示部分内容]"
                    )
                    break

                for file_name in files:
                    if file_name.startswith("._") or file_name.startswith("__MACOSX"):
                        continue

                    if files_processed >= max_files:
                        break

                    file_path_full = os.path.join(root, file_name)
                    normalized_path = os.path.normpath(file_path_full)
                    if not normalized_path.startswith(temp_dir):
                        content.append(f"[警告: 跳过不安全路径: {file_name}]")
                        continue

                    try:
                        file_size = os.path.getsize(file_path_full)
                        if total_size + file_size > max_size:
                            content.append("\n[警告: 达到总大小限制，跳过剩余文件]")
                            break

                        if file_size < 1024 * 1024:  # 1MB
                            try:
                                with open(
                                    file_path_full,
                                    "r",
                                    encoding="utf-8",
                                    errors="ignore",
                                ) as f:
                                    text = f.read(10240)
                                    if text.strip():
                                        content.append(
                                            f"File: {file_name}\n{text[:1024]}..."
                                        )
                                        files_processed += 1
                                        total_size += file_size
                            except (UnicodeDecodeError, IOError):
                                pass
                    except OSError:
                        content.append(f"[警告: 无法读取文件: {file_name}]")

            return {
                "content": "\n\n".join(content),
                "metadata": {"file_type": Path(file_path).suffix[1:]},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(
                Path(file_path).suffix[1:], f"压缩文件处理失败: {str(e)} - 跳过"
            )
        finally:
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"清理临时目录失败: {temp_dir} - {str(e)}")

    def _extract_doc_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOC文件内容，使用专用的Word实例，改进COM对象管理"""
        logger.info(f"处理DOC文件: {file_path}")
        max_retries = 3
        timeout = 10  # 默认超时时间(秒)

        for attempt in range(max_retries):
            try:
                # 设置超时标志和结果变量
                timeout_occurred = [False]
                result = [None]
                exception = [None]
                processing_completed = [False]

                # 主处理函数
                def process_doc():
                    try:
                        # 初始化COM
                        pythoncom.CoInitialize()
                        word_app = None
                        doc = None

                        try:
                            # 创建Word应用实例
                            word_app = win32com.client.Dispatch("Word.Application")
                            word_app.Visible = False
                            word_app.DisplayAlerts = False

                            # 尝试打开文档
                            abs_path = os.path.abspath(file_path)
                            doc = word_app.Documents.Open(
                                abs_path,
                                ReadOnly=True,
                                PasswordDocument="",
                                Visible=False,
                                NoEncodingDialog=True,
                                OpenAndRepair=True,
                            )

                            # 提取文本内容
                            content = doc.Range().Text

                            # 设置结果
                            result[0] = {
                                "content": content,
                                "metadata": {
                                    "file_type": "doc",
                                    "extractor": "win32com",
                                },
                                "error": None
                                if content.strip()
                                else "未提取到任何文本内容",
                            }

                        except Exception as e:
                            # 捕获处理异常
                            exception[0] = e
                        finally:
                            # 确保资源释放
                            try:
                                if doc:
                                    doc.Close(SaveChanges=False)
                            except:
                                pass
                            try:
                                if word_app:
                                    word_app.Quit()
                            except:
                                pass
                            try:
                                pythoncom.CoUninitialize()
                            except:
                                pass

                            # 标记处理完成
                            processing_completed[0] = True
                    except Exception as e:
                        # 捕获总体异常
                        exception[0] = e
                        processing_completed[0] = True

                # 超时处理函数
                def handle_timeout():
                    if not processing_completed[0]:
                        timeout_occurred[0] = True
                        logger.warning(f"处理文件超时 ({timeout}秒): {file_path}")
                        self._force_close_office_processes("WINWORD.EXE")

                # 创建并启动处理线程
                process_thread = threading.Thread(target=process_doc)
                process_thread.daemon = True
                process_thread.start()

                # 创建并启动超时定时器
                timer = threading.Timer(timeout, handle_timeout)
                timer.daemon = True
                timer.start()

                # 等待处理完成或超时
                process_thread.join(timeout + 2)  # 给超时处理留出额外时间
                timer.cancel()  # 取消定时器

                # 检查处理结果
                if timeout_occurred[0]:
                    logger.error(f"处理DOC文件超时: {file_path}")
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": f"处理超时 ({timeout}秒)",
                    }
                # 处理结果检查部分的修改
                elif exception[0]:
                    # 处理过程中发生异常
                    error_str = str(exception[0])

                    # 以下错误类型直接跳过，不再重试
                    skip_retry_errors = [
                        "(-2147352567, '发生意外。', (0, 'Microsoft Word', '命令失败'",  # Word命令失败
                        "(-2147023174, '服务器执行操作时发生错误'",  # 服务器错误
                        "RPC_E_SERVERFAULT",  # RPC服务器错误
                        "COMError",  # 一般COM错误
                        "AutomationException",  # 自动化异常
                    ]

                    # 检查是否是需要直接跳过的错误
                    if any(skip_err in error_str for skip_err in skip_retry_errors):
                        logger.warning(
                            f"检测到无需重试的错误，直接跳过: {error_str} - 文件: {file_path}"
                        )
                        return {
                            "content": "",
                            "metadata": {"file_type": "doc", "extractor": "win32com"},
                            "error": f"Office错误，跳过处理: {error_str}",
                        }

                    # 如果不是跳过类型的错误，且未达到最大重试次数，则重试
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 失败: {error_str}, 重试中..."
                        )
                        time.sleep(2)
                        continue

                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": f"Win32COM 处理失败: {error_str}",
                    }

                elif result[0]:
                    # 处理成功
                    return result[0]
                else:
                    # 未知错误
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 出现未知错误, 重试中..."
                        )
                        time.sleep(2)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": "处理过程遇到未知错误",
                    }

            except Exception as e:
                logger.error(f"处理DOC文件异常: {str(e)}")
                if attempt == max_retries - 1:
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": f"处理失败: {str(e)}",
                    }

        # 清理可能残留的进程
        self._force_close_office_processes("WINWORD.EXE")
        return {
            "content": "",
            "metadata": {"file_type": "doc", "extractor": "win32com"},
            "error": "所有处理尝试均失败",
        }

    def _extract_ppt_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPT文件内容，处理Visible属性错误和多级回退策略"""
        logger.info(f"处理PPT文件: {file_path}")
        max_retries = 3
        timeout = 10  # 默认超时时间(秒)

        for attempt in range(max_retries):
            try:
                # 设置超时标志和结果变量
                timeout_occurred = [False]
                result = [None]
                exception = [None]
                processing_completed = [False]

                # 主处理函数
                def process_ppt():
                    try:
                        # 初始化COM
                        pythoncom.CoInitialize()
                        ppt_app = None
                        presentation = None

                        try:
                            # 创建PowerPoint应用实例
                            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                            # 不设置Visible属性，避免某些环境下的错误

                            # 打开演示文稿
                            abs_path = os.path.abspath(file_path)
                            presentation = ppt_app.Presentations.Open(
                                abs_path, WithWindow=False, ReadOnly=True
                            )

                            # 提取文本内容
                            text_content = []
                            slide_count = presentation.Slides.Count

                            for i in range(1, slide_count + 1):
                                try:
                                    slide = presentation.Slides.Item(i)
                                    slide_text = []

                                    for shape_idx in range(1, slide.Shapes.Count + 1):
                                        try:
                                            shape = slide.Shapes.Item(shape_idx)
                                            if shape.HasTextFrame:
                                                text_frame = shape.TextFrame
                                                if (
                                                    hasattr(text_frame, "TextRange")
                                                    and text_frame.TextRange.Text.strip()
                                                ):
                                                    slide_text.append(
                                                        text_frame.TextRange.Text.strip()
                                                    )
                                        except:
                                            continue

                                    if slide_text:
                                        text_content.append(
                                            f"Slide {i}:\n" + "\n".join(slide_text)
                                        )
                                except:
                                    continue

                            content = "\n\n".join(text_content)

                            # 设置结果
                            result[0] = {
                                "content": content,
                                "metadata": {
                                    "file_type": "ppt",
                                    "extractor": "win32com",
                                    "slide_count": slide_count,
                                },
                                "error": None
                                if content.strip()
                                else "未提取到任何文本内容",
                            }

                        except Exception as e:
                            # 捕获处理异常
                            exception[0] = e
                        finally:
                            # 确保资源释放
                            try:
                                if presentation:
                                    presentation.Close()
                            except:
                                pass
                            try:
                                if ppt_app:
                                    ppt_app.Quit()
                            except:
                                pass
                            try:
                                pythoncom.CoUninitialize()
                            except:
                                pass

                            # 标记处理完成
                            processing_completed[0] = True
                    except Exception as e:
                        # 捕获总体异常
                        exception[0] = e
                        processing_completed[0] = True

                # 超时处理函数
                def handle_timeout():
                    if not processing_completed[0]:
                        timeout_occurred[0] = True
                        logger.warning(f"处理文件超时 ({timeout}秒): {file_path}")
                        self._force_close_office_processes("POWERPNT.EXE")

                # 创建并启动处理线程
                process_thread = threading.Thread(target=process_ppt)
                process_thread.daemon = True
                process_thread.start()

                # 创建并启动超时定时器
                timer = threading.Timer(timeout, handle_timeout)
                timer.daemon = True
                timer.start()

                # 等待处理完成或超时
                process_thread.join(timeout + 2)  # 给超时处理留出额外时间
                timer.cancel()  # 取消定时器

                # 检查处理结果
                if timeout_occurred[0]:
                    logger.error(f"处理PPT文件超时: {file_path}")
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": f"处理超时 ({timeout}秒)",
                    }
                elif exception[0]:
                    # 处理过程中发生异常
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 失败: {str(exception[0])}, 重试中..."
                        )
                        time.sleep(2)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": f"Win32COM 处理失败: {str(exception[0])}",
                    }
                elif result[0]:
                    # 处理成功
                    return result[0]
                else:
                    # 未知错误
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 出现未知错误, 重试中..."
                        )
                        time.sleep(2)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": "处理过程遇到未知错误",
                    }

            except Exception as e:
                logger.error(f"处理PPT文件异常: {str(e)}")
                if attempt == max_retries - 1:
                    return {
                        "content": "",
                        "metadata": {"file_type": "ppt", "extractor": "win32com"},
                        "error": f"处理失败: {str(e)}",
                    }

        # 清理可能残留的进程
        self._force_close_office_processes("POWERPNT.EXE")
        return {
            "content": "",
            "metadata": {"file_type": "ppt", "extractor": "win32com"},
            "error": "所有处理尝试均失败",
        }

    def _force_close_office_processes(self, process_name: str) -> None:
        """强制关闭指定Office进程"""
        try:
            for proc in psutil.process_iter(["pid", "name"]):
                try:
                    if proc.info["name"].upper() == process_name:
                        logger.info(
                            f"关闭Office进程: {proc.info['name']} (PID: {proc.info['pid']})"
                        )

                        # 尝试优雅终止
                        proc.terminate()

                        # 等待进程终止
                        try:
                            proc.wait(timeout=3)
                        except psutil.TimeoutExpired:
                            # 强制终止
                            proc.kill()
                except:
                    continue
        except Exception as e:
            logger.debug(f"关闭Office进程失败: {e}")

    def _extract_binary_text_content(
        self, file_path: str, mime_type: str
    ) -> Dict[str, Any]:
        """从二进制文件中提取可能的文本内容，作为最后的回退方案"""
        try:
            with open(file_path, "rb") as f:
                binary_data = f.read()

            # 尝试不同编码
            potential_text = ""
            encodings = ["utf-8", "latin1", "cp1252", "gb18030", "big5"]

            for encoding in encodings:
                try:
                    decoded = binary_data.decode(encoding, errors="ignore")
                    # 检查解码是否产生有意义的内容
                    if len(decoded.strip()) > len(potential_text.strip()):
                        potential_text = decoded
                except:
                    continue

            # 只提取可打印ASCII字符和常见空白
            printable_text = ""
            for char in potential_text:
                if char.isprintable() or char in " \t\n\r":
                    printable_text += char

            # 移除连续空白
            cleaned_text = re.sub(r"\s+", " ", printable_text).strip()

            # 移除二进制垃圾（随机字符序列）
            final_text = re.sub(
                r"[^\w\s.,;:!?(){}\[\]\'\"<>@#$%^&*+=\-_\\|/]", "", cleaned_text
            )

            # 移除非常短的片段
            words = [w for w in final_text.split(" ") if len(w) > 1]
            result_text = " ".join(words)

            if len(result_text) < 50:  # 如果提取的文本不够多
                return self._create_error_result(
                    mime_type, f"二进制内容提取失败: 未能提取到足够的文本内容"
                )

            return {
                "content": result_text,
                "metadata": {"file_type": mime_type, "extractor": "binary"},
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(mime_type, f"二进制内容提取失败: {str(e)}")



    def extract_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """提取文件内容，包含多级失败回退策略，增强容错性，基于MIME类型而非扩展名"""

        # 将所有操作包装在一个大的try-except块中，确保不会崩溃
        try:
            if not Path(file_path).exists():
                return {
                    "error": f"文件不存在: {file_path}",
                    "content": "",
                    "metadata": {},
                }

            if Path(file_path).stat().st_size == 0:
                return self._create_empty_result(mime_type)

            # 跳过不支持的文件类型
            skip_mime_types = [
                "image/jpeg",
                "image/png",
                "image/gif",
                "image/bmp",
                "audio/mpeg",
                "audio/wav",
                "video/mp4",
                "video/x-msvideo",
                "video/quicktime",
                "font/ttf",
                "font/otf",
                "font/woff",
                "font/woff2",
                "application/x-executable",
            ]

            # 对于dwg文件特别处理
            if mime_type == "image/vnd.dwg":
                return self._create_error_result(mime_type, "DWG图纸文件, 跳过处理")

            # 对于明确是二进制格式且不是文档/压缩文件的特殊处理
            if mime_type in skip_mime_types:
                return self._create_error_result(
                    mime_type, f"不支持的MIME类型: {mime_type}, 跳过处理"
                )

            # 检查文件大小，过大的文件可能会导致处理缓慢
            file_size = os.path.getsize(file_path)
            size_mb = file_size / (1024 * 1024)
            if size_mb > 50:  # 大于50MB的文件
                logger.warning(
                    f"文件过大 ({size_mb:.2f} MB)，可能影响处理性能: {file_path}"
                )

            # 根据MIME类型选择提取器
            # 优先使用基于MIME类型的提取策略，而不是文件扩展名
            if mime_type in self.extractors:
                try:
                    logger.info(f"使用MIME类型 '{mime_type}' 处理: {file_path}")
                    result = self.extractors[mime_type](file_path)
                    if not result.get("error"):
                        return result

                    # 如果主要提取器失败，记录错误并继续尝试备用方法
                    logger.warning(
                        f"主要提取器 '{mime_type}' 失败: {result.get('error')} - 尝试备用方法"
                    )
                except Exception as e:
                    logger.warning(
                        f"主要提取器 '{mime_type}' 发生异常: {str(e)} - 尝试备用方法"
                    )
            else:
                logger.info(f"没有找到MIME类型 '{mime_type}' 的专用提取器，尝试备用方法")

            # 如果无法识别MIME类型或主要提取器失败，尝试使用MarkItDown
            if self.md is not None:
                try:
                    logger.info(f"尝试使用MarkItDown处理: {file_path}")

                    # 设置超时，防止MarkItDown卡住
                    import threading
                    import signal

                    timeout_occurred = [False]
                    markdown_result = [None]
                    markdown_exception = [None]

                    # 处理函数
                    def process_markdown():
                        try:
                            result = self.md.convert(file_path)
                            markdown_result[0] = {
                                "content": result.text_content,
                                "metadata": {
                                    "file_type": mime_type,
                                    "extractor": "markitdown_fallback",
                                },
                                "error": None,
                            }
                        except Exception as e:
                            markdown_exception[0] = e

                    # 创建并启动处理线程
                    markdown_thread = threading.Thread(target=process_markdown)
                    markdown_thread.daemon = True
                    markdown_thread.start()

                    # 设置最长5秒的超时
                    markdown_thread.join(5)

                    if markdown_thread.is_alive():
                        # 如果线程仍在运行，表示超时
                        timeout_occurred[0] = True
                        logger.warning(f"MarkItDown处理超时: {file_path}")
                    elif markdown_exception[0]:
                        # 如果有异常，记录错误
                        logger.warning(
                            f"MarkItDown处理失败: {str(markdown_exception[0])} - 文件: {file_path}"
                        )
                    elif markdown_result[0]:
                        # 如果处理成功，返回结果
                        return markdown_result[0]

                    # 如果未成功处理，继续尝试其他方法
                    logger.info(f"MarkItDown处理不成功，继续尝试其他方法: {file_path}")

                except Exception as e:
                    logger.warning(
                        f"MarkItDown备用提取尝试失败: {str(e)} - 文件: {file_path}"
                    )

            # 基于文件签名的内容推断
            # 如果文件是Office格式，但MIME类型没有匹配到提取器
            header_hex = self.detector.read_file_header(file_path, 16) if hasattr(self.detector, 'read_file_header') else ""
            
            # OLE格式 (DOC/XLS/PPT)
            if header_hex.startswith("D0CF11E0A1B11AE1"):
                try:
                    # 尝试使用OLE分析仪来确定具体类型
                    if hasattr(self.detector, '_analyze_ole_file'):
                        file_type, specific_mime = self.detector._analyze_ole_file(file_path)
                        if specific_mime and specific_mime in self.extractors:
                            logger.info(f"通过OLE分析确定文件类型: {file_type}，使用对应提取器")
                            result = self.extractors[specific_mime](file_path)
                            if not result.get("error"):
                                return result
                except Exception as e:
                    logger.warning(f"OLE格式分析失败: {str(e)}")
            
            # ZIP/OOXML格式 (DOCX/XLSX/PPTX)
            elif header_hex.startswith("504B0304"):
                try:
                    # 尝试使用OOXML分析仪来确定具体类型
                    if hasattr(self.detector, '_analyze_ooxml_file'):
                        file_type, specific_mime = self.detector._analyze_ooxml_file(file_path)
                        if specific_mime and specific_mime in self.extractors:
                            logger.info(f"通过OOXML分析确定文件类型: {file_type}，使用对应提取器")
                            result = self.extractors[specific_mime](file_path)
                            if not result.get("error"):
                                return result
                except Exception as e:
                    logger.warning(f"OOXML格式分析失败: {str(e)}")

            # 最后尝试二进制文本提取作为最后的手段
            try:
                logger.info(f"尝试二进制文本提取: {file_path}")
                result = self._extract_binary_text_content(file_path, mime_type)
                if not result.get("error") and result.get("content"):
                    return result
            except Exception as e:
                logger.warning(f"二进制提取失败: {str(e)} - 文件: {file_path}")

            # 如果所有方法都失败，返回错误
            return self._create_error_result(
                mime_type, f"所有内容提取方法都失败，无法处理文件类型: {mime_type}"
            )

        except Exception as e:
            # 总体异常处理，确保函数永远不会崩溃
            logger.error(f"提取内容时发生严重错误: {str(e)} - 文件: {file_path}")
            return self._create_error_result(
                mime_type, f"提取过程发生严重错误: {str(e)}"
            )

    def __del__(self):
        """对象销毁时清理资源"""
        try:
            self._cleanup_word_app()
            pythoncom.CoUninitialize()
        except Exception as e:
            logger.debug(f"清理资源时出错: {e}")


class SensitiveChecker:
    """敏感内容检查器，使用正则表达式替代 Aho-Corasick"""

    def __init__(self, config_path: str = "sensitive_config.yaml"):
        """初始化敏感词配置"""
        self.config = self._load_config(config_path)
        self.all_keywords = self.config.get("security_marks", []) + [
            kw
            for cat in self.config.get("sensitive_patterns", {}).values()
            for kw in cat.get("keywords", [])
        ]
        escaped_keywords = [re.escape(kw) for kw in self.all_keywords]
        self.keyword_pattern = re.compile("|".join(escaped_keywords))

    def _load_config(self, config_path: str) -> Dict:
        """加载敏感词配置文件"""
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return yaml.safe_load(f)
        except Exception as e:
            logger.error(f"加载敏感词配置失败: {e}")
            return {}

    def check_content(self, text: str) -> List[Tuple[str, List[int]]]:
        """检查文本中的敏感词"""
        keyword_matches = {}
        for match in self.keyword_pattern.finditer(text or ""):
            keyword = match.group()
            if keyword not in keyword_matches:
                keyword_matches[keyword] = []
            keyword_matches[keyword].append(match.start())
        keyword_results = list(keyword_matches.items())

        structured_results = []
        for pattern, weight in self.config.get("structured_patterns", {}).items():
            matches = list(re.finditer(pattern, text or ""))
            if matches:
                positions = [m.start() for m in matches]
                structured_results.append((pattern, positions))

        number_results = []
        for pattern in self.config.get("number_patterns", []):
            matches = list(re.finditer(pattern, text or ""))
            if matches:
                positions = [m.start() for m in matches]
                number_results.append((pattern, positions))

        return keyword_results + structured_results + number_results


class ResultExporter:
    """增强版处理结果导出器，添加文本内容到Excel"""

    def export_to_json(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 JSON 文件"""
        export_data = [
            {
                "file_path": r.file_path,
                "mime_type": r.mime_type,
                "content_preview": self._get_content_preview(r, 200),
                "sensitive_words": [
                    {"word": w, "positions": p} for w, p in r.sensitive_words
                ],
                "error": r.error,
                "processing_time": r.processing_time,
            }
            for r in results
        ]
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)

    def export_to_excel(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 Excel 文件，包含文本内容"""
        try:
            # 创建主工作表数据
            summary_data = [
                {
                    "文件路径": r.file_path,
                    "文件类型": r.mime_type,
                    "敏感词统计": "; ".join(
                        [f"{w}({len(p)}次)" for w, p in r.sensitive_words]
                    ),
                    "处理时间(秒)": round(r.processing_time, 3),
                    "错误信息": r.error or "",
                    "内容预览": self._get_content_preview(r, 500),
                }
                for r in results
            ]

            # 创建内容工作表数据
            content_data = [
                {
                    "文件路径": r.file_path,
                    "文件类型": r.mime_type,
                    "内容": self._get_full_content(r),
                }
                for r in results
            ]

            # 创建Excel工作簿并写入两个工作表
            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                # 摘要工作表
                pd.DataFrame(summary_data).to_excel(
                    writer, sheet_name="处理结果", index=False
                )
                # 内容工作表
                pd.DataFrame(content_data).to_excel(
                    writer, sheet_name="文本内容", index=False
                )

            logger.info(f"结果已导出到Excel（含文本内容）: {output_path}")
        except Exception as e:
            logger.error(f"导出到Excel失败: {e}")
            # 尝试导出简化版本，不包含内容
            try:
                simplified_path = output_path.replace(".xlsx", "_简化版.xlsx")
                # 创建简化版本
                simple_data = [
                    {
                        "文件路径": r.file_path,
                        "文件类型": r.mime_type,
                        "敏感词统计": "; ".join(
                            [f"{w}({len(p)}次)" for w, p in r.sensitive_words]
                        ),
                        "处理时间(秒)": round(r.processing_time, 3),
                        "错误信息": r.error or "",
                    }
                    for r in results
                ]
                pd.DataFrame(simple_data).to_excel(
                    simplified_path, index=False, engine="openpyxl"
                )
                logger.info(f"已导出简化版Excel: {simplified_path}")
            except Exception as e2:
                logger.error(f"导出简化版Excel也失败: {e2}")

    def _get_content_preview(
        self, result: ProcessingResult, max_length: int = 200
    ) -> str:
        """从处理结果中获取内容预览"""
        try:
            if isinstance(result.content, dict):
                content = result.content.get("content", "")
                if isinstance(content, str):
                    return content[:max_length] + (
                        "..." if len(content) > max_length else ""
                    )
            return ""
        except Exception:
            return ""

    def _get_full_content(self, result: ProcessingResult) -> str:
        """从处理结果中获取完整内容"""
        try:
            if isinstance(result.content, dict):
                content = result.content.get("content", "")
                if isinstance(content, str):
                    return content
            return ""
        except Exception:
            return ""


class ResultMonitor:
    """结果监控器，用于实时输出处理进度和结果"""

    def __init__(self, output_csv: str = "processing_results.csv"):
        self.output_csv = output_csv
        self._init_csv()

    def _init_csv(self):
        """初始化 CSV 文件"""
        with open(self.output_csv, "w", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow(["文件名", "识别类型"])

    def record_result(self, result: ProcessingResult):
        """记录单个处理结果"""
        file_name = os.path.basename(result.file_path)
        sensitive_words_count = len(result.sensitive_words)
        recognition_type = "敏感" if sensitive_words_count > 0 else "常规"

        with open(self.output_csv, "a", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow([file_name, recognition_type])

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 处理文件: {result.file_path}")
        print(f"  类型: {result.mime_type}")
        print(f"  识别类型: {recognition_type}")
        if sensitive_words_count > 0:
            print(f"  发现敏感词: {'; '.join([f'{w}({len(p)}次)' for w, p in result.sensitive_words])}")
        if result.error:
            print(f"  错误: {result.error}")
        print("-" * 80)


class FileProcessor:
    """优化后的文件处理器主类，使用基于签名的文件类型检测"""

    def __init__(
        self,
        config_path: str = "sensitive_config.yaml",
        monitor_output: str = "processing_results.csv",
        chunk_size: int = 1000,
        max_workers: Optional[int] = None,
        is_windows: bool = True,
        verbose: bool = False
    ):
        self.detector = EnhancedFileSignatureDetector(verbose=verbose)
        self.extractor = ContentExtractor(detector=self.detector, is_windows=is_windows)
        self.checker = SensitiveChecker(config_path)
        self.exporter = ResultExporter()
        self.monitor = ResultMonitor(monitor_output)
        self.chunk_size = chunk_size
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self._mime_cache = {}
        self._file_size_cache = {}
        self.verbose = verbose

    def _scan_directory(self, directory: str) -> Iterator[List[str]]:
        """使用生成器分批扫描目录文件，确保正确关闭scandir迭代器"""
        current_chunk = []

        try:
            with os.scandir(directory) as dir_iter:
                for entry in dir_iter:
                    try:
                        if entry.is_file(follow_symlinks=False):
                            # 使用更先进的文件过滤逻辑，基于文件名而非扩展名
                            if entry.name.startswith(".") or entry.name.startswith("~$"):
                                continue
                            if bool(re.match(r"^\[\d+\].+", entry.name)):
                                continue
                            
                            # 收集文件
                            current_chunk.append(entry.path)
                            if len(current_chunk) >= self.chunk_size:
                                yield current_chunk
                                current_chunk = []
                        elif entry.is_dir(follow_symlinks=False):
                            for sub_chunk in self._scan_directory(entry.path):
                                yield sub_chunk
                    except (PermissionError, OSError) as e:
                        logger.warning(f"访问文件/目录出错 {entry.path}: {e}")
                        continue
        except (PermissionError, OSError) as e:
            logger.warning(f"访问目录出错 {directory}: {e}")

        if current_chunk:
            yield current_chunk

    def _preload_file_info(self, file_paths: List[str]):
        """预加载文件信息以减少 I/O 操作"""
        for file_path in file_paths:
            try:
                stat = os.stat(file_path)
                self._file_size_cache[file_path] = stat.st_size
                
                # 使用基于签名的文件类型检测 - 修改此处以统一存储格式
                file_type_info = self.detector.detect_file_type(file_path)
                
                # 安全处理返回值并统一存储格式
                if isinstance(file_type_info, tuple):
                    if len(file_type_info) >= 2:
                        # 标准返回值，存储MIME类型
                        mime_type = file_type_info[1]
                    else:
                        # 非标准元组，使用第一个元素
                        mime_type = file_type_info[0] if file_type_info else "application/octet-stream"
                else:
                    # 非元组返回值，直接使用
                    mime_type = file_type_info

                if self.verbose:
                    logger.info(f"检测文件类型: {file_path} -> {mime_type}")
                    
                self._mime_cache[file_path] = mime_type
            except Exception as e:
                logger.warning(f"预加载文件信息失败 {file_path}: {e}")

    def _process_file_batch(self, file_paths: List[str]) -> List[ProcessingResult]:
        """处理一批文件，Office文件串行处理，其他文件并行处理"""
        results = []
        self._preload_file_info(file_paths)

        # 按类型对文件进行分类
        office_files = []
        other_files = []

        for fp in file_paths:
            mime_type = self._mime_cache.get(fp)
            
            # 基于MIME类型来确定是否是需要串行处理的Office文件
            # 不再尝试解包mime_type
            if mime_type in (
                "application/msword",
                "application/vnd.ms-powerpoint",
                "application/vnd.ms-excel"
            ):
                office_files.append(fp)
            else:
                other_files.append(fp)

        # 串行处理Office文件
        if office_files:
            logger.info(f"串行处理 {len(office_files)} 个Office文件")
            for file_path in office_files:
                try:
                    result = self.process_file(file_path)
                    results.append(result)
                    self.monitor.record_result(result)
                except Exception as e:
                    logger.error(f"处理Office文件异常 {file_path}: {e}")
                    # 添加错误记录而不是跳过
                    error_result = ProcessingResult(
                        file_path=file_path,
                        mime_type="unknown",
                        content={"content": "", "error": str(e)},
                        sensitive_words=[],
                        error=f"处理异常: {str(e)}",
                        processing_time=0.0
                    )
                    results.append(error_result)
                    self.monitor.record_result(error_result)

        # 并行处理其他文件
        if other_files:
            current_workers = max(1, min(len(other_files), self.max_workers))
            logger.info(
                f"并行处理 {len(other_files)} 个非Office文件 (工作线程: {current_workers})"
            )

            with ThreadPoolExecutor(max_workers=current_workers) as executor:
                future_to_file = {
                    executor.submit(self.process_file, fp): fp for fp in other_files
                }
                for future in as_completed(future_to_file):
                    file_path = future_to_file[future]
                    try:
                        result = future.result()
                        results.append(result)
                        self.monitor.record_result(result)
                    except Exception as e:
                        logger.error(f"处理文件异常 {file_path}: {e}")

        # 当缓存过大时清理
        if len(self._mime_cache) > 10000:
            self._mime_cache.clear()
            self._file_size_cache.clear()

        return results

    def process_directory(self, directory: str) -> List[ProcessingResult]:
        """处理目录下的所有文件"""
        results = []
        total_files = len(self.detector.get_all_files(directory))
        completed = 0

        print(f"\n开始处理目录: {directory}")
        print(f"共发现 {total_files} 个文件")
        print("=" * 80)

        try:
            for file_chunk in self._scan_directory(directory):
                chunk_results = self._process_file_batch(file_chunk)
                results.extend(chunk_results)
                completed += len(file_chunk)
                print(
                    f"\n已完成: {completed}/{total_files} ({completed / total_files * 100:.1f}%)"
                )
        except KeyboardInterrupt:
            print("\n用户中断处理...")
        except Exception as e:
            logger.error(f"处理目录时发生错误: {e}")
        finally:
            if results:
                try:
                    self.exporter.export_to_json(results, "interrupted_results.json")
                    self.exporter.export_to_excel(results, "interrupted_results.xlsx")
                except Exception as e:
                    logger.error(f"导出中断结果失败: {e}")

        print(f"\n处理完成！共处理 {len(results)}/{total_files} 个文件")
        print(f"详细结果已保存至: {self.monitor.output_csv}")
        print("=" * 80)
        return results

    def process_file(self, file_path: str) -> ProcessingResult:
        """处理单个文件，基于文件签名确定类型，然后应用适当的提取器"""
        start_time = time.time()
        try:
            # 使用基于签名的文件类型检测 - 修改此处以防止解包错误
            mime_type = None
            cached_info = self._mime_cache.get(file_path)
            
            if cached_info is not None:
                # 从缓存中获取MIME类型
                if isinstance(cached_info, tuple):
                    if len(cached_info) >= 2:
                        # 如果是一个至少有两个元素的元组，取第二个元素作为MIME类型
                        mime_type = cached_info[1]
                    else:
                        # 元组但少于两个元素，使用第一个元素
                        mime_type = cached_info[0] if cached_info else "application/octet-stream"
                else:
                    # 如果不是元组，直接使用
                    mime_type = cached_info
            else:
                # 缓存中没有，使用detector获取
                file_type_info = self.detector.detect_file_type(file_path)
                
                # 安全处理不同格式的返回值
                if isinstance(file_type_info, tuple):
                    if len(file_type_info) >= 2:
                        mime_type = file_type_info[1]
                    else:
                        mime_type = file_type_info[0] if file_type_info else "application/octet-stream"
                else:
                    mime_type = file_type_info
            
            # 确保mime_type不为None
            if mime_type is None:
                mime_type = "application/octet-stream"
            
            # 检查文件是否为空
            file_size = self._file_size_cache.get(file_path, 0) or os.path.getsize(file_path)
            if file_size == 0:
                logger.info(f"空文件，跳过: {file_path}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={"content": "", "is_empty": True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time,
                )

            # 检查文件大小，过大的文件可能会导致处理缓慢
            size_mb = file_size / (1024 * 1024)
            if size_mb > 50:  # 大于50MB的文件
                logger.warning(
                    f"文件过大 ({size_mb:.2f} MB)，可能影响处理性能: {file_path}"
                )

            # 使用容错增强版的提取内容方法处理文件
            try:
                content = self.extractor.extract_content(file_path, mime_type)
            except Exception as extract_error:
                # 检查是否是加密工作簿错误
                error_msg = str(extract_error)
                if "Workbook is encrypted" in error_msg:
                    logger.warning(f"检测到加密工作簿: {file_path}")
                    # 将加密工作簿直接标记为敏感文件
                    return ProcessingResult(
                        file_path=file_path,
                        mime_type=mime_type,
                        content={
                            "content": "[加密文件]",
                            "metadata": {"file_type": "加密工作簿", "is_encrypted": True},
                        },
                        sensitive_words=[("文件已加密", [0])],  # 添加一个伪敏感词条目，确保被标记为敏感
                        error="文件已加密，无法读取内容",
                        processing_time=time.time() - start_time,
                    )
                
                # 其他错误按原方式处理
                logger.error(
                    f"提取内容时发生严重错误: {error_msg} - 文件: {file_path}"
                )
                content = {
                    "content": "",
                    "metadata": {"file_type": mime_type},
                    "error": f"提取内容时发生严重错误: {error_msg}",
                }

            # 检查提取结果是否包含加密提示
            if content.get("error") and "Workbook is encrypted" in content.get("error"):
                logger.warning(f"检测到加密工作簿: {file_path}")
                # 将加密工作簿直接标记为敏感文件
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={
                        "content": "[加密文件]",
                        "metadata": {"file_type": "加密工作簿", "is_encrypted": True},
                    },
                    sensitive_words=[("文件已加密", [0])],  # 添加一个伪敏感词条目，确保被标记为敏感
                    error="文件已加密，无法读取内容",
                    processing_time=time.time() - start_time,
                )

            # 检查提取是否失败
            if content.get("error") and "UnsupportedFormatException" in content.get(
                "error"
            ):
                # 如果是不支持的格式，标记为跳过
                logger.info(
                    f"不支持的文件格式，标记为跳过: {file_path} (type: {mime_type})"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={
                        "content": "",
                        "metadata": {"file_type": "unsupported"},
                        "skipped": True,
                    },
                    sensitive_words=[],
                    error=content.get("error"),
                    processing_time=time.time() - start_time,
                )
            elif content.get("error"):
                # 再次检查错误信息中是否包含关于加密的提示
                if "password" in content.get("error").lower() or "encrypted" in content.get("error").lower():
                    logger.warning(f"检测到可能加密的文件: {file_path}")
                    return ProcessingResult(
                        file_path=file_path,
                        mime_type=mime_type,
                        content={
                            "content": "[可能是加密文件]",
                            "metadata": {"file_type": "可能加密文件", "is_encrypted": True},
                        },
                        sensitive_words=[("文件可能已加密", [0])],
                        error=content.get("error"),
                        processing_time=time.time() - start_time,
                    )
                
                logger.warning(
                    f"提取内容失败: {file_path} - MIME: {mime_type} - 错误: {content.get('error')}"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content=content,
                    sensitive_words=[],
                    error=content.get("error"),
                    processing_time=time.time() - start_time,
                )

            # 敏感内容检查 - 包装在try-except中，确保不会因为敏感词检查失败而中断
            try:
                sensitive_result = self.checker.check_content(
                    content.get("content", "")
                )
                sensitive_words = sensitive_result
            except Exception as check_error:
                logger.error(f"敏感内容检查失败: {file_path} - {str(check_error)}")
                sensitive_words = []

            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type,
                content=content,
                sensitive_words=sensitive_words,
                error=None,
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            error_msg = str(e)
            # 再次检查整体异常中是否有加密相关错误
            if "Workbook is encrypted" in error_msg or "password" in error_msg.lower() or "encrypted" in error_msg.lower():
                logger.warning(f"处理过程中检测到加密文件: {file_path} - {error_msg}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type if "mime_type" in locals() else "unknown",
                    content={
                        "content": "[加密文件]",
                        "metadata": {"file_type": "加密文件", "is_encrypted": True},
                    },
                    sensitive_words=[("文件已加密", [0])],
                    error=f"文件已加密: {error_msg}",
                    processing_time=time.time() - start_time,
                )
                
            logger.error(f"处理文件失败: {error_msg} - {file_path}")
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type if "mime_type" in locals() else "unknown",
                content={
                    "content": "",
                    "metadata": {"file_type": "regular"},
                    "skipped": True,
                },
                sensitive_words=[],
                error=f"处理文件失败: {error_msg}",
                processing_time=time.time() - start_time,
            )


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="增强型文件敏感内容检测工具 - 基于文件签名而非扩展名")
    parser.add_argument("path", help="要处理的文件或目录路径")
    parser.add_argument(
        "--config", default="sensitive_config.yaml", help="敏感词配置文件路径"
    )
    parser.add_argument(
        "--output", default="results", help="输出结果文件名(不含扩展名)"
    )
    parser.add_argument(
        "--chunk-size", type=int, default=1000, help="每批处理的文件数量"
    )
    parser.add_argument("--workers", type=int, default=None, help="最大工作线程数")
    parser.add_argument(
        "--no-windows",
        action="store_true",
        help="指定非Windows平台，禁用win32com相关功能",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="启用详细日志输出",
    )
    parser.add_argument(
        "--repair",
        action="store_true",
        help="尝试修复损坏的文件签名",
    )

    args = parser.parse_args()

    try:
        # 初始化文件处理器，使用基于签名的检测方法
        processor = FileProcessor(
            config_path=args.config,
            monitor_output=f"{args.output}_processing.csv",
            chunk_size=args.chunk_size,
            max_workers=args.workers,
            is_windows=not args.no_windows,
            verbose=args.verbose
        )

        results = []
        if Path(args.path).is_file():
            print(f"处理单个文件: {args.path}")
            results = [processor.process_file(args.path)]
            processor.monitor.record_result(results[0])
        else:
            print(f"处理目录: {args.path}")
            results = processor.process_directory(args.path)

        # 输出结果统计
        total = len(results)
        error_count = sum(1 for r in results if r.error)
        sensitive_count = sum(1 for r in results if r.sensitive_words)
        
        print(f"\n处理统计:")
        print(f"总文件数: {total}")
        print(f"敏感文件: {sensitive_count}")
        print(f"处理错误: {error_count}")
        
        # 按MIME类型统计
        mime_counts = {}
        for r in results:
            if r.mime_type not in mime_counts:
                mime_counts[r.mime_type] = 0
            mime_counts[r.mime_type] += 1
            
        print("\n文件类型统计:")
        for mime_type, count in sorted(mime_counts.items(), key=lambda x: x[1], reverse=True):
            if count > 0:
                print(f"  {mime_type}: {count} 个文件")

        # 导出结果
        processor.exporter.export_to_json(results, f"{args.output}.json")
        processor.exporter.export_to_excel(results, f"{args.output}.xlsx")

        print(f"\n处理完成，结果已导出到:")
        print(f"  {args.output}.json")
        print(f"  {args.output}.xlsx")
        print(f"  {args.output}_processing.csv")
        
    except KeyboardInterrupt:
        print("\n用户中断程序执行")
        sys.exit(1)
    except Exception as e:
        logger.error(f"程序执行出错: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
