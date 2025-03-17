import json
import os
import filetype
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime

import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
from typing import List, Dict, Tuple
from io import BytesIO


# 功能：生成data.db（原始content）

class db_info:
    def detect_file_type(filename: str) -> str:
        """检测文件类型"""
        try:
            kind = filetype.guess(filename)
            if kind is None:
                return "未知"
            return kind.mime
        except Exception as e:
            print(f"无法检测文件 {filename} 的类型: {e}")
            return "未知"

    def calculate_md5(filename: str) -> str:
        """计算文件的 MD5 值"""
        hash_md5 = hashlib.md5()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def extract_text_from_file(filename: str, file_type: str) -> str:
        """根据文件类型提取文本内容"""
        try:
            # 【加了文本框的：处理失败】处理 Word 文档 (.docx)
            if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(filename)
                return "\n".join([para.text for para in doc.paragraphs])

            # 处理 PDF 文件
            elif file_type == "application/pdf":
                text = ""
                with pdfplumber.open(filename) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or "" + "\n"
                return text.strip()

            # 处理纯文本文件 (.txt)
            elif file_type == "text/plain":
                with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()

            # 【不行】处理zip文件
            elif file_type == "application/zip":
                with zipfile.ZipFile(filename, 'r') as zip_ref:
                    content = ""
                    for file_info in zip_ref.infolist():
                        with zip_ref.open(file_info) as file:
                            content += f"Content of {file_info.filename}:\n{file.read().decode('utf-8')}\n"
                    return content

            # 【不行】处理 .xlsx 文件
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                workbook = load_workbook(filename=filename)
                content = ""
                for sheet in workbook.sheetnames:
                    content += f"Sheet name: {sheet}\n"
                    worksheet = workbook[sheet]
                    for row in worksheet.iter_rows(values_only=True):
                        content += "\t".join([str(cell) for cell in row]) + "\n"
                return content

            # 【部分不行】处理旧版 .xls 文件
            elif file_type == "application/vnd.ms-excel":
                workbook = xlrd.open_workbook(filename)
                content = ""
                for sheet in workbook.sheets():
                    content += f"Sheet name: {sheet.name}\n"
                    for row in range(sheet.nrows):
                        content += "\t".join([str(cell) for cell in sheet.row_values(row)]) + "\n"
                return content

            # 【部分不行】处理 .pptx 文件
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":

                prs = Presentation(filename)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content

            # 【不行】处理旧版 .ppt 文件
            elif file_type == "application/vnd.ms-powerpoint":
                ole = olefile.OleFileIO(filename)
                data = ole.openstream("PowerPoint Document").read()

                # 使用 BytesIO 来模拟文件对象
                file_like_object = BytesIO(data)
                prs = Presentation(file_like_object)

                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content

            # 【加了文本框的：处理失败】处理 .rtf 文件
            elif file_type == "application/rtf":
                output = pypandoc.convert_file(filename, 'plain', format='rtf')
                return output

            # 【不行】处理 .doc 文件
            elif file_type == "application/msword":
                doc = Document(filename)
                content = "\n".join([para.text for para in doc.paragraphs])
                return content

            # 其他文件类型
            else:
                return "不支持的文件类型"

        except Exception as e:
            print(f"无法从文件 {filename} 中提取文本: {e}")
            return "提取失败"

    def get_file_metadata(filename: str):
        """获取文件的元数据（大小和创建时间）"""
        file_stat = os.stat(filename)
        return {
            "size": file_stat.st_size,  # 文件大小（字节）
            "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")  # 创建时间
        }

    def create_db_and_save_data(directory: str,S_or_R: int):
        """创建数据库并保存文件信息"""
        db_name = "data.db"
        # 连接到 SQLite 数据库（如果不存在则创建）
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()

        # 创建表
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL
            );
        """)

        # 遍历目录下的所有文件
        for idx, filename in enumerate(os.listdir(directory)):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path) :  # 确保是文件
                print(f"正在处理文件: {filename}")

                # 获取文件信息

                file_type = db_info.detect_file_type(file_path)
                md5 = db_info.calculate_md5(file_path)
                content = db_info.extract_text_from_file(file_path, file_type)
                metadata = db_info.get_file_metadata(file_path)


                # 插入数据到数据库

                error = ['提取失败','不支持的文件类型']
                flag = 0
                for e in error:
                    if e in content:
                        flag = 1 # 如果存在error情况则 则不插入数据库

                if flag == 0 and (content != ''):
                    cursor.execute("""
                        INSERT INTO files (md5,filename , file_type, file_size, created_time, content,is_sensitive)
                        VALUES (?, ?, ?, ?, ?, ?, ?)
                    """, (md5,filename, file_type, metadata["size"], metadata["created_time"], content,S_or_R))

        # 删除md5重复数据
        try:

            # 步骤1: 找出所有重复的 md5 值
            cursor.execute('''
                    SELECT md5
                    FROM files
                    GROUP BY md5
                    HAVING COUNT(*) > 1
                ''')
            duplicate_md5_list = [row[0] for row in cursor.fetchall()]

            # 步骤2: 删除所有重复的 md5 对应行
            if duplicate_md5_list:
                # 使用参数化查询防止 SQL 注入
                placeholders = ', '.join(['?'] * len(duplicate_md5_list))
                delete_sql = f'''
                        DELETE FROM files
                        WHERE md5 IN ({placeholders})
                    '''
                cursor.execute(delete_sql, duplicate_md5_list)
                deleted_rows = cursor.rowcount
            else:
                deleted_rows = 0

            # 提交事务
            conn.commit()
            print(f"已删除 {deleted_rows} 条重复数据")

        except sqlite3.Error as e:
            # 回滚事务
            conn.rollback()
            print(f"操作失败: {e}")

        # 提交事务并关闭连接
        conn.commit()
        conn.close()
        print(f"文件信息已保存到数据库: {db_name}")

# 示例用法
if __name__ == "__main__":

    # 敏感文件位置
    #path_directory = r"D:\PyCharm Community Edition 2024.3.1.1\my_program\西湖创新项目\样本\敏感文件/"
    path_directory = input("检测的文件位置（如：D:\样本\敏感文件/）：")
    S_or_R = int(input("敏感文件为1，常规文件标记为0。当前文件标记为："))

    # 将数据最终导入到db中
    db_info.create_db_and_save_data(path_directory,S_or_R)