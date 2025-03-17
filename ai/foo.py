import json
import os
import filetype
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
from typing import List, Dict, Tuple
from io import BytesIO

class db_info:
    @staticmethod
    def detect_file_type(filename: str) -> str:
        """检测文件类型"""
        try:
            kind = filetype.guess(filename)
            if kind is None:
                return "未知"
            return kind.mime
        except Exception as e:
            print(f"无法检测文件 {filename} 的类型: {e}")
            return "未知"

    @staticmethod
    def calculate_md5(filename: str) -> str:
        """计算文件的 MD5 值"""
        hash_md5 = hashlib.md5()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    @staticmethod
    def extract_text_from_file(filename: str, file_type: str) -> str:
        """根据文件类型提取文本内容"""
        try:
            if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(filename)
                return "\n".join([para.text for para in doc.paragraphs])
            elif file_type == "application/pdf":
                text = ""
                with pdfplumber.open(filename) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or "" + "\n"
                return text.strip()
            elif file_type == "text/plain":
                with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            elif file_type == "application/zip":
                with zipfile.ZipFile(filename, 'r') as zip_ref:
                    content = ""
                    for file_info in zip_ref.infolist():
                        with zip_ref.open(file_info) as file:
                            content += f"Content of {file_info.filename}:\n{file.read().decode('utf-8')}\n"
                    return content
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                workbook = load_workbook(filename=filename)
                content = ""
                for sheet in workbook.sheetnames:
                    content += f"Sheet name: {sheet}\n"
                    worksheet = workbook[sheet]
                    for row in worksheet.iter_rows(values_only=True):
                        content += "\t".join([str(cell) for cell in row]) + "\n"
                return content
            elif file_type == "application/vnd.ms-excel":
                workbook = xlrd.open_workbook(filename)
                content = ""
                for sheet in workbook.sheets():
                    content += f"Sheet name: {sheet.name}\n"
                    for row in range(sheet.nrows):
                        content += "\t".join([str(cell) for cell in sheet.row_values(row)]) + "\n"
                return content
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(filename)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content
            elif file_type == "application/vnd.ms-powerpoint":
                ole = olefile.OleFileIO(filename)
                data = ole.openstream("PowerPoint Document").read()
                file_like_object = BytesIO(data)
                prs = Presentation(file_like_object)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content
            elif file_type == "application/rtf":
                output = pypandoc.convert_file(filename, 'plain', format='rtf')
                return output
            elif file_type == "application/msword":
                doc = Document(filename)
                content = "\n".join([para.text for para in doc.paragraphs])
                return content
            else:
                return "不支持的文件类型"
        except Exception as e:
            print(f"无法从文件 {filename} 中提取文本: {e}")
            return "提取失败"

    @staticmethod
    def get_file_metadata(filename: str):
        """获取文件的元数据（大小和创建时间）"""
        file_stat = os.stat(filename)
        return {
            "size": file_stat.st_size,  # 文件大小（字节）
            "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")  # 创建时间
        }

    @staticmethod
    def create_db_and_save_data(directory: str, S_or_R: int):
        """创建数据库并保存文件信息"""
        db_name = "data.db"
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL
            );
        """)

        files_data = []

        for idx, filename in enumerate(os.listdir(directory)):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                print(f"正在处理文件: {filename}")
                file_type = db_info.detect_file_type(file_path)
                md5 = db_info.calculate_md5(file_path)
                content = db_info.extract_text_from_file(file_path, file_type)
                metadata = db_info.get_file_metadata(file_path)
                error = ['提取失败', '不支持的文件类型']
                if content not in error and content:
                    files_data.append((md5, filename, file_type, metadata["size"], metadata["created_time"], content, S_or_R))

        if files_data:
            cursor.executemany("""
                INSERT INTO files (md5, filename, file_type, file_size, created_time, content, is_sensitive)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """, files_data)

        try:
            cursor.execute('''
                DELETE FROM files
                WHERE id NOT IN (
                    SELECT MIN(id)
                    FROM files
                    GROUP BY md5
                )
            ''')
            conn.commit()
            print(f"已删除 {cursor.rowcount} 条重复数据")
        except sqlite3.Error as e:
            conn.rollback()
            print(f"操作失败: {e}")

        conn.commit()
        conn.close()
        print(f"文件信息已保存到数据库: {db_name}")

if __name__ == "__main__":
    path_directory = input("检测的文件位置（如：D:\\样本\\敏感文件\\）：")
    S_or_R = int(input("敏感文件为1，常规文件标记为0。当前文件标记为："))
    db_info.create_db_and_save_data(path_directory, S_or_R)