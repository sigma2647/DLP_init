import json
import os
import filetype
import hashlib
import sqlite3
import pandas as pd
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
import pypandoc
from docx import Document
import logging
from io import BytesIO

# 配置日志记录
logging.basicConfig(
    level=logging.INFO,
    filename='file_processing.log',
    filemode='a',
    format='%(asctime)s - %(levelname)s - %(message)s'
)

class db_info:
    @staticmethod
    def detect_file_type(filename: str) -> str:
        """检测文件类型"""
        try:
            kind = filetype.guess(filename)
            return kind.mime if kind else "未知"
        except Exception as e:
            logging.error(f"无法检测文件 {filename} 的类型: {e}")
            return "未知"

    @staticmethod
    def calculate_md5(filename: str) -> str:
        """计算文件的 MD5 值"""
        hash_md5 = hashlib.md5()
        try:
            with open(filename, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logging.error(f"无法计算文件 {filename} 的 MD5: {e}")
            return ""

    @staticmethod
    def extract_text_from_file(filename: str, file_type: str) -> str:
        """根据文件类型提取文本内容"""
        # 预检查文件是否可访问
        try:
            with open(filename, "rb") as f:
                f.read(1)
        except Exception as e:
            logging.error(f"文件 {filename} 不可访问: {e}")
            return "提取失败"

        try:
            # 处理 Word 文档 (.docx)
            if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(filename)
                return "\n".join([para.text for para in doc.paragraphs])

            # 处理 PDF 文件
            elif file_type == "application/pdf":
                text = ""
                with pdfplumber.open(filename) as pdf:
                    for page in pdf.pages:
                        text += (page.extract_text() or "") + "\n"
                return text.strip()

            # 处理纯文本文件 (.txt)
            elif file_type == "text/plain":
                with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()

            # 处理 ZIP 文件
            elif file_type == "application/zip":
                with zipfile.ZipFile(filename, 'r') as zip_ref:
                    content = ""
                    for file_info in zip_ref.infolist():
                        try:
                            with zip_ref.open(file_info) as file:
                                content += f"Content of {file_info.filename}:\n{file.read().decode('utf-8')}\n"
                        except Exception as e:
                            logging.warning(f"无法读取 ZIP 内文件 {file_info.filename}: {e}")
                    return content

            # 处理 .xlsx 文件
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                df = pd.read_excel(filename, engine='openpyxl', nrows=1000)
                df.dropna(how='all', inplace=True)
                return df.to_string(index=False) if not df.empty else "表格为空"

            # 处理 .xls 文件
            elif file_type == "application/vnd.ms-excel":
                df = pd.read_excel(filename, engine='xlrd', nrows=1000)
                df.dropna(how='all', inplace=True)
                return df.to_string(index=False) if not df.empty else "表格为空"

            # 处理 .pptx 文件
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(filename)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content

            # 处理旧版 .ppt 文件
            elif file_type == "application/vnd.ms-powerpoint":
                ole = olefile.OleFileIO(filename)
                data = ole.openstream("PowerPoint Document").read()
                file_like_object = BytesIO(data)
                prs = Presentation(file_like_object)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content

            # 处理 .rtf 文件
            elif file_type == "application/rtf":
                output = pypandoc.convert_file(filename, 'plain', format='rtf')
                return output

            # 处理 .doc 文件（需要额外工具支持，如 antiword 或转换为 docx）
            elif file_type == "application/msword":
                logging.warning(f"{filename} 是 .doc 文件，建议转换为 .docx 后处理")
                return "不支持的文件类型"

            else:
                return "不支持的文件类型"

        except Exception as e:
            logging.error(f"无法从文件 {filename} 中提取文本: {e}")
            return "提取失败"

    @staticmethod
    def get_file_metadata(filename: str) -> dict:
        """获取文件的元数据（大小和创建时间）"""
        try:
            file_stat = os.stat(filename)
            return {
                "size": file_stat.st_size,
                "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")
            }
        except Exception as e:
            logging.error(f"无法获取文件 {filename} 的元数据: {e}")
            return {"size": 0, "created_time": "未知"}

    @staticmethod
    def create_db_and_save_data(directory: str, S_or_R: int):
        """创建数据库并保存文件信息"""
        db_name = "data.db"
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()

        # 创建表
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL
            );
        """)

        # 遍历目录下的所有文件
        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path) and not filename.startswith("~$"):  # 忽略临时文件
                logging.info(f"正在处理文件: {filename}")

                # 获取文件信息
                file_type = db_info.detect_file_type(file_path)
                md5 = db_info.calculate_md5(file_path)
                content = db_info.extract_text_from_file(file_path, file_type)
                metadata = db_info.get_file_metadata(file_path)

                # 插入数据到数据库
                error_conditions = ['提取失败', '不支持的文件类型']
                if content and content not in error_conditions:
                    cursor.execute("""
                        INSERT INTO files (md5, filename, file_type, file_size, created_time, content, is_sensitive)
                        VALUES (?, ?, ?, ?, ?, ?, ?)
                    """, (md5, filename, file_type, metadata["size"], metadata["created_time"], content, S_or_R))

        # 删除重复的 MD5 数据
        try:
            cursor.execute('''
                SELECT md5
                FROM files
                GROUP BY md5
                HAVING COUNT(*) > 1
            ''')
            duplicate_md5_list = [row[0] for row in cursor.fetchall()]
            if duplicate_md5_list:
                placeholders = ', '.join(['?'] * len(duplicate_md5_list))
                cursor.execute(f'''
                    DELETE FROM files
                    WHERE md5 IN ({placeholders})
                ''', duplicate_md5_list)
                logging.info(f"已删除 {cursor.rowcount} 条重复数据")
            conn.commit()
        except sqlite3.Error as e:
            conn.rollback()
            logging.error(f"删除重复数据失败: {e}")

        conn.close()
        logging.info(f"文件信息已保存到数据库: {db_name}")

if __name__ == "__main__":
    path_directory = input("检测的文件位置（如：D:\样本\敏感文件/）：")
    S_or_R = int(input("敏感文件为1，常规文件标记为0。当前文件标记为："))
    db_info.create_db_and_save_data(path_directory, S_or_R)
