import os
import pandas as pd
import xlrd
import openpyxl
import tempfile
import traceback
import argparse
import binascii
import zipfile
from collections import OrderedDict

# 可以在这里添加更多的导入，用于处理其他类型的文档
# 例如，用于处理Word文档
# import docx  # 需要安装python-docx库

# 用于处理PPT文档
# import pptx  # 需要安装python-pptx库

def detect_file_signature(file_path):
    """检测文件签名来识别文件类型"""
    # 文件签名字典（十六进制 -> 文件类型）
    signatures = OrderedDict([
        # Office文档格式
        (b'504B030414000600', 'Office Open XML (.xlsx)'),  # XLSX专用签名
        (b'504B0304140008', 'Office Open XML (.xlsx)'),    # 另一种XLSX签名
        (b'504B030414000600080000002100', 'Office Open XML (.xlsx)'),  # XLSX另一变体
        
        (b'504B030414000600080000002100E2', 'PowerPoint (.pptx)'),  # PPTX专用签名
        (b'504B0304140006000800000021', 'PowerPoint (.pptx)'),      # 另一种PPTX签名格式
        
        (b'504B0304140006000800000021006', 'Word (.docx)'),  # DOCX专用签名
        (b'504B0304140006000800', 'Word (.docx)'),          # 另一种DOCX签名
        
        (b'D0CF11E0A1B11AE1', 'Microsoft Office Binary File Format'),  # 通用Office二进制格式
        
        # 更具体的Office二进制格式识别（需要检查文件内部更多信息）
        (b'D0CF11E0A1B11AE100', 'Microsoft Excel (.xls)'),  # XLS
        (b'D0CF11E0A1B11AE1000000', 'Microsoft Word (.doc)'),  # DOC
        (b'D0CF11E0A1B11AE10000', 'Microsoft PowerPoint (.ppt)'),  # PPT
        
        # 其他常见文件格式
        (b'25504446', 'PDF (.pdf)'),
        (b'504B0304', 'ZIP Archive (.zip)'),
        (b'89504E47', 'PNG (.png)'),
        (b'47494638', 'GIF (.gif)'),
        (b'FFD8FFE0', 'JPEG (.jpg, .jpeg)'),
        (b'FFD8FFE1', 'JPEG with EXIF (.jpg, .jpeg)'),
        (b'FFD8FFE2', 'JPEG with EXIF (.jpg, .jpeg)'),
        (b'FFD8FFE8', 'JPEG (.jpg, .jpeg)'),
        (b'424D', 'Bitmap (.bmp)'),
        (b'377ABCAF271C', 'Microsoft Access Database (.mdb)'),
        (b'7B5C727466', 'Rich Text Format (.rtf)'),
        (b'CAFEBABE', 'Java Class File (.class)'),
        (b'4D546864', 'MIDI File (.mid)'),
        (b'1A45DFA3', 'Matroska Media File (.mkv)'),
        (b'3026B2758E66CF', 'Microsoft Word 97-2003 (.doc)'),
        (b'52617221', 'RAR Archive (.rar)'),
        (b'1F8B08', 'GZIP (.gz)'),
        (b'FD377A585A00', 'XZ Archive (.xz)'),
        (b'4F676753', 'Ogg Vorbis (.ogg)'),
        (b'38425053', 'Photoshop Document (.psd)'),
        (b'213C617263', 'Internet Archive (.arc)'),
        (b'3C3F786D6C', 'XML (.xml)'),
        (b'0A0D0D0A', 'PCAP Packet Capture (.pcap)'),
        (b'000001BA', 'MPEG-PS (.mpg, .mpeg)'),
        (b'000001B3', 'MPEG-1/2 Video (.mpg, .mpeg)'),
        (b'494433', 'MP3 (.mp3)'),
        (b'EFBBBF', 'UTF-8 encoding text with BOM'),
        (b'FFFE', 'UTF-16 (LE) encoding text with BOM'),
        (b'FEFF', 'UTF-16 (BE) encoding text with BOM'),
        (b'0000FEFF', 'UTF-32 (BE) encoding text with BOM'),
        (b'FFFE0000', 'UTF-32 (LE) encoding text with BOM'),
    ])
    
    try:
        with open(file_path, 'rb') as f:
            # 读取文件开头的32字节，用于更精确的文件类型检测
            header_bytes = f.read(32)
            header_hex = binascii.hexlify(header_bytes).upper()
            
            # 进一步分析Office二进制格式
            if header_hex.startswith(b'D0CF11E0A1B11AE1'):
                # 尝试确定是XLS、DOC还是PPT
                try:
                    # 读取更多数据来查找特定标记
                    f.seek(0)
                    content = f.read(10000)  # 读取前10KB用于分析
                    
                    # 查找特定字符串
                    if b'Microsoft Excel' in content or b'Worksheet' in content or b'Spreadsheet' in content:
                        return 'Microsoft Excel (.xls)'
                    elif b'Microsoft Word' in content or b'Document' in content or b'MSWordDoc' in content:
                        return 'Microsoft Word (.doc)'
                    elif b'Microsoft PowerPoint' in content or b'Presentation' in content or b'PowerPoint' in content:
                        return 'Microsoft PowerPoint (.ppt)'
                except:
                    # 如果无法确定具体类型，返回通用Office格式
                    return 'Microsoft Office Binary File Format'
            
            # 查找ZIP-based Office格式的特定标记
            if header_hex.startswith(b'504B0304'):
                try:
                    # 判断是否是Office Open XML格式
                    import zipfile
                    with zipfile.ZipFile(file_path) as zf:
                        file_list = zf.namelist()
                        
                        # 检查是否含有特定Office格式的文件
                        if any('xl/' in f for f in file_list) or any('workbook.xml' in f for f in file_list):
                            return 'Microsoft Excel (.xlsx)'
                        elif any('word/' in f for f in file_list) or any('document.xml' in f for f in file_list):
                            return 'Microsoft Word (.docx)'
                        elif any('ppt/' in f for f in file_list) or any('presentation.xml' in f for f in file_list):
                            return 'Microsoft PowerPoint (.pptx)'
                except:
                    # 如果无法打开ZIP或判断内部结构，尝试使用常规签名检测
                    pass
            
            # 标准签名检测
            for signature, file_type in signatures.items():
                if header_hex.startswith(signature):
                    return file_type
            
            # 如果未能通过签名识别，尝试获取文件扩展名
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext:
                ext_map = {
                    '.xlsx': 'Microsoft Excel (.xlsx)',
                    '.xls': 'Microsoft Excel (.xls)',
                    '.docx': 'Microsoft Word (.docx)',
                    '.doc': 'Microsoft Word (.doc)',
                    '.pptx': 'Microsoft PowerPoint (.pptx)',
                    '.ppt': 'Microsoft PowerPoint (.ppt)',
                    '.pdf': 'PDF (.pdf)',
                    '.zip': 'ZIP Archive (.zip)',
                    '.txt': 'Text File (.txt)',
                    '.csv': 'CSV File (.csv)',
                }
                if file_ext in ext_map:
                    return f"{ext_map[file_ext]} (基于文件扩展名)"
            
            print("未能通过文件签名识别文件类型")
            return None
    except Exception as e:
        print(f"文件签名检测出错: {str(e)}")
        return None

def extract_excel_content(file_path, output_dir=None, methods=None):
    """根据文件类型使用相应的方法提取内容"""
    print(f"原始文件路径: {file_path}")
    print(f"原始文件大小: {os.path.getsize(file_path)} 字节")
    
    # 首先使用文件签名检测文件类型
    file_type = detect_file_signature(file_path)
    if file_type:
        print(f"文件类型检测结果: {file_type}")
    
    # Set output directory
    if not output_dir:
        output_dir = os.path.dirname(file_path)
    
    # 根据文件类型选择提取方法
    if methods is None:
        if file_type:
            if 'Excel (.xls)' in file_type:
                # 对于XLS文件，使用XLRD和Pandas
                methods = [extract_with_xlrd, extract_with_pandas]
            elif 'Excel (.xlsx)' in file_type:
                # 对于XLSX文件，使用OpenPyXL和Pandas
                methods = [extract_with_openpyxl, extract_with_pandas]
            elif 'Word (.doc' in file_type or 'Word (.docx' in file_type:
                # 对于Word文件，先尝试二进制提取
                methods = [extract_with_binary_read]
                # 还可以添加专门的Word文档提取方法
            elif 'PowerPoint (.ppt' in file_type or 'PowerPoint (.pptx' in file_type:
                # 对于PPT文件，先尝试二进制提取
                methods = [extract_with_binary_read]
                # 还可以添加专门的PPT提取方法
            else:
                # 对于未识别或其他类型文件，尝试所有方法
                methods = [
                    extract_with_pandas,
                    extract_with_xlrd,
                    extract_with_openpyxl,
                    extract_with_binary_read
                ]
        else:
            # 如果无法识别文件类型，尝试所有方法
            methods = [
                extract_with_pandas,
                extract_with_xlrd,
                extract_with_openpyxl,
                extract_with_binary_read
            ]
    
    results = []
    
    for method in methods:
        try:
            print(f"\n尝试使用 {method.__name__} 方法提取内容...")
            content = method(file_path)
            if content and content.strip():
                print(f"使用 {method.__name__} 成功提取了 {len(content)} 字符")
                results.append((method.__name__, content))
            else:
                print(f"使用 {method.__name__} 没有提取到内容")
        except Exception as e:
            print(f"使用 {method.__name__} 提取时出错: {str(e)}")
            traceback.print_exc()
    
    # Return results from all successful methods
    return results, output_dir
    
    results = []
    
    for method in methods:
        try:
            print(f"\n尝试使用 {method.__name__} 方法提取内容...")
            content = method(file_path)
            if content and content.strip():
                print(f"使用 {method.__name__} 成功提取了 {len(content)} 字符")
                results.append((method.__name__, content))
            else:
                print(f"使用 {method.__name__} 没有提取到内容")
        except Exception as e:
            print(f"使用 {method.__name__} 提取时出错: {str(e)}")
            traceback.print_exc()
    
    # Return results from all successful methods
    return results, output_dir

def extract_with_pandas(file_path):
    """使用pandas提取Excel内容"""
    try:
        # 尝试不同的引擎
        engines = ['openpyxl', 'xlrd']
        content = ""
        
        for engine in engines:
            try:
                print(f"尝试使用pandas的{engine}引擎...")
                df = pd.read_excel(file_path, engine=engine)
                sheet_content = df.to_string(index=False)
                if sheet_content:
                    content += f"Sheet (使用 {engine}):\n{sheet_content}\n\n"
            except Exception as e:
                print(f"使用pandas的{engine}引擎失败: {str(e)}")
        
        return content
    except Exception as e:
        print(f"Pandas提取失败: {str(e)}")
        return ""

def extract_with_xlrd(file_path):
    """使用xlrd直接提取Excel内容（适用于.xls格式）"""
    try:
        workbook = xlrd.open_workbook(file_path, formatting_info=False)
        content = []
        
        for sheet_idx in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_idx)
            sheet_content = f"表格名称: {sheet.name}\n"
            
            # 提取表头
            headers = []
            for col in range(sheet.ncols):
                headers.append(str(sheet.cell_value(0, col)))
            
            sheet_content += "表头: " + ", ".join(headers) + "\n\n"
            
            # 提取数据
            for row in range(1, sheet.nrows):
                row_data = []
                for col in range(sheet.ncols):
                    cell_value = sheet.cell_value(row, col)
                    cell_type = sheet.cell_type(row, col)
                    
                    # 处理日期类型
                    if cell_type == xlrd.XL_CELL_DATE:
                        try:
                            date_tuple = xlrd.xldate_as_tuple(cell_value, workbook.datemode)
                            cell_value = f"{date_tuple[0]}-{date_tuple[1]}-{date_tuple[2]}"
                        except:
                            pass
                    
                    row_data.append(str(cell_value))
                
                sheet_content += " | ".join(row_data) + "\n"
            
            content.append(sheet_content)
        
        return "\n\n".join(content)
    except Exception as e:
        print(f"XLRD提取失败: {str(e)}")
        return ""

def extract_with_openpyxl(file_path):
    """使用openpyxl提取Excel内容（适用于.xlsx格式）"""
    try:
        # 创建临时文件
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
        temp_file.close()
        
        # 复制原始文件到临时文件
        with open(file_path, 'rb') as src, open(temp_file.name, 'wb') as dst:
            dst.write(src.read())
        
        try:
            workbook = openpyxl.load_workbook(temp_file.name, read_only=True, data_only=True)
            content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = f"表格名称: {sheet_name}\n"
                
                rows = list(sheet.rows)
                if not rows:
                    continue
                
                # 提取表头
                headers = [str(cell.value) for cell in rows[0]]
                sheet_content += "表头: " + ", ".join(headers) + "\n\n"
                
                # 提取数据
                for row in rows[1:]:
                    row_data = [str(cell.value if cell.value is not None else "") for cell in row]
                    sheet_content += " | ".join(row_data) + "\n"
                
                content.append(sheet_content)
            
            return "\n\n".join(content)
        finally:
            # 删除临时文件
            try:
                os.unlink(temp_file.name)
            except:
                pass
    except Exception as e:
        print(f"OpenPyXL提取失败: {str(e)}")
        return ""

def extract_with_binary_read(file_path):
    """以二进制方式读取文件并提取可读文本"""
    try:
        with open(file_path, 'rb') as file:
            content = file.read()
        
        # 尝试以不同的编码解析文件
        encodings = ['utf-8', 'latin-1', 'ascii', 'gbk', 'gb2312', 'gb18030']
        all_decoded_text = ""
        
        for encoding in encodings:
            try:
                decoded = content.decode(encoding, errors='replace')
                
                # 过滤出可打印字符和基本控制字符
                printable_text = ''.join(c for c in decoded if c.isprintable() or c in ['\n', '\t', '\r', ' '])
                
                # 查找可能的单元格值
                # 文档中常见的字符串片段
                potential_cells = []
                lines = printable_text.split('\n')
                
                for line in lines:
                    # 跳过太短的行
                    if len(line.strip()) < 3:
                        continue
                    
                    # 跳过大量重复字符的行（可能是格式化字符）
                    if len(set(line)) < len(line) / 5:
                        continue
                    
                    potential_cells.append(line)
                
                cell_content = '\n'.join(potential_cells)
                
                # 只有当找到的文本超过当前总文本时才添加
                if len(cell_content) > len(all_decoded_text):
                    all_decoded_text = cell_content
                    print(f"使用编码 {encoding} 提取了 {len(cell_content)} 字符")
            except Exception as e:
                print(f"编码 {encoding} 解析失败: {str(e)}")
        
        return all_decoded_text
    except Exception as e:
        print(f"二进制读取失败: {str(e)}")
        return ""

# 可以添加其他文档类型的提取方法
"""
# 示例: Word文档提取方法
def extract_with_docx(file_path):
    '''使用python-docx提取.docx文档内容'''
    try:
        import docx
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        tables = []
        
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            tables.append("\n".join(table_text))
        
        content = "\n\n".join(paragraphs) + "\n\n表格内容:\n" + "\n\n".join(tables)
        return content
    except Exception as e:
        print(f"DOCX提取失败: {str(e)}")
        return ""

# 示例: PowerPoint提取方法
def extract_with_pptx(file_path):
    '''使用python-pptx提取.pptx文档内容'''
    try:
        import pptx
        prs = pptx.Presentation(file_path)
        
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = [f"===== 幻灯片 {i+1} ====="]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))
        
        content = "\n\n".join(slides_text)
        return content
    except Exception as e:
        print(f"PPTX提取失败: {str(e)}")
        return ""
"""

def main():
    # 设置命令行参数解析
    parser = argparse.ArgumentParser(description='文档内容提取工具 - 根据文件类型智能选择提取方法')
    parser.add_argument('file_path', help='要提取内容的文件路径')
    parser.add_argument('-o', '--output-dir', dest='output_dir', help='输出目录路径，默认为文件所在目录')
    parser.add_argument('-q', '--quiet', action='store_true', help='安静模式，减少输出信息')
    parser.add_argument('-v', '--verbose', action='store_true', help='详细模式，显示更多调试信息')
    parser.add_argument('--methods', nargs='+', choices=['pandas', 'xlrd', 'openpyxl', 'binary'], 
                        help='指定要使用的提取方法，默认根据文件类型自动选择')
    parser.add_argument('--preview-length', type=int, default=500, 
                        help='内容预览的字符长度，默认为500')
    parser.add_argument('--signature-only', action='store_true',
                        help='仅检测文件签名而不尝试提取内容')
    
    args = parser.parse_args()
    
    # 设置日志级别
    if args.quiet:
        import logging
        logging.basicConfig(level=logging.ERROR)
    elif args.verbose:
        import logging
        logging.basicConfig(level=logging.DEBUG)
    
    # 检查文件是否存在
    if not os.path.exists(args.file_path):
        print(f"错误: 文件不存在: {args.file_path}")
        return 1
    
    # 如果只需检测签名
    if args.signature_only:
        file_type = detect_file_signature(args.file_path)
        if file_type:
            print(f"文件类型（根据签名）: {file_type}")
            return 0
        else:
            print("未能通过文件签名识别文件类型")
            # 尝试使用系统file命令
            try:
                import subprocess
                process = subprocess.Popen(['file', args.file_path], stdout=subprocess.PIPE)
                output, error = process.communicate()
                if output:
                    print(f"系统file命令检测结果: {output.decode('utf-8')}")
            except:
                pass
            return 1
    
    # 设置输出目录
    output_dir = args.output_dir if args.output_dir else os.path.dirname(args.file_path)
    if not os.path.exists(output_dir):
        try:
            os.makedirs(output_dir)
            print(f"已创建输出目录: {output_dir}")
        except Exception as e:
            print(f"创建输出目录失败: {str(e)}")
            output_dir = os.path.dirname(args.file_path)
    
    # 根据参数选择要使用的方法
    if args.methods:
        selected_methods = []
        method_map = {
            'pandas': extract_with_pandas,
            'xlrd': extract_with_xlrd,
            'openpyxl': extract_with_openpyxl, 
            'binary': extract_with_binary_read
        }
        for method_name in args.methods:
            if method_name in method_map:
                selected_methods.append(method_map[method_name])
        methods = selected_methods if selected_methods else None
    else:
        methods = None
    
    print("正在尝试提取文档内容...")
    results, output_dir = extract_excel_content(args.file_path, output_dir, methods)
    
    if results:
        print("\n" + "=" * 50)
        print(f"成功使用 {len(results)} 种方法提取内容")
        
        # 保存所有成功提取的内容
        for i, (method_name, content) in enumerate(results):
            base_filename = os.path.basename(args.file_path)
            # 获取检测到的文件类型
            file_type = detect_file_signature(args.file_path)
            file_type_str = file_type.replace(' ', '_').replace('(', '').replace(')', '').replace(',', '').replace('.', '')
            
            output_file = os.path.join(output_dir, f"{base_filename}_{file_type_str}_content_{i}_{method_name}.txt")
            with open(output_file, 'w', encoding='utf-8', errors='replace') as f:
                f.write(content)
            print(f"已将使用 {method_name} 提取的内容保存到: {output_file}")
            
            # 显示内容预览
            preview_length = args.preview_length
            preview = content[:preview_length] + ("..." if len(content) > preview_length else "")
            print(f"\n{method_name} 提取的内容预览:\n{'-' * 50}\n{preview}\n{'-' * 50}")
    else:
        # 如果所有方法都失败，再次尝试使用文件签名检测详细类型
        print("\n所有提取方法都失败，尝试更详细的文件类型检测...")
        
        # 再次使用文件签名检测
        file_type = detect_file_signature(args.file_path)
        if file_type:
            print(f"文件签名检测结果: {file_type}")
        
        # 尝试使用系统file命令检测文件类型
        try:
            import subprocess
            process = subprocess.Popen(['file', args.file_path], stdout=subprocess.PIPE)
            output, error = process.communicate()
            if output:
                print(f"系统file命令检测结果: {output.decode('utf-8')}")
        except:
            pass
        
        # 尝试提取更多的十六进制文件头信息
        try:
            with open(args.file_path, 'rb') as f:
                # 读取文件开头的32字节
                header_bytes = f.read(32)
                header_hex = binascii.hexlify(header_bytes).upper()
                print(f"文件头32字节的十六进制表示: {header_hex.decode()}")
        except Exception as e:
            print(f"读取文件头失败: {str(e)}")
            
        print("\n未能成功提取Excel内容，请尝试使用专业文件恢复工具")
        return 1
    
    return 0

if __name__ == "__main__":
    exit(main())
