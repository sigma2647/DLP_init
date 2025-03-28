import magic
import json
import os
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
from typing import List, Dict, Tuple
from io import BytesIO

class db_info:
    @staticmethod
    def detect_file_type(filename: str) -> str:
        """检测文件类型"""
        try:
            mime = magic.Magic(mime=True)
            file_type = mime.from_file(filename)
            return file_type
        except Exception as e:
            print(f"无法检测文件 {filename} 的类型: {e}")
            return "未知"

    @staticmethod
    def calculate_md5(filename: str) -> str:
        """计算文件的 MD5 值"""
        hash_md5 = hashlib.md5()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    @staticmethod
    def extract_text_from_file(filename: str, file_type: str) -> str:
        """根据文件类型提取文本内容"""
        try:
            if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(filename)
                return "\n".join([para.text for para in doc.paragraphs])
            elif file_type == "application/pdf":
                text = ""
                with pdfplumber.open(filename) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or "" + "\n"
                return text.strip()
            elif file_type == "text/plain":
                with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            elif file_type == "application/zip":
                with zipfile.ZipFile(filename, 'r') as zip_ref:
                    content = ""
                    for file_info in zip_ref.infolist():
                        with zip_ref.open(file_info) as file:
                            content += f"Content of {file_info.filename}:\n{file.read().decode('utf-8')}\n"
                    return content
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                workbook = load_workbook(filename=filename)
                content = ""
                for sheet in workbook.sheetnames:
                    content += f"Sheet name: {sheet}\n"
                    worksheet = workbook[sheet]
                    for row in worksheet.iter_rows(values_only=True):
                        content += "\t".join([str(cell) for cell in row]) + "\n"
                return content
            elif file_type == "application/vnd.ms-excel":
                workbook = xlrd.open_workbook(filename)
                content = ""
                for sheet in workbook.sheets():
                    content += f"Sheet name: {sheet.name}\n"
                    for row in range(sheet.nrows):
                        content += "\t".join([str(cell) for cell in sheet.row_values(row)]) + "\n"
                return content
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                prs = Presentation(filename)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content
            elif file_type == "application/vnd.ms-powerpoint":
                ole = olefile.OleFileIO(filename)
                data = ole.openstream("PowerPoint Document").read()
                file_like_object = BytesIO(data)
                prs = Presentation(file_like_object)
                content = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                return content
            elif file_type == "application/rtf":
                output = pypandoc.convert_file(filename, 'plain', format='rtf')
                return output
            elif file_type == "application/msword":
                doc = Document(filename)
                content = "\n".join([para.text for para in doc.paragraphs])
                return content
            else:
                # 使用 markitdown 库提取内容
                from markitdown import MarkItDown
                content = MarkItDown().convert(filename)
                if content and len(content.strip()) > 0:
                    return content
        except Exception as e:
            print(f"无法从文件 {filename} 中提取文本: {e}")
            return "提取失败"

    @staticmethod
    def get_file_metadata(filename: str):
        """获取文件的元数据（大小和创建时间）"""
        file_stat = os.stat(filename)
        return {
            "size": file_stat.st_size,  # 文件大小（字节）
            "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")  # 创建时间
        }

    @staticmethod
    def create_db_and_save_data(directory: str, S_or_R: int):
        """创建数据库并保存文件信息"""
        db_name = "data.db"
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL
            );
        """)

        files_data = []

        for idx, filename in enumerate(os.listdir(directory)):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                print(f"正在处理文件: {filename}")
                file_type = db_info.detect_file_type(file_path)
                md5 = db_info.calculate_md5(file_path)
                content = db_info.extract_text_from_file(file_path, file_type)
                metadata = db_info.get_file_metadata(file_path)
                error = ['提取失败', '不支持的文件类型']
                if content not in error and content:
                    files_data.append((md5, filename, file_type, metadata["size"], metadata["created_time"], content, S_or_R))

        if files_data:
            cursor.executemany("""
                INSERT INTO files (md5, filename, file_type, file_size, created_time, content, is_sensitive)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """, files_data)

        try:
            cursor.execute('''
                DELETE FROM files
                WHERE id NOT IN (
                    SELECT MIN(id)
                    FROM files
                    GROUP BY md5
                )
            ''')
            conn.commit()
            print(f"已删除 {cursor.rowcount} 条重复数据")
        except sqlite3.Error as e:
            conn.rollback()
            print(f"操作失败: {e}")

        conn.commit()
        conn.close()
        print(f"文件信息已保存到数据库: {db_name}")

##### 数据库字段content初步格式化 ######
class format_DBcontent:
    def process_content(self, content, file_type):
        """根据文件类型处理内容"""
        if not content:
            return ""

        # XML处理（application/zip类型）
        if file_type == "application/zip":
            return self._process_xml(content)

        # 其他类型按CSS处理
        return self._process_css(content)

    def _process_xml(self, text):
        """处理XML内容：移除标签并提取文本"""
        try:
            # 尝试解析XML
            root = ElementTree.fromstring(text)
            elements = [elem.text.strip() for elem in root.iter() if elem.text]
            cleaned = " ".join(elements)
        except ElementTree.ParseError:
            # 解析失败时使用正则清理
            cleaned = re.sub(r"<[^>]+>", "", text)

        # 处理HTML转义字符并最终清理
        return self._final_clean(html.unescape(cleaned))

    def _process_css(self, text):
        """处理CSS内容：移除注释和特殊结构"""
        # 移除CSS注释
        text = re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)
        # 移除URL引用
        text = re.sub(r"url$[^)]*$", "", text, flags=re.IGNORECASE)
        # 移除CSS变量
        text = re.sub(r"var$--[^)]*$", "", text)
        return self._final_clean(text)

    def _final_clean(self, text):
        """最终清理：保留有效的中英文词汇"""
        # 转换为小写
        text = text.lower()
        # 匹配中文汉字和基本拉丁字母
        words = re.findall(r"[\u4e00-\u9fa5a-zA-Z]+", text)
        return " ".join(words).strip()

    def process_database(self, source_db, target_db):
        """处理数据库迁移"""
        try:
            # 数据库连接
            src_conn = sqlite3.connect(source_db)
            src_cur = src_conn.cursor()
            dst_conn = sqlite3.connect(target_db)
            dst_cur = dst_conn.cursor()

            # 复制表结构（使用IF NOT EXISTS避免表已存在的错误）
            src_cur.execute("SELECT sql FROM sqlite_master WHERE type='table' AND name='files'")
            create_table_sql = src_cur.fetchone()[0]
            # 在CREATE TABLE后添加IF NOT EXISTS
            create_table_sql = create_table_sql.replace("CREATE TABLE", "CREATE TABLE IF NOT EXISTS")
            dst_cur.execute(create_table_sql)

            # 获取字段索引
            src_cur.execute("PRAGMA table_info(files)")
            columns = [col[1] for col in src_cur.fetchall()]
            content_idx = columns.index("content")
            type_idx = columns.index("file_type")

            # 先清空目标表（如果需要）
            dst_cur.execute("DELETE FROM files")
            dst_conn.commit()

            # 批量处理数据
            src_cur.execute("SELECT * FROM files")
            while True:
                rows = src_cur.fetchmany(1000)
                if not rows:
                    break

                processed = []
                for row in rows:
                    row = list(row)
                    content = row[content_idx]
                    file_type = row[type_idx]

                    # 处理内容
                    row[content_idx] = self.process_content(content, file_type) if content else ""
                    processed.append(row)

                # 批量插入
                dst_cur.executemany(
                    f"INSERT INTO files VALUES ({','.join(['?'] * len(columns))})",
                    processed
                )
                dst_conn.commit()

            print(f"数据库处理完成：{target_db}")

        except Exception as e:
            print(f"处理失败：{str(e)}")
            dst_conn.rollback()
        finally:
            src_conn.close()
            dst_conn.close()

if __name__ == "__main__":
    path_directory = input("检测的文件位置（如：D:\\样本\\敏感文件\\）：")
    S_or_R = int(input("敏感文件为1，常规文件标记为0。当前文件标记为："))
    db_info.create_db_and_save_data(path_directory, S_or_R)

    ###### 数据库content内容初步格式化 ######
    fdb = format_DBcontent()
    fdb.process_database('data.db', 'new_data.db')





