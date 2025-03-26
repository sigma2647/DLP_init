import json
import os
import filetype
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime

import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
from typing import List, Dict, Tuple
from io import BytesIO
from xml.etree import ElementTree
import html
from docx import Document
import magic
import json
import os
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
from typing import List, Dict, Tuple
from io import BytesIO
import re
import win32com.client
import os
import re
from pythoncom import com_error


# 功能：生成data.db（原始content）

class db_info:
    @staticmethod
    def extract_text_from_doc_win32(filename):
        """
        使用win32 COM接口从DOC文件中提取文本
        需要安装：pip install pywin32
        """
        try:

            # 创建Word应用对象
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False  # 不显示Word应用程序
            word.DisplayAlerts = False  # 不显示任何警告

            try:
                # 获取绝对路径
                abs_path = os.path.abspath(filename)

                # 打开文档
                doc = word.Documents.Open(abs_path, ReadOnly=True, Visible=False)

                # 提取文本
                text = doc.Content.Text

                # 关闭文档
                doc.Close(SaveChanges=False)

                # 清理文本 - 移除多余的空白字符等
                text = re.sub(r'\r+', '\n', text)  # 替换回车符为换行符
                text = re.sub(r'\n\s*\n', '\n\n', text)  # 删除空行中的空白字符

                return text

            except com_error as e:
                error_msg = str(e)
                if "is protected" in error_msg.lower() or "password" in error_msg.lower():
                    return f"提取失败: Word文件已加密 - {e}"
                else:
                    return f"提取失败: Word COM错误 - {e}"
            finally:
                # 总是确保关闭Word应用
                word.Quit()
        except Exception as e:
            return f"提取失败: 无法使用Win32提取DOC - {e}"

    @staticmethod
    def detect_file_type(filename: str) -> Tuple[str, list]:
        """改进的文件类型检测方法，不依赖扩展名，使用多种检测方法，修正优先级逻辑"""
        try:
            # 尝试多种方法检测文件类型
            detected_types = []

            # 1. 使用魔术字节库检测
            mime = magic.Magic(mime=True)
            detected_mime = mime.from_file(filename)
            detected_types.append({"method": "magic_bytes", "type": detected_mime})

            # 2. 尝试从文件名获取提示
            filename_lower = filename.lower()
            if "_ppt" in filename_lower or filename_lower.endswith(".ppt") or filename_lower.endswith(".pptx"):
                detected_types.append({"method": "filename_hint", "type": "powerpoint_hint"})
            elif "_doc" in filename_lower or filename_lower.endswith(".doc") or filename_lower.endswith(".docx"):
                detected_types.append({"method": "filename_hint", "type": "word_hint"})
            elif "_xls" in filename_lower or filename_lower.endswith(".xls") or filename_lower.endswith(".xlsx"):
                detected_types.append({"method": "filename_hint", "type": "excel_hint"})

            # 3. 尝试读取文件头部以进行更精确的Office文档识别
            try:
                with open(filename, 'rb') as f:
                    header = f.read(8)  # 读取文件头部字节

                    # 检测Office文档
                    if header.startswith(b'\xD0\xCF\x11\xE0'):  # OLE2 Compound Document Header
                        detected_types.append({"method": "file_header", "type": "ms_compound_document"})
                    elif header.startswith(b'PK\x03\x04'):  # ZIP header (used by modern Office formats)
                        detected_types.append({"method": "file_header", "type": "zip_based_document"})
            except Exception as e:
                print(f"读取文件头部失败: {e}")

            # 4. 尝试使用特定库检测
            # 尝试PowerPoint
            try:
                if olefile.isOleFile(filename):
                    ole = olefile.OleFileIO(filename)
                    if ole.exists('PowerPoint Document'):
                        detected_types.append({"method": "ole_specific", "type": "application/vnd.ms-powerpoint"})
                        print(f"检测到OLE格式的PPT文件: {filename}")
                    elif ole.exists('WordDocument'):
                        detected_types.append({"method": "ole_specific", "type": "application/msword"})
                        print(f"检测到OLE格式的DOC文件: {filename}")
            except Exception as e:
                print(f"OLE文件检测失败: {e}")

            # 尝试docx
            try:
                Document(filename)
                detected_types.append({"method": "library_specific",
                                       "type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"})
            except Exception as e:
                if "not a zip file" not in str(e).lower() and "invalid header" not in str(e).lower():
                    print(f"docx检测失败: {e}")

            # 尝试 xlsx
            try:
                load_workbook(filename, read_only=True)
                detected_types.append({"method": "library_specific",
                                       "type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"})
            except Exception as e:
                if "not a zip file" not in str(e).lower() and "invalid header" not in str(e).lower():
                    print(f"xlsx检测失败: {e}")

                    # 如果检测到"加密的工作簿"错误，可能是误识别的docx文件
                    if "workbook is encrypted" in str(e).lower():
                        # 再次尝试将其视为docx
                        try:
                            with open(filename, 'rb') as f:
                                header = f.read(50)
                                if b'word/document.xml' in header or b'word/' in header:
                                    detected_types.append({"method": "content_hint", "type": "possible_docx"})
                        except:
                            pass

            # 尝试 PDF
            try:
                with pdfplumber.open(filename) as pdf:
                    if pdf.pages:
                        detected_types.append({"method": "library_specific", "type": "application/pdf"})
            except Exception as e:
                if "not a pdf file" not in str(e).lower():
                    print(f"pdf检测失败: {e}")

            # 尝试使用Win32检测
            try:
                # 尝试作为DOC文件
                win32_text = db_info.extract_text_from_doc_win32(filename)
                if win32_text and not win32_text.startswith("提取失败"):
                    detected_types.append({"method": "win32_specific", "type": "possible_msword"})
                    print(f"Win32成功读取文件(可能是DOC): {filename}")
            except Exception as e:
                print(f"Win32文件检测失败: {e}")

            # 打印检测到的所有类型以便调试
            print(f"文件 {filename} 的检测结果: {detected_types}")

            # ===== 修改的优先级逻辑 =====
            # 根据多种检测方法和提示确定最可能的文件类型

            # 1. 检查魔术字节检测结果是否明确（非通用值）
            non_generic_mime_types = [
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/pdf",
                "text/plain"
            ]

            # 2. 基于文件名提示调整优先级
            has_ppt_hint = any(
                t["method"] == "filename_hint" and t["type"] == "powerpoint_hint" for t in detected_types)
            has_doc_hint = any(t["method"] == "filename_hint" and t["type"] == "word_hint" for t in detected_types)
            has_xls_hint = any(t["method"] == "filename_hint" and t["type"] == "excel_hint" for t in detected_types)

            # 3. 检查OLE特定检测
            ole_ppt = any(
                t["method"] == "ole_specific" and t["type"] == "application/vnd.ms-powerpoint" for t in detected_types)
            ole_doc = any(t["method"] == "ole_specific" and t["type"] == "application/msword" for t in detected_types)

            # 基于以上信息确定文件类型
            final_type = detected_mime  # 默认使用magic库结果

            # 如果magic检测结果明确，优先使用
            if detected_mime in non_generic_mime_types:
                final_type = detected_mime
                print(f"使用明确的MIME类型: {final_type}")

                # 文件名提示比magic有更高优先级（仅当有强烈提示时）
                if has_ppt_hint and "powerpoint" not in final_type.lower() and "presentation" not in final_type.lower():
                    if ole_ppt:
                        final_type = "application/vnd.ms-powerpoint"
                        print(f"基于文件名和OLE检测覆盖为PPT: {filename}")
                    elif detected_mime == "application/msword":  # 特别处理被误识别为Word的PPT
                        final_type = "application/vnd.ms-powerpoint"
                        print(f"基于文件名提示修正被误判为Word的PPT: {filename}")

                # 对于被Win32识别为Word但magic表明是PPT的文件，保持为PPT
                if "powerpoint" in detected_mime.lower() and any(
                        t["method"] == "win32_specific" for t in detected_types):
                    print(f"保持PowerPoint类型，忽略Win32识别结果: {filename}")
                    final_type = detected_mime

            # 如果magic检测为通用格式，使用其他方法
            else:
                # 使用文件名提示
                if has_ppt_hint:
                    final_type = "application/vnd.ms-powerpoint"
                    print(f"基于文件名提示确定为PPT: {filename}")
                elif has_doc_hint:
                    final_type = "application/msword"
                    print(f"基于文件名提示确定为DOC: {filename}")
                elif has_xls_hint:
                    final_type = "application/vnd.ms-excel"
                    print(f"基于文件名提示确定为XLS: {filename}")
                # 使用OLE检测
                elif ole_ppt:
                    final_type = "application/vnd.ms-powerpoint"
                    print(f"基于OLE检测确定为PPT: {filename}")
                elif ole_doc:
                    final_type = "application/msword"
                    print(f"基于OLE检测确定为DOC: {filename}")
                # 使用库特定检测
                elif any(t["method"] == "library_specific" for t in detected_types):
                    lib_type = next(t["type"] for t in detected_types if t["method"] == "library_specific")
                    final_type = lib_type
                    print(f"基于库特定检测确定为: {final_type}")
                # 特殊情况：处理可能被误识别为加密Excel的docx
                elif any("workbook is encrypted" in str(t).lower() for t in detected_types) and \
                        any(t["method"] == "content_hint" and t["type"] == "possible_docx" for t in detected_types):
                    final_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    print(f"处理被误判为加密Excel的DOCX: {filename}")

            print(f"最终确定的文件类型: {final_type}")
            return final_type, detected_types

        except Exception as e:
            print(f"无法检测文件 {filename} 的类型: {e}")
            return "未知", [{"method": "error", "type": str(e)}]

    @staticmethod
    def extract_text_from_file(filename: str) -> str:
        """根据检测到的文件类型提取文本内容，增强处理无扩展名文件的能力"""
        try:
            # 使用增强的检测方法获取文件类型
            file_type, detection_info = db_info.detect_file_type(filename)

            # 处理Word文档 (DOCX)
            if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                try:
                    doc = Document(filename)
                    return "\n".join([para.text for para in doc.paragraphs])
                except Exception as e:
                    error_msg = str(e)
                    print(f"DOCX处理失败: {error_msg}")

                    # 尝试多种方法处理可能被误识别的文件
                    if "is encrypted" in error_msg.lower() or "workbook is encrypted" in error_msg.lower():
                        # 尝试其他方法读取
                        try:
                            # 1. 尝试使用pypandoc
                            try:
                                output = pypandoc.convert_file(filename, 'plain', format='docx')
                                if output:
                                    return output
                            except Exception as pe:
                                print(f"Pypandoc转换失败: {pe}")

                            # 2. 尝试使用python-docx的另一种方式
                            try:
                                from docx import Document
                                with open(filename, 'rb') as f:
                                    doc = Document(f)
                                    return "\n".join([para.text for para in doc.paragraphs])
                            except Exception as de:
                                print(f"Document再次尝试失败: {de}")

                            # 3. 尝试读取ZIP结构中的XML文件
                            try:
                                import zipfile
                                from lxml import etree
                                with zipfile.ZipFile(filename, 'r') as zip_ref:
                                    doc_xml = zip_ref.read('word/document.xml')
                                    root = etree.fromstring(doc_xml)
                                    ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                                    paragraphs = root.xpath('//w:p', namespaces=ns)
                                    text = '\n'.join([''.join([run.text for run in para.xpath('.//w:t', namespaces=ns)])
                                                      for para in paragraphs])
                                    return text
                            except Exception as ze:
                                print(f"ZIP XML提取失败: {ze}")

                        except Exception as compound_error:
                            print(f"所有替代方法都失败: {compound_error}")
                            return f"提取失败: 文档似乎是DOCX但无法提取 - 尝试了多种方法都失败"

                    # 尝试使用 MarkItDown
                    try:
                        from markitdown import MarkItDown
                        content = MarkItDown().convert(filename)
                        if content and len(content.strip()) > 0:
                            return content
                    except Exception as md_error:
                        print(f"MarkItDown处理失败: {md_error}")

                    return f"提取失败: Word文件 {filename} 处理错误 - {e}"

            # 处理PDF文件
            elif file_type == "application/pdf":
                try:
                    text = ""
                    with pdfplumber.open(filename) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                    return text.strip()
                except Exception as e:
                    if "password" in str(e).lower():
                        return f"提取失败: PDF文件已加密 - {e}"
                    return f"提取失败: PDF文件 {filename} 处理错误 - {e}"

            # 处理纯文本文件
            elif file_type == "text/plain":
                try:
                    # 尝试多种编码方式
                    encodings = ['utf-8', 'latin-1', 'gbk', 'gb2312', 'big5']
                    for encoding in encodings:
                        try:
                            with open(filename, "r", encoding=encoding) as f:
                                return f.read()
                        except UnicodeDecodeError:
                            continue

                    # 如果所有编码都失败，尝试使用markitdown
                    try:
                        from markitdown import MarkItDown
                        return MarkItDown().convert(filename)
                    except Exception as md_error:
                        print(f"MarkItDown处理纯文本失败: {md_error}")
                        return f"提取失败: 无法以任何已知编码读取文本文件"

                except Exception as e:
                    print(f"无法以任何编码读取文本文件 {filename}: {e}")
                    return f"提取失败: {e}"

            # 处理ZIP文件
            elif file_type == "application/zip" or file_type == "application/x-zip-compressed":
                try:
                    with zipfile.ZipFile(filename, 'r') as zip_ref:
                        content = ""
                        for file_info in zip_ref.infolist():
                            try:
                                if file_info.filename.endswith('/'):  # 跳过目录
                                    continue
                                with zip_ref.open(file_info) as file:
                                    content += f"Content of {file_info.filename}:\n{file.read().decode('utf-8', errors='replace')}\n"
                            except Exception as e:
                                content += f"Content of {file_info.filename}: 无法解码内容 - {str(e)}\n"
                        return content
                except Exception as e:
                    if "password" in str(e).lower() or "encrypted" in str(e).lower():
                        return f"提取失败: ZIP文件已加密 - {e}"
                    return f"提取失败: ZIP文件 {filename} 处理错误 - {e}"

            # 处理Excel文件 (XLSX)
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                try:
                    # 首先尝试使用MarkItDown
                    try:
                        from markitdown import MarkItDown
                        content = MarkItDown().convert(filename)
                        if content and len(content.strip()) > 0:
                            return content
                    except Exception as md_error:
                        print(f"MarkItDown处理Excel失败: {md_error}")

                    # 如果MarkItDown失败，回退到openpyxl
                    workbook = load_workbook(filename=filename, read_only=True, data_only=True)
                    content = ""
                    for sheet_name in workbook.sheetnames:
                        content += f"Sheet name: {sheet_name}\n"
                        sheet = workbook[sheet_name]
                        for row in sheet.iter_rows(values_only=True):
                            content += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                    return content
                except Exception as e:
                    if "is encrypted" in str(e).lower() or "password" in str(e).lower():
                        return f"提取失败: Excel文件已加密 - {e}"
                    return f"提取失败: Excel文件 {filename} 处理错误 - {e}"

            # 处理Excel文件 (XLS)
            elif file_type == "application/vnd.ms-excel":
                try:
                    # 首先尝试使用MarkItDown
                    try:
                        from markitdown import MarkItDown
                        content = MarkItDown().convert(filename)
                        if content and len(content.strip()) > 0:
                            return content
                    except Exception as md_error:
                        print(f"MarkItDown处理XLS失败: {md_error}")

                    # 如果MarkItDown失败，回退到xlrd
                    workbook = xlrd.open_workbook(filename)
                    content = ""
                    for sheet in workbook.sheets():
                        content += f"Sheet name: {sheet.name}\n"
                        for row in range(sheet.nrows):
                            content += "\t".join([str(cell) if cell else "" for cell in sheet.row_values(row)]) + "\n"
                    return content
                except Exception as e:
                    if "workbook is encrypted" in str(e).lower():
                        # 如果报错是"workbook is encrypted"，尝试作为docx处理
                        print("检测到可能误判的加密Excel错误，尝试作为DOCX处理")
                        try:
                            doc = Document(filename)
                            return "\n".join([para.text for para in doc.paragraphs])
                        except Exception as doc_e:
                            print(f"作为DOCX重试失败: {doc_e}")

                    if "password" in str(e).lower() or "encrypted" in str(e).lower():
                        return f"提取失败: Excel文件已加密 - {e}"
                    return f"提取失败: Excel文件 {filename} 处理错误 - {e}"

            # 处理PowerPoint文件 (PPTX)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                try:
                    prs = Presentation(filename)
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                    return content
                except Exception as e:
                    if "password" in str(e).lower() or "encrypted" in str(e).lower():
                        return f"提取失败: PowerPoint文件已加密 - {e}"
                    return f"提取失败: PowerPoint文件 {filename} 处理错误 - {e}"

            # 处理PowerPoint文件 (PPT)
            elif file_type == "application/vnd.ms-powerpoint":
                try:
                    ole = olefile.OleFileIO(filename)
                    data = ole.openstream("PowerPoint Document").read()
                    file_like_object = BytesIO(data)
                    prs = Presentation(file_like_object)
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                    return content
                except Exception as e:
                    if "password" in str(e).lower() or "encrypted" in str(e).lower():
                        return f"提取失败: PowerPoint文件已加密 - {e}"
                    return f"提取失败: PowerPoint文件 {filename} 处理错误 - {e}"

            # 处理RTF文件
            elif file_type == "application/rtf":
                try:
                    output = pypandoc.convert_file(filename, 'plain', format='rtf')
                    return output
                except Exception as e:
                    return f"提取失败: RTF文件 {filename} 处理错误 - {e}"

            # 处理Word文件 (DOC)
            elif file_type == "application/msword" or "ms_compound_document" in str(detection_info):
                try:
                    # 首先尝试使用Win32方法
                    win32_text = db_info.extract_text_from_doc_win32(filename)
                    if win32_text and not win32_text.startswith("提取失败"):
                        return win32_text

                    # 如果Win32方法失败，尝试使用python-docx
                    print("Win32方法提取DOC失败，尝试使用python-docx")
                    doc = Document(filename)
                    content = "\n".join([para.text for para in doc.paragraphs])
                    return content
                except Exception as e:
                    if "password" in str(e).lower() or "encrypted" in str(e).lower() or "protected" in str(e).lower():
                        return f"提取失败: Word文件已加密 - {e}"

                    # 尝试使用 MarkItDown
                    try:
                        from markitdown import MarkItDown
                        content = MarkItDown().convert(filename)
                        if content and len(content.strip()) > 0:
                            return content
                    except Exception as md_error:
                        print(f"MarkItDown处理DOC失败: {md_error}")

                    # 最后尝试使用OLE方法提取
                    try:
                        import olefile
                        if olefile.isOleFile(filename):
                            ole = olefile.OleFileIO(filename)
                            if ole.exists('WordDocument'):
                                stream = ole.openstream('WordDocument')
                                data = stream.read()
                                return f"提取了{len(data)}字节的原始DOC数据，需要进一步解析"
                    except Exception as ole_error:
                        print(f"OLE提取DOC失败: {ole_error}")

                    return f"提取失败: Word文件 {filename} 处理错误 - 原始错误: {e}"

            # 对于无法识别的文件类型，尝试使用多种方法
            else:
                # 首先尝试作为docx处理（很多无扩展名文件实际上是docx）
                try:
                    doc = Document(filename)
                    extracted_text = "\n".join([para.text for para in doc.paragraphs])
                    if extracted_text.strip():
                        print(f"成功将未知类型文件作为DOCX处理: {filename}")
                        return extracted_text
                except Exception as docx_e:
                    print(f"作为DOCX处理失败: {docx_e}")

                # 尝试使用 MarkItDown
                try:
                    from markitdown import MarkItDown
                    content = MarkItDown().convert(filename)
                    if content and len(content.strip()) > 0:
                        return content
                except Exception as md_error:
                    print(f"MarkItDown处理未知类型失败: {md_error}")

                # 其他尝试...
                # 1. 尝试使用pypandoc
                try:
                    output = pypandoc.convert_file(filename, 'plain', format='docx')
                    if output:
                        return output
                except Exception as pe:
                    print(f"Pypandoc转换未知类型失败: {pe}")

                return f"提取失败: 未知文件类型 {file_type} - 尝试了多种方法都失败"
        except Exception as e:
            error_msg = f"提取失败: {e}"
            print(f"无法从文件 {filename} 中提取文本: {e}")
            return error_msg

    @staticmethod
    def calculate_md5(filename: str) -> str:
        """计算文件的 MD5 值"""
        hash_md5 = hashlib.md5()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    @staticmethod
    def get_file_metadata(filename: str):
        """获取文件的元数据（大小和创建时间）"""
        file_stat = os.stat(filename)
        return {
            "size": file_stat.st_size,  # 文件大小（字节）
            "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")  # 创建时间
        }

    @staticmethod
    def process_file(filename: str) -> Tuple[str, str, Dict, str, str]:
        """处理单个文件，返回所有必要信息"""
        print(f"正在处理文件: {filename}")
        file_type, detection_info = db_info.detect_file_type(filename)
        md5 = db_info.calculate_md5(filename)
        content = db_info.extract_text_from_file(filename)
        metadata = db_info.get_file_metadata(filename)

        # 检查内容是否提取失败
        extraction_status = "成功"
        if isinstance(content, str) and content.startswith("提取失败"):
            extraction_status = content
            content = ""  # 清空内容，但保留错误信息

        return file_type, content, metadata, extraction_status, md5

    @staticmethod
    def create_db_and_save_data(directory: str, S_or_R: int):
        """创建数据库并保存文件信息，包括提取失败的文件"""
        db_name = "data.db"
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()

        # Check if extraction_status column exists and add it if it doesn't
        try:
            cursor.execute("SELECT extraction_status FROM files LIMIT 1")
        except sqlite3.OperationalError:
            # Column doesn't exist, add it
            try:
                cursor.execute("ALTER TABLE files ADD COLUMN extraction_status TEXT DEFAULT '成功'")
                conn.commit()
                print("添加 extraction_status 列到数据库表")
            except sqlite3.OperationalError:
                # Table might not exist yet
                pass

        # Create table if it doesn't exist
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL,
                extraction_status TEXT DEFAULT '成功'
            );
        """)
        conn.commit()

        files_data = []

        for idx, filename in enumerate(os.listdir(directory)):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                try:
                    file_type, content, metadata, extraction_status, md5 = db_info.process_file(file_path)

                    # 添加所有文件，包括提取失败的
                    files_data.append((md5, filename, file_type, metadata["size"], metadata["created_time"], content,
                                       S_or_R, extraction_status))
                except Exception as e:
                    print(f"处理文件 {filename} 时出错: {e}")
                    # 即使处理过程出错，也添加文件记录
                    try:
                        md5 = db_info.calculate_md5(file_path)
                        metadata = db_info.get_file_metadata(file_path)
                        files_data.append((
                                          md5, filename, "未知", metadata["size"], metadata["created_time"], "", S_or_R,
                                          f"处理失败: {e}"))
                    except Exception as inner_e:
                        print(f"无法获取文件 {filename} 的基本信息: {inner_e}")

        if files_data:
            cursor.executemany("""
                INSERT INTO files (md5, filename, file_type, file_size, created_time, content, is_sensitive, extraction_status)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """, files_data)

        try:
            cursor.execute('''
                DELETE FROM files
                WHERE id NOT IN (
                    SELECT MIN(id)
                    FROM files
                    GROUP BY md5
                )
            ''')
            conn.commit()
            print(f"已删除 {cursor.rowcount} 条重复数据")
        except sqlite3.Error as e:
            conn.rollback()
            print(f"操作失败: {e}")

        conn.commit()
        conn.close()
        print(f"文件信息已保存到数据库: {db_name}")


###### 数据库字段content初步格式化 ######
class format_DBcontent:
    def process_content(self, content, file_type):
        """根据文件类型处理内容"""
        if not content:
            return ""

        # XML处理（application/zip类型）
        if file_type == "application/zip":
            return self._process_xml(content)

        # 其他类型按CSS处理
        return self._process_css(content)

    def _process_xml(self, text):
        """处理XML内容：移除标签并提取文本"""
        try:
            # 尝试解析XML
            root = ElementTree.fromstring(text)
            elements = [elem.text.strip() for elem in root.iter() if elem.text]
            cleaned = " ".join(elements)
        except ElementTree.ParseError:
            # 解析失败时使用正则清理
            cleaned = re.sub(r"<[^>]+>", "", text)

        # 处理HTML转义字符并最终清理
        return self._final_clean(html.unescape(cleaned))

    def _process_css(self, text):
        """处理CSS内容：移除注释和特殊结构"""
        # 移除CSS注释
        text = re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)
        # 移除URL引用
        text = re.sub(r"url$[^)]*$", "", text, flags=re.IGNORECASE)
        # 移除CSS变量
        text = re.sub(r"var$--[^)]*$", "", text)
        return self._final_clean(text)

    def _final_clean(self, text):
        """最终清理：保留有效的中英文词汇"""
        # 转换为小写
        text = text.lower()
        # 匹配中文汉字和基本拉丁字母
        words = re.findall(r"[\u4e00-\u9fa5a-zA-Z]+", text)
        return " ".join(words).strip()

    def process_database(self, source_db, target_db):
        """处理数据库迁移"""
        try:
            # 数据库连接
            src_conn = sqlite3.connect(source_db)
            src_cur = src_conn.cursor()
            dst_conn = sqlite3.connect(target_db)
            dst_cur = dst_conn.cursor()

            # 复制表结构
            src_cur.execute("SELECT sql FROM sqlite_master WHERE type='table' AND name='files'")
            dst_cur.execute(src_cur.fetchone()[0])

            # 获取字段索引
            src_cur.execute("PRAGMA table_info(files)")
            columns = [col[1] for col in src_cur.fetchall()]
            content_idx = columns.index("content")
            type_idx = columns.index("file_type")

            # 批量处理数据
            src_cur.execute("SELECT * FROM files")
            while True:
                rows = src_cur.fetchmany(1000)
                if not rows:
                    break

                processed = []
                for row in rows:
                    row = list(row)
                    content = row[content_idx]
                    file_type = row[type_idx]

                    # 处理内容
                    row[content_idx] = self.process_content(content, file_type) if content else ""
                    processed.append(row)

                # 批量插入
                dst_cur.executemany(
                    f"INSERT INTO files VALUES ({','.join(['?']*len(columns))})",
                    processed
                )
                dst_conn.commit()

            print(f"数据库处理完成：{target_db}")

        except Exception as e:
            print(f"处理失败：{str(e)}")
            dst_conn.rollback()
        finally:
            src_conn.close()
            dst_conn.close()

# 示例用法
if __name__ == "__main__":

    # 敏感文件位置
    #path_directory = r"D:\PyCharm Community Edition 2024.3.1.1\my_program\西湖创新项目\样本\敏感文件/"
    path_directory = input("检测的文件位置（如：D:\样本\敏感文件/）：")
    S_or_R = int(input("敏感文件为1，常规文件标记为0。当前文件标记为："))
    #
    # # 将数据最终导入到db中
    # db_info.create_db_and_save_data(path_directory,S_or_R)


    ###### 数据库content内容初步格式化 ######
    fdb = format_DBcontent()
    fdb.process_database('data.db', 'new_data.db')