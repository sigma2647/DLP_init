import magic
import json
import os
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
import binascii
from typing import List, Dict, Tuple, Optional
from io import BytesIO
from pathlib import Path
import logging
import html
from xml.etree import ElementTree

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class FileAnalyzer:
    """
    Enhanced file signature detection and content extraction
    """
    
    # Common file signatures (magic numbers) and their corresponding information
    FILE_SIGNATURES = {
        # Office documents
        "504B0304": {
            "mime_types": ["application/zip"],
            "possible_formats": ["ZIP", "DOCX", "XLSX", "PPTX", "ODT"],
            "extensions": [".zip", ".docx", ".xlsx", ".pptx", ".odt"],
        },
        "D0CF11E0A1B11AE1": {
            "mime_types": [
                "application/x-ole-storage",
                "application/msword",
                "application/vnd.ms-excel",
                "application/vnd.ms-powerpoint",
            ],
            "possible_formats": ["DOC", "XLS", "PPT", "OLE Compound Document"],
            "extensions": [".doc", ".xls", ".ppt", ".msg"],
        },
        "25504446": {
            "mime_types": ["application/pdf"],
            "possible_formats": ["PDF"],
            "extensions": [".pdf"],
        },
        # Compressed files
        "526172211A07": {
            "mime_types": ["application/x-rar"],
            "possible_formats": ["RAR Archive"],
            "extensions": [".rar"],
        },
        "377ABCAF": {
            "mime_types": ["application/x-7z-compressed"],
            "possible_formats": ["7-Zip Archive"],
            "extensions": [".7z"],
        },
        "1F8B08": {
            "mime_types": ["application/gzip"],
            "possible_formats": ["GZIP Archive"],
            "extensions": [".gz"],
        },
        # Image files
        "FFD8FFE0": {
            "mime_types": ["image/jpeg"],
            "possible_formats": ["JPEG Image"],
            "extensions": [".jpg", ".jpeg"],
        },
        "89504E470D0A1A0A": {
            "mime_types": ["image/png"],
            "possible_formats": ["PNG Image"],
            "extensions": [".png"],
        },
        "47494638": {
            "mime_types": ["image/gif"],
            "possible_formats": ["GIF Image"],
            "extensions": [".gif"],
        },
    }

    # Office specific format detection markers
    OFFICE_MARKERS = {
        # DOCX markers (inside ZIP)
        "word/document.xml": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        # XLSX markers (inside ZIP)
        "xl/workbook.xml": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        # PPTX markers (inside ZIP)
        "ppt/presentation.xml": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        # DOC markers (inside OLE)
        "WordDocument": "application/msword",
        # XLS markers (inside OLE)
        "Workbook": "application/vnd.ms-excel",
        # PPT markers (inside OLE)
        "PowerPoint Document": "application/vnd.ms-powerpoint",
    }

    # File extension to MIME type mapping (for fallback)
    MIME_TYPES = {
        ".txt": "text/plain",
        ".csv": "text/csv",
        ".xml": "text/xml",
        ".html": "text/html",
        ".htm": "text/html",
        ".json": "application/json",
        ".yaml": "application/yaml",
        ".yml": "application/yaml",
        ".md": "text/markdown",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".zip": "application/zip",
        ".rar": "application/x-rar",
        ".7z": "application/x-7z-compressed",
        ".tar": "application/x-tar",
        ".gz": "application/gzip",
        ".bz2": "application/x-bzip2",
        ".pdf": "application/pdf",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".svg": "image/svg+xml",
        ".webp": "image/webp",
        ".mp3": "audio/mpeg",
        ".mp4": "video/mp4",
        ".avi": "video/x-msvideo",
        ".mov": "video/quicktime",
        ".wav": "audio/wav",
        ".bin": "application/octet-stream",
        ".exe": "application/x-msdownload",
        ".dll": "application/x-msdownload",
        ".rtf": "application/rtf",
        ".dot": "application/msword",
        ".pps": "application/vnd.ms-powerpoint",
        ".xlt": "application/vnd.ms-excel",
    }

    # Extensions to skip processing
    SKIP_EXTENSIONS = {
        ".dwg", ".mp3", ".wav", ".mp4", ".avi", ".mkv", ".flv",
        ".mov", ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif",
        ".tiff", ".webp", ".exe", ".dll", ".so", ".class", ".pyc", 
        ".pyd", ".wasm", ".ttf", ".otf", ".woff", ".woff2", ".eot",
    }

    def __init__(self):
        """Initialize the file analyzer with magic library"""
        try:
            self.mime = magic.Magic(mime=True)
            logger.info("Initialized magic library for MIME type detection")
        except ImportError:
            logger.warning("python-magic library not installed, some detection features will be limited")
            self.mime = None
        except Exception as e:
            logger.error(f"Error initializing magic library: {e}")
            self.mime = None

        # Reverse mapping
        self.MIME_TO_EXT = {mime: ext for ext, mime in self.MIME_TYPES.items()}

    def read_file_header(self, filename: str, bytes_to_read: int = 32) -> str:
        """
        Read the first N bytes of a file and return as hex string
        
        Args:
            filename: File path
            bytes_to_read: Number of bytes to read
            
        Returns:
            Uppercase hex string of the file header
        """
        try:
            with open(filename, "rb") as f:
                header = f.read(bytes_to_read)
            return binascii.hexlify(header).decode("ascii").upper()
        except Exception as e:
            logger.error(f"Error reading file header: {filename} - {e}")
            return ""

    def process_content(self, content: str, file_type: str) -> str:
        """
        Process content based on file type
        
        Args:
            content: Raw text content
            file_type: MIME type of the file
            
        Returns:
            Processed text
        """
        if not content or content in ["Unsupported file type", "Text extraction failed"]:
            return ""
            
        # XML processing (for application/zip type)
        if file_type == "application/zip" or file_type.endswith("+xml") or file_type == "text/xml":
            return self._process_xml(content)
            
        # Other types processed as text
        return self._process_text(content)
        
    def _process_xml(self, text: str) -> str:
        """
        Process XML content: remove tags and extract text
        
        Args:
            text: XML content
            
        Returns:
            Cleaned text
        """
        try:
            # Try to parse as XML by wrapping in root element
            try:
                root = ElementTree.fromstring(f"<root>{text}</root>")
                elements = [elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()]
                cleaned = " ".join(elements)
            except ElementTree.ParseError:
                # If parsing fails, use regex cleaning
                cleaned = re.sub(r"<[^>]+>", " ", text)
                
            # Handle HTML escape characters and final cleanup
            return self._final_clean(html.unescape(cleaned))
        except Exception as e:
            logger.error(f"Error processing XML content: {e}")
            # Fallback
            return self._final_clean(text)
        
    def _process_text(self, text: str) -> str:
        """
        Process text content: remove special structures and unnecessary parts
        
        Args:
            text: Text content
            
        Returns:
            Cleaned text
        """
        try:
            # Remove CSS comments
            text = re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)
            # Remove URL references
            text = re.sub(r"url\([^)]*\)", "", text, flags=re.IGNORECASE)
            # Remove CSS variables
            text = re.sub(r"var\(--[^)]*\)", "", text)
            # Remove non-printable characters
            text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
            # Final cleanup
            return self._final_clean(text)
        except Exception as e:
            logger.error(f"Error processing text content: {e}")
            # Fallback
            return text.strip()
        
    def _final_clean(self, text: str) -> str:
        """
        Final cleanup: preserve valid Chinese and Latin words
        
        Args:
            text: Text to clean
            
        Returns:
            Cleaned text
        """
        try:
            # Convert to lowercase
            text = text.lower()
            # Match Chinese characters and basic Latin letters
            words = re.findall(r"[\u4e00-\u9fa5a-zA-Z0-9]+", text)
            # Join with spaces
            return " ".join(words).strip()
        except Exception as e:
            logger.error(f"Error in final text cleaning: {e}")
            # Return original if cleaning fails
            return text.strip()
    
    def create_db_and_save_data(self, directory: str, S_or_R: int):
        """
        Create database and save file information
        
        Args:
            directory: Directory to scan
            S_or_R: Sensitivity flag (1 for sensitive, 0 for regular)
        """
        db_name = "enhanced_data.db"
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()
        
        # Create table with improved structure
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                processed_content TEXT,
                is_sensitive INTEGER NOT NULL,
                processing_time REAL,
                error_message TEXT
            );
        """)
        
        files_data = []
        
        # Get list of files in directory
        for idx, filename in enumerate(os.listdir(directory)):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                logger.info(f"Processing file {idx+1}: {filename}")
                start_time = datetime.now()
                error_message = ""
                
                try:
                    # Detect file type
                    file_type = self.detect_file_type(file_path)
                    # Calculate MD5
                    md5 = self.calculate_md5(file_path)
                    # Get metadata
                    metadata = self.get_file_metadata(file_path)
                    
                    # Skip certain file types
                    ext = Path(filename).suffix.lower()
                    if ext in self.SKIP_EXTENSIONS:
                        logger.info(f"Skipping binary file: {filename}")
                        continue
                        
                    # Extract text content
                    content = self.extract_text_from_file(file_path, file_type)
                    
                    # Process content
                    processed_content = self.process_content(content, file_type)
                    
                    # Skip empty or failed extraction
                    if not content or content in ["Unsupported file type", "Text extraction failed"]:
                        error_message = f"Failed to extract text from {filename}"
                        logger.warning(error_message)
                    
                    # Record processing time
                    end_time = datetime.now()
                    processing_time = (end_time - start_time).total_seconds()
            normalized_path = os.path.normpath(os.path.abspath(filename))
            return self.mime.from_file(normalized_path)
        except Exception as e:
            logger.error(f"Magic MIME type detection failed: {filename} - {e}")
            return ""

    def inspect_zip_content(self, filename: str) -> Optional[str]:
        """
        Check ZIP file contents to determine if it's an OOXML document
        
        Args:
            filename: ZIP file path
            
        Returns:
            MIME type if determined, otherwise None
        """
        try:
            with zipfile.ZipFile(filename) as zip_file:
                file_list = zip_file.namelist()
                
                # Check for DOCX
                if any(name.startswith("word/document.xml") for name in file_list):
                    return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                # Check for XLSX
                elif any(name.startswith("xl/workbook.xml") for name in file_list):
                    return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                # Check for PPTX
                elif any(name.startswith("ppt/presentation.xml") for name in file_list):
                    return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    
                # If no specific markers found, it's a generic ZIP
                return "application/zip"
        except Exception as e:
            logger.warning(f"ZIP content inspection failed: {filename} - {e}")
            return None

    def inspect_ole_content(self, filename: str) -> Optional[str]:
        """
        Check OLE file contents to determine specific Office document type
        
        Args:
            filename: OLE file path
            
        Returns:
            MIME type if determined, otherwise None
        """
        try:
            # First, quick determination by extension
            ext = Path(filename).suffix.lower()
            if ext == ".doc":
                return "application/msword"
            elif ext == ".xls":
                return "application/vnd.ms-excel"
            elif ext == ".ppt":
                return "application/vnd.ms-powerpoint"
                
            # If extension doesn't help, try deep analysis
            # Read a large chunk to search for markers
            with open(filename, "rb") as f:
                content = f.read(16384)  # Read 16KB
                
            content_str = content.decode("latin1", errors="ignore")
            
            # Check for Word markers
            if "WordDocument" in content_str or "Microsoft Word" in content_str:
                return "application/msword"
                
            # Check for Excel markers
            if "Workbook" in content_str or "Microsoft Excel" in content_str:
                return "application/vnd.ms-excel"
                
            # Check for PowerPoint markers
            if "PowerPoint Document" in content_str or "Microsoft PowerPoint" in content_str:
                return "application/vnd.ms-powerpoint"
                
            # If no specific markers found, it's a generic OLE file
            return "application/x-ole-storage"
        except Exception as e:
            logger.warning(f"OLE content inspection failed: {filename} - {e}")
            return None

    def detect_file_type(self, filename: str) -> str:
        """
        Comprehensive file type detection, returning MIME type
        
        Args:
            filename: File path
            
        Returns:
            Detected MIME type
        """
        try:
            # Basic file information
            file_size = os.path.getsize(filename)
            file_extension = Path(filename).suffix.lower()
            expected_mime = self.MIME_TYPES.get(file_extension, "application/octet-stream")
            
            # Skip very small or empty files
            if file_size == 0:
                return "application/octet-stream"
                
            # Read file header (first 32 bytes)
            header_hex = self.read_file_header(filename, 32)
            
            # Get magic MIME type (if available)
            magic_mime = self.get_magic_mime_type(filename) if self.mime else ""
            
            # Detection through file header signature
            
            # Check for ZIP signature
            if header_hex.startswith("504B0304"):  # ZIP signature
                zip_mime = self.inspect_zip_content(filename)
                if zip_mime:
                    return zip_mime
                    
            # Check for OLE signature
            elif header_hex.startswith("D0CF11E0A1B11AE1"):  # OLE signature
                ole_mime = self.inspect_ole_content(filename)
                if ole_mime:
                    return ole_mime
                    
            # Check for PDF signature
            elif header_hex.startswith("25504446"):
                return "application/pdf"
                
            # If signature doesn't identify, try using magic library
            if magic_mime and magic_mime != "application/octet-stream":
                return magic_mime
                
            # Finally, try using extension
            if file_extension and expected_mime != "application/octet-stream":
                return expected_mime
                
            # If everything else fails, return binary stream type
            return "application/octet-stream"
            
        except Exception as e:
            logger.error(f"File type detection failed: {filename} - {e}")
            return "application/octet-stream"

    def calculate_md5(self, filename: str) -> str:
        """Calculate MD5 hash of a file"""
        hash_md5 = hashlib.md5()
        with open(filename, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def get_file_metadata(self, filename: str) -> Dict:
        """
        Get file metadata (size and creation time)
        
        Args:
            filename: File path
            
        Returns:
            Dictionary with file metadata
        """
        file_stat = os.stat(filename)
        return {
            "size": file_stat.st_size,  # File size in bytes
            "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")  # Creation time
        }
        
    def extract_text_from_file(self, filename: str, file_type: str) -> str:
        """
        Extract text content from file based on its type
        
        Args:
            filename: File path
            file_type: MIME type of the file
            
        Returns:
            Extracted text content
        """
        try:
            # Office documents - DOCX
            if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                try:
                    doc = Document(filename)
                    paragraphs = [para.text for para in doc.paragraphs if para.text]
                    return "\n".join(paragraphs)
                except Exception as e:
                    logger.error(f"Error extracting text from DOCX: {filename} - {e}")
                    # Fallback for problematic DOCX files
                    try:
                        return pypandoc.convert_file(filename, 'plain', format='docx')
                    except Exception as e2:
                        logger.error(f"Fallback extraction failed for DOCX: {filename} - {e2}")
                        
            # PDF documents
            elif file_type == "application/pdf":
                try:
                    text = ""
                    with pdfplumber.open(filename) as pdf:
                        for page in pdf.pages:
                            extracted = page.extract_text()
                            if extracted:
                                text += extracted + "\n"
                    return text.strip()
                except Exception as e:
                    logger.error(f"Error extracting text from PDF: {filename} - {e}")
                    
            # Plain text
            elif file_type == "text/plain" or file_type.startswith("text/"):
                try:
                    # Try different encodings
                    for encoding in ['utf-8', 'latin1', 'cp1252']:
                        try:
                            with open(filename, "r", encoding=encoding) as f:
                                return f.read()
                        except UnicodeDecodeError:
                            continue
                    # Fallback to binary read with errors ignored
                    with open(filename, "r", encoding='utf-8', errors="ignore") as f:
                        return f.read()
                except Exception as e:
                    logger.error(f"Error reading text file: {filename} - {e}")
                    
            # ZIP archives
            elif file_type == "application/zip":
                try:
                    with zipfile.ZipFile(filename, 'r') as zip_ref:
                        content = []
                        # Only process text files within ZIP
                        for file_info in zip_ref.infolist():
                            if file_info.file_size < 1024*1024:  # Skip files larger than 1MB
                                try:
                                    if not any(file_info.filename.endswith(ext) for ext in ['.jpg', '.png', '.gif', '.exe']):
                                        with zip_ref.open(file_info) as file:
                                            text = file.read().decode('utf-8', errors='ignore')
                                            if len(text.strip()) > 0:
                                                content.append(f"Content of {file_info.filename}:\n{text}")
                                except:
                                    continue
                        return "\n\n".join(content)
                except Exception as e:
                    logger.error(f"Error extracting content from ZIP: {filename} - {e}")
                    
            # Excel files - XLSX
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                try:
                    workbook = load_workbook(filename=filename, read_only=True, data_only=True)
                    content = []
                    for sheet_name in workbook.sheetnames:
                        worksheet = workbook[sheet_name]
                        sheet_content = [f"Sheet name: {sheet_name}"]
                        for row in worksheet.iter_rows(values_only=True):
                            if any(cell is not None for cell in row):
                                sheet_content.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
                        content.append("\n".join(sheet_content))
                    return "\n\n".join(content)
                except Exception as e:
                    logger.error(f"Error extracting content from XLSX: {filename} - {e}")
                    
            # Excel files - XLS
            elif file_type == "application/vnd.ms-excel":
                try:
                    workbook = xlrd.open_workbook(filename)
                    content = []
                    for sheet in workbook.sheets():
                        sheet_content = [f"Sheet name: {sheet.name}"]
                        for row in range(sheet.nrows):
                            if any(cell for cell in sheet.row_values(row)):
                                sheet_content.append("\t".join([str(cell) for cell in sheet.row_values(row)]))
                        content.append("\n".join(sheet_content))
                    return "\n\n".join(content)
                except Exception as e:
                    logger.error(f"Error extracting content from XLS: {filename} - {e}")
                    
            # PowerPoint - PPTX
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                try:
                    prs = Presentation(filename)
                    content = []
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text)
                        if slide_text:
                            content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
                    return "\n\n".join(content)
                except Exception as e:
                    logger.error(f"Error extracting content from PPTX: {filename} - {e}")
                    
            # PowerPoint - PPT
            elif file_type == "application/vnd.ms-powerpoint":
                try:
                    # Try using pypandoc to convert PPT to text
                    return pypandoc.convert_file(filename, 'plain', format='ppt')
                except Exception as e:
                    logger.error(f"Error extracting content from PPT: {filename} - {e}")
                    
            # Word documents - DOC
            elif file_type == "application/msword":
                try:
                    # Try using docx library first
                    try:
                        doc = Document(filename)
                        return "\n".join([para.text for para in doc.paragraphs if para.text])
                    except:
                        # Fallback to pypandoc
                        return pypandoc.convert_file(filename, 'plain', format='doc')
                except Exception as e:
                    logger.error(f"Error extracting content from DOC: {filename} - {e}")
                    
            # RTF files
            elif file_type == "application/rtf":
                try:
                    output = pypandoc.convert_file(filename, 'plain', format='rtf')
                    return output
                except Exception as e:
                    logger.error(f"Error extracting content from RTF: {filename} - {e}")
                    
            # Fallback for unrecognized types
            else:
                # Try reading as text first
                try:
                    with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                        # Check if content looks like text
                        if len(content) > 0 and len(re.findall(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', content)) < len(content) * 0.01:
                            return content
                except:
                    pass
                    
                # Try using markitdown if available
                try:
                    from markitdown import MarkItDown
                    content = MarkItDown().convert(filename)
                    if content and len(content.strip()) > 0:
                        return content
                except:
                    pass
                    
            return "Unsupported file type"
            
        except Exception as e:
            logger.error(f"Error extracting text from file: {filename} - {e}")
            return "Text extraction failed"
            
    def process_content(self, content: str, file_type: str) -> str:
        """
        Process content based on file type
        
        Args:
            content: Raw text content
            file_type: MIME type of the file
            
        Returns:
            Processed text
        """
        if not content or content in ["Unsupported file type", "Text extraction failed"]:
            return ""
            
        # XML processing (for application/zip type)
        if file_type == "application/zip" or file_type.endswith("+xml") or file_type == "text/xml":
            return self._process_xml(content)
            
        # Other types processed as text
        return self._process_text(content)
        
    def _process_xml(self, text: str) -> str:
        """
        Process XML content: remove tags and extract text
        
        Args:
            text: XML content
            
        Returns:
            Cleaned text
        """
        try:
            # Try to parse as XML by wrapping in root element
            try:
                root = ElementTree.fromstring(f"<root>{text}</root>")
                elements = [elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()]
                cleaned = " ".join(elements)
            except ElementTree.ParseError:
                # If parsing fails, use regex cleaning
                cleaned = re.sub(r"<[^>]+>", " ", text)
                
            # Handle HTML escape characters and final cleanup
            return self._final_clean(html.unescape(cleaned))
        except Exception as e:
            logger.error(f"Error processing XML content: {e}")
            # Fallback
            return self._final_clean(text)
        
    def _process_text(self, text: str) -> str:
        """
        Process text content: remove special structures and unnecessary parts
        
        Args:
            text: Text content
            
        Returns:
            Cleaned text
        """
        try:
            # Remove CSS comments
            text = re.sub(r"/\*.*?\*/", "", text, flags=re.DOTALL)
            # Remove URL references
            text = re.sub(r"url\([^)]*\)", "", text, flags=re.IGNORECASE)
            # Remove CSS variables
            text = re.sub(r"var\(--[^)]*\)", "", text)
            # Remove non-printable characters
            text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
            # Final cleanup
            return self._final_clean(text)
        except Exception as e:
            logger.error(f"Error processing text content: {e}")
            # Fallback
            return text.strip()
        
    def _final_clean(self, text: str) -> str:
        """
        Final cleanup: preserve valid Chinese and Latin words
        
        Args:
            text: Text to clean
            
        Returns:
            Cleaned text
        """
        try:
            # Convert to lowercase
            text = text.lower()
            # Match Chinese characters and basic Latin letters
            words = re.findall(r"[\u4e00-\u9fa5a-zA-Z0-9]+", text)
            # Join with spaces
            return " ".join(words).strip()
        except Exception as e:
            logger.error(f"Error in final text cleaning: {e}")
            # Return original if cleaning fails
            return text.strip()
    
    def create_db_and_save_data(self, directory: str, S_or_R: int):
        """
        Create database and save file information
        
        Args:
            directory: Directory to scan
            S_or_R: Sensitivity flag (1 for sensitive, 0 for regular)
        """
        db_name = "enhanced_data.db"
        conn = sqlite3.connect(db_name)
        cursor = conn.cursor()
        
        # Create table with improved structure
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                processed_content TEXT,
                is_sensitive INTEGER NOT NULL,
                processing_time REAL,
                error_message TEXT
            );
        """)
        
        files_data = []
        
        # Get list of files in directory
        for idx, filename in enumerate(os.listdir(directory)):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                logger.info(f"Processing file {idx+1}: {filename}")
                start_time = datetime.now()
                error_message = ""
                
                try:
                    # Detect file type
                    file_type = self.detect_file_type(file_path)
                    # Calculate MD5
                    md5 = self.calculate_md5(file_path)
                    # Get metadata
                    metadata = self.get_file_metadata(file_path)
                    
                    # Skip certain file types
                    ext = Path(filename).suffix.lower()
                    if ext in self.SKIP_EXTENSIONS:
                        logger.info(f"Skipping binary file: {filename}")
                        continue
                        
                    # Extract text content
                    content = self.extract_text_from_file(file_path, file_type)
                    
                    # Process content
                    processed_content = self.process_content(content, file_type)
                    
                    # Skip empty or failed extraction
                    if not content or content in ["Unsupported file type", "Text extraction failed"]:
                        error_message = f"Failed to extract text from {filename}"
                        logger.warning(error_message)
                    
                    # Record processing time
                    end_time = datetime.now()
                    processing_time = (end_time - start_time).total_seconds()
                    
                    # Add to dataset
                    files_data.append((
                        md5, filename, file_type, metadata["size"], 
                        metadata["created_time"], content, processed_content,
                        S_or_R, processing_time, error_message
                    ))
                    
                except Exception as e:
                    error_message = str(e)
                    logger.error(f"Error processing file {filename}: {error_message}")
        
        # Insert data in batches
        if files_data:
            cursor.executemany("""
                INSERT INTO files (
                    md5, filename, file_type, file_size, created_time, 
                    content, processed_content, is_sensitive, processing_time, error_message
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, files_data)
            
        # Remove duplicates based on MD5
        try:
            cursor.execute('''
                DELETE FROM files
                WHERE id NOT IN (
                    SELECT MIN(id)
                    FROM files
                    GROUP BY md5
                )
            ''')
            conn.commit()
            logger.info(f"Removed {cursor.rowcount} duplicate records")
        except sqlite3.Error as e:
            conn.rollback()
            logger.error(f"Failed to remove duplicates: {e}")
            
        # Log summary statistics
        cursor.execute("SELECT COUNT(*) FROM files")
        total_files = cursor.fetchone()[0]
        
        cursor.execute("SELECT COUNT(*) FROM files WHERE is_sensitive=1")
        sensitive_files = cursor.fetchone()[0]
        
        cursor.execute("SELECT COUNT(*) FROM files WHERE error_message != ''")
        error_files = cursor.fetchone()[0]
        
        logger.info(f"Database created: {db_name}")
        logger.info(f"Total files processed: {total_files}")
        logger.info(f"Sensitive files: {sensitive_files}")
        logger.info(f"Files with errors: {error_files}")
        
        conn.commit()
        conn.close()
        logger.info(f"File information saved to database: {db_name}")
        
    def process_database(self, source_db, target_db):
        """
        Process database and improve content formatting
        
        Args:
            source_db: Source database path
            target_db: Target database path
        """
        try:
            # Connect to databases
            src_conn = sqlite3.connect(source_db)
            src_cur = src_conn.cursor()
            dst_conn = sqlite3.connect(target_db)
            dst_cur = dst_conn.cursor()
            
            # Get schema from source database
            src_cur.execute("SELECT sql FROM sqlite_master WHERE type='table' AND name='files'")
            table_sql = src_cur.fetchone()[0]
            
            # Create table in target database
            dst_cur.execute(table_sql)
            
            # Get field indices
            src_cur.execute("PRAGMA table_info(files)")
            columns = [col[1] for col in src_cur.fetchall()]
            content_idx = columns.index("content") if "content" in columns else -1
            type_idx = columns.index("file_type") if "file_type" in columns else -1
            
            # If no content or file_type column, copy as is
            if content_idx == -1 or type_idx == -1:
                src_cur.execute("SELECT * FROM files")
                rows = src_cur.fetchall()
                placeholders = ",".join(["?"] * len(columns))
                dst_cur.executemany(f"INSERT INTO files VALUES ({placeholders})", rows)
                dst_conn.commit()
                logger.info(f"Copied {len(rows)} rows without processing (missing required columns)")
                return
                
            # Process in batches
            batch_size = 1000
            processed_count = 0
            
            src_cur.execute("SELECT COUNT(*) FROM files")
            total_count = src_cur.fetchone()[0]
            
            src_cur.execute("SELECT * FROM files")
            while True:
                rows = src_cur.fetchmany(batch_size)
                if not rows:
                    break
                    
                processed_rows = []
                for row in rows:
                    row_list = list(row)
                    
                    # Process content if available
                    if content_idx >= 0 and type_idx >= 0:
                        content = row_list[content_idx]
                        file_type = row_list[type_idx]
                        
                        # Find processed_content index
                        processed_idx = columns.index("processed_content") if "processed_content" in columns else -1
                        
                        # Process content if column exists
                        if processed_idx >= 0 and content and isinstance(content, str):
                            row_list[processed_idx] = self.process_content(content, file_type)
                            
                    processed_rows.append(tuple(row_list))
                    
                # Insert processed batch
                placeholders = ",".join(["?"] * len(columns))
                dst_cur.executemany(f"INSERT INTO files VALUES ({placeholders})", processed_rows)
                dst_conn.commit()
                
                processed_count += len(processed_rows)
                logger.info(f"Processed {processed_count}/{total_count} rows")
                
            logger.info(f"Database processing completed: {target_db}")
            
        except Exception as e:
            logger.error(f"Database processing failed: {str(e)}")
            dst_conn.rollback()
        finally:
            src_conn.close()
            dst_conn.close()

def main():
    """Main function to run the enhanced file analyzer"""
    # Configure logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers=[
            logging.FileHandler("file_analyzer.log"),
            logging.StreamHandler()
        ]
    )
    logger = logging.getLogger(__name__)
    
    print("Enhanced File Type Detection and Content Extraction")
    print("==================================================")
    
    try:
        path_directory = input("Directory to scan (e.g., D:\\samples\\sensitive_files\\): ")
        S_or_R = int(input("Mark as sensitive (1) or regular (0) files: "))
        
        # Instantiate the file analyzer
        analyzer = FileAnalyzer()
        
        # Create database and save data
        analyzer.create_db_and_save_data(path_directory, S_or_R)
        
        # Ask if user wants to run content processing
        process_content = input("Process content for improved analysis? (y/n): ").lower() == 'y'
        if process_content:
            analyzer.process_database('enhanced_data.db', 'processed_data.db')
            print("Content processing completed. Results saved to processed_data.db")
        
    except Exception as e:
        logger.error(f"Error in main execution: {str(e)}")
        print(f"An error occurred: {str(e)}")
        
if __name__ == "__main__":
    main()
