import magic
import json
import os
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
from typing import List, Dict, Tuple, Optional
from io import BytesIO
import xml.etree.ElementTree as ElementTree
import html
import logging
import binascii
import struct

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("file_extraction.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class FileSignatureDetector:
    """
    增强型文件签名检测器，使用文件头部签名、扩展名和内容特征识别文件类型
    """
    # 常见文件签名（魔术数字）及其对应信息
    FILE_SIGNATURES = {
        # Office文档
        "504B0304": {
            "mime_types": ["application/zip"],
            "possible_formats": ["ZIP", "DOCX", "XLSX", "PPTX", "ODT"],
            "extensions": [".zip", ".docx", ".xlsx", ".pptx", ".odt"],
        },
        "D0CF11E0A1B11AE1": {
            "mime_types": [
                "application/x-ole-storage",
                "application/msword",
                "application/vnd.ms-excel",
                "application/vnd.ms-powerpoint",
            ],
            "possible_formats": ["DOC", "XLS", "PPT", "OLE Compound Document"],
            "extensions": [".doc", ".xls", ".ppt", ".msg"],
        },
        "25504446": {
            "mime_types": ["application/pdf"],
            "possible_formats": ["PDF"],
            "extensions": [".pdf"],
        },
        # 压缩文件
        "526172211A07": {
            "mime_types": ["application/x-rar-compressed"],
            "possible_formats": ["RAR"],
            "extensions": [".rar"],
        },
        "1F8B08": {
            "mime_types": ["application/gzip"],
            "possible_formats": ["GZIP"],
            "extensions": [".gz", ".tgz"],
        },
        "425A68": {
            "mime_types": ["application/x-bzip2"],
            "possible_formats": ["BZIP2"],
            "extensions": [".bz2"],
        },
        "377ABCAF271C": {
            "mime_types": ["application/x-7z-compressed"],
            "possible_formats": ["7ZIP"],
            "extensions": [".7z"],
        },
        # 图像文件
        "FFD8FFE0": {
            "mime_types": ["image/jpeg"],
            "possible_formats": ["JPEG/JFIF"],
            "extensions": [".jpg", ".jpeg"],
        },
        "FFD8FFE1": {
            "mime_types": ["image/jpeg"],
            "possible_formats": ["JPEG/Exif"],
            "extensions": [".jpg", ".jpeg"],
        },
        "89504E47": {
            "mime_types": ["image/png"],
            "possible_formats": ["PNG"],
            "extensions": [".png"],
        },
        "47494638": {
            "mime_types": ["image/gif"],
            "possible_formats": ["GIF"],
            "extensions": [".gif"],
        },
        "424D": {
            "mime_types": ["image/bmp"],
            "possible_formats": ["BMP"],
            "extensions": [".bmp"],
        },
        # 文本文件 - 通常没有明确的签名，但可以有BOM
        "EFBBBF": {
            "mime_types": ["text/plain"],
            "possible_formats": ["UTF-8 Text"],
            "extensions": [".txt"],
        },
        "FEFF": {
            "mime_types": ["text/plain"],
            "possible_formats": ["UTF-16 Text"],
            "extensions": [".txt"],
        },
        # 其他常见格式
        "4D5A": {
            "mime_types": ["application/x-msdownload"],
            "possible_formats": ["EXE", "DLL"],
            "extensions": [".exe", ".dll"],
        },
        "CAFEBABE": {
            "mime_types": ["application/java-archive"],
            "possible_formats": ["Java Class"],
            "extensions": [".class"],
        },
        "7573746172": {
            "mime_types": ["application/x-tar"],
            "possible_formats": ["TAR"],
            "extensions": [".tar"],
        }
    }

    # Office专用的格式检测标记
    OFFICE_MARKERS = {
        # DOCX标记(ZIP内部)
        "word/document.xml": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        # XLSX标记(ZIP内部)
        "xl/workbook.xml": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        # PPTX标记(ZIP内部)
        "ppt/presentation.xml": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        # DOC标记(OLE内部)
        "WordDocument": "application/msword",
        # XLS标记(OLE内部)
        "Workbook": "application/vnd.ms-excel",
        # PPT标记(OLE内部)
        "PowerPoint Document": "application/vnd.ms-powerpoint",
    }

    @staticmethod
    def read_file_header(filename: str, num_bytes: int = 16) -> str:
        """读取文件头部并返回十六进制表示"""
        try:
            with open(filename, 'rb') as f:
                header = f.read(num_bytes)
            return binascii.hexlify(header).decode('utf-8').upper()
        except Exception as e:
            logger.error(f"读取文件头部失败 {filename}: {e}")
            return ""

    @staticmethod
    def detect_mime_type(filename: str) -> str:
        """通过文件签名和内部结构检测MIME类型"""
        # 1. 先通过文件头部签名检测
        header = FileSignatureDetector.read_file_header(filename)
        if not header:
            return "unknown"

        # 检查是否匹配任何已知签名
        for signature, info in FileSignatureDetector.FILE_SIGNATURES.items():
            if header.startswith(signature):
                # 进一步检测ZIP文件内容来识别Office文档类型
                if signature == "504B0304":  # ZIP格式
                    office_type = FileSignatureDetector.check_office_zip(filename)
                    if office_type:
                        return office_type
                # 进一步检测OLE文件内容来识别Office文档类型
                elif signature == "D0CF11E0A1B11AE1":  # OLE格式
                    office_type = FileSignatureDetector.check_office_ole(filename)
                    if office_type:
                        return office_type
                
                # 返回第一个匹配的MIME类型
                return info["mime_types"][0]
        
        # 2. 通过扩展名辅助判断
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".docx":
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif ext == ".xlsx":
            return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif ext == ".pptx":
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif ext == ".doc":
            return "application/msword"
        elif ext == ".xls":
            return "application/vnd.ms-excel"
        elif ext == ".ppt":
            return "application/vnd.ms-powerpoint"
        elif ext == ".pdf":
            return "application/pdf"
        elif ext == ".txt":
            return "text/plain"
        
        # 3. 最后尝试使用magic库
        try:
            mime = magic.Magic(mime=True)
            return mime.from_file(filename)
        except Exception as e:
            logger.error(f"使用magic库检测文件类型失败 {filename}: {e}")
            return "unknown"

    @staticmethod
    def check_office_zip(filename: str) -> Optional[str]:
        """检查ZIP文件内是否包含Office文档的特征文件"""
        try:
            with zipfile.ZipFile(filename, 'r') as zip_file:
                file_list = zip_file.namelist()
                
                for marker, mime_type in FileSignatureDetector.OFFICE_MARKERS.items():
                    if marker in file_list:
                        return mime_type
                
                # 检查更复杂的路径匹配
                for file_path in file_list:
                    if "word/document.xml" in file_path:
                        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    elif "xl/workbook.xml" in file_path:
                        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    elif "ppt/presentation.xml" in file_path:
                        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            
            # 没有找到Office标记，可能是普通ZIP文件
            return "application/zip"
        except Exception as e:
            logger.warning(f"检查ZIP文件内容失败 {filename}: {e}")
            return None

    @staticmethod
    def check_office_ole(filename: str) -> Optional[str]:
        """检查OLE文件内是否包含Office文档的特征"""
        try:
            if not olefile.isOleFile(filename):
                return None
                
            ole = olefile.OleFile(filename)
            streams = ole.listdir()
            
            # 扁平化流列表
            stream_paths = []
            for stream in streams:
                if isinstance(stream, list):
                    stream_paths.append('/'.join(stream))
                else:
                    stream_paths.append(stream)
            
            # 检查Word文档标记
            if "WordDocument" in stream_paths:
                return "application/msword"
            
            # 检查Excel工作簿标记
            if "Workbook" in stream_paths:
                return "application/vnd.ms-excel"
            
            # 检查PowerPoint演示文稿标记
            if "PowerPoint Document" in stream_paths:
                return "application/vnd.ms-powerpoint"
            
            # 默认返回OLE存储格式
            return "application/x-ole-storage"
        except Exception as e:
            logger.warning(f"检查OLE文件内容失败 {filename}: {e}")
            return None

class FileProcessor:
    @staticmethod
    def calculate_md5(filename: str) -> str:
        """计算文件的MD5哈希值"""
        hash_md5 = hashlib.md5()
        try:
            with open(filename, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logger.error(f"计算MD5失败 {filename}: {e}")
            return ""

    @staticmethod
    def get_metadata(filename: str) -> Dict:
        """获取文件元数据（大小和创建时间）"""
        try:
            file_stat = os.stat(filename)
            return {
                "size": file_stat.st_size,
                "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")
            }
        except Exception as e:
            logger.error(f"获取元数据失败 {filename}: {e}")
            return {"size": 0, "created_time": ""}

class TextExtractor:
    # 不支持提取的文件类型列表
    UNSUPPORTED_TYPES = [
        "application/x-msdownload",
        "application/x-executable",
        "application/x-dosexec",
        "application/x-sharedlib",
        "image/jpeg",
        "image/png",
        "image/gif",
        "image/bmp",
        "image/tiff",
        "audio/",
        "video/",
        "application/octet-stream"
    ]
    
    @staticmethod
    def is_supported_type(file_type: str) -> bool:
        """检查文件类型是否支持文本提取"""
        return not any(file_type.startswith(unsupported) for unsupported in TextExtractor.UNSUPPORTED_TYPES)

    @staticmethod
    def extract_from_docx(filename: str) -> str:
        """从DOCX文件提取文本"""
        try:
            doc = Document(filename)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"从DOCX提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_pdf(filename: str) -> str:
        """从PDF文件提取文本"""
        try:
            text = ""
            with pdfplumber.open(filename) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"从PDF提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_txt(filename: str) -> str:
        """从纯文本文件提取文本"""
        try:
            with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"从TXT提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_zip(filename: str) -> str:
        """从ZIP文件提取文本，排除二进制内容"""
        try:
            with zipfile.ZipFile(filename, 'r') as zip_ref:
                content = ""
                # 只处理最多10个文本文件，避免处理过多内容
                text_files_processed = 0
                for file_info in zip_ref.infolist():
                    # 跳过目录和二进制文件
                    if file_info.filename.endswith('/') or not TextExtractor._is_likely_text_file(file_info.filename):
                        continue
                        
                    try:
                        with zip_ref.open(file_info) as file:
                            file_content = file.read(8192)  # 只读取前8KB判断是否为文本
                            if TextExtractor._is_binary_content(file_content):
                                continue
                                
                            # 重新打开并读取完整内容
                            with zip_ref.open(file_info) as full_file:
                                try:
                                    text = full_file.read().decode('utf-8', errors='ignore')
                                    if text.strip():
                                        content += f"===== {file_info.filename} =====\n"
                                        content += text[:4096] + "\n"  # 限制每个文件最多4KB
                                        text_files_processed += 1
                                        if text_files_processed >= 10:
                                            content += "...[更多文件已省略]...\n"
                                            break
                                except:
                                    continue
                    except Exception as inner_e:
                        continue
                return content
        except Exception as e:
            logger.error(f"从ZIP提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def _is_likely_text_file(filename: str) -> bool:
        """判断文件名是否可能为文本文件"""
        text_extensions = ['.txt', '.xml', '.json', '.csv', '.md', '.html', '.htm', '.css', '.js', '.java', '.py', '.c', '.cpp', '.h', '.php', '.log', '.ini', '.cfg', '.conf']
        return any(filename.lower().endswith(ext) for ext in text_extensions)

    @staticmethod
    def _is_binary_content(content: bytes) -> bool:
        """判断内容是否为二进制"""
        # 检查是否包含空字节，通常文本文件不包含
        if b'\x00' in content:
            return True
            
        # 检查非ASCII字符比例
        non_ascii = sum(1 for b in content if b > 127)
        if len(content) > 0 and non_ascii / len(content) > 0.3:  # 如果非ASCII字符超过30%，可能是二进制
            return True
            
        return False

    @staticmethod
    def extract_from_xlsx(filename: str) -> str:
        """从XLSX文件提取文本"""
        try:
            workbook = load_workbook(filename=filename, read_only=True, data_only=True)
            content = ""
            for sheet_name in workbook.sheetnames:
                content += f"===== 工作表: {sheet_name} =====\n"
                worksheet = workbook[sheet_name]
                # 限制处理的行数和列数
                max_rows = min(worksheet.max_row, 500)
                max_cols = min(worksheet.max_column, 20)
                for row in range(1, max_rows + 1):
                    row_values = []
                    for col in range(1, max_cols + 1):
                        cell = worksheet.cell(row=row, column=col)
                        row_values.append(str(cell.value) if cell.value is not None else "")
                    if any(v.strip() for v in row_values):  # 只添加非空行
                        content += "\t".join(row_values) + "\n"
            return content
        except Exception as e:
            logger.error(f"从XLSX提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_xls(filename: str) -> str:
        """从XLS文件提取文本"""
        try:
            workbook = xlrd.open_workbook(filename, logfile=open(os.devnull, 'w'), on_demand=True)
            content = ""
            for sheet in workbook.sheets():
                content += f"===== 工作表: {sheet.name} =====\n"
                # 限制处理的行数
                max_rows = min(sheet.nrows, 500)
                for row in range(max_rows):
                    try:
                        row_values = [str(cell.value) if cell.ctype != xlrd.XL_CELL_EMPTY else "" for cell in sheet.row(row)]
                        if any(v.strip() for v in row_values):  # 只添加非空行
                            content += "\t".join(row_values) + "\n"
                    except Exception:
                        continue
            return content
        except Exception as e:
            logger.error(f"从XLS提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_pptx(filename: str) -> str:
        """从PPTX文件提取文本"""
        try:
            prs = Presentation(filename)
            content = ""
            for i, slide in enumerate(prs.slides):
                content += f"===== 幻灯片 {i+1} =====\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"
            return content
        except Exception as e:
            logger.error(f"从PPTX提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_rtf(filename: str) -> str:
        """从RTF文件提取文本"""
        try:
            output = pypandoc.convert_file(filename, 'plain', format='rtf')
            return output
        except Exception as e:
            logger.error(f"从RTF提取文本失败 {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_file(filename: str, file_type: str) -> str:
        """根据文件类型提取文本"""
        # 检查是否为不支持的文件类型
        if not TextExtractor.is_supported_type(file_type):
            logger.info(f"不支持的文件类型: {file_type} - {filename}")
            return ""
            
        # 根据MIME类型提取文本
        if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return TextExtractor.extract_from_docx(filename)
        elif file_type == "application/pdf":
            return TextExtractor.extract_from_pdf(filename)
        elif file_type == "text/plain":
            return TextExtractor.extract_from_txt(filename)
        elif file_type == "application/zip":
            return TextExtractor.extract_from_zip(filename)
        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return TextExtractor.extract_from_xlsx(filename)
        elif file_type == "application/vnd.ms-excel":
            return TextExtractor.extract_from_xls(filename)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return TextExtractor.extract_from_pptx(filename)
        elif file_type == "application/rtf":
            return TextExtractor.extract_from_rtf(filename)
        elif file_type == "application/msword":
            # 尝试使用Document，可能适用于某些DOC文件
            try:
                return TextExtractor.extract_from_docx(filename)
            except:
                logger.warning(f"无法使用docx库从DOC文件提取文本 {filename}")
                return ""
        
        # 对于无法识别的类型，尝试基于扩展名
        ext = os.path.splitext(filename)[1].lower()
        if ext in ['.txt', '.csv', '.log', '.md', '.xml', '.json', '.html', '.htm', '.css', '.js', '.java', '.py', '.c', '.cpp', '.h', '.php']:
            return TextExtractor.extract_from_txt(filename)
            
        logger.warning(f"未知文件类型: {file_type} (扩展名: {ext}) - {filename}")
        return ""

class ContentFormatter:
    @staticmethod
    def process_content(content, file_type):
        """根据文件类型处理内容"""
        if not content:
            return ""

        # 对于ZIP和XML内容特殊处理
        if file_type == "application/zip" or file_type.endswith("+xml"):
            return ContentFormatter._process_xml(content)
        
        # 其他内容处理
        return ContentFormatter._process_general_text(content)

    @staticmethod
    def _process_xml(text):
        """处理XML内容：移除标签并提取文本"""
        # 尝试清理XML标签
        try:
            cleaned = re.sub(r"<[^>]+>", " ", text)
            cleaned = html.unescape(cleaned)
        except:
            cleaned = text
            
        return ContentFormatter._final_clean(cleaned)

    @staticmethod
    def _process_general_text(text):
        """处理一般文本内容"""
        # 转换为小写并移除多余空白
        text = text.lower()
        text = re.sub(r'\s+', ' ', text)
        return ContentFormatter._final_clean(text)

    @staticmethod
    def _final_clean(text):
        """最终清理：保留有效的中英文词汇"""
        # 匹配中文汉字和英文字母
        words = re.findall(r"[\u4e00-\u9fa5a-zA-Z]+", text)
        result = " ".join(words).strip()
        
        # 限制长度，避免内容过大
        if len(result) > 100000:  # 限制为约10万字符
            result = result[:100000] + "..."
            
        return result

class DatabaseManager:
    def __init__(self, db_name="data.db"):
        self.db_name = db_name
        self.conn = None
        self.cursor = None
        
    def connect(self):
        """连接数据库"""
        self.conn = sqlite3.connect(self.db_name)
        self.cursor = self.conn.cursor()
        return self
        
    def close(self):
        """关闭数据库连接"""
        if self.conn:
            self.conn.close()
            
    def create_tables(self):
        """创建数据库表"""
        self.cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL
            );
        """)
        self.conn.commit()
        
    def save_file_data(self, files_data):
        """保存文件数据到数据库"""
        if not files_data:
            return
            
        self.cursor.executemany("""
            INSERT INTO files (md5, filename, file_type, file_size, created_time, content, is_sensitive)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, files_data)
        self.conn.commit()
        
    def remove_duplicates(self):
        """根据MD5移除重复记录"""
        try:
            self.cursor.execute('''
                DELETE FROM files
                WHERE id NOT IN (
                    SELECT MIN(id)
                    FROM files
                    GROUP BY md5
                )
            ''')
            self.conn.commit()
            logger.info(f"已移除 {self.cursor.rowcount} 条重复记录")
        except sqlite3.Error as e:
            self.conn.rollback()
            logger.error(f"移除重复记录失败: {e}")
            
    def process_and_migrate(self, target_db):
        """处理并迁移数据到新数据库"""
        formatter = ContentFormatter()
        target_db_mgr = DatabaseManager(target_db)
        target_db_mgr.connect()
        target_db_mgr.create_tables()
        
        # 获取所有记录
        self.cursor.execute("SELECT * FROM files")
        batch_size = 100
        processed_count = 0
        
        while True:
            rows = self.cursor.fetchmany(batch_size)
            if not rows:
                break
                
            processed_data = []
            for row in rows:
                row_id, md5, filename, file_type, file_size, created_time, content, is_sensitive = row
                
                # 处理内容
                processed_content = formatter.process_content(content, file_type) if content else ""
                
                processed_data.append((md5, filename, file_type, file_size, created_time, processed_content, is_sensitive))
                
            # 保存处理后的数据
            target_db_mgr.save_file_data(processed_data)
            processed_count += len(processed_data)
            logger.info(f"已处理 {processed_count} 条记录")
            
        target_db_mgr.close()
        logger.info(f"数据库迁移完成: {target_db}")

class FileSystem:
    @staticmethod
    def process_directory(directory, is_sensitive):
        """处理目录中的所有文件"""
        if not os.path.exists(directory):
            logger.error(f"目录不存在: {directory}")
            return
            
        # 创建数据库管理器
        db_mgr = DatabaseManager()
        db_mgr.connect()
        db_mgr.create_tables()
        
        files_data = []
        file_count = 0
        skipped_count = 0
        
        # 处理每个文件
        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if not os.path.isfile(file_path):
                continue
                
            file_count += 1
            logger.info(f"正在处理文件 {file_count}: {filename}")
            
            # 获取文件信息
            file_type = FileSignatureDetector.detect_mime_type(file_path)
            
            # 跳过不支持的文件类型
            if not TextExtractor.is_supported_type(file_type):
                logger.info(f"跳过不支持的文件类型: {file_type} - {filename}")
                skipped_count += 1
                continue
                
            md5 = FileProcessor.calculate_md5(file_path)
            metadata = FileProcessor.get_metadata(file_path)
            
            if not md5:
                logger.warning(f"跳过文件 {filename}: 计算MD5失败")
                skipped_count += 1
                continue
                
            # 提取文本内容
            content = TextExtractor.extract_from_file(file_path, file_type)
            
            # 跳过没有可提取内容的文件
            if not content:
                logger.info(f"跳过文件 {filename}: 无可提取内容")
                skipped_count += 1
                continue
                
            files_data.append((
                md5, 
                filename, 
                file_type, 
                metadata["size"], 
                metadata["created_time"], 
                content, 
                1 if is_sensitive else 0
            ))
            
            # 批量保存以避免内存问题
            if len(files_data) >= 50:
                db_mgr.save_file_data(files_data)
                logger.info(f"已保存50个文件到数据库")
                files_data = []
        
        # 保存剩余文件
        if files_data:
            db_mgr.save_file_data(files_data)
            
        # 移除重复项
        db_mgr.remove_duplicates()
        
        # 关闭连接
        db_mgr.close()
        logger.info(f"处理完成 {file_count} 个文件, 跳过 {skipped_count} 个文件. 数据已保存到数据库: {db_mgr.db_name}")

def main():
    """主函数"""
    print("文件处理与数据库存储系统")
    print("====================")
    
    # 获取目录和敏感度标记
    path_directory = input("要扫描的文件目录路径 (例如: D:\\样本\\敏感文件\\): ")
    try:
        sensitivity_input = input("标记为敏感文件 (1) 或普通文件 (0): ")
        is_sensitive = int(sensitivity_input)
        if is_sensitive not in [0, 1]:
            raise ValueError("输入必须为0或1")
    except ValueError as e:
        print(f"输入无效. 使用默认值0 (普通文件): {e}")
        is_sensitive = 0
    
    # 处理目录
    FileSystem.process_directory(path_directory, is_sensitive)
    
    # 格式化数据库内容
    perform_formatting = input("是否格式化数据库内容? (y/n): ").lower() == 'y'
    if perform_formatting:
        db_mgr = DatabaseManager()
        db_mgr.connect()
        db_mgr.process_and_migrate('new_data.db')
        db_mgr.close()
        print("数据库内容已格式化并保存到 new_data.db")

if __name__ == "__main__":
    main()
