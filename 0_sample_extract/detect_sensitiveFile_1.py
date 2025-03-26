import magic
import json
import os
import hashlib
import sqlite3
import openpyxl
import pdfplumber
from datetime import datetime
import zipfile
import olefile
from pptx import Presentation
from openpyxl import load_workbook
import xlrd
import pypandoc
from docx import Document
import re
from typing import List, Dict, Tuple
from io import BytesIO
import xml.etree.ElementTree as ElementTree
import html
import logging

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("file_extraction.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

class FileTypeDetector:
    @staticmethod
    def detect_type(filename: str) -> str:
        """Detect file type using magic library"""
        try:
            mime = magic.Magic(mime=True)
            file_type = mime.from_file(filename)
            return file_type
        except Exception as e:
            logger.error(f"Failed to detect file type for {filename}: {e}")
            return "unknown"
            
    @staticmethod
    def get_extension(filename: str) -> str:
        """Get file extension"""
        return os.path.splitext(filename)[1].lower()

class FileProcessor:
    @staticmethod
    def calculate_md5(filename: str) -> str:
        """Calculate MD5 hash of a file"""
        hash_md5 = hashlib.md5()
        try:
            with open(filename, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception as e:
            logger.error(f"Failed to calculate MD5 for {filename}: {e}")
            return ""

    @staticmethod
    def get_metadata(filename: str) -> Dict:
        """Get file metadata (size and creation time)"""
        try:
            file_stat = os.stat(filename)
            return {
                "size": file_stat.st_size,
                "created_time": datetime.fromtimestamp(file_stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")
            }
        except Exception as e:
            logger.error(f"Failed to get metadata for {filename}: {e}")
            return {"size": 0, "created_time": ""}

class TextExtractor:
    @staticmethod
    def extract_from_docx(filename: str) -> str:
        """Extract text from DOCX file"""
        try:
            doc = Document(filename)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_pdf(filename: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with pdfplumber.open(filename) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Failed to extract text from PDF {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_txt(filename: str) -> str:
        """Extract text from plain text file"""
        try:
            with open(filename, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to extract text from TXT {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_zip(filename: str) -> str:
        """Extract text from ZIP file"""
        try:
            with zipfile.ZipFile(filename, 'r') as zip_ref:
                content = ""
                for file_info in zip_ref.infolist():
                    try:
                        with zip_ref.open(file_info) as file:
                            content += f"Content of {file_info.filename}:\n"
                            try:
                                content += file.read().decode('utf-8', errors='ignore') + "\n"
                            except:
                                content += "[Binary content]\n"
                    except Exception as inner_e:
                        logger.warning(f"Could not read {file_info.filename} in ZIP: {inner_e}")
                return content
        except Exception as e:
            logger.error(f"Failed to extract text from ZIP {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_xlsx(filename: str) -> str:
        """Extract text from XLSX file"""
        try:
            workbook = load_workbook(filename=filename, read_only=True, data_only=True)
            content = ""
            for sheet_name in workbook.sheetnames:
                content += f"Sheet name: {sheet_name}\n"
                worksheet = workbook[sheet_name]
                for row in worksheet.iter_rows(values_only=True):
                    content += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            return content
        except Exception as e:
            logger.error(f"Failed to extract text from XLSX {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_xls(filename: str) -> str:
        """Extract text from XLS file"""
        try:
            workbook = xlrd.open_workbook(filename, logfile=open(os.devnull, 'w'), on_demand=True)
            content = ""
            for sheet in workbook.sheets():
                content += f"Sheet name: {sheet.name}\n"
                for row in range(min(sheet.nrows, 1000)):  # Limit to 1000 rows for performance
                    try:
                        content += "\t".join([str(cell) for cell in sheet.row_values(row)]) + "\n"
                    except Exception as row_e:
                        logger.warning(f"Error processing row {row} in sheet {sheet.name}: {row_e}")
            return content
        except Exception as e:
            logger.error(f"Failed to extract text from XLS {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_pptx(filename: str) -> str:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(filename)
            content = ""
            for i, slide in enumerate(prs.slides):
                content += f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"
            return content
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_ppt(filename: str) -> str:
        """Extract text from PPT file"""
        try:
            # Try using olefile first
            if olefile.isOleFile(filename):
                ole = olefile.OleFileIO(filename)
                content = "PowerPoint content (limited extraction)"
                return content
            return ""
        except Exception as e:
            logger.error(f"Failed to extract text from PPT {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_rtf(filename: str) -> str:
        """Extract text from RTF file"""
        try:
            output = pypandoc.convert_file(filename, 'plain', format='rtf')
            return output
        except Exception as e:
            logger.error(f"Failed to extract text from RTF {filename}: {e}")
            return ""

    @staticmethod
    def extract_from_file(filename: str, file_type: str) -> str:
        """Extract text from file based on file type"""
        # Get file extension as a backup
        ext = os.path.splitext(filename)[1].lower()
        
        # First try based on MIME type
        if file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return TextExtractor.extract_from_docx(filename)
        elif file_type == "application/pdf":
            return TextExtractor.extract_from_pdf(filename)
        elif file_type == "text/plain" or ext in ['.txt', '.csv', '.log', '.md']:
            return TextExtractor.extract_from_txt(filename)
        elif file_type == "application/zip" or ext == '.zip':
            return TextExtractor.extract_from_zip(filename)
        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or ext == '.xlsx':
            return TextExtractor.extract_from_xlsx(filename)
        elif file_type == "application/vnd.ms-excel" or ext == '.xls':
            return TextExtractor.extract_from_xls(filename)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or ext == '.pptx':
            return TextExtractor.extract_from_pptx(filename)
        elif file_type == "application/vnd.ms-powerpoint" or ext == '.ppt':
            return TextExtractor.extract_from_ppt(filename)
        elif file_type == "application/rtf" or ext == '.rtf':
            return TextExtractor.extract_from_rtf(filename)
        elif file_type == "application/msword" or ext == '.doc':
            # Try to use Document for .doc files, may work for some
            try:
                return TextExtractor.extract_from_docx(filename)
            except:
                logger.warning(f"Could not extract text from DOC file {filename} using docx. Fallback methods failed.")
                return "Text extraction not supported for this file type"
        
        # Try extraction based on extension if MIME type didn't work
        if ext in ['.docx', '.doc']:
            try:
                return TextExtractor.extract_from_docx(filename)
            except:
                pass
        elif ext == '.pdf':
            return TextExtractor.extract_from_pdf(filename)
        elif ext in ['.xlsx', '.xls']:
            try:
                return TextExtractor.extract_from_xlsx(filename)
            except:
                try:
                    return TextExtractor.extract_from_xls(filename)
                except:
                    pass
        
        # If we get here, we couldn't extract text
        logger.warning(f"Unsupported file type for text extraction: {file_type} (ext: {ext})")
        return "Text extraction not supported for this file type"

class ContentFormatter:
    @staticmethod
    def process_content(content, file_type):
        """Process content based on file type"""
        if not content or content == "Text extraction not supported for this file type":
            return ""

        # XML handling (for ZIP files content)
        if file_type == "application/zip":
            return ContentFormatter._process_xml(content)
        
        # Other content processing
        return ContentFormatter._process_general_text(content)

    @staticmethod
    def _process_xml(text):
        """Process XML content: remove tags and extract text"""
        try:
            # Try parsing as XML
            root = ElementTree.fromstring(text)
            elements = [elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()]
            cleaned = " ".join(elements)
        except ElementTree.ParseError:
            # If parsing fails, use regex to clean
            cleaned = re.sub(r"<[^>]+>", "", text)

        # Process HTML escape characters and final cleaning
        return ContentFormatter._final_clean(html.unescape(cleaned))

    @staticmethod
    def _process_general_text(text):
        """Process general text content"""
        # Convert to lowercase
        text = text.lower()
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        return ContentFormatter._final_clean(text)

    @staticmethod
    def _final_clean(text):
        """Final cleaning: preserve valid Chinese and English words"""
        # Match Chinese characters and Latin letters
        words = re.findall(r"[\u4e00-\u9fa5a-zA-Z]+", text)
        return " ".join(words).strip()

class DatabaseManager:
    def __init__(self, db_name="data.db"):
        self.db_name = db_name
        self.conn = None
        self.cursor = None
        
    def connect(self):
        """Connect to the database"""
        self.conn = sqlite3.connect(self.db_name)
        self.cursor = self.conn.cursor()
        return self
        
    def close(self):
        """Close the database connection"""
        if self.conn:
            self.conn.close()
            
    def create_tables(self):
        """Create database tables"""
        self.cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                md5 TEXT NOT NULL,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                file_size INTEGER NOT NULL,
                created_time TEXT NOT NULL,
                content TEXT,
                is_sensitive INTEGER NOT NULL
            );
        """)
        self.conn.commit()
        
    def save_file_data(self, files_data):
        """Save file data to database"""
        if not files_data:
            return
            
        self.cursor.executemany("""
            INSERT INTO files (md5, filename, file_type, file_size, created_time, content, is_sensitive)
            VALUES (?, ?, ?, ?, ?, ?, ?)
        """, files_data)
        self.conn.commit()
        
    def remove_duplicates(self):
        """Remove duplicate entries based on MD5"""
        try:
            self.cursor.execute('''
                DELETE FROM files
                WHERE id NOT IN (
                    SELECT MIN(id)
                    FROM files
                    GROUP BY md5
                )
            ''')
            self.conn.commit()
            logger.info(f"Removed {self.cursor.rowcount} duplicate records")
        except sqlite3.Error as e:
            self.conn.rollback()
            logger.error(f"Failed to remove duplicates: {e}")
            
    def process_and_migrate(self, target_db):
        """Process and migrate data to a new database"""
        formatter = ContentFormatter()
        target_db_mgr = DatabaseManager(target_db)
        target_db_mgr.connect()
        target_db_mgr.create_tables()
        
        # Get all records
        self.cursor.execute("SELECT * FROM files")
        batch_size = 100
        processed_count = 0
        
        while True:
            rows = self.cursor.fetchmany(batch_size)
            if not rows:
                break
                
            processed_data = []
            for row in rows:
                row_id, md5, filename, file_type, file_size, created_time, content, is_sensitive = row
                
                # Process content
                processed_content = formatter.process_content(content, file_type) if content else ""
                
                processed_data.append((md5, filename, file_type, file_size, created_time, processed_content, is_sensitive))
                
            # Save processed data
            target_db_mgr.save_file_data(processed_data)
            processed_count += len(processed_data)
            logger.info(f"Processed {processed_count} records so far")
            
        target_db_mgr.close()
        logger.info(f"Database migration complete: {target_db}")

class FileSystem:
    @staticmethod
    def process_directory(directory, is_sensitive):
        """Process all files in a directory"""
        if not os.path.exists(directory):
            logger.error(f"Directory does not exist: {directory}")
            return
            
        # Create database manager
        db_mgr = DatabaseManager()
        db_mgr.connect()
        db_mgr.create_tables()
        
        files_data = []
        file_count = 0
        
        # Process each file
        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if not os.path.isfile(file_path):
                continue
                
            file_count += 1
            logger.info(f"Processing file {file_count}: {filename}")
            
            # Get file information
            file_type = FileTypeDetector.detect_type(file_path)
            md5 = FileProcessor.calculate_md5(file_path)
            metadata = FileProcessor.get_metadata(file_path)
            
            if not md5:
                logger.warning(f"Skipping file {filename}: Failed to calculate MD5")
                continue
                
            # Extract text content
            content = TextExtractor.extract_from_file(file_path, file_type)
            
            # Skip files with no extractable content
            if content and content != "Text extraction not supported for this file type":
                files_data.append((
                    md5, 
                    filename, 
                    file_type, 
                    metadata["size"], 
                    metadata["created_time"], 
                    content, 
                    1 if is_sensitive else 0
                ))
                
            # Save in batches to avoid memory issues
            if len(files_data) >= 50:
                db_mgr.save_file_data(files_data)
                files_data = []
        
        # Save any remaining files
        if files_data:
            db_mgr.save_file_data(files_data)
            
        # Remove duplicates
        db_mgr.remove_duplicates()
        
        # Close connection
        db_mgr.close()
        logger.info(f"Processed {file_count} files. Data saved to database: {db_mgr.db_name}")

def main():
    """Main entry point"""
    print("File Processing and Database Storage")
    print("====================================")
    
    # Get directory and sensitivity flag
    path_directory = input("File directory path to scan (e.g. D:\\samples\\sensitive\\): ")
    try:
        sensitivity_input = input("Mark as sensitive files (1) or regular files (0): ")
        is_sensitive = int(sensitivity_input)
        if is_sensitive not in [0, 1]:
            raise ValueError("Input must be 0 or 1")
    except ValueError as e:
        print(f"Invalid input. Using default value 0 (regular files): {e}")
        is_sensitive = 0
    
    # Process directory
    FileSystem.process_directory(path_directory, is_sensitive)
    
    # Format database content
    perform_formatting = input("Format database content? (y/n): ").lower() == 'y'
    if perform_formatting:
        db_mgr = DatabaseManager()
        db_mgr.connect()
        db_mgr.process_and_migrate('new_data.db')
        db_mgr.close()
        print("Database content formatted and saved to new_data.db")

if __name__ == "__main__":
    main()
