import os
import yaml
import magic
import time
import xlrd
import logging
import threading
import uuid
import shutil
from pathlib import Path
from typing import Optional, List, Dict, Tuple, Any, Iterator
from concurrent.futures import ThreadPoolExecutor, as_completed, TimeoutError
from markitdown import MarkItDown
from dataclasses import dataclass, field
import re
from pdfminer.high_level import extract_text as pdf_extract_text
import pandas as pd
import openpyxl
import csv
import json
from datetime import datetime
import argparse
import sys
import win32com.client
import pythoncom
import psutil
import patoolib
from pptx import Presentation

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('sensitive_detector.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class ProcessingResult:
    """数据类，用于存储文件处理结果"""
    file_path: str
    mime_type: str
    content: Dict[str, Any]
    sensitive_words: List[Tuple[str, List[int]]] = field(default_factory=list)
    error: Optional[str] = None
    processing_time: float = 0.0

class FileTypeDetector:
    """文件类型检测器，使用 python-magic 库，增强OLE文件识别"""
    
    MIME_TYPES = {
        ".txt": "text/plain", ".csv": "text/csv", ".xml": "text/xml", ".html": "text/html",
        ".htm": "text/html", ".json": "application/json", ".yaml": "application/yaml",
        ".yml": "application/yaml", ".md": "text/markdown", ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".zip": "application/zip", ".rar": "application/x-rar", ".7z": "application/x-7z-compressed",
        ".tar": "application/x-tar", ".gz": "application/gzip", ".bz2": "application/x-bzip2",
        ".pdf": "application/pdf", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
        ".gif": "image/gif", ".bmp": "image/bmp", ".svg": "image/svg+xml", ".webp": "image/webp",
        ".mp3": "audio/mpeg", ".mp4": "video/mp4", ".avi": "video/x-msvideo", ".mov": "video/quicktime",
        ".wav": "audio/wav", ".bin": "application/octet-stream", ".exe": "application/x-msdownload",
        ".dll": "application/x-msdownload",
        # 增加一些容易混淆的文件类型
        ".rtf": "application/rtf", ".dot": "application/msword", 
        ".pps": "application/vnd.ms-powerpoint", ".xlt": "application/vnd.ms-excel",
    }
    MIME_TO_EXT = {mime: ext for ext, mime in MIME_TYPES.items()}
    SKIP_EXTENSIONS = {'.dwg', '.mp3', '.wav', '.mp4', '.avi', '.mkv', '.flv', '.mov'}

    def __init__(self):
        """初始化文件类型检测器"""
        self.mime = magic.Magic(mime=True)
    
    def get_all_files(self, directory: str) -> List[str]:
        """递归获取目录中的所有文件，忽略 ~ 开头或内部流文件及跳过文件"""
        files = []
        for f in Path(directory).rglob('*'):
            if f.is_file() and not f.name.startswith('~$') and not self._is_internal_stream(f.name):
                ext = f.suffix.lower()
                if ext not in self.SKIP_EXTENSIONS:
                    files.append(str(f))
        return files
    
    def _is_internal_stream(self, filename: str) -> bool:
        """判断文件是否为 Office 内部流文件（如 [5]SummaryInformation）"""
        return bool(re.match(r'^\[\d+\].+', filename))
    
    def get_file_info(self, file_path: str) -> Dict:
        """获取文件信息，包括 MIME 类型和文件头"""
        try:
            mime_type = self.detect_file_type(file_path)
            with open(file_path, "rb") as f:
                file_header = f.read(16).hex().upper()
            file_extension = Path(file_path).suffix.lower()
            return {
                "file_path": file_path,
                "mime_type": mime_type,
                "file_extension": file_extension,
                "file_header": file_header,
                "size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"获取文件信息失败: {e} - {file_path}")
            return {"file_path": file_path, "mime_type": "unknown", "error": str(e)}
    
    def _detect_ole_file_type(self, file_path: str) -> str:
        """更精确地检测OLE文件类型，区分DOC/XLS/PPT"""
        try:
            # 读取文件头判断是否为OLE
            with open(file_path, 'rb') as f:
                header = f.read(8)
                if not header.startswith(b'\xD0\xCF\x11\xE0'):
                    return "unknown"  # 不是OLE文件
            
            # 获取扩展名
            ext = Path(file_path).suffix.lower()
            if ext in ['.doc', '.docx', '.rtf', '.dot']:
                logger.info(f"基于扩展名判断为Word文档: {file_path}")
                return "application/msword"
            elif ext in ['.xls', '.xlsx', '.xlt', '.xlm']:
                logger.info(f"基于扩展名判断为Excel工作簿: {file_path}")
                return "application/vnd.ms-excel"
            elif ext in ['.ppt', '.pptx', '.pps']:
                logger.info(f"基于扩展名判断为PowerPoint演示文稿: {file_path}")
                return "application/vnd.ms-powerpoint"
            
            # 无有效扩展名或扩展名不明确，尝试特征检测
            file_size = os.path.getsize(file_path)
            
            # 尝试特征检测
            with open(file_path, 'rb') as f:
                content = f.read(min(16384, file_size))  # 读取前16KB或文件全部内容
                
                # 搜索特征字符串
                content_str = content.decode('latin1', errors='ignore')
                
                # PowerPoint特征
                ppt_markers = ['PowerPoint Document', 'PP97_DUALSTORAGE', 'Current User', 
                              'SlideListWithText', 'Pictures', 'SlidePersist']
                if any(marker in content_str for marker in ppt_markers):
                    logger.info(f"基于特征检测为PowerPoint文件: {file_path}")
                    return "application/vnd.ms-powerpoint"
                
                # Excel特征
                excel_markers = ['Workbook', 'Excel', 'Worksheet', 'Microsoft Excel', 
                                'Sheet', 'PivotTable', 'PivotCache']
                if any(marker in content_str for marker in excel_markers):
                    logger.info(f"基于特征检测为Excel文件: {file_path}")
                    return "application/vnd.ms-excel"
                
                # Word特征
                word_markers = ['Word.Document', 'MSWordDoc', 'Microsoft Word', 
                               'Normal.dot', 'WordDocument', 'Document Summary']
                if any(marker in content_str for marker in word_markers):
                    logger.info(f"基于特征检测为Word文件: {file_path}")
                    return "application/msword"
                
                # 基于文件大小和一些概率特征的推测
                if 'Slide' in content_str or 'PowerPoint' in content_str:
                    return "application/vnd.ms-powerpoint"
                
                if 'Excel' in content_str or 'Worksheet' in content_str or 'Workbook' in content_str:
                    return "application/vnd.ms-excel"
                
                if file_size < 100*1024:  # 小于100KB
                    if content_str.count('Table') > 5 or content_str.count('Cell') > 10:
                        logger.info(f"基于内容特征可能是Excel文件: {file_path}")
                        return "application/vnd.ms-excel"  # 可能是Excel
                
                logger.info(f"无法精确判断OLE文件类型，默认为Word文档: {file_path}")
                # 默认假设是Word文档(最常见)
                return "application/msword"
        except Exception as e:
            logger.warning(f"OLE文件类型检测失败: {file_path} - {str(e)}")
            # 返回通用OLE类型
            return "application/x-ole-storage"
    
    def detect_file_type(self, file_path: str) -> str:
        """检测文件 MIME 类型，优化文件类型检测，增强OLE文件识别"""
        try:
            # 规范化路径
            normalized_path = os.path.normpath(file_path.encode('utf-8', errors='replace').decode('utf-8'))
            
            # 获取文件扩展名
            ext = Path(file_path).suffix.lower()
            
            # 通过扩展名快速确定常见Office文件
            if ext == '.docx':
                return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif ext == '.xlsx':
                return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif ext == '.pptx':
                return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif ext == '.doc':
                return 'application/msword'
            elif ext == '.xls':
                return 'application/vnd.ms-excel'
            elif ext == '.ppt':
                return 'application/vnd.ms-powerpoint'
            
            # 使用magic库获取MIME类型
            mime_type = self.mime.from_file(normalized_path)
            
            # 检查文件头部以确认类型
            try:
                with open(file_path, 'rb') as f:
                    header = f.read(8)
                    
                    # 检查Office Open XML格式(ZIP-based)
                    if header.startswith(b'PK\x03\x04'):
                        try:
                            import zipfile
                            with zipfile.ZipFile(file_path) as zf:
                                file_list = zf.namelist()
                                if 'word/document.xml' in file_list:
                                    return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                                elif 'xl/workbook.xml' in file_list:
                                    return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                                elif 'ppt/presentation.xml' in file_list:
                                    return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        except Exception as zip_error:
                            logger.debug(f"ZIP文件内容检查失败: {file_path} - {str(zip_error)}")
                        
                        # 如果无法通过内容确定，但有扩展名，使用扩展名映射
                        if ext in self.MIME_TYPES:
                            return self.MIME_TYPES[ext]
                        return 'application/zip'
                    
                    # 检查OLE复合文档(DOC/XLS/PPT)
                    elif header.startswith(b'\xD0\xCF\x11\xE0'):
                        # 使用增强的OLE文件类型检测
                        ole_type = self._detect_ole_file_type(file_path)
                        if ole_type != "unknown":
                            return ole_type
                        
                        # 如果OLE检测不确定，尝试使用扩展名
                        if ext == '.doc':
                            return 'application/msword'
                        elif ext == '.xls':
                            return 'application/vnd.ms-excel'
                        elif ext == '.ppt':
                            return 'application/vnd.ms-powerpoint'
                        
                        # 最后返回通用OLE类型或基于magic的类型
                        return 'application/x-ole-storage' if mime_type == 'application/octet-stream' else mime_type
                    
                    # 检查PDF文件
                    elif header.startswith(b'%PDF'):
                        return 'application/pdf'
            except Exception as header_error:
                logger.debug(f"文件头检查失败: {file_path} - {str(header_error)}")
            
            # 如果magic库返回的不是通用类型，使用它
            if mime_type not in ('application/octet-stream', 'text/plain'):
                return mime_type
            
            # 最后尝试通过扩展名确定类型
            return self.MIME_TYPES.get(ext, 'application/octet-stream')
            
        except Exception as e:
            logger.error(f"文件类型检测失败 {file_path}: {e}")
            # 如果一切都失败了，尝试用扩展名，否则返回二进制流类型
            return self.MIME_TYPES.get(Path(file_path).suffix.lower(), "application/octet-stream")
    
    def get_mime_by_extension(self, extension: str) -> str:
        """通过文件扩展名获取MIME类型"""
        if not extension.startswith('.'):
            extension = '.' + extension
        return self.MIME_TYPES.get(extension.lower(), "application/octet-stream")
    
    def get_extension_by_mime(self, mime_type: str) -> str:
        """通过MIME类型获取文件扩展名"""
        return self.MIME_TO_EXT.get(mime_type, ".bin")

class ContentExtractor:
    """优化后的文件内容提取器"""
    
    MIME_TYPE = {
        'TEXT': 'text/plain',
        'CSV': 'text/csv',
        'PDF': 'application/pdf',
        'MARKDOWN': 'text/markdown',
        'DOCX': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'XLSX': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'XLS': 'application/vnd.ms-excel',
        'DOC': 'application/msword',
        'PPT': 'application/vnd.ms-powerpoint',
        'PPTX': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'ZIP': 'application/zip',
        'RAR': 'application/x-rar',
        'SEVENZ': 'application/x-7z-compressed'
    }
    
    def __init__(self, detector, is_windows: bool = True):
        self.detector = detector
        try:
            self.md = MarkItDown()
        except Exception as e:
            self.md = None
            logger.error(f"初始化 MarkItDown 失败: {str(e)}")
        self.is_windows = is_windows
        self.word_app = None
        self.word_lock = threading.Lock()
        
        self.extractors = {
            self.MIME_TYPE['TEXT']: self._extract_text_content,
            self.MIME_TYPE['CSV']: self._extract_csv_content,
            self.MIME_TYPE['PDF']: self._extract_pdf_content,
            self.MIME_TYPE['MARKDOWN']: self._extract_markdown_content,
            self.MIME_TYPE['ZIP']: self._extract_archive_content,
            self.MIME_TYPE['RAR']: self._extract_archive_content,
            self.MIME_TYPE['SEVENZ']: self._extract_archive_content,
            self.MIME_TYPE['DOCX']: self._extract_docx_content,
            self.MIME_TYPE['XLSX']: self._extract_xlsx_content,
            self.MIME_TYPE['PPTX']: self._extract_pptx_content,
        }
        
        if self.is_windows:
            # 尝试初始化Word应用
            if self._init_word_app():
                self.extractors.update({
                    self.MIME_TYPE['DOC']: self._extract_doc_content,
                })
            else:
                logger.warning("Word应用程序初始化失败，.doc文件将使用备用方法处理")
            
            # PowerPoint不需要持久应用实例
            self.extractors.update({
                self.MIME_TYPE['PPT']: self._extract_ppt_content,
                self.MIME_TYPE['XLS']: self._extract_xls_content,
            })
    
    def _init_word_app(self):
        """初始化Word应用程序实例"""
        if not self.is_windows:
            return False
        
        with self.word_lock:
            if self.word_app is None:
                try:
                    # 强制关闭任何现有Word进程
                    self._force_close_office_processes('WINWORD.EXE')
                    
                    # 初始化当前线程的COM
                    pythoncom.CoInitialize()
                    
                    # 创建新的Word应用实例
                    self.word_app = win32com.client.Dispatch("Word.Application")
                    self.word_app.Visible = False
                    self.word_app.DisplayAlerts = False
                    
                    logger.info("成功初始化Word应用程序")
                    return True
                except Exception as e:
                    logger.error(f"初始化Word应用程序失败: {str(e)}")
                    self.word_app = None
                    
                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass
                    
                    return False
            return True
    
    def _cleanup_word_app(self):
        """清理 Word 应用程序实例"""
        with self.word_lock:
            if self.word_app is not None:
                try:
                    self.word_app.Quit()
                except:
                    pass
                self._force_close_office_processes('WINWORD.EXE')
                self.word_app = None
                pythoncom.CoUninitialize()
                logger.info("成功清理 Word 应用程序")
    
    def _create_empty_result(self, file_type: str) -> Dict[str, Any]:
        """创建空结果对象"""
        return {
            'content': '',
            'metadata': {'file_type': file_type},
            'error': None,
            'is_empty': True
        }
    
    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果对象"""
        return {
            'content': '',
            'metadata': {'file_type': file_type},
            'error': error_msg
        }
    
    def _get_file_header(self, file_path: str, size: int = 16) -> str:
        """获取文件头部字节的十六进制表示"""
        try:
            with open(file_path, 'rb') as f:
                return f.read(size).hex().upper()
        except Exception as e:
            logger.error(f"无法读取文件头: {file_path} - {e}")
            return ""
    
    def _precheck_file_type(self, file_path: str, detected_mime: str) -> Tuple[bool, str]:
        """预检查文件类型和格式"""
        try:
            ext = Path(file_path).suffix.lower()
            expected_mime = self.detector.MIME_TYPES.get(ext, detected_mime)
            
            if detected_mime != expected_mime:
                logger.warning(f"检测到的 MIME 类型 {detected_mime} 与基于扩展名推断的 {expected_mime} 不匹配，文件: {file_path}")
            
            with open(file_path, 'rb') as f:
                header = f.read(8)
                
                if expected_mime == self.MIME_TYPE['PPT'] and not header.startswith(b'\xD0\xCF\x11\xE0'):
                    logger.warning(f"文件头不符合 .ppt 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE['PPTX'] and not header.startswith(b'PK\x03\x04'):
                    logger.warning(f"文件头不符合 .pptx 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE['XLSX'] and not header.startswith(b'PK\x03\x04'):
                    logger.warning(f"文件头不符合 .xls 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE['DOC'] and not header.startswith(b'\xD0\xCF\x11\xE0'):
                    logger.warning(f"文件头不符合 .doc 格式，文件: {file_path}")
                
            return True, ""
        except Exception as e:
            return False, f"预检查失败: {str(e)}"

    def _is_valid_xlsx(self, file_path: str) -> Tuple[bool, str]:
        """检查 .xlsx 文件是否有效"""
        try:
            with open(file_path, 'rb') as f:
                header = f.read(4)
                if not header.startswith(b'PK\x03\x04'):
                    return False, "文件头不符合 .xlsx 格式"
            
            wb = openpyxl.load_workbook(file_path, read_only=True)
            ws = wb.active
            has_data = False
            
            for row in ws.iter_rows(max_rows=10, values_only=True):
                if any(cell is not None for cell in row):
                    has_data = True
                    break
            
            wb.close()
            
            if not has_data:
                return False, "文件为空或无有效数据"
            
            return True, ""
        except openpyxl.utils.exceptions.InvalidFileException as e:
            return False, f"文件格式无效或损坏: {str(e)}"
        except Exception as e:
            return False, f"文件检查失败: {str(e)}"
    
    def _extract_text_content(self, file_path: str) -> Dict[str, Any]:
        """提取文本文件内容"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                'metadata': {'file_type': 'text'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('text', f"读取文本文件失败: {str(e)}")
    
    def _extract_csv_content(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容"""
        try:
            encodings = ['utf-8', 'latin1', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, engine='python')
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                df = pd.read_csv(file_path, encoding='utf-8', errors='ignore')
            
            if len(df) > 10000:
                df = df.head(10000)
                notice = "\n\n[注意: 文件过大，仅显示前10000行]"
            else:
                notice = ""
                
            return {
                'content': df.to_string(index=False) + notice,
                'metadata': {'file_type': 'csv'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('csv', f"CSV处理失败: {str(e)}")
    
    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文件内容"""
        try:
            text = pdf_extract_text(file_path)
            if len(text) > 100000:
                text = text[:100000] + "\n\n[注意: 文件过大，内容已截断]"
            return {
                'content': text,
                'metadata': {'file_type': 'pdf'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('pdf', f"PDF处理失败: {str(e)}")
    
    def _extract_markdown_content(self, file_path: str) -> Dict[str, Any]:
        """提取Markdown文件内容"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1024000)  # 限制最大读取1MB
            return {
                'content': content[:1000000] + ("\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""),
                'metadata': {'file_type': 'markdown'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result('markdown', f"Markdown文件处理失败: {str(e)}")
    
    def _extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOCX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result('docx', "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                'content': result.text_content,
                'metadata': {'file_type': 'docx', 'extractor': 'markitdown'},
                'error': None
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理DOCX失败: {file_path} - {str(e)}")
            return self._create_error_result('docx', f"DOCX处理失败: {str(e)}")
    
    def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLSX文件内容"""
        try:
            sheets = pd.read_excel(file_path, sheet_name=None)
            content = []
            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': 'excel'},
                'error': None
            }
        except ImportError as e:
            return self._create_error_result('excel', f"缺少必要库: {str(e)}")
        except ValueError as e:
            return self._create_error_result('excel', f"文件格式错误: {str(e)}")
        except Exception as e:
            return self._create_error_result('excel', f"Excel处理失败: {str(e)}，类型: {type(e).__name__}")
    
    def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPTX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result('pptx', "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            return {
                'content': result.text_content,
                'metadata': {'file_type': 'pptx', 'extractor': 'markitdown'},
                'error': None
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理PPTX失败: {file_path} - {str(e)}")
            if self.is_windows:
                try:
                    prs = Presentation(file_path)
                    text_content = []
                    
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        
                        if slide_text:
                            text_content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
                    
                    content = "\n\n".join(text_content)
                    return {
                        'content': content,
                        'metadata': {'file_type': 'pptx', 'extractor': 'python-pptx'},
                        'error': None
                    }
                except Exception as e2:
                    return self._create_error_result('pptx', f"所有提取方法失败: MarkItDown: {str(e)}; python-pptx: {str(e2)}")
            return self._create_error_result('pptx', f"PPTX处理失败: {str(e)}")
    
    def _extract_xls_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLS文件内容，增强错误处理和多级回退策略"""
        logger.info(f'处理XLS文件: {file_path}')
        
        # 策略1: 尝试使用pandas通用方法
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
            content = []
            
            for sheet_name, df in sheets.items():
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)
                
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': 'excel', 'extractor': 'pandas-xlrd'},
                'error': None
            }
        except xlrd.biffh.XLRDError as xlrd_error:
            # 特殊处理XLRD的"Can't find workbook in OLE2 compound document"错误
            logger.warning(f"XLRD无法解析XLS文件: {file_path} - {str(xlrd_error)}")
            # 继续尝试其他方法
        except Exception as pandas_error:
            logger.warning(f"Pandas处理XLS失败: {file_path} - {str(pandas_error)}")
            # 继续尝试其他方法
        
        # 策略2: 如果在Windows上，尝试使用COM接口
        if self.is_windows:
            try:
                pythoncom.CoInitialize()
                excel_app = win32com.client.Dispatch("Excel.Application")
                excel_app.Visible = False
                excel_app.DisplayAlerts = False
                
                workbook = excel_app.Workbooks.Open(
                    os.path.abspath(file_path),
                    ReadOnly=True,
                    UpdateLinks=False
                )
                
                content = []
                for sheet_index in range(1, workbook.Sheets.Count + 1):
                    try:
                        worksheet = workbook.Sheets(sheet_index)
                        sheet_name = worksheet.Name
                        
                        # 获取已使用范围
                        used_range = worksheet.UsedRange
                        if used_range.Rows.Count > 0 and used_range.Columns.Count > 0:
                            # 创建二维数组存储数据
                            data = []
                            for row in range(1, min(5001, used_range.Rows.Count + 1)):
                                row_data = []
                                for col in range(1, used_range.Columns.Count + 1):
                                    try:
                                        cell_value = worksheet.Cells(row, col).Value
                                        row_data.append(str(cell_value) if cell_value is not None else "")
                                    except:
                                        row_data.append("")
                                data.append(row_data)
                            
                            # 格式化为类似表格的文本
                            if data:
                                max_lengths = [max(len(str(row[i])) for row in data if i < len(row)) for i in range(len(data[0]))]
                                formatted_rows = []
                                
                                for row in data:
                                    formatted_row = [str(val).ljust(max_lengths[i]) for i, val in enumerate(row) if i < len(max_lengths)]
                                    formatted_rows.append(" | ".join(formatted_row))
                                
                                sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(formatted_rows)
                                if used_range.Rows.Count > 5000:
                                    sheet_text += "\n[注意: 表格过大，仅显示前5000行]"
                                    
                                content.append(sheet_text)
                    except Exception as sheet_error:
                        logger.warning(f"处理工作表 {sheet_index} 时出错: {str(sheet_error)}")
                        continue
                
                workbook.Close(SaveChanges=False)
                excel_app.Quit()
                
                return {
                    'content': '\n\n'.join(content),
                    'metadata': {'file_type': 'excel', 'extractor': 'win32com'},
                    'error': None
                }
            except Exception as com_error:
                logger.warning(f"使用COM接口处理XLS失败: {file_path} - {str(com_error)}")
            finally:
                try:
                    if 'workbook' in locals() and workbook is not None:
                        workbook.Close(SaveChanges=False)
                except:
                    pass
                    
                try:
                    if 'excel_app' in locals() and excel_app is not None:
                        excel_app.Quit()
                except:
                    pass
                    
                self._force_close_office_processes('EXCEL.EXE')
                pythoncom.CoUninitialize()
        
        # 策略3: 尝试openpyxl (虽然通常用于xlsx，但有时也能处理某些xls)
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                data = []
                row_count = 0
                
                for row in sheet.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else "" for cell in row])
                    row_count += 1
                    if row_count >= 5000:
                        break
                
                if data:
                    max_lengths = [max(len(str(row[i])) for row in data if i < len(row)) for i in range(len(data[0]))]
                    formatted_rows = []
                    
                    for row in data:
                        formatted_row = [str(val).ljust(max_lengths[i]) for i, val in enumerate(row) if i < len(max_lengths)]
                        formatted_rows.append(" | ".join(formatted_row))
                    
                    sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(formatted_rows)
                    if row_count >= 5000:
                        sheet_text += "\n[注意: 表格过大，仅显示前5000行]"
                        
                    content.append(sheet_text)
            
            workbook.close()
            
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': 'excel', 'extractor': 'openpyxl'},
                'error': None
            }
        except Exception as openpyxl_error:
            logger.warning(f"使用openpyxl处理XLS失败: {file_path} - {str(openpyxl_error)}")
        
        # 最后策略: 尝试提取二进制文本
        try:
            binary_result = self._extract_binary_text_content(file_path, 'application/vnd.ms-excel')
            if not binary_result.get('error'):
                return binary_result
        except Exception as e:
            pass
        
        # 所有方法失败，返回错误
        return self._create_error_result(
            'excel', 
            "无法提取XLS文件内容: 所有已知方法均已尝试并失败"
        )
    
    def _extract_archive_content(self, file_path: str) -> Dict[str, Any]:
        """提取压缩文件内容，使用安全的临时目录和路径验证"""
        temp_dir = os.path.join(
            os.path.dirname(file_path), 
            f"temp_extract_{uuid.uuid4().hex}"
        )
        
        try:
            os.makedirs(temp_dir, exist_ok=True)
            patoolib.extract_archive(file_path, outdir=temp_dir)
            
            content = []
            max_files = 100
            max_size = 10 * 1024 * 1024  # 10MB
            total_size = 0
            files_processed = 0
            
            for root, _, files in os.walk(temp_dir):
                if files_processed >= max_files:
                    content.append(f"\n[警告: 超过最大处理文件数 {max_files}, 仅展示部分内容]")
                    break
                
                for file_name in files:
                    if file_name.startswith('._') or file_name.startswith('__MACOSX'):
                        continue
                    
                    if files_processed >= max_files:
                        break
                    
                    file_path_full = os.path.join(root, file_name)
                    normalized_path = os.path.normpath(file_path_full)
                    if not normalized_path.startswith(temp_dir):
                        content.append(f"[警告: 跳过不安全路径: {file_name}]")
                        continue
                    
                    try:
                        file_size = os.path.getsize(file_path_full)
                        if total_size + file_size > max_size:
                            content.append("\n[警告: 达到总大小限制，跳过剩余文件]")
                            break
                            
                        if file_size < 1024 * 1024:  # 1MB
                            try:
                                with open(file_path_full, 'r', encoding='utf-8', errors='ignore') as f:
                                    text = f.read(10240)
                                    if text.strip():
                                        content.append(f"File: {file_name}\n{text[:1024]}...")
                                        files_processed += 1
                                        total_size += file_size
                            except (UnicodeDecodeError, IOError):
                                pass
                    except OSError:
                        content.append(f"[警告: 无法读取文件: {file_name}]")
            
            return {
                'content': '\n\n'.join(content),
                'metadata': {'file_type': Path(file_path).suffix[1:]},
                'error': None
            }
        except patoolib.util.PatoolError as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过（请确保安装 7-Zip 或 unrar）"
            )
        except Exception as e:
            return self._create_error_result(
                Path(file_path).suffix[1:],
                f"压缩文件处理失败: {str(e)} - 跳过"
            )
        finally:
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"清理临时目录失败: {temp_dir} - {str(e)}")
    
    def _extract_doc_content(self, file_path: str, timeout_seconds: int = 10) -> Dict[str, Any]:
        """提取DOC文件内容，增加超时处理，超时后跳过处理"""
        logger.info(f'处理DOC文件: {file_path}')
        
        # 使用线程池执行文档处理，支持超时
        with ThreadPoolExecutor(max_workers=1) as executor:
            # 提交任务到线程池
            future = executor.submit(self._doc_content_worker, file_path)
            
            try:
                # 等待任务完成，设置超时时间
                result = future.result(timeout=timeout_seconds)
                return result
            except TimeoutError:
                # 超时处理
                logger.warning(f"处理DOC文件超时 ({timeout_seconds}秒): {file_path}")
                
                # 确保停止任务执行
                future.cancel()
                
                # 强制清理资源
                self._force_close_office_processes('WINWORD.EXE')
                
                # 返回跳过标记
                return {
                    'content': '',
                    'metadata': {'file_type': 'doc', 'extractor': 'timeout'},
                    'skipped': True,
                    'error': f"DOC文件处理超时 ({timeout_seconds}秒)，可能是编码问题或文件损坏，已跳过"
                }

    def _doc_content_worker(self, file_path: str) -> Dict[str, Any]:
        """实际执行DOC文件内容提取的工作函数 (在线程中执行)"""
        word_app = None
        doc = None
        max_retries = 3
        
        try:
            abs_path = os.path.abspath(file_path)
            pythoncom.CoInitialize()  # 初始化当前线程的COM
            
            # 创建新的Word应用实例
            word_app = win32com.client.Dispatch("Word.Application")
            word_app.Visible = False
            word_app.DisplayAlerts = False
            
            for attempt in range(max_retries):
                try:
                    # 尝试打开文档，启用修复选项
                    doc = word_app.Documents.Open(
                        abs_path,
                        ReadOnly=True,
                        PasswordDocument="",
                        Visible=False,
                        NoEncodingDialog=True,
                        OpenAndRepair=True  # 尝试修复损坏文件
                    )
                    
                    # 提取文本内容
                    content = doc.Range().Text
                    
                    # 检查是否成功提取到内容
                    if not content.strip():
                        logger.warning(f"DOC文件未提取到文本内容: {abs_path}")
                    
                    return {
                        'content': content,
                        'metadata': {'file_type': 'doc', 'extractor': 'win32com'},
                        'error': None if content.strip() else "未提取到任何文本内容"
                    }
                except Exception as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"第 {attempt + 1} 次尝试处理 {file_path} 失败: {str(e)}, 重试中...")
                        # 清理打开的文档后重试
                        try:
                            if doc:
                                doc.Close(SaveChanges=False)
                                doc = None
                        except:
                            pass
                        
                        time.sleep(2)  # 重试前等待更长时间
                        continue
                        
                    logger.error(f"Win32COM 处理失败: {str(e)} - 文件: {file_path}")
                    return self._create_error_result('doc', f"Win32COM 处理失败: {str(e)}")
            
        except Exception as e:
            logger.error(f"初始化 Word 应用程序失败: {str(e)}")
            return self._create_error_result('doc', f"初始化 Word 应用程序失败: {str(e)}")
        
        finally:
            # 确保正确清理资源
            try:
                if doc:
                    doc.Close(SaveChanges=False)
            except Exception as close_err:
                logger.debug(f"关闭文档时出错: {str(close_err)}")
            
            try:
                if word_app:
                    word_app.Quit()
            except Exception as quit_err:
                logger.debug(f"退出Word应用程序时出错: {str(quit_err)}")
            
            # 强制关闭任何残留进程
            self._force_close_office_processes('WINWORD.EXE')
            
            # 清理COM
            try:
                pythoncom.CoUninitialize()
            except Exception as com_err:
                logger.debug(f"COM清理时出错: {str(com_err)}")

    def _force_close_office_processes(self, process_name: str = 'WINWORD.EXE') -> None:
        """强制关闭指定Office进程方法，更健壮的实现"""
        try:
            process_found = False
            for proc in psutil.process_iter(['pid', 'name']):
                try:
                    if proc.info['name'].upper() == process_name:
                        process_found = True
                        logger.info(f"正在关闭进程: {proc.info['name']} (PID: {proc.info['pid']})")
                        
                        # 先尝试优雅终止
                        proc.terminate()
                        
                        # 等待最多3秒让进程终止
                        try:
                            proc.wait(timeout=3)
                        except psutil.TimeoutExpired:
                            # 如果优雅终止失败，强制结束
                            logger.warning(f"强制结束进程: {proc.info['name']} (PID: {proc.info['pid']})")
                            proc.kill()
                            
                            # 等待进程被杀死
                            try:
                                proc.wait(timeout=3)
                            except:
                                logger.error(f"无法终止进程: {proc.info['name']} (PID: {proc.info['pid']})")
                        
                        if not proc.is_running():
                            logger.info(f"成功关闭进程: {proc.info['name']} (PID: {proc.info['pid']})")
                except (psutil.NoSuchProcess, psutil.AccessDenied) as e:
                    logger.debug(f"关闭进程时出错: {e}")
                    continue
            
            if not process_found:
                logger.debug(f"未找到 {process_name} 进程，无需关闭")
        except Exception as e:
            logger.warning(f"关闭进程失败: {e}")

    def _extract_binary_text_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """从二进制文件中提取可能的文本内容，作为最后的回退方案"""
        try:
            with open(file_path, 'rb') as f:
                binary_data = f.read()
            
            # 尝试不同编码
            potential_text = ""
            encodings = ['utf-8', 'latin1', 'cp1252', 'gb18030', 'big5']
            
            for encoding in encodings:
                try:
                    decoded = binary_data.decode(encoding, errors='ignore')
                    # 检查解码是否产生有意义的内容
                    if len(decoded.strip()) > len(potential_text.strip()):
                        potential_text = decoded
                except:
                    continue
            
            # 只提取可打印ASCII字符和常见空白
            printable_text = ""
            for char in potential_text:
                if char.isprintable() or char in ' \t\n\r':
                    printable_text += char
            
            # 移除连续空白
            cleaned_text = re.sub(r'\s+', ' ', printable_text).strip()
            
            # 移除二进制垃圾（随机字符序列）
            final_text = re.sub(r'[^\w\s.,;:!?(){}\[\]\'\"<>@#$%^&*+=\-_\\|/]', '', cleaned_text)
            
            # 移除非常短的片段
            words = [w for w in final_text.split(' ') if len(w) > 1]
            result_text = ' '.join(words)
            
            if len(result_text) < 50:  # 如果提取的文本不够多
                return self._create_error_result(
                    mime_type, 
                    f"二进制内容提取失败: 未能提取到足够的文本内容"
                )
            
            return {
                'content': result_text,
                'metadata': {'file_type': mime_type, 'extractor': 'binary'},
                'error': None
            }
        except Exception as e:
            return self._create_error_result(mime_type, f"二进制内容提取失败: {str(e)}")

    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果字典"""
        return {
            'content': None,
            'metadata': {'file_type': file_type, 'extractor': 'win32com'},
            'error': error_msg
        }

    def _cleanup_word_app(self) -> None:
        """清理 Word 应用程序"""
        if hasattr(self, 'word_app') and self.word_app:
            try:
                self.word_app.Quit()
            except Exception as e:
                logger.debug(f"关闭 Word 应用程序时出错: {e}")
                self._force_close_office_processes('WINWORD.EXE')

    def __del__(self):
        """对象销毁时清理资源"""
        try:
            self._cleanup_word_app()
            pythoncom.CoUninitialize()
        except Exception as e:
            logger.debug(f"清理资源时出错: {e}")
    
    # 优化PowerPoint PPT文件提取器
    def _extract_ppt_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPT文件内容，处理Visible属性错误和多级回退策略"""
        logger.info(f'处理PPT文件: {file_path}')
        ppt_app = None
        presentation = None
        
        # 策略1: 使用修改后的COM接口
        try:
            abs_path = os.path.abspath(file_path)
            pythoncom.CoInitialize()
            
            # 创建PowerPoint应用实例 - 不设置Visible属性
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            
            # 尝试设置不可见，但捕获可能的错误
            try:
                ppt_app.Visible = False  # 有些环境下这会失败
            except Exception as visible_error:
                logger.warning(f"无法设置PowerPoint为不可见模式: {str(visible_error)}，继续处理")
                # 继续处理，接受PowerPoint可能会短暂可见
            
            # 尝试关闭警告对话框
            try:
                ppt_app.DisplayAlerts = False
            except:
                pass
                
            # 打开演示文稿
            presentation = ppt_app.Presentations.Open(
                abs_path, 
                WithWindow=False,  # 尝试不显示窗口打开
                ReadOnly=True      # 只读模式
            )
            
            # 提取所有幻灯片的文本
            text_content = []
            slide_count = 0
            
            try:
                # 获取幻灯片总数
                slide_count = presentation.Slides.Count
                logger.info(f"PPT文件包含 {slide_count} 张幻灯片: {abs_path}")
                
                # 遍历幻灯片
                for i in range(1, slide_count + 1):
                    try:
                        slide = presentation.Slides.Item(i)
                        slide_text = []
                        
                        # 遍历形状
                        for shape_idx in range(1, slide.Shapes.Count + 1):
                            try:
                                shape = slide.Shapes.Item(shape_idx)
                                # 检查是否有文本框架
                                if shape.HasTextFrame:
                                    text_frame = shape.TextFrame
                                    if hasattr(text_frame, "TextRange") and text_frame.TextRange.Text.strip():
                                        slide_text.append(text_frame.TextRange.Text.strip())
                            except Exception as shape_err:
                                # 记录但继续处理下一个形状
                                continue
                        
                        # 只有当幻灯片包含文本时才添加
                        if slide_text:
                            text_content.append(f"Slide {i}:\n" + "\n".join(slide_text))
                    except Exception as slide_err:
                        logger.warning(f"处理幻灯片 {i} 时出错: {str(slide_err)}")
                        continue
            except Exception as slides_err:
                logger.error(f"处理幻灯片集合时出错: {str(slides_err)}")
                # 继续使用已提取的部分内容
            
            # 检查是否提取到内容
            content = "\n\n".join(text_content)
            if not content.strip() and slide_count > 0:
                logger.warning(f"未能从PPT提取到文本内容，但文件包含 {slide_count} 张幻灯片")
                content = f"[注意：此PPT文件包含 {slide_count} 张幻灯片，但未能提取到文本内容]"
            
            return {
                'content': content,
                'metadata': {'file_type': 'ppt', 'extractor': 'win32com', 'slide_count': slide_count},
                'error': None if content.strip() else "未提取到任何文本内容"
            }
        except Exception as e:
            logger.error(f"Win32COM处理PPT失败: {file_path} - {str(e)}")
            # 继续尝试其他方法
        finally:
            # 确保资源释放
            try:
                if presentation:
                    presentation.Close()
            except:
                pass
            
            try:
                if ppt_app:
                    ppt_app.Quit()
            except:
                pass
            
            # 强制关闭任何残留进程
            self._force_close_office_processes('POWERPNT.EXE')
            
            try:
                pythoncom.CoUninitialize()
            except:
                pass
        
        # 策略2: 如果有MarkItDown且COM接口失败，尝试使用MarkItDown
        if self.md is not None:
            try:
                result = self.md.convert(file_path)
                return {
                    'content': result.text_content,
                    'metadata': {'file_type': 'ppt', 'extractor': 'markitdown_fallback'},
                    'error': None
                }
            except Exception as md_error:
                logger.warning(f"MarkItDown处理PPT失败: {file_path} - {str(md_error)}")
                # 继续尝试其他方法
        
        # 策略3: 尝试使用python-pptx (通常用于pptx，但可能能处理某些ppt)
        try:
            from pptx import Presentation
            try:
                prs = Presentation(file_path)
                text_content = []
                
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        text_content.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
                
                content = "\n\n".join(text_content)
                
                if content.strip():
                    return {
                        'content': content,
                        'metadata': {'file_type': 'ppt', 'extractor': 'python-pptx'},
                        'error': None
                    }
            except Exception as pptx_error:
                logger.warning(f"python-pptx处理PPT失败: {file_path} - {str(pptx_error)}")
                # 继续尝试其他方法
        except ImportError:
            logger.debug("python-pptx不可用，跳过此方法")
        
        # 策略4: 尝试使用antiword (适用于某些二进制文档)
        try:
            import subprocess
            try:
                # 检查antiword是否可用
                antiword_path = shutil.which('antiword')
                if antiword_path:
                    result = subprocess.run(
                        [antiword_path, file_path], 
                        capture_output=True, 
                        text=True,
                        timeout=30
                    )
                    
                    if result.returncode == 0 and result.stdout.strip():
                        return {
                            'content': result.stdout,
                            'metadata': {'file_type': 'ppt', 'extractor': 'antiword'},
                            'error': None
                        }
            except Exception as antiword_error:
                logger.debug(f"antiword处理失败: {file_path} - {str(antiword_error)}")
                # 继续尝试其他方法
        except ImportError:
            logger.debug("subprocess不可用，跳过antiword方法")
        
        # 最后策略: 尝试二进制文本提取
        try:
            binary_result = self._extract_binary_text_content(file_path, 'application/vnd.ms-powerpoint')
            if not binary_result.get('error') and binary_result.get('content'):
                return binary_result
        except Exception as e:
            pass
        
        # 所有方法都失败
        return self._create_error_result(
            'ppt', 
            "无法提取PPT文件内容: 所有已知方法均已尝试并失败"
        )
    
    def extract_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """提取文件内容，增强错误处理以防止程序中断"""
        try:
            # 检查文件是否存在
            if not Path(file_path).exists():
                return {'error': f'文件不存在: {file_path}', 'content': '', 'metadata': {}}
            
            # 检查文件是否为空
            if Path(file_path).stat().st_size == 0:
                return self._create_empty_result(mime_type)
            
            # 检查文件扩展名
            ext = Path(file_path).suffix.lower()
            
            # 定义需要跳过的MIME类型
            skip_mime_types = [
                'image/jpeg', 'image/png', 'image/gif', 'image/bmp',  # 图片文件
                'audio/mpeg', 'audio/wav',                            # 音频文件
                'video/mp4', 'video/x-msvideo', 'video/quicktime',    # 视频文件
                'font/ttf', 'font/otf', 'font/woff', 'font/woff2',    # 字体文件
                'application/x-executable'                             # 可执行文件
            ]
            
            # 特殊处理DWG文件
            if mime_type == 'image/vnd.dwg':
                return self._create_error_result(mime_type, "DWG图纸文件, 跳过处理")
            
            # 处理明确不支持的MIME类型
            if mime_type in skip_mime_types:
                return self._create_error_result(mime_type, f"不支持的MIME类型: {mime_type}, 跳过处理")
            
            # ===== 核心内容提取逻辑开始 =====
            
            # 1. 使用已知的提取器尝试处理文件
            if mime_type in self.extractors:
                try:
                    content = self.extractors[mime_type](file_path)
                    if not content.get('error'):
                        return content
                    # 如果提取器返回错误，记录但继续尝试其他方法
                    logger.warning(f"主提取器 {mime_type} 失败: {content.get('error')} - 尝试备用方法")
                except Exception as e:
                    # 捕获所有异常，确保程序不会中断
                    logger.warning(f"主提取器 {mime_type} 抛出异常: {str(e)} - 尝试备用方法")
            
            # 2. 尝试对无扩展名文件特殊处理
            if not ext:
                try:
                    # 读取文件头
                    with open(file_path, 'rb') as f:
                        header = f.read(8)
                    
                    # 检查是否为Office Open XML文件 (docx, xlsx, pptx), 它们以PK\03\04开头(ZIP格式)
                    if header.startswith(b'PK\x03\x04'):
                        try:
                            import zipfile
                            with zipfile.ZipFile(file_path) as zf:
                                file_list = zf.namelist()
                                
                                # 检查是否为Word文档
                                if any('word/document.xml' in name for name in file_list):
                                    logger.info(f"无扩展名文件检测为Word文档: {file_path}")
                                    if self.md is not None:
                                        try:
                                            result = self.md.convert(file_path)
                                            return {
                                                'content': result.text_content,
                                                'metadata': {'file_type': 'docx', 'converter': 'markitdown'},
                                                'error': None
                                            }
                                        except Exception as e:
                                            logger.warning(f"MarkItDown转换Word文档失败: {file_path} - {str(e)}")
                                
                                # 检查是否为Excel工作簿
                                elif any('xl/workbook.xml' in name for name in file_list):
                                    logger.info(f"无扩展名文件检测为Excel工作簿: {file_path}")
                                    if 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in self.extractors:
                                        return self.extractors['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'](file_path)
                                
                                # 检查是否为PowerPoint演示文稿
                                elif any('ppt/presentation.xml' in name for name in file_list):
                                    logger.info(f"无扩展名文件检测为PowerPoint演示文稿: {file_path}")
                                    if self.md is not None:
                                        try:
                                            result = self.md.convert(file_path)
                                            return {
                                                'content': result.text_content,
                                                'metadata': {'file_type': 'pptx', 'converter': 'markitdown'},
                                                'error': None
                                            }
                                        except Exception as e:
                                            logger.warning(f"MarkItDown转换PowerPoint文档失败: {file_path} - {str(e)}")
                                
                                # 其他ZIP文件 - 尝试解压缩处理
                                if 'application/zip' in self.extractors:
                                    return self.extractors['application/zip'](file_path)
                        except Exception as zip_error:
                            logger.warning(f"尝试处理ZIP格式文件失败: {file_path} - {str(zip_error)}")
                    
                    # 检查是否为旧版Office文件 (doc, xls, ppt)
                    elif header.startswith(b'\xD0\xCF\x11\xE0'):
                        # 尝试使用默认的DOC提取器
                        if 'application/msword' in self.extractors:
                            try:
                                logger.info(f"尝试使用DOC提取器处理OLE文件: {file_path}")
                                return self.extractors['application/msword'](file_path)
                            except Exception as doc_error:
                                logger.warning(f"使用DOC提取器处理OLE文件失败: {str(doc_error)}")
                    
                    # 检查是否为PDF文件
                    elif header.startswith(b'%PDF'):
                        logger.info(f"无扩展名文件检测为PDF: {file_path}")
                        if 'application/pdf' in self.extractors:
                            return self.extractors['application/pdf'](file_path)
                    
                    # 尝试作为文本处理
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            sample = f.read(1024)
                            
                        # 如果前1024字节是可显示文本，尝试作为文本文件处理
                        if all(ord(c) < 128 and (c.isprintable() or c in '\n\r\t') for c in sample):
                            logger.info(f"无扩展名文件检测为文本: {file_path}")
                            return self._extract_text_content(file_path)
                    except Exception as text_error:
                        logger.debug(f"尝试作为文本处理失败: {file_path} - {str(text_error)}")
                
                except Exception as e:
                    logger.warning(f"检测无扩展名文件类型出错: {file_path} - {str(e)}")
            
            # 3. 尝试使用MarkItDown
            if self.md is not None:
                try:
                    result = self.md.convert(file_path)
                    return {
                        'content': result.text_content,
                        'metadata': {'file_type': mime_type, 'extractor': 'markitdown'},
                        'error': None
                    }
                except Exception as md_error:
                    # 捕获MarkItDown错误，但不中断处理
                    error_msg = str(md_error)
                    logger.warning(f"MarkItDown转换失败: {file_path} - {error_msg}")
                    
                    # 检查是否为已知的XLRD错误
                    if "XLRDError" in error_msg and "OLE2 compound document" in error_msg:
                        logger.info(f"检测到无法处理的OLE文档，将尝试使用DOC处理器: {file_path}")
                        # 尝试使用DOC处理器处理
                        if 'application/msword' in self.extractors:
                            try:
                                return self.extractors['application/msword'](file_path)
                            except Exception as doc_error:
                                logger.warning(f"使用DOC处理器处理OLE文件失败: {str(doc_error)}")
                    
                    # 如果看起来像文本文件，尝试直接读取
                    try:
                        with open(file_path, 'rb') as f:
                            sample = f.read(1024)
                        
                        is_likely_text = True
                        for byte in sample:
                            # 检查是否包含非ASCII控制字符
                            if byte < 32 and byte not in (9, 10, 13):
                                is_likely_text = False
                                break
                        
                        if is_likely_text:
                            logger.info(f"尝试作为文本处理: {file_path}")
                            return self._extract_text_content(file_path)
                    except Exception as text_error:
                        logger.debug(f"尝试作为文本处理失败: {str(text_error)}")
            
            # 4. 尝试作为二进制文件提取文本
            try:
                logger.info(f"尝试从二进制文件中提取文本: {file_path}")
                with open(file_path, 'rb') as f:
                    data = f.read()
                
                # 尝试不同的编码解码
                for encoding in ['utf-8', 'latin1', 'cp1252', 'gb18030']:
                    try:
                        text = data.decode(encoding, errors='ignore')
                        # 如果解码出来的文本看起来合理
                        if len(text) > 100 and text.count(' ') > len(text) / 20:
                            return {
                                'content': text[:100000],  # 限制内容长度
                                'metadata': {'file_type': 'binary', 'encoding': encoding},
                                'error': None
                            }
                    except Exception:
                        continue
            except Exception as bin_error:
                logger.debug(f"二进制提取失败: {str(bin_error)}")
            
            # 5. 最后的回退策略 - 标记为无法处理但不抛出异常
            logger.warning(f"所有方法均无法处理文件: {file_path}")
            return {
                'content': f"[无法识别的文件格式: {mime_type}]",
                'metadata': {'file_type': 'unknown', 'mime_type': mime_type},
                'skipped': True,
                'error': f"所有可用方法均无法提取内容，文件可能为未知格式或损坏"
            }
            
        except Exception as e:
            # 捕获所有可能的异常，确保程序不会中断
            logger.error(f"提取内容时发生意外错误: {file_path} - {str(e)}")
            return {
                'content': '',
                'metadata': {'file_type': 'error'},
                'skipped': True,
                'error': f"内容提取异常: {str(e)}"
            }

class SensitiveChecker:
    """敏感内容检查器，使用正则表达式替代 Aho-Corasick"""
    
    def __init__(self, config_path: str = "sensitive_config.yaml"):
        """初始化敏感词配置"""
        self.config = self._load_config(config_path)
        self.all_keywords = self.config.get('security_marks', []) + \
                           [kw for cat in self.config.get('sensitive_patterns', {}).values() 
                            for kw in cat.get('keywords', [])]
        escaped_keywords = [re.escape(kw) for kw in self.all_keywords]
        self.keyword_pattern = re.compile('|'.join(escaped_keywords))
    
    def _load_config(self, config_path: str) -> Dict:
        """加载敏感词配置文件"""
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                return yaml.safe_load(f)
        except Exception as e:
            logger.error(f"加载敏感词配置失败: {e}")
            return {}
    
    def check_content(self, text: str) -> List[Tuple[str, List[int]]]:
        """检查文本中的敏感词"""
        keyword_matches = {}
        for match in self.keyword_pattern.finditer(text or ''):
            keyword = match.group()
            if keyword not in keyword_matches:
                keyword_matches[keyword] = []
            keyword_matches[keyword].append(match.start())
        keyword_results = list(keyword_matches.items())
        
        structured_results = []
        for pattern, weight in self.config.get('structured_patterns', {}).items():
            matches = list(re.finditer(pattern, text or ''))
            if matches:
                positions = [m.start() for m in matches]
                structured_results.append((pattern, positions))
        
        number_results = []
        for pattern in self.config.get('number_patterns', []):
            matches = list(re.finditer(pattern, text or ''))
            if matches:
                positions = [m.start() for m in matches]
                number_results.append((pattern, positions))
        
        return keyword_results + structured_results + number_results

class ResultExporter:
    """处理结果导出器"""
    
    def export_to_json(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 JSON 文件"""
        export_data = [
            {
                'file_path': r.file_path,
                'mime_type': r.mime_type,
                'content_preview': r.content.get('content', '')[:200],
                'sensitive_words': [{'word': w, 'positions': p} for w, p in r.sensitive_words],
                'error': r.error,
                'processing_time': r.processing_time
            } for r in results
        ]
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)
    
    def export_to_excel(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 Excel 文件"""
        data = [
            {
                '文件路径': r.file_path,
                '文件类型': r.mime_type,
                '敏感词统计': '; '.join([f"{w}({len(p)}次)" for w, p in r.sensitive_words]),
                '处理时间(秒)': round(r.processing_time, 3),
                '错误信息': r.error or ''
            } for r in results
        ]
        df = pd.DataFrame(data)
        df.to_excel(output_path, index=False, engine='openpyxl')

class ResultMonitor:
    """结果监控器，用于实时输出处理进度和结果"""
    
    def __init__(self, output_csv: str = "processing_results.csv"):
        self.output_csv = output_csv
        self._init_csv()
    
    def _init_csv(self):
        """初始化 CSV 文件"""
        with open(self.output_csv, 'w', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([
                '处理时间', '文件路径', '文件类型', '文件大小(bytes)', '是否为空',
                '敏感词数量', '敏感词列表', '处理时长(秒)', '状态', '错误信息'
            ])
    
    def record_result(self, result: ProcessingResult):
        """记录单个处理结果"""
        try:
            file_size = Path(result.file_path).stat().st_size
        except:
            file_size = 0
        
        sensitive_words_count = len(result.sensitive_words)
        sensitive_words_list = '; '.join([f"{w}({len(p)}处)" for w, p in result.sensitive_words])
        
        status = '失败' if result.error else ('空文件' if result.content.get('is_empty') else 
                                              '已跳过' if result.content.get('skipped') else '成功')
        
        with open(self.output_csv, 'a', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerow([
                datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                result.file_path,
                result.mime_type,
                file_size,
                'Yes' if result.content.get('is_empty') else 'No',
                sensitive_words_count,
                sensitive_words_list,
                round(result.processing_time, 3),
                status,
                result.error or ''
            ])
        
        print(f"[{datetime.now().strftime('%H:%M:%S')}] 处理文件: {result.file_path}")
        print(f"  类型: {result.mime_type}")
        print(f"  大小: {file_size} bytes")
        print(f"  状态: {status}")
        if sensitive_words_count > 0:
            print(f"  发现敏感词: {sensitive_words_list}")
        if result.error:
            print(f"  错误: {result.error}")
        print("-" * 80)

class FileProcessor:
    """优化后的文件处理器主类"""
    
    def __init__(self, config_path: str = "sensitive_config.yaml", 
                 monitor_output: str = "processing_results.csv",
                 chunk_size: int = 1000,
                 max_workers: Optional[int] = None,
                 is_windows: bool = True):
        self.detector = FileTypeDetector()
        self.extractor = ContentExtractor(detector=self.detector, is_windows=is_windows)
        self.checker = SensitiveChecker(config_path)
        self.exporter = ResultExporter()
        self.monitor = ResultMonitor(monitor_output)
        self.chunk_size = chunk_size
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self._mime_cache = {}
        self._file_size_cache = {}
    
    def _scan_directory(self, directory: str) -> Iterator[List[str]]:
        """使用生成器分批扫描目录文件，确保正确关闭scandir迭代器"""
        current_chunk = []
        
        try:
            with os.scandir(directory) as dir_iter:
                for entry in dir_iter:
                    try:
                        if entry.is_file(follow_symlinks=False) and not entry.name.startswith('~$') and not self.detector._is_internal_stream(entry.name):
                            ext = Path(entry.path).suffix.lower()
                            if ext not in self.detector.SKIP_EXTENSIONS:
                                current_chunk.append(entry.path)
                                if len(current_chunk) >= self.chunk_size:
                                    yield current_chunk
                                    current_chunk = []
                        elif entry.is_dir(follow_symlinks=False):
                            for sub_chunk in self._scan_directory(entry.path):
                                yield sub_chunk
                    except (PermissionError, OSError) as e:
                        logger.warning(f"访问文件/目录出错 {entry.path}: {e}")
                        continue
        except (PermissionError, OSError) as e:
            logger.warning(f"访问目录出错 {directory}: {e}")
        
        if current_chunk:
            yield current_chunk
    
    def _preload_file_info(self, file_paths: List[str]):
        """预加载文件信息以减少 I/O 操作"""
        for file_path in file_paths:
            try:
                stat = os.stat(file_path)
                self._file_size_cache[file_path] = stat.st_size
                mime_type = self.detector.detect_file_type(file_path)
                self._mime_cache[file_path] = mime_type
            except Exception as e:
                logger.warning(f"预加载文件信息失败 {file_path}: {e}")
    
    def _process_file_batch(self, file_paths: List[str]) -> List[ProcessingResult]:
        """处理一批文件，Office文件串行处理，其他文件并行处理"""
        results = []
        self._preload_file_info(file_paths)
        
        # 按类型对文件进行分类
        office_files = []
        other_files = []
        
        for fp in file_paths:
            mime_type = self._mime_cache.get(fp)
            ext = Path(fp).suffix.lower()
            
            # 检查是否为需要串行处理的Office文件
            if mime_type in (self.extractor.MIME_TYPE['DOC'], self.extractor.MIME_TYPE['PPT']) or \
               ext in ('.doc', '.ppt', '.xls'):
                office_files.append(fp)
            else:
                other_files.append(fp)
        
        # 串行处理Office文件
        if office_files:
            logger.info(f"串行处理 {len(office_files)} 个Office文件")
            for file_path in office_files:
                try:
                    result = self.process_file(file_path)
                    results.append(result)
                    self.monitor.record_result(result)
                except Exception as e:
                    logger.error(f"处理Office文件异常 {file_path}: {e}")
        
        # 并行处理其他文件
        if other_files:
            current_workers = max(1, min(len(other_files), self.max_workers))
            logger.info(f"并行处理 {len(other_files)} 个非Office文件 (工作线程: {current_workers})")
            
            with ThreadPoolExecutor(max_workers=current_workers) as executor:
                future_to_file = {executor.submit(self.process_file, fp): fp for fp in other_files}
                for future in as_completed(future_to_file):
                    file_path = future_to_file[future]
                    try:
                        result = future.result()
                        results.append(result)
                        self.monitor.record_result(result)
                    except Exception as e:
                        logger.error(f"处理文件异常 {file_path}: {e}")
        
        # 当缓存过大时清理
        if len(self._mime_cache) > 10000:
            self._mime_cache.clear()
            self._file_size_cache.clear()
        
        return results
    
    def process_directory(self, directory: str) -> List[ProcessingResult]:
        """处理目录下的所有文件"""
        results = []
        total_files = len(self.detector.get_all_files(directory))
        completed = 0
        
        print(f"\n开始处理目录: {directory}")
        print(f"共发现 {total_files} 个文件（已忽略 ~$ 开头、内部流文件及 .dwg/音频视频文件）")
        print("=" * 80)
        
        try:
            for file_chunk in self._scan_directory(directory):
                chunk_results = self._process_file_batch(file_chunk)
                results.extend(chunk_results)
                completed += len(file_chunk)
                print(f"\n已完成: {completed}/{total_files} ({completed/total_files*100:.1f}%)")
        except KeyboardInterrupt:
            print("\n用户中断处理...")
        except Exception as e:
            logger.error(f"处理目录时发生错误: {e}")
        finally:
            if results:
                try:
                    self.exporter.export_to_json(results, "interrupted_results.json")
                    self.exporter.export_to_excel(results, "interrupted_results.xlsx")
                except Exception as e:
                    logger.error(f"导出中断结果失败: {e}")
        
        print(f"\n处理完成！共处理 {len(results)}/{total_files} 个文件")
        print(f"详细结果已保存至: {self.monitor.output_csv}")
        print("=" * 80)
        return results
    

    def process_file(self, file_path: str) -> ProcessingResult:
        start_time = time.time()
        try:
            mime_type = self._mime_cache.get(file_path) or self.detector.detect_file_type(file_path)
            ext = Path(file_path).suffix.lower()
            
            # 检查文件头部以辅助判断文件类型
            file_header = ""
            try:
                with open(file_path, 'rb') as f:
                    file_header = f.read(16).hex().upper()
            except Exception:
                pass
            
            # 扩展不支持处理的MIME类型列表
            skip_mime_types = [
                'image/vnd.dwg',                    # DWG图纸文件
                'application/octet-stream',         # 二进制文件
                'application/x-msdownload',         # 可执行文件
                'application/font-sfnt',            # 字体文件
                'font/ttf',                         # TTF字体
                'font/otf',                         # OTF字体
                'font/woff',                        # WOFF字体
                'font/woff2',                       # WOFF2字体
                'text/css',                         # CSS文件
                'application/encrypted',
                'text/javascript',                  # JS文件
                'image/jpeg', 'image/png', 'image/gif', 'image/bmp', 'image/webp', 'image/svg+xml',  # 图片文件
                'audio/mpeg', 'audio/wav',          # 音频文件
                'video/mp4', 'video/x-msvideo', 'video/quicktime'  # 视频文件
            ]
            
            # 扩展不支持的MIME类型前缀
            skip_mime_prefixes = ['font/', 'image/', 'audio/', 'video/']
            
            # 定义常规文件的条件 - 不再主要依赖文件扩展名
            regular_file_conditions = (
                mime_type in skip_mime_types or
                any(mime_type.startswith(prefix) for prefix in skip_mime_prefixes) or
                (ext and ext in self.detector.SKIP_EXTENSIONS)
            )
            
            if regular_file_conditions:
                logger.info(f"识别为无需处理的文件类型，跳过内容提取: {file_path} (type: {mime_type}, header: {file_header})")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={'content': '', 'metadata': {'file_type': 'regular', 'file_header': file_header}, 'skipped': True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time
                )
            
            # 检查文件是否为空
            file_size = self._file_size_cache.get(file_path, 0) or os.path.getsize(file_path)
            if file_size == 0:
                logger.info(f"空文件，跳过: {file_path}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={'is_empty': True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time
                )
            
            # 处理内容的其余部分保持不变
            content = self.extractor.extract_content(file_path, mime_type)
            
            # 检查提取是否失败
            if content.get('error') and "UnsupportedFormatException" in content.get('error'):
                # 如果是不支持的格式，标记为跳过
                logger.info(f"不支持的文件格式，标记为跳过: {file_path} (type: {mime_type})")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={'content': '', 'metadata': {'file_type': 'unsupported'}, 'skipped': True},
                    sensitive_words=[],
                    error=content.get('error'),
                    processing_time=time.time() - start_time
                )
            elif content.get('error'):
                logger.warning(f"提取内容失败: {file_path} - MIME: {mime_type} - 错误: {content.get('error')}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content=content,
                    sensitive_words=[],
                    error=content.get('error'),
                    processing_time=time.time() - start_time
                )
            
            sensitive_result = self.checker.check_content(content.get('content', ''))
            sensitive_words = sensitive_result
            
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type,
                content=content,
                sensitive_words=sensitive_words,
                error=None,
                processing_time=time.time() - start_time
            )
        except Exception as e:
            error_msg = f"处理文件失败: {str(e)}"
            logger.error(f"{error_msg} - {file_path}")
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type if 'mime_type' in locals() else "unknown",
                content={'content': '', 'metadata': {'file_type': 'regular'}, 'skipped': True},
                sensitive_words=[],
                error=error_msg,
                processing_time=time.time() - start_time
            )

def main():
    """主函数"""
    parser = argparse.ArgumentParser(description='文件敏感内容检测工具')
    parser.add_argument('path', help='要处理的文件或目录路径')
    parser.add_argument('--config', default='sensitive_config.yaml', help='敏感词配置文件路径')
    parser.add_argument('--output', default='results', help='输出结果文件名(不含扩展名)')
    parser.add_argument('--chunk-size', type=int, default=1000, help='每批处理的文件数量')
    parser.add_argument('--workers', type=int, default=None, help='最大工作线程数')
    parser.add_argument('--no-windows', action='store_true', help='指定非Windows平台，禁用win32com相关功能')
    
    args = parser.parse_args()
    
    try:
        processor = FileProcessor(
            config_path=args.config,
            monitor_output=f"{args.output}_processing.csv",
            chunk_size=args.chunk_size,
            max_workers=args.workers,
            is_windows=not args.no_windows
        )
        
        results = []
        if Path(args.path).is_file():
            results = [processor.process_file(args.path)]
        else:
            results = processor.process_directory(args.path)
        
        processor.exporter.export_to_json(results, f"{args.output}.json")
        processor.exporter.export_to_excel(results, f"{args.output}.xlsx")
        
        logger.info(f"处理完成，共处理 {len(results)} 个文件")
    except KeyboardInterrupt:
        print("\n用户中断程序执行")
        sys.exit(1)
    except Exception as e:
        logger.error(f"程序执行出错: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
