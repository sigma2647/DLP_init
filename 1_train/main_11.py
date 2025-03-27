o
class FileSignatureDetector:
    """
    增强型文件签名检测器，使用文件头部签名、扩展名和内容特征识别文件类型
    """

    def __init__(self, detector, is_windows: bool = True, verbose=False, cache_size=1024, max_workers=None):
        """
        初始化文件签名检测器，使用python-magic和基于特征的文件类型识别
        
        Args:
            detector: 父类探测器的引用
            is_windows (bool): 是否为Windows平台
            verbose (bool): 是否显示详细日志
            cache_size (int): 缓存最近处理的文件数量
            max_workers (int): 多线程处理的最大线程数，默认为CPU核心数
        """
        self.verbose = verbose
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self.is_windows = is_windows
        self.detector = detector
        
        # 设置日志级别
        if verbose:
            logger.setLevel(logging.DEBUG)
        
        # 初始化python-magic
        try:
            self.mime = magic.Magic(mime=True)
            logger.info("初始化magic库用于MIME类型检测")
        except ImportError:
            logger.warning("未安装python-magic库，部分文件检测功能将受限")
            self.mime = None
        except Exception as e:
            logger.error(f"初始化magic库出错: {e}")
            self.mime = None
        
        # 反向映射
        self.MIME_TO_EXT = {mime: ext for ext, mime in self.MIME_TYPES.items()}
        
        # 编译正则表达式以提高性能
        self.ole_pattern = re.compile(r'(WordDocument|Workbook|PowerPoint Document|Microsoft (?:Word|Excel|PowerPoint))', re.IGNORECASE)
        self.excel_cell_pattern = re.compile(r'Cell[^a-zA-Z0-9]', re.IGNORECASE)
        self.excel_sheet_pattern = re.compile(r'Sheet\d', re.IGNORECASE)
        self.internal_stream_pattern = re.compile(r'^\[\d+\].+|^~\$|^__SRP_\d+$|^__substg1\.0_|^__attach_version1\.0_|^__properties_version1\.0$|^MsoDataStore|Thumbs\.db$|\.tmp$')
        
        # 文件签名字典（十六进制 -> 文件类型）
        self.signatures = OrderedDict([
            # 办公文档格式
            (b'504B0304', 'ZIP格式 (可能是Office文档/APK/JAR等)'),
            (b'D0CF11E0A1B11AE1', 'Microsoft复合文档二进制格式 (DOC/XLS/PPT等)'),
            
            # 图像文件格式 - 增强PNG识别
            (b'FFD8FF', 'JPEG图像'),
            (b'89504E470D0A1A0A', 'PNG图像'),  # 后续会通过_analyze_png_structure进一步分析
            (b'47494638', 'GIF图像'),
            (b'424D', 'BMP图像'),
            (b'49492A00', 'TIFF图像 (小端)'),
            (b'4D4D002A', 'TIFF图像 (大端)'),
            (b'52494646', 'WEBP/WAV/AVI (RIFF格式)'),
            (b'38425053', 'Photoshop文档 (PSD)'),
            
            # 音视频文件格式 - 增强视频识别
            (b'494433', 'MP3音频 (ID3标签)'),
            (b'FFFB', 'MP3音频'),
            (b'FFF3', 'MP3音频'),
            (b'FFF2', 'MP3音频'),
            (b'664C6143', 'FLAC音频'),
            (b'52494646', 'WAV音频 (RIFF标记)'),
            (b'4F676753', 'OGG音频/视频'),
            (b'1A45DFA3', 'Matroska视频 (MKV)'),
            (b'000001BA', 'MPEG传输流 (TS/MTS)'),
            (b'000001B3', 'MPEG视频 (MPG/MPEG)'),
            (b'6674797069736F6D', 'MP4视频'),
            (b'4654595049534F4D', 'MP4视频'),
            (b'52494646', 'AVI视频 (RIFF标记)'),
            (b'3026B2758E66CF11', 'ASF/WMV/WMA'),
            (b'464C5601', 'FLV视频'),
            (b'55525330', 'Video URT'),
            (b'00000020667479706D70', 'MP4视频 (高级格式)'),
            (b'667479704D344120', 'M4A音频容器'),
            (b'2E524D46', 'RealMedia文件'),
            (b'2E524D5600000012', 'RealVideo文件'),
            (b'2E4B4445', 'KDE分区表'),
            
            # 压缩和存档格式
            (b'504B0304', 'ZIP压缩包'),
            (b'526172211A0700', 'RAR压缩包 (v1.5+)'),
            (b'526172211A070100', 'RAR压缩包 (v5.0+)'),
            (b'377ABCAF271C', 'RAR压缩包 (旧版本)'),
            (b'1F8B08', 'GZIP压缩包'),
            (b'1F9D', 'COMPRESS压缩数据'),
            (b'425A68', 'BZIP2压缩包'),
            (b'FD377A585A00', 'XZ压缩包'),
            (b'04224D18', '7-ZIP压缩包'),
            (b'213C617263683E0A', 'DEB安装包'),
            
            # 可执行文件和二进制格式
            (b'4D5A', 'Windows可执行文件 (EXE/DLL/OCX)'),
            (b'7F454C46', 'ELF可执行文件 (Linux/Unix)'),
            (b'CAFEBABE', 'Java类文件/Mach-O Fat文件'),
            (b'FEEDFACE', 'Mach-O (32位)'),
            (b'FEEDFACF', 'Mach-O (64位)'),
            (b'CEFAEDFE', 'Mach-O (逆序, 32位)'),
            (b'CFFAEDFE', 'Mach-O (逆序, 64位)'),
            (b'DEY\x01', 'DEX (Dalvik可执行文件, Android)'),
            (b'FEEDFEED', 'JKS Java密钥存储'),
            
            # 数据库和索引格式
            (b'53514C69746520666F726D6174203300', 'SQLite数据库'),
            (b'4D7953514C', 'MySQL数据库'),
            (b'1F8B08', 'GZ数据库备份'),
            (b'5249464644', 'Webm格式'),
            
            # 文档和富文本格式
            (b'25504446', 'PDF文档'),
            (b'7B5C72746631', 'RTF文档'),
            (b'0A2524454E44', 'UNIX mailbox'),
            
            # 字体格式
            (b'4F54544F', 'OpenType字体'),
            (b'00010000', 'TrueType字体'),
            (b'774F4646', 'WOFF字体'),
            (b'774F4632', 'WOFF2字体'),
            
            # 系统和容器格式
            (b'4C000000011402', 'Windows链接文件 (LNK)'),
            (b'49536328', 'Windows安装程序包 (CAB)'),
            (b'0000000C6A502020', 'JPEG2000图像'),
            (b'EDABEEDB', 'RPM包管理器包'),
            (b'53494D504C4520', 'FITS科学图像格式'),
            (b'414331303031', 'GPG/PGP加密文件'),
            
            # 开发相关格式
            (b'3C3F786D6C', 'XML文件'),
            (b'3C68746D6C', 'HTML文件'),
            (b'3C21444F43', 'HTML文件 (带DOCTYPE)'),
            (b'255044462D', 'PDF文档 (带版本)'),
            (b'2321', 'Shell脚本'),
            (b'7061636B', 'Python已编译代码'),
            (b'4C01', 'ELF文件对象'),
            
            # 文本编码标记
            (b'EFBBBF', 'UTF-8带BOM文本'),
            (b'FEFF', 'UTF-16 (BE)带BOM文本'),
            (b'FFFE', 'UTF-16 (LE)带BOM文本'),
            (b'FFFE0000', 'UTF-32 (LE)带BOM文本'),
            (b'0000FEFF', 'UTF-32 (BE)带BOM文本'),
            
            # 虚拟机和容器格式
            (b'4B444D', 'VMDK虚拟机磁盘'),
            (b'7F454C46', 'ELF可执行文件'),
            (b'23204469736B2044', 'VDI虚拟机磁盘'),
            (b'636F6E6563746978', 'VHD虚拟机磁盘'),
            (b'213C617263683E', 'Linux ar归档'),
            (b'52656365697665', 'BitTorrent部分下载'),
        ])
        
        # 使用LRU缓存装饰器缓存文件检测结果
        self._detect = lru_cache(maxsize=cache_size)(self._detect_impl)
        
        # 文件类型统计
        self.stats = defaultdict(int)
        
        # 创建文件扩展名到MIME类型的映射
        self._init_extension_mapping()
        
        # 文件签名缓存
        self.signature_cache = {}

    def _is_internal_stream(self, file_name: str) -> bool:
        """
        检查文件是否为内部流文件（如OLE文件中的那些）
        
        Args:
            file_name: 要检查的文件名
            
        Returns:
            如果文件是内部流，则为True，否则为False
        """
        return bool(self.internal_stream_pattern.search(file_name))

    def read_file_header(self, file_path: str, bytes_to_read: int = 32) -> str:
        """
        读取文件的前N个字节并返回十六进制字符串

        Args:
            file_path: 文件路径
            bytes_to_read: 要读取的字节数

        Returns:
            文件头的大写十六进制字符串
        """
        try:
            with open(file_path, "rb") as f:
                header = f.read(bytes_to_read)
            return binascii.hexlify(header).decode("ascii").upper()
        except Exception as e:
            logger.error(f"读取文件头出错: {file_path} - {e}")
            return ""

    def get_magic_mime_type(self, file_path: str) -> str:
        """
        使用python-magic检测MIME类型

        Args:
            file_path: 文件路径

        Returns:
            MIME类型字符串，失败时返回空字符串
        """
        if not self.mime:
            return ""

        try:
            # 处理非ASCII文件名或含空格的路径
            normalized_path = os.path.normpath(os.path.abspath(file_path))
            return self.mime.from_file(normalized_path)
        except Exception as e:
            logger.error(f"Magic MIME类型检测失败: {file_path} - {e}")
            return ""

    def inspect_zip_content(self, file_path: str) -> Optional[str]:
        """
        检查ZIP文件内容以确定是否为OOXML文档

        Args:
            file_path: ZIP文件路径

        Returns:
            如果确定，返回MIME类型，否则返回None
        """
        try:
            with zipfile.ZipFile(file_path) as zip_file:
                file_list = set(zip_file.namelist())  # 使用集合提高查找效率

                # 检查DOCX
                if "word/document.xml" in file_list:
                    return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                # 检查XLSX
                elif "xl/workbook.xml" in file_list:
                    return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                # 检查PPTX
                elif "ppt/presentation.xml" in file_list:
                    return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                # 检查ODT (OpenDocument Text)
                elif "mimetype" in file_list and "content.xml" in file_list:
                    try:
                        mimetype = zip_file.read("mimetype").decode('utf-8').strip()
                        if mimetype == "application/vnd.oasis.opendocument.text":
                            return "application/vnd.oasis.opendocument.text"
                        elif mimetype == "application/vnd.oasis.opendocument.spreadsheet":
                            return "application/vnd.oasis.opendocument.spreadsheet"
                        elif mimetype == "application/vnd.oasis.opendocument.presentation":
                            return "application/vnd.oasis.opendocument.presentation"
                    except:
                        pass

                # 如果没有找到特定标记，它是通用ZIP
                return "application/zip"
        except zipfile.BadZipFile:
            logger.warning(f"无效的ZIP文件: {file_path}")
            return None
        except Exception as e:
            logger.warning(f"ZIP内容检查失败: {file_path} - {e}")
            return None

    def inspect_ole_content(self, file_path: str) -> Optional[str]:
        """
        检查OLE文件内容以确定是否为特定Office文档类型

        Args:
            file_path: OLE文件路径

        Returns:
            如果确定，返回MIME类型，否则返回None
        """
        try:
            # 首先，通过扩展名快速判断
            ext = Path(file_path).suffix.lower()
            if ext == ".doc":
                return "application/msword"
            elif ext == ".xls":
                return "application/vnd.ms-excel"
            elif ext == ".ppt":
                return "application/vnd.ms-powerpoint"

            # 如果扩展名不能帮助，尝试深度分析
            # 读取较大的块以搜索标记
            with open(file_path, "rb") as f:
                content = f.read(16384)  # 读取16KB

            content_str = content.decode("latin1", errors="ignore")

            # 使用预编译正则表达式检查特定标记
            if self.ole_pattern.search(content_str):
                # 检查Word标记
                if "WordDocument" in content_str or "Microsoft Word" in content_str:
                    return "application/msword"
                
                # 检查Excel标记
                if "Workbook" in content_str or "Microsoft Excel" in content_str:
                    return "application/vnd.ms-excel"
                
                # 检查单元格和工作表模式（Excel特有）
                if self.excel_cell_pattern.search(content_str) and self.excel_sheet_pattern.search(content_str):
                    return "application/vnd.ms-excel"
                
                # 检查PowerPoint标记
                if "PowerPoint Document" in content_str or "Microsoft PowerPoint" in content_str:
                    return "application/vnd.ms-powerpoint"

            # 检查二进制模式，查找Excel特有的标记
            if b"Standard Jet DB" in content or b"Excel.Sheet" in content:
                return "application/vnd.ms-excel"

            # 如果没有找到特定标记，它是通用OLE文件
            return "application/x-ole-storage"
        except Exception as e:
            logger.warning(f"OLE内容检查失败: {file_path} - {e}")
            return None

    def detect_file_type(self, file_path: str) -> str:
        """
        综合检测文件类型，返回MIME类型

        Args:
            file_path: 文件路径

        Returns:
            检测到的MIME类型
        """
        try:
            # 基本文件信息
            file_size = os.path.getsize(file_path)
            file_extension = Path(file_path).suffix.lower()
            expected_mime = self.MIME_TYPES.get(
                file_extension, "application/octet-stream"
            )

            # 跳过非常小或空的文件
            if file_size == 0:
                return "application/octet-stream"

            # 缓存机制 - 使用文件路径和修改时间作为键
            file_mtime = os.path.getmtime(file_path)
            cache_key = f"{file_path}:{file_mtime}"
            
            # 可以在这里添加缓存查找逻辑
            # if cache_key in self.mime_cache:
            #    return self.mime_cache[cache_key]

            # 读取文件头部（前32字节）
            header_hex = self.read_file_header(file_path, 32)

            # 获取magic MIME类型（如果可用）
            magic_mime = self.get_magic_mime_type(file_path) if self.mime else ""

            # 通过文件头部签名检测
            detected_mime = "application/octet-stream"

            # 检查PDF签名
            if header_hex.startswith("25504446"):
                return "application/pdf"
            
            # 检查图像格式
            elif header_hex.startswith("FFD8FF"):
                return "image/jpeg"
            elif header_hex.startswith("89504E47"):
                return "image/png"
            elif header_hex.startswith("47494638"):
                return "image/gif"
            elif header_hex.startswith("424D"):
                return "image/bmp"
            
            # 检查音频/视频
            elif header_hex.startswith("494433") or header_hex.startswith("FFFB"):
                return "audio/mpeg"
            elif header_hex.startswith("52494646") and "57415645" in header_hex:
                return "audio/wav"
            elif header_hex.startswith("00000020667479704D3441"):
                return "video/mp4"
            
            # 检查ZIP签名
            elif header_hex.startswith("504B0304"):  # ZIP签名
                zip_mime = self.inspect_zip_content(file_path)
                if zip_mime:
                    return zip_mime

            # 检查OLE签名
            elif header_hex.startswith("D0CF11E0A1B11AE1"):  # OLE签名
                ole_mime = self.inspect_ole_content(file_path)
                if ole_mime:
                    return ole_mime

            # 检查压缩文件
            elif header_hex.startswith("526172"):  # RAR
                return "application/x-rar"
            elif header_hex.startswith("377ABCAF"):  # 7-Zip
                return "application/x-7z-compressed"
            elif header_hex.startswith("1F8B08"):  # GZIP
                return "application/gzip"

            # 如果通过签名无法识别，尝试使用magic库
            if magic_mime and magic_mime != "application/octet-stream":
                return magic_mime

            # 检查文本文件 - 尝试读取并解码前4KB
            if file_size < 1024 * 1024 * 5:  # 小于5MB的文件
                try:
                    with open(file_path, "rb") as f:
                        content = f.read(4096)
                    
                    # 尝试以UTF-8解码
                    try:
                        content.decode('utf-8')
                        
                        # 检查具体的文本格式
                        text_content = content.decode('utf-8')
                        
                        # 检查JSON
                        if file_extension == '.json' or (text_content.strip().startswith('{') and text_content.strip().endswith('}')):
                            return "application/json"
                        
                        # 检查XML
                        if file_extension == '.xml' or (text_content.strip().startswith('<') and text_content.strip().endswith('>')):
                            return "text/xml"
                        
                        # 检查CSV
                        if file_extension == '.csv' or ',' in text_content and '\n' in text_content:
                            comma_count = text_content.count(',')
                            newline_count = text_content.count('\n')
                            if comma_count > 0 and newline_count > 0 and comma_count / newline_count > 1:
                                return "text/csv"
                        
                        # 检查Markdown
                        if file_extension == '.md' or '##' in text_content or '```' in text_content:
                            return "text/markdown"
                        
                        # 检查YAML
                        if file_extension in ['.yml', '.yaml'] or ':' in text_content and '\n' in text_content:
                            return "application/yaml"
                        
                        # 默认为纯文本
                        return "text/plain"
                    except UnicodeDecodeError:
                        pass
                except:
                    pass

            # 最后尝试使用扩展名
            if file_extension and expected_mime != "application/octet-stream":
                return expected_mime

            # 如果一切都失败了，返回二进制流类型
            return "application/octet-stream"

        except Exception as e:
            logger.error(f"文件类型检测失败: {file_path} - {e}")
            return "application/octet-stream"

    def get_all_files(self, directory: str) -> List[str]:
        """
        递归获取目录中的所有文件，过滤掉某些类型

        Args:
            directory: 要扫描的目录路径

        Returns:
            文件路径列表
        """
        files = []
        try:
            # 更高效的文件遍历
            for root, dirs, file_names in os.walk(directory):
                # 跳过隐藏目录
                dirs[:] = [d for d in dirs if not d.startswith('.')]
                
                for file_name in file_names:
                    # 跳过隐藏文件和Office临时文件
                    if file_name.startswith(".") or file_name.startswith("~$"):
                        continue

                    # 跳过Office内部流文件
                    if self._is_internal_stream(file_name):
                        continue

                    file_path = os.path.join(root, file_name)
                    
                    # 跳过已知的二进制格式
                    ext = Path(file_path).suffix.lower()
                    if ext in self.SKIP_EXTENSIONS:
                        continue
                    
                    files.append(file_path)
        except Exception as e:
            logger.error(f"遍历目录失败: {directory} - {e}")
        
        return files

    def check_file_signature(self, file_path: str) -> Dict:
        """
        检查文件签名并返回详细信息

        Args:
            file_path: 文件路径

        Returns:
            包含详细文件签名信息的字典
        """
        try:
            mime_type = self.detect_file_type(file_path)

            # 获取文件头
            header_hex = self.read_file_header(file_path, 32)

            # 获取更多元数据
            file_size = os.path.getsize(file_path)
            creation_time = os.path.getctime(file_path)
            modification_time = os.path.getmtime(file_path)
            
            return {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_size": file_size,
                "extension": Path(file_path).suffix.lower(),
                "mime_type": mime_type,
                "header_hex": header_hex,
                "created_at": datetime.fromtimestamp(creation_time).strftime("%Y-%m-%d %H:%M:%S"),
                "modified_at": datetime.fromtimestamp(modification_time).strftime("%Y-%m-%d %H:%M:%S"),
            }
        except Exception as e:
            logger.error(f"检查文件签名失败: {file_path} - {e}")
            return {
                "file_path": file_path,
                "error": str(e)
            }


class ContentExtractor:
    """优化后的文件内容提取器"""

    MIME_TYPE = {
        "TEXT": "text/plain",
        "CSV": "text/csv",
        "HTML": "text/html",
        "XML": "text/xml",
        "JSON": "application/json",
        "YAML": "application/yaml",
        "MARKDOWN": "text/markdown",
        "PDF": "application/pdf",
        "DOC": "application/msword",
        "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "XLS": "application/vnd.ms-excel",
        "XLSX": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "PPT": "application/vnd.ms-powerpoint",
        "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ZIP": "application/zip",
        "RAR": "application/x-rar",
        "SEVENZ": "application/x-7z-compressed",
        "TAR": "application/x-tar",
        "GZIP": "application/gzip",
        "BZIP2": "application/x-bzip2",
    }

    def __init__(self, detector, is_windows: bool = True):
        """
        初始化内容提取器
        
        Args:
            detector: 文件类型检测器实例
            is_windows: 是否为Windows系统
        """
        self.detector = detector
        try:
            self.md = MarkItDown()
        except Exception as e:
            self.md = None
            logger.error(f"初始化 MarkItDown 失败: {str(e)}")
        self.is_windows = is_windows
        self.word_app = None
        self.word_lock = threading.Lock()

        self.extractors = {
            self.MIME_TYPE["TEXT"]: self._extract_text_content,
            self.MIME_TYPE["CSV"]: self._extract_csv_content,
            self.MIME_TYPE["PDF"]: self._extract_pdf_content,
            self.MIME_TYPE["MARKDOWN"]: self._extract_markdown_content,
            self.MIME_TYPE["ZIP"]: self._extract_archive_content,
            self.MIME_TYPE["RAR"]: self._extract_archive_content,
            self.MIME_TYPE["SEVENZ"]: self._extract_archive_content,
            self.MIME_TYPE["DOCX"]: self._extract_docx_content,
            self.MIME_TYPE["XLSX"]: self._extract_xlsx_content,
            self.MIME_TYPE["PPTX"]: self._extract_pptx_content,
            self.MIME_TYPE["HTML"]: self._extract_text_content,
            self.MIME_TYPE["XML"]: self._extract_text_content,
            self.MIME_TYPE["JSON"]: self._extract_text_content,
            self.MIME_TYPE["YAML"]: self._extract_text_content,
        }

        if self.is_windows:
            # 尝试初始化Word应用
            if self._init_word_app():
                self.extractors.update({
                    self.MIME_TYPE["DOC"]: self._extract_doc_content,
                })
            else:
                logger.warning("Word应用程序初始化失败，.doc文件将使用备用方法处理")

            # PowerPoint不需要持久应用实例
            self.extractors.update({
                self.MIME_TYPE["PPT"]: self._extract_ppt_content,
                self.MIME_TYPE["XLS"]: self._extract_xls_content,
            })

    def _init_word_app(self):
        """初始化Word应用程序实例"""
        if not self.is_windows:
            return False

        with self.word_lock:
            if self.word_app is None:
                try:
                    # 强制关闭任何现有Word进程
                    self._force_close_office_processes("WINWORD.EXE")

                    # 初始化当前线程的COM
                    pythoncom.CoInitialize()

                    # 创建新的Word应用实例
                    self.word_app = win32com.client.Dispatch("Word.Application")
                    self.word_app.Visible = False
                    self.word_app.DisplayAlerts = False

                    logger.info("成功初始化Word应用程序")
                    return True
                except Exception as e:
                    logger.error(f"初始化Word应用程序失败: {str(e)}")
                    self.word_app = None

                    try:
                        pythoncom.CoUninitialize()
                    except:
                        pass

                    return False
            return True

    def _cleanup_word_app(self):
        """清理 Word 应用程序实例"""
        with self.word_lock:
            if self.word_app is not None:
                try:
                    self.word_app.Quit()
                except:
                    pass
                self._force_close_office_processes("WINWORD.EXE")
                self.word_app = None
                pythoncom.CoUninitialize()
                logger.info("成功清理 Word 应用程序")

    def _create_empty_result(self, file_type: str) -> Dict[str, Any]:
        """创建空结果对象"""
        return {
            "content": "",
            "metadata": {"file_type": file_type},
            "error": None,
            "is_empty": True,
        }

    def _create_error_result(self, file_type: str, error_msg: str) -> Dict[str, Any]:
        """创建错误结果对象"""
        return {"content": "", "metadata": {"file_type": file_type}, "error": error_msg}

    def _get_file_header(self, file_path: str, size: int = 16) -> str:
        """获取文件头部字节的十六进制表示"""
        try:
            with open(file_path, "rb") as f:
                return f.read(size).hex().upper()
        except Exception as e:
            logger.error(f"无法读取文件头: {file_path} - {e}")
            return ""

    def _precheck_file_type(self, file_path: str, detected_mime: str) -> Tuple[bool, str]:
        """预检查文件类型和格式"""
        try:
            ext = Path(file_path).suffix.lower()
            expected_mime = self.detector.MIME_TYPES.get(ext, detected_mime)

            if detected_mime != expected_mime:
                logger.warning(
                    f"检测到的 MIME 类型 {detected_mime} 与基于扩展名推断的 {expected_mime} 不匹配，文件: {file_path}"
                )

            with open(file_path, "rb") as f:
                header = f.read(8)

                if expected_mime == self.MIME_TYPE["PPT"] and not header.startswith(
                    b"\xd0\xcf\x11\xe0"
                ):
                    logger.warning(f"文件头不符合 .ppt 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE["PPTX"] and not header.startswith(
                    b"PK\x03\x04"
                ):
                    logger.warning(f"文件头不符合 .pptx 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE["XLSX"] and not header.startswith(
                    b"PK\x03\x04"
                ):
                    logger.warning(f"文件头不符合 .xlsx 格式，文件: {file_path}")
                elif expected_mime == self.MIME_TYPE["DOC"] and not header.startswith(
                    b"\xd0\xcf\x11\xe0"
                ):
                    logger.warning(f"文件头不符合 .doc 格式，文件: {file_path}")

            return True, ""
        except Exception as e:
            return False, f"预检查失败: {str(e)}"

    def _is_valid_xlsx(self, file_path: str) -> Tuple[bool, str]:
        """检查 .xlsx 文件是否有效"""
        try:
            with open(file_path, "rb") as f:
                header = f.read(4)
                if not header.startswith(b"PK\x03\x04"):
                    return False, "文件头不符合 .xlsx 格式"

            wb = openpyxl.load_workbook(file_path, read_only=True)
            ws = wb.active
            has_data = False

            for row in ws.iter_rows(max_rows=10, values_only=True):
                if any(cell is not None for cell in row):
                    has_data = True
                    break

            wb.close()

            if not has_data:
                return False, "文件为空或无有效数据"

            return True, ""
        except openpyxl.utils.exceptions.InvalidFileException as e:
            return False, f"文件格式无效或损坏: {str(e)}"
        except Exception as e:
            return False, f"文件检查失败: {str(e)}"

    def _extract_text_content(self, file_path: str) -> Dict[str, Any]:
        """提取文本文件内容"""
        try:
            encodings = ["utf-8", "latin1", "cp1252", "gbk", "gb18030"]
            content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read(1024000)  # 限制最大读取1MB
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                # 如果所有编码都失败，使用二进制模式读取并强制解码
                with open(file_path, "rb") as f:
                    binary_content = f.read(1024000)
                content = binary_content.decode("utf-8", errors="ignore")
                
            return {
                "content": content[:1000000]
                + (
                    "\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""
                ),
                "metadata": {
                    "file_type": "text",
                    "encoding": encoding if 'encoding' in locals() else "utf-8(forced)",
                    "size": len(content)
                },
                "error": None,
            }
        except Exception as e:
            return self._create_error_result("text", f"读取文本文件失败: {str(e)}")

    def _extract_csv_content(self, file_path: str) -> Dict[str, Any]:
        """提取CSV文件内容"""
        try:
            # 检测分隔符
            encodings = ["utf-8", "latin1", "cp1252", "gbk", "gb18030"]
            df = None
            encoding_used = None
            delimiter_used = None
            
            # 先尝试不同的编码
            for encoding in encodings:
                try:
                    # 尝试猜测分隔符
                    with open(file_path, 'r', encoding=encoding) as f:
                        sample = f.read(4096)
                    
                    # 计算可能的分隔符
                    delimiters = [',', ';', '\t', '|']
                    delimiter_counts = {d: sample.count(d) for d in delimiters}
                    likely_delimiter = max(delimiter_counts, key=delimiter_counts.get)
                    
                    # 如果最可能的分隔符次数太少，使用逗号作为默认值
                    if delimiter_counts[likely_delimiter] < 5:
                        likely_delimiter = ','
                    
                    # 尝试用检测到的分隔符解析
                    df = pd.read_csv(file_path, encoding=encoding, sep=likely_delimiter, engine="python")
                    encoding_used = encoding
                    delimiter_used = likely_delimiter
                    break
                except UnicodeDecodeError:
                    continue
                except pd.errors.ParserError:
                    # 如果解析错误，可能是分隔符检测错误，尝试其他分隔符
                    for d in delimiters:
                        if d != likely_delimiter:
                            try:
                                df = pd.read_csv(file_path, encoding=encoding, sep=d, engine="python")
                                encoding_used = encoding
                                delimiter_used = d
                                break
                            except:
                                continue
                    if df is not None:
                        break
                except Exception:
                    continue

            # 如果所有尝试都失败，使用更健壮的方法
            if df is None:
                try:
                    # 使用Python CSV模块直接读取，支持更多异常情况
                    rows = []
                    
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as csvfile:
                        # 先尝试自动检测方言
                        dialect = csv.Sniffer().sniff(csvfile.read(4096))
                        csvfile.seek(0)
                        
                        reader = csv.reader(csvfile, dialect)
                        for row in reader:
                            rows.append(row)
                    
                    # 转换为DataFrame
                    if rows:
                        headers = rows[0]
                        data = rows[1:]
                        df = pd.DataFrame(data, columns=headers)
                        delimiter_used = dialect.delimiter
                        encoding_used = 'utf-8(forced)'
                except Exception as e:
                    return self._create_error_result("csv", f"CSV处理失败，无法检测格式: {str(e)}")

            # 最后保底尝试
            if df is None:
                df = pd.read_csv(file_path, encoding='utf-8', errors='ignore')
                encoding_used = 'utf-8(forced)'
                delimiter_used = ','

            # 限制大小
            if len(df) > 10000:
                df = df.head(10000)
                notice = "\n\n[注意: 文件过大，仅显示前10000行]"
            else:
                notice = ""

            # 分析CSV结构
            column_types = {col: str(df[col].dtype) for col in df.columns}
            num_numeric_cols = sum(1 for dtype in column_types.values() if 'int' in dtype or 'float' in dtype)
            
            return {
                "content": df.to_string(index=False) + notice,
                "metadata": {
                    "file_type": "csv",
                    "rows": len(df),
                    "columns": len(df.columns),
                    "encoding": encoding_used,
                    "delimiter": delimiter_used,
                    "column_types": column_types,
                    "numeric_columns": num_numeric_cols
                },
                "error": None,
            }
        except Exception as e:
            return self._create_error_result("csv", f"CSV处理失败: {str(e)}")

    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """提取PDF文件内容"""
        try:
            start_time = time.time()
            text = pdf_extract_text(file_path)
            extraction_time = time.time() - start_time
            
            # 提取元数据（如果可能）
            metadata = {}
            try:
                with open(file_path, 'rb') as f:
                    parser = PDFParser(f)
                    doc = PDFDocument(parser)
                    metadata = {
                        'pages': len(list(PDFPage.create_pages(doc))),
                        'encrypted': doc.encryption is not None
                    }
                    
                    if 'Title' in doc.info[0]:
                        metadata['title'] = doc.info[0]['Title']
                    if 'Author' in doc.info[0]:
                        metadata['author'] = doc.info[0]['Author']
                    if 'CreationDate' in doc.info[0]:
                        metadata['creation_date'] = doc.info[0]['CreationDate']
            except Exception as meta_err:
                logger.warning(f"PDF元数据提取失败: {str(meta_err)}")
                # 尝试简单估计页数
                try:
                    metadata['pages'] = text.count('\f') + 1  # 估计页数通过换页符
                except:
                    pass
            
            if len(text) > 100000:
                text = text[:100000] + "\n\n[注意: 文件过大，内容已截断]"
                
            # 检测PDF是否有文本层
            has_text = len(text.strip()) > 0
            if not has_text:
                logger.warning(f"PDF可能仅包含扫描图像，没有提取到文本: {file_path}")
                
            return {
                "content": text, 
                "metadata": {
                    "file_type": "pdf", 
                    "extraction_time": f"{extraction_time:.2f}秒",
                    "has_text_layer": has_text,
                    **metadata
                }, 
                "error": None
            }
        except Exception as e:
            return self._create_error_result("pdf", f"PDF处理失败: {str(e)}")

    def _extract_markdown_content(self, file_path: str) -> Dict[str, Any]:
        """提取Markdown文件内容"""
        try:
            encodings = ["utf-8", "latin1", "cp1252", "gbk", "gb18030"]
            content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read(1024000)  # 限制最大读取1MB
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                # 如果所有编码都失败，使用二进制模式读取并强制解码
                with open(file_path, "rb") as f:
                    binary_content = f.read(1024000)
                content = binary_content.decode("utf-8", errors="ignore")
            
            # 简单分析Markdown结构
            headers = len(re.findall(r'^#{1,6}\s+', content, re.MULTILINE))
            code_blocks = len(re.findall(r'```', content)) // 2
            links = len(re.findall(r'\[.*?\]\(.*?\)', content))
            images = len(re.findall(r'!\[.*?\]\(.*?\)', content))
            
            return {
                "content": content[:1000000]
                + (
                    "\n\n[注意: 文件过大，内容已截断]" if len(content) > 1000000 else ""
                ),
                "metadata": {
                    "file_type": "markdown",
                    "encoding": encoding if 'encoding' in locals() else "utf-8(forced)",
                    "headers": headers,
                    "code_blocks": code_blocks,
                    "links": links,
                    "images": images
                },
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(
                "markdown", f"Markdown文件处理失败: {str(e)}"
            )

    def _extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOCX文件内容，优先使用markitdown"""
        if self.md is None:
            try:
                # 备用方法：使用python-docx
                import docx
                doc = docx.Document(file_path)
                
                # 提取文本
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                
                # 提取元数据
                metadata = {
                    "file_type": "docx",
                    "extractor": "python-docx",
                    "paragraphs": len(doc.paragraphs),
                    "sections": len(doc.sections),
                    "tables": len(doc.tables)
                }
                
                # 尝试提取文档属性
                try:
                    core_props = doc.core_properties
                    if core_props.title:
                        metadata["title"] = core_props.title
                    if core_props.author:
                        metadata["author"] = core_props.author
                    if core_props.created:
                        metadata["created"] = str(core_props.created)
                    if core_props.modified:
                        metadata["modified"] = str(core_props.modified)
                except:
                    pass
                
                return {
                    "content": "\n\n".join(full_text),
                    "metadata": metadata,
                    "error": None
                }
            except Exception as docx_err:
                return self._create_error_result("docx", f"DOCX处理失败 (python-docx): {str(docx_err)}")
        
        try:
            result = self.md.convert(file_path)
            # 提取更多元数据
            metadata = {
                "file_type": "docx", 
                "extractor": "markitdown"
            }
            
            # 如果能访问到文档结构，提取更多信息
            if hasattr(result, 'structure'):
                metadata.update({
                    "paragraphs": result.structure.get('paragraphs', 0),
                    "tables": result.structure.get('tables', 0),
                    "images": result.structure.get('images', 0),
                    "sections": result.structure.get('sections', 0)
                })
            
            return {
                "content": result.text_content,
                "metadata": metadata,
                "error": None,
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理DOCX失败: {file_path} - {str(e)}")
            # 尝试备用方法
            try:
                # 使用python-docx
                import docx
                doc = docx.Document(file_path)
                
                # 提取文本
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                
                return {
                    "content": "\n\n".join(full_text),
                    "metadata": {
                        "file_type": "docx", 
                        "extractor": "python-docx",
                        "paragraphs": len(doc.paragraphs)
                    },
                    "error": None
                }
            except Exception as docx_err:
                return self._create_error_result("docx", f"DOCX处理失败: {str(e)}; 备用方法也失败: {str(docx_err)}")

    def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLSX文件内容"""
        try:
            # 检查文件有效性
            is_valid, error_msg = self._is_valid_xlsx(file_path)
            if not is_valid:
                return self._create_error_result("excel", error_msg)
                
            # 使用pandas读取所有工作表
            sheets = pd.read_excel(file_path, sheet_name=None)
            content = []
            total_rows = 0
            total_cells = 0
            sheet_stats = {}
            
            for sheet_name, df in sheets.items():
                # 计算统计数据
                sheet_rows = len(df)
                total_rows += sheet_rows
                sheet_cells = sheet_rows * len(df.columns) if sheet_rows > 0 else 0
                total_cells += sheet_cells
                
                # 记录工作表统计
                sheet_stats[sheet_name] = {
                    "rows": sheet_rows,
                    "columns": len(df.columns),
                    "cells": sheet_cells
                }
                
                # 限制大工作表的输出
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                    
                content.append(sheet_content)
                
            # 尝试提取元数据
            metadata = {
                "file_type": "excel",
                "sheets": len(sheets),
                "total_rows": total_rows,
                "total_cells": total_cells,
                "sheet_stats": sheet_stats
            }
            
            # 尝试获取更多元数据
            try:
                wb = openpyxl.load_workbook(file_path, read_only=True)
                if wb.properties.title:
                    metadata["title"] = wb.properties.title
                if wb.properties.creator:
                    metadata["creator"] = wb.properties.creator
                if wb.properties.created:
                    metadata["created"] = str(wb.properties.created)
                wb.close()
            except:
                pass
                
            return {
                "content": "\n\n".join(content),
                "metadata": metadata,
                "error": None,
            }
        except ImportError as e:
            return self._create_error_result("excel", f"缺少必要库: {str(e)}")
        except ValueError as e:
            return self._create_error_result("excel", f"文件格式错误: {str(e)}")
        except Exception as e:
            return self._create_error_result(
                "excel", f"Excel处理失败: {str(e)}，类型: {type(e).__name__}"
            )

    def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """提取PPTX文件内容，优先使用markitdown"""
        if self.md is None:
            return self._create_error_result("pptx", "MarkItDown 不可用")
        try:
            result = self.md.convert(file_path)
            
            # 提取更多元数据
            metadata = {
                "file_type": "pptx", 
                "extractor": "markitdown"
            }
            
            # 如果能访问到演示文稿结构，提取更多信息
            if hasattr(result, 'structure'):
                metadata.update({
                    "slides": result.structure.get('slides', 0),
                    "images": result.structure.get('images', 0),
                    "charts": result.structure.get('charts', 0),
                    "tables": result.structure.get('tables', 0)
                })
                
            # 检查是否成功提取了内容
            if not result.text_content or result.text_content.strip() == "":
                raise ValueError("MarkItDown无法提取任何文本内容")
                
            content_length = len(result.text_content)
            if content_length > 100000:
                truncated_content = result.text_content[:100000] + "\n\n[注意: 内容过长，已截断]"
                return {
                    "content": truncated_content,
                    "metadata": {
                        **metadata,
                        "original_length": content_length,
                        "truncated": True
                    },
                    "error": None,
                }
            
            return {
                "content": result.text_content,
                "metadata": metadata,
                "error": None,
            }
        except Exception as e:
            logger.warning(f"MarkItDown处理PPTX失败: {file_path} - {str(e)}")
            
            # 尝试使用python-pptx作为备用方法
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text_content = []
                slide_metadata = []

                # 提取每张幻灯片的内容和元数据
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    shape_count = 0
                    text_shape_count = 0
                    image_count = 0
                    chart_count = 0
                    table_count = 0
                    
                    # 遍历幻灯片中的所有形状
                    for shape in slide.shapes:
                        shape_count += 1
                        
                        # 提取文本内容
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                            text_shape_count += 1
                        
                        # 识别图像
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            image_count += 1
                        
                        # 识别图表
                        if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART
                            chart_count += 1
                        
                        # 识别表格
                        if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                            table_count += 1
                            if hasattr(shape, 'table'):
                                # 尝试提取表格内容
                                table_data = []
                                for row in shape.table.rows:
                                    row_data = []
                                    for cell in row.cells:
                                        if cell.text_frame.text:
                                            row_data.append(cell.text_frame.text.strip())
                                        else:
                                            row_data.append("")
                                    if any(cell for cell in row_data):  # 确保行不是空的
                                        table_data.append(" | ".join(row_data))
                                
                                if table_data:
                                    slide_text.append("表格内容:")
                                    slide_text.append("\n".join(table_data))

                    # 将幻灯片文本添加到内容列表
                    if slide_text:
                        text_content.append(
                            f"幻灯片 {i + 1}:\n" + "\n".join(slide_text)
                        )
                    
                    # 记录幻灯片元数据
                    slide_metadata.append({
                        "slide_number": i + 1,
                        "shapes": shape_count,
                        "text_shapes": text_shape_count,
                        "images": image_count,
                        "charts": chart_count,
                        "tables": table_count,
                        "has_content": len(slide_text) > 0
                    })

                content = "\n\n".join(text_content)
                
                # 创建完整的元数据
                detailed_metadata = {
                    "file_type": "pptx",
                    "extractor": "python-pptx",
                    "slides": len(prs.slides),
                    "slides_with_content": sum(1 for slide in slide_metadata if slide["has_content"]),
                    "total_shapes": sum(slide["shapes"] for slide in slide_metadata),
                    "total_images": sum(slide["images"] for slide in slide_metadata),
                    "total_charts": sum(slide["charts"] for slide in slide_metadata),
                    "total_tables": sum(slide["tables"] for slide in slide_metadata),
                    "slide_details": slide_metadata
                }
                
                # 获取作者和标题（如果可用）
                try:
                    if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
                        detailed_metadata["author"] = prs.core_properties.author
                    if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                        detailed_metadata["title"] = prs.core_properties.title
                except:
                    pass
                    
                return {
                    "content": content,
                    "metadata": detailed_metadata,
                    "error": None,
                }
            except ImportError:
                logger.error("Python-pptx库未安装，无法提取PPTX")
                return self._create_error_result("pptx", f"MarkItDown处理失败，且未安装python-pptx: {str(e)}")
            except Exception as pptx_error:
                logger.error(f"Python-pptx处理失败: {str(pptx_error)}")
                return self._create_error_result(
                    "pptx",
                    f"所有提取方法均失败: MarkItDown: {str(e)}; python-pptx: {str(pptx_error)}",
                )


    def _extract_xls_content(self, file_path: str) -> Dict[str, Any]:
        """提取XLS文件内容，增强错误处理和多级回退策略"""
        logger.info(f"处理XLS文件: {file_path}")

        # 策略1: 尝试使用pandas通用方法
        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine="xlrd")
            content = []
            total_rows = 0
            total_cells = 0
            sheet_stats = {}

            for sheet_name, df in sheets.items():
                # 计算和记录统计数据
                sheet_rows = len(df)
                total_rows += sheet_rows
                sheet_cells = sheet_rows * len(df.columns) if sheet_rows > 0 else 0
                total_cells += sheet_cells
                
                sheet_stats[sheet_name] = {
                    "rows": sheet_rows,
                    "columns": len(df.columns),
                    "cells": sheet_cells
                }
                
                # 限制输出大小
                if len(df) > 5000:
                    df = df.head(5000)
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}\n[注意: 表格过大，仅显示前5000行]"
                else:
                    sheet_content = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                content.append(sheet_content)

            return {
                "content": "\n\n".join(content),
                "metadata": {
                    "file_type": "excel", 
                    "extractor": "pandas-xlrd",
                    "sheets": len(sheets),
                    "total_rows": total_rows,
                    "total_cells": total_cells,
                    "sheet_stats": sheet_stats
                },
                "error": None,
            }
        except xlrd.biffh.XLRDError as xlrd_error:
            # 特殊处理XLRD的"Can't find workbook in OLE2 compound document"错误
            logger.warning(f"XLRD无法解析XLS文件: {file_path} - {str(xlrd_error)}")
            # 继续尝试其他方法
        except Exception as pandas_error:
            logger.warning(f"Pandas处理XLS失败: {file_path} - {str(pandas_error)}")
            # 继续尝试其他方法

        # 策略2: 如果在Windows上，尝试使用COM接口
        if self.is_windows:
            try:
                pythoncom.CoInitialize()
                excel_app = win32com.client.Dispatch("Excel.Application")
                excel_app.Visible = False
                excel_app.DisplayAlerts = False
                
                # 设置超时保护
                timeout_occurred = [False]
                excel_result = [None]
                
                def process_excel():
                    try:
                        workbook = excel_app.Workbooks.Open(
                            os.path.abspath(file_path), ReadOnly=True, UpdateLinks=False
                        )

                        content = []
                        sheet_stats = {}
                        total_rows = 0
                        total_cells = 0
                        
                        for sheet_index in range(1, workbook.Sheets.Count + 1):
                            try:
                                worksheet = workbook.Sheets(sheet_index)
                                sheet_name = worksheet.Name

                                # 获取已使用范围
                                used_range = worksheet.UsedRange
                                if used_range.Rows.Count > 0 and used_range.Columns.Count > 0:
                                    # 记录工作表统计信息
                                    sheet_rows = used_range.Rows.Count
                                    sheet_cols = used_range.Columns.Count
                                    total_rows += sheet_rows
                                    total_cells += sheet_rows * sheet_cols
                                    
                                    sheet_stats[sheet_name] = {
                                        "rows": sheet_rows,
                                        "columns": sheet_cols,
                                        "cells": sheet_rows * sheet_cols
                                    }
                                    
                                    # 创建二维数组存储数据
                                    data = []
                                    max_display_rows = min(5001, used_range.Rows.Count)
                                    for row in range(1, max_display_rows):
                                        row_data = []
                                        for col in range(1, used_range.Columns.Count + 1):
                                            try:
                                                cell_value = worksheet.Cells(row, col).Value
                                                row_data.append(
                                                    str(cell_value)
                                                    if cell_value is not None
                                                    else ""
                                                )
                                            except:
                                                row_data.append("")
                                        data.append(row_data)

                                    # 格式化为类似表格的文本
                                    if data:
                                        max_lengths = [
                                            max(
                                                len(str(row[i])) for row in data if i < len(row)
                                            )
                                            for i in range(len(data[0]))
                                        ]
                                        formatted_rows = []

                                        for row in data:
                                            formatted_row = [
                                                str(val).ljust(max_lengths[i])
                                                for i, val in enumerate(row)
                                                if i < len(max_lengths)
                                            ]
                                            formatted_rows.append(" | ".join(formatted_row))

                                        sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(
                                            formatted_rows
                                        )
                                        if used_range.Rows.Count > 5000:
                                            sheet_text += "\n[注意: 表格过大，仅显示前5000行]"

                                        content.append(sheet_text)
                            except Exception as sheet_error:
                                logger.warning(
                                    f"处理工作表 {sheet_index} 时出错: {str(sheet_error)}"
                                )
                                continue

                        workbook.Close(SaveChanges=False)
                        
                        # 设置结果
                        excel_result[0] = {
                            "content": "\n\n".join(content),
                            "metadata": {
                                "file_type": "excel", 
                                "extractor": "win32com",
                                "sheets": workbook.Sheets.Count,
                                "total_rows": total_rows,
                                "total_cells": total_cells,
                                "sheet_stats": sheet_stats
                            },
                            "error": None,
                        }
                    except Exception as e:
                        logger.error(f"COM处理Excel过程中出错: {str(e)}")
                        excel_result[0] = {
                            "content": "",
                            "metadata": {"file_type": "excel", "extractor": "win32com"},
                            "error": f"COM处理失败: {str(e)}",
                        }
                
                # 超时处理函数
                def handle_timeout():
                    timeout_occurred[0] = True
                    logger.warning(f"Excel COM处理超时: {file_path}")
                    try:
                        excel_app.Quit()
                    except:
                        pass
                    self._force_close_office_processes("EXCEL.EXE")
                
                # 设置超时
                timeout = 30  # 30秒超时
                timer = threading.Timer(timeout, handle_timeout)
                process_thread = threading.Thread(target=process_excel)
                
                # 启动线程和计时器
                timer.daemon = True
                process_thread.daemon = True
                timer.start()
                process_thread.start()
                
                # 等待处理完成
                process_thread.join(timeout + 5)  # 额外5秒缓冲
                timer.cancel()
                
                # 检查结果
                if timeout_occurred[0]:
                    return self._create_error_result("excel", f"COM处理超时({timeout}秒)")
                
                if excel_result[0]:
                    return excel_result[0]
                
                return self._create_error_result("excel", "COM处理Excel未返回结果")
            except Exception as com_error:
                logger.warning(
                    f"使用COM接口处理XLS失败: {file_path} - {str(com_error)}"
                )
            finally:
                try:
                    if 'workbook' in locals() and workbook is not None:
                        workbook.Close(SaveChanges=False)
                except:
                    pass

                try:
                    if 'excel_app' in locals() and excel_app is not None:
                        excel_app.Quit()
                except:
                    pass

                self._force_close_office_processes("EXCEL.EXE")
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass

        # 策略3: 尝试openpyxl (虽然通常用于xlsx，但有时也能处理某些xls)
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

            content = []
            sheet_stats = {}
            total_rows = 0
            total_cells = 0

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                data = []
                row_count = 0

                for row in sheet.iter_rows(values_only=True):
                    data.append([str(cell) if cell is not None else "" for cell in row])
                    row_count += 1
                    if row_count >= 5000:
                        break

                # 更新统计信息
                total_rows += row_count
                columns = len(data[0]) if data and len(data) > 0 else 0
                total_cells += row_count * columns
                
                sheet_stats[sheet_name] = {
                    "rows": row_count,
                    "columns": columns,
                    "cells": row_count * columns
                }

                if data:
                    max_lengths = [
                        max(len(str(row[i])) for row in data if i < len(row))
                        for i in range(len(data[0]))
                    ]
                    formatted_rows = []

                    for row in data:
                        formatted_row = [
                            str(val).ljust(max_lengths[i])
                            for i, val in enumerate(row)
                            if i < len(max_lengths)
                        ]
                        formatted_rows.append(" | ".join(formatted_row))

                    sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(formatted_rows)
                    if row_count >= 5000:
                        sheet_text += "\n[注意: 表格过大，仅显示前5000行]"

                    content.append(sheet_text)

            workbook.close()

            return {
                "content": "\n\n".join(content),
                "metadata": {
                    "file_type": "excel", 
                    "extractor": "openpyxl",
                    "sheets": len(workbook.sheetnames),
                    "total_rows": total_rows,
                    "total_cells": total_cells,
                    "sheet_stats": sheet_stats
                },
                "error": None,
            }
        except Exception as openpyxl_error:
            logger.warning(
                f"使用openpyxl处理XLS失败: {file_path} - {str(openpyxl_error)}"
            )

        # 最后策略: 尝试提取二进制文本
        try:
            binary_result = self._extract_binary_text_content(
                file_path, "application/vnd.ms-excel"
            )
            if not binary_result.get("error"):
                return binary_result
        except Exception as e:
            logger.error(f"二进制文本提取失败: {str(e)}")

        # 所有方法失败，返回错误
        return self._create_error_result(
            "excel", "无法提取XLS文件内容: 所有已知方法均已尝试并失败"
        )


    def _extract_doc_content(self, file_path: str) -> Dict[str, Any]:
        """提取DOC文件内容，使用专用的Word实例，改进COM对象管理"""
        logger.info(f"处理DOC文件: {file_path}")
        max_retries = 3
        timeout = 30  # 默认超时时间(秒)

        for attempt in range(max_retries):
            try:
                # 设置超时标志和结果变量
                timeout_occurred = [False]
                result = [None]
                exception = [None]
                processing_completed = [False]

                # 主处理函数
                def process_doc():
                    try:
                        # 初始化COM
                        pythoncom.CoInitialize()
                        word_app = None
                        doc = None

                        try:
                            # 创建Word应用实例
                            word_app = win32com.client.Dispatch("Word.Application")
                            word_app.Visible = False
                            word_app.DisplayAlerts = False

                            # 打开文档
                            abs_path = os.path.abspath(file_path)
                            doc = word_app.Documents.Open(
                                abs_path, ReadOnly=True, Visible=False, 
                                AddToRecentFiles=False, PasswordDocument=""
                            )

                            # 提取文本内容
                            content = doc.Content.Text
                            
                            # 提取文档属性和元数据
                            metadata = {
                                "file_type": "doc",
                                "extractor": "win32com",
                            }
                            
                            # 尝试提取文档统计信息
                            try:
                                metadata["pages"] = doc.ComputeStatistics(2)  # wdStatisticPages
                                metadata["paragraphs"] = doc.ComputeStatistics(3)  # wdStatisticParagraphs
                                metadata["words"] = doc.ComputeStatistics(0)  # wdStatisticWords
                                metadata["characters"] = doc.ComputeStatistics(1)  # wdStatisticCharacters
                            except:
                                pass
                                
                            # 尝试提取文档属性
                            try:
                                if doc.BuiltInDocumentProperties("Title").Value:
                                    metadata["title"] = doc.BuiltInDocumentProperties("Title").Value
                                if doc.BuiltInDocumentProperties("Author").Value:
                                    metadata["author"] = doc.BuiltInDocumentProperties("Author").Value
                                if doc.BuiltInDocumentProperties("Subject").Value:
                                    metadata["subject"] = doc.BuiltInDocumentProperties("Subject").Value
                            except:
                                pass

                            # 设置结果
                            result[0] = {
                                "content": content,
                                "metadata": metadata,
                                "error": None
                            }

                        except Exception as e:
                            # 捕获处理异常
                            exception[0] = e
                        finally:
                            # 确保资源释放
                            try:
                                if doc:
                                    doc.Close(SaveChanges=False)
                            except:
                                pass
                            try:
                                if word_app:
                                    word_app.Quit()
                            except:
                                pass
                            try:
                                pythoncom.CoUninitialize()
                            except:
                                pass

                            # 标记处理完成
                            processing_completed[0] = True
                    except Exception as e:
                        # 捕获总体异常
                        exception[0] = e
                        processing_completed[0] = True

                # 超时处理函数
                def handle_timeout():
                    if not processing_completed[0]:
                        timeout_occurred[0] = True
                        logger.warning(f"处理文件超时 ({timeout}秒): {file_path}")
                        self._force_close_office_processes("WINWORD.EXE")

                # 创建并启动处理线程
                process_thread = threading.Thread(target=process_doc)
                process_thread.daemon = True
                process_thread.start()

                # 创建并启动超时定时器
                timer = threading.Timer(timeout, handle_timeout)
                timer.daemon = True
                timer.start()

                # 等待处理完成或超时
                process_thread.join(timeout + 2)  # 给超时处理留出额外时间
                timer.cancel()  # 取消定时器

                # 检查处理结果
                if timeout_occurred[0]:
                    logger.error(f"处理DOC文件超时: {file_path}")
                    if attempt < max_retries - 1:
                        logger.warning(f"将在 {attempt + 1} 秒后重试...")
                        time.sleep(attempt + 1)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": f"处理超时 ({timeout}秒)",
                    }
                elif exception[0]:
                    # 处理过程中发生异常
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 失败: {str(exception[0])}, 重试中..."
                        )
                        time.sleep(2)
                        continue
                        
                    # 所有重试都失败，尝试备用方法
                    try:
                        # 尝试使用antiword（如果安装）
                        antiword_result = self._extract_doc_with_antiword(file_path)
                        if antiword_result:
                            return antiword_result
                    except:
                        pass
                        
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": f"Win32COM 处理失败: {str(exception[0])}",
                    }
                elif result[0]:
                    # 处理成功
                    return result[0]
                else:
                    # 未知错误
                    if attempt < max_retries - 1:
                        logger.warning(
                            f"第 {attempt + 1} 次尝试处理 {file_path} 出现未知错误, 重试中..."
                        )
                        time.sleep(2)
                        continue
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": "处理过程遇到未知错误",
                    }

            except Exception as e:
                logger.error(f"处理DOC文件异常: {str(e)}")
                if attempt == max_retries - 1:
                    # 尝试备用方法
                    try:
                        # 尝试使用antiword（如果安装）
                        antiword_result = self._extract_doc_with_antiword(file_path)
                        if antiword_result:
                            return antiword_result
                    except:
                        pass
                        
                    return {
                        "content": "",
                        "metadata": {"file_type": "doc", "extractor": "win32com"},
                        "error": f"处理失败: {str(e)}",
                    }

        # 清理可能残留的进程
        self._force_close_office_processes("WINWORD.EXE")
        
        # 尝试最后的备用方法：使用二进制文本提取
        try:
            binary_result = self._extract_binary_text_content(
                file_path, "application/msword"
            )
            if not binary_result.get("error"):
                return binary_result
        except:
            pass
            
        return {
            "content": "",
            "metadata": {"file_type": "doc", "extractor": "win32com"},
            "error": "所有处理尝试均失败",
        }

    def _extract_doc_with_antiword(self, file_path: str) -> Dict[str, Any]:
        """
        使用antiword工具提取DOC文件内容（备用方法）
        
        Args:
            file_path: DOC文件路径
            
        Returns:
            包含内容和元数据的字典，如果失败则返回None
        """
        try:
            # 检查antiword是否可用
            subprocess.run(["antiword", "-h"], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
            
            # 使用antiword提取文本
            result = subprocess.run(
                ["antiword", file_path], 
                stdout=subprocess.PIPE, 
                stderr=subprocess.PIPE,
                timeout=10,
                check=False
            )
            
            if result.returncode == 0:
                content = result.stdout.decode('utf-8', errors='ignore')
                return {
                    "content": content,
                    "metadata": {
                        "file_type": "doc", 
                        "extractor": "antiword"
                    },
                    "error": None
                }
            else:
                error = result.stderr.decode('utf-8', errors='ignore')
                logger.warning(f"Antiword处理失败: {error}")
                return None
        except FileNotFoundError:
            logger.warning("Antiword工具未安装")
            return None
        except subprocess.TimeoutExpired:
            logger.warning("Antiword处理超时")
            return None
        except Exception as e:
            logger.warning(f"使用Antiword处理失败: {str(e)}")
            return None


    def _extract_archive_content(self, file_path: str) -> Dict[str, Any]:
        """提取压缩文件内容，使用安全的临时目录和路径验证"""
        # 使用系统临时目录和随机UUID创建安全的临时目录
        temp_dir = os.path.join(
            tempfile.gettempdir(), f"extract_{uuid.uuid4().hex}"
        )

        try:
            # 创建临时目录
            os.makedirs(temp_dir, exist_ok=True)
            archive_type = Path(file_path).suffix.lower()[1:]  # 去掉点号
            archive_size = os.path.getsize(file_path) / (1024 * 1024)  # 大小(MB)
            
            # 检查归档大小，防止解压炸弹
            if archive_size > 100:  # 100MB
                return self._create_error_result(
                    archive_type, f"压缩文件过大 ({archive_size:.2f} MB)，跳过解压，防止解压炸弹"
                )
            
            # 尝试使用合适的工具解压缩
            extract_start = time.time()
            try:
                # 首先尝试使用zipfile处理.zip文件
                if file_path.lower().endswith('.zip') or zipfile.is_zipfile(file_path):
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        # 安全检查路径，防止路径遍历攻击
                        safe_members = []
                        for member in zip_ref.namelist():
                            target_path = os.path.normpath(os.path.join(temp_dir, member))
                            if not target_path.startswith(os.path.normpath(temp_dir)):
                                logger.warning(f"跳过不安全路径: {member}")
                                continue
                            safe_members.append(member)
                        
                        # 仅解压安全文件
                        for member in safe_members[:1000]:  # 限制最多解压1000个文件
                            try:
                                zip_ref.extract(member, temp_dir)
                            except Exception as e:
                                logger.warning(f"解压文件 {member} 失败: {str(e)}")
                        
                    extractor = "zipfile"
                    
                # 然后尝试使用7z命令行工具
                elif shutil.which('7z'):
                    result = subprocess.run(
                        ['7z', 'x', '-o' + temp_dir, '-y', file_path],
                        capture_output=True,
                        text=True,
                        timeout=60  # 设置超时
                    )
                    if result.returncode == 0:
                        extractor = "7z"
                    else:
                        raise Exception(f"7z解压失败: {result.stderr}")
                        
                # 最后尝试使用patoolib
                else:
                    import patoolib
                    patoolib.extract_archive(file_path, outdir=temp_dir)
                    extractor = "patoolib"
                    
            except Exception as extract_error:
                return self._create_error_result(
                    archive_type, f"解压缩失败: {str(extract_error)} - 跳过（请确保安装 7-Zip 或 unrar）"
                )
            
            extract_time = time.time() - extract_start
            logger.info(f"解压 {file_path} 完成，耗时 {extract_time:.2f} 秒")
                
            # 处理提取的文件
            content = []
            max_files = 100
            max_size = 10 * 1024 * 1024  # 10MB
            total_size = 0
            files_processed = 0
            file_stats = {
                "text_files": 0,
                "binary_files": 0,
                "skipped_files": 0,
                "total_files": 0,
                "max_depth": 0,
                "extensions": {}
            }
            
            # 定义有用的文本文件扩展名
            text_extensions = {
                '.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.yaml', 
                '.yml', '.cfg', '.ini', '.conf', '.log', '.properties', '.sql',
                '.js', '.py', '.java', '.c', '.cpp', '.h', '.cs', '.php', '.rb'
            }
            
            # 定义列出归档文件结构的函数
            def get_archive_structure(start_dir, max_depth=3):
                structure = []
                
                for root, dirs, files in os.walk(start_dir):
                    # 计算当前深度
                    rel_path = os.path.relpath(root, start_dir)
                    depth = 0 if rel_path == '.' else rel_path.count(os.sep) + 1
                    
                    # 更新最大深度统计
                    file_stats["max_depth"] = max(file_stats["max_depth"], depth)
                    
                    # 如果超出最大深度，则不再继续
                    if depth > max_depth:
                        continue
                    
                    # 构建文件夹结构信息
                    if depth > 0:  # 跳过根目录
                        dir_name = os.path.basename(root)
                        structure.append(f"{'  ' * (depth-1)}📁 {dir_name}/")
                    
                    # 添加文件信息
                    for file in sorted(files):
                        file_stats["total_files"] += 1
                        ext = os.path.splitext(file)[1].lower()
                        
                        # 统计扩展名
                        if ext in file_stats["extensions"]:
                            file_stats["extensions"][ext] += 1
                        else:
                            file_stats["extensions"][ext] = 1
                        
                        # 限制深度和文件数量
                        if depth <= max_depth and len(structure) < 100:
                            structure.append(f"{'  ' * depth}📄 {file}")
                
                # 如果文件太多，添加提示
                if file_stats["total_files"] > 100:
                    structure.append(f"\n[注意: 仅显示部分文件，共 {file_stats['total_files']} 个文件]")
                    
                return structure
            
            # 获取并添加归档结构
            archive_structure = get_archive_structure(temp_dir)
            content.append("📦 归档文件结构:")
            content.append("\n".join(archive_structure))
            content.append("")  # 空行分隔

            # 提取文件内容样本
            for root, _, files in os.walk(temp_dir):
                if files_processed >= max_files:
                    content.append(f"\n[警告: 超过最大处理文件数 {max_files}, 仅展示部分内容]")
                    break

                for file_name in sorted(files):
                    # 跳过MacOS元数据
                    if file_name.startswith("._") or file_name.startswith("__MACOSX"):
                        file_stats["skipped_files"] += 1
                        continue

                    if files_processed >= max_files:
                        break

                    file_path_full = os.path.join(root, file_name)
                    normalized_path = os.path.normpath(file_path_full)
                    
                    # 安全检查
                    if not normalized_path.startswith(temp_dir):
                        content.append(f"[警告: 跳过不安全路径: {file_name}]")
                        file_stats["skipped_files"] += 1
                        continue

                    try:
                        file_size = os.path.getsize(file_path_full)
                        ext = os.path.splitext(file_name)[1].lower()
                        
                        # 检查限制
                        if total_size + file_size > max_size:
                            content.append("\n[警告: 达到总大小限制，跳过剩余文件]")
                            break

                        # 只处理文本文件且大小适中
                        if ext in text_extensions and file_size < 100 * 1024:  # 100KB
                            try:
                                with open(file_path_full, "r", encoding="utf-8", errors="ignore") as f:
                                    text = f.read(5120)  # 最多读取5KB
                                    if text.strip():
                                        # 获取文件相对路径
                                        rel_path = os.path.relpath(file_path_full, temp_dir)
                                        content.append(f"📄 文件: {rel_path}")
                                        
                                        # 检测并限制文本长度
                                        if len(text) > 1000:
                                            text = text[:1000] + "...[内容已截断]"
                                        
                                        content.append(f"```\n{text}\n```")
                                        files_processed += 1
                                        total_size += file_size
                                        file_stats["text_files"] += 1
                            except Exception as text_err:
                                file_stats["binary_files"] += 1
                        else:
                            file_stats["binary_files"] += 1
                    except OSError:
                        content.append(f"[警告: 无法读取文件: {file_name}]")
                        file_stats["skipped_files"] += 1

            # 添加统计信息
            content.append("\n📊 归档统计信息:")
            content.append(f"- 文件总数: {file_stats['total_files']}")
            content.append(f"- 文件夹最大深度: {file_stats['max_depth']}")
            content.append(f"- 已处理文本文件: {file_stats['text_files']}")
            content.append(f"- 二进制或跳过的文件: {file_stats['binary_files'] + file_stats['skipped_files']}")
            
            # 显示文件类型分布
            if file_stats["extensions"]:
                content.append("\n📁 文件类型分布:")
                for ext, count in sorted(file_stats["extensions"].items(), key=lambda x: x[1], reverse=True)[:10]:
                    if ext:
                        content.append(f"- {ext}: {count} 个文件")
                
                if len(file_stats["extensions"]) > 10:
                    content.append(f"- 其他: {sum(count for ext, count in file_stats['extensions'].items() if ext not in dict(sorted(file_stats['extensions'].items(), key=lambda x: x[1], reverse=True)[:10]))} 个文件")

            return {
                "content": "\n\n".join(content),
                "metadata": {
                    "file_type": archive_type,
                    "extractor": extractor,
                    "extraction_time": f"{extract_time:.2f}秒",
                    "archive_size_mb": f"{archive_size:.2f}",
                    "files_processed": files_processed,
                    "total_files": file_stats["total_files"],
                    "text_files": file_stats["text_files"],
                    "binary_files": file_stats["binary_files"],
                    "skipped_files": file_stats["skipped_files"]
                },
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(
                Path(file_path).suffix[1:], f"压缩文件处理失败: {str(e)} - 跳过"
            )
        finally:
            try:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)
            except Exception as e:
                logger.warning(f"清理临时目录失败: {temp_dir} - {str(e)}")



    def _extract_binary_text_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """从二进制文件中提取可能的文本内容，作为最后的回退方案"""
        try:
            # 读取二进制数据，限制大小
            file_size = os.path.getsize(file_path)
            max_size = 1024 * 1024  # 限制为1MB
            read_size = min(file_size, max_size)
            
            with open(file_path, "rb") as f:
                binary_data = f.read(read_size)

            # 尝试多种编码，按可能性排序
            encodings = ["utf-8", "latin1", "cp1252", "gb18030", "big5", "utf-16", "utf-16-le", "utf-16-be"]
            
            # 保存最佳解码结果
            best_decoded = ""
            best_encoding = ""
            best_valid_ratio = 0.0
            
            # 尝试不同编码并评估结果质量
            for encoding in encodings:
                try:
                    decoded = binary_data.decode(encoding, errors="ignore")
                    
                    # 计算可打印字符比例
                    total_chars = len(decoded)
                    if total_chars == 0:
                        continue
                        
                    printable_chars = sum(1 for char in decoded if char.isprintable() or char in " \t\n\r")
                    valid_ratio = printable_chars / total_chars
                    
                    # 如果这个编码产生了更好的结果
                    if valid_ratio > best_valid_ratio:
                        best_valid_ratio = valid_ratio
                        best_decoded = decoded
                        best_encoding = encoding
                        
                    # 如果结果非常好，可以提前结束
                    if valid_ratio > 0.95:
                        break
                except:
                    continue
            
            # 如果没有找到好的解码，返回错误
            if best_valid_ratio < 0.5:
                return self._create_error_result(
                    mime_type, f"二进制内容提取失败: 未能提取到有效的文本内容"
                )
            
            # 处理解码后的文本
            cleaned_text = best_decoded
            
            # 1. 移除控制字符，仅保留可打印字符和基本空白
            printable_text = ""
            for char in cleaned_text:
                if char.isprintable() or char in " \t\n\r":
                    printable_text += char
                    
            # 2. 移除连续空白
            cleaned_text = re.sub(r"\s+", " ", printable_text).strip()
            
            # 3. 智能分段：按照可能的段落分隔符分段
            paragraphs = re.split(r"\n\s*\n|\r\n\s*\r\n|\r\s*\r", cleaned_text)
            paragraphs = [p.strip() for p in paragraphs if p.strip()]
            
            # 4. 移除过短的片段和明显是二进制垃圾的片段
            min_para_length = 20  # 最小段落长度
            clean_paragraphs = []
            
            for para in paragraphs:
                # 跳过太短的段落
                if len(para) < min_para_length:
                    continue
                    
                # 跳过看起来像乱码的段落（包含太多特殊字符）
                special_char_ratio = sum(1 for c in para if not c.isalnum() and c not in " .,;:!?-\"'()[]{}<>/\\") / len(para)
                if special_char_ratio > 0.3:  # 如果特殊字符比例过高
                    continue
                    
                clean_paragraphs.append(para)
            
            # 限制结果大小
            max_content_length = 50000
            if sum(len(p) for p in clean_paragraphs) > max_content_length:
                # 选择最具代表性的段落
                result_paragraphs = []
                current_length = 0
                
                # 优先选择较长的段落
                for para in sorted(clean_paragraphs, key=len, reverse=True):
                    if current_length + len(para) + 2 <= max_content_length:  # +2 for newlines
                        result_paragraphs.append(para)
                        current_length += len(para) + 2
                    else:
                        break
                        
                final_text = "\n\n".join(result_paragraphs)
                truncated = True
            else:
                final_text = "\n\n".join(clean_paragraphs)
                truncated = False
                
            # 如果提取的文本不够多
            if len(final_text) < 100:
                return self._create_error_result(
                    mime_type, f"二进制内容提取失败: 未能提取到足够的文本内容"
                )

            return {
                "content": final_text,
                "metadata": {
                    "file_type": mime_type, 
                    "extractor": "binary",
                    "encoding": best_encoding,
                    "quality_ratio": f"{best_valid_ratio:.2f}",
                    "original_size": file_size,
                    "processed_size": read_size,
                    "paragraphs": len(clean_paragraphs),
                    "truncated": truncated
                },
                "error": None,
            }
        except Exception as e:
            return self._create_error_result(mime_type, f"二进制内容提取失败: {str(e)}")



    def extract_content(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        """提取文件内容，包含多级失败回退策略，增强容错性"""

        # 将所有操作包装在一个大的try-except块中，确保不会崩溃
        try:
            # 检查文件路径
            if not os.path.exists(file_path):
                return {
                    "error": f"文件不存在: {file_path}",
                    "content": "",
                    "metadata": {},
                }

            # 检查文件大小
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                return self._create_empty_result(mime_type)
                
            # 检查文件权限
            if not os.access(file_path, os.R_OK):
                return {
                    "error": f"无权限读取文件: {file_path}",
                    "content": "",
                    "metadata": {"file_type": mime_type},
                }

            # 获取扩展名
            ext = Path(file_path).suffix.lower()

            # 检查文件大小，记录大文件
            size_mb = file_size / (1024 * 1024)
            if size_mb > 50:  # 大于50MB的文件
                logger.warning(f"文件过大 ({size_mb:.2f} MB)，可能影响处理性能: {file_path}")
                
                # 对于非常大的文件(>200MB)，直接跳过
                if size_mb > 200:
                    return {
                        "error": f"文件过大 ({size_mb:.2f} MB)，跳过处理",
                        "content": "",
                        "metadata": {"file_type": mime_type, "size_mb": f"{size_mb:.2f}"},
                    }

            # 跳过不支持的MIME类型
            skip_mime_types = [
                "image/vnd.dwg",  # DWG图纸文件
                "application/x-msdownload",  # 可执行文件
                "application/font-sfnt",  # 字体文件
                "font/ttf",  # TTF字体
                "font/otf",  # OTF字体
                "font/woff",  # WOFF字体
                "font/woff2",  # WOFF2字体
                "text/css",  # CSS文件
                "application/encrypted",
                "text/javascript",  # JS文件
                "image/jpeg",
                "image/png",
                "image/gif",
                "image/bmp",
                "image/webp",
                "image/svg+xml",  # 图片文件
                "audio/mpeg",
                "audio/wav",  # 音频文件
                "video/mp4",
                "video/x-msvideo",
                "video/quicktime",  # 视频文件
            ]

            # 跳过不支持的MIME类型前缀
            skip_mime_prefixes = ["font/", "image/", "audio/", "video/"]
            
            # 检查是否为不需要处理的文件类型
            if mime_type in skip_mime_types or any(mime_type.startswith(prefix) for prefix in skip_mime_prefixes):
                logger.info(f"识别为无需处理的文件类型，跳过内容提取: {file_path} (类型: {mime_type})")
                return {
                    "content": "",
                    "metadata": {"file_type": mime_type, "skipped": True},
                    "error": None
                }

            # 检测文件实际类型
            start_time = time.time()
            
            # 尝试使用注册的提取器提取内容
            if mime_type in self.extractors:
                try:
                    logger.info(f"使用主要提取器 '{mime_type}' 处理: {file_path}")
                    content = self.extractors[mime_type](file_path)
                    
                    # 添加处理时间到元数据
                    if isinstance(content, dict) and 'metadata' in content:
                        content['metadata']['processing_time'] = f"{time.time() - start_time:.2f} 秒"
                    
                    return content
                except Exception as e:
                    logger.error(f"主要提取器失败: {str(e)}")
                    # 继续尝试备用提取器
            
            # 如果没有注册的提取器，尝试基于MIME类型的通用处理
            if mime_type.startswith("text/"):
                # 所有文本类型
                return self._extract_text_content(file_path)
            elif mime_type == "application/json":
                return self._extract_text_content(file_path)
            elif mime_type == "application/xml" or mime_type == "text/xml":
                return self._extract_text_content(file_path)
            elif mime_type == "application/pdf":
                return self._extract_pdf_content(file_path)
            elif mime_type in ["application/zip", "application/x-zip-compressed"]:
                return self._extract_archive_content(file_path)
            elif mime_type in ["application/x-rar", "application/vnd.rar"]:
                return self._extract_archive_content(file_path)
            elif mime_type == "application/x-7z-compressed":
                return self._extract_archive_content(file_path)
            
            # 基于扩展名的备用处理
            if ext == ".txt" or ext == ".log" or ext == ".ini" or ext == ".cfg":
                return self._extract_text_content(file_path)
            elif ext == ".csv":
                return self._extract_csv_content(file_path)
            elif ext == ".json":
                return self._extract_text_content(file_path)
            elif ext == ".xml":
                return self._extract_text_content(file_path)
            elif ext == ".md" or ext == ".markdown":
                return self._extract_markdown_content(file_path)
            elif ext == ".pdf":
                return self._extract_pdf_content(file_path)
            elif ext in [".zip", ".jar", ".war", ".ear"]:
                return self._extract_archive_content(file_path)
            elif ext in [".rar"]:
                return self._extract_archive_content(file_path)
            elif ext in [".7z"]:
                return self._extract_archive_content(file_path)
            elif ext in [".gz", ".bz2", ".xz", ".tar", ".tgz", ".tbz"]:
                return self._extract_archive_content(file_path)
            
            # 最后的回退：尝试提取二进制文本内容
            logger.warning(f"没有找到适合 {mime_type} 的提取器，尝试二进制提取: {file_path}")
            return self._extract_binary_text_content(file_path, mime_type)
            
        except Exception as e:
            # 总体异常处理，确保函数永远不会崩溃
            logger.error(f"提取内容时发生严重错误: {str(e)} - 文件: {file_path}")
            return self._create_error_result(
                mime_type, f"提取过程发生严重错误: {str(e)}"
            )

    def _force_close_office_processes(self, process_name: str) -> None:
        """强制关闭指定Office进程"""
        try:
            for proc in psutil.process_iter(["pid", "name"]):
                try:
                    proc_name = proc.info["name"].upper()
                    if proc_name == process_name or process_name in proc_name:
                        logger.info(
                            f"关闭Office进程: {proc.info['name']} (PID: {proc.info['pid']})"
                        )

                        # 尝试优雅终止
                        proc.terminate()

                        # 等待进程终止
                        try:
                            proc.wait(timeout=3)
                        except psutil.TimeoutExpired:
                            # 强制终止
                            proc.kill()
                except:
                    continue
        except Exception as e:
            logger.debug(f"关闭Office进程失败: {e}")

    def __del__(self):
        """对象销毁时清理资源"""
        try:
            self._cleanup_word_app()
            try:
                pythoncom.CoUninitialize()
            except:
                pass
        except Exception as e:
            logger.debug(f"清理资源时出错: {e}")


class SensitiveChecker:
    """敏感内容检查器，根据提供的参数使用BERT模型或YAML配置"""

    def __init__(self, config_path="sensitive_config.yaml", model_path="best_model.pth", use_model=False, cache_size=1000, threshold=0.7):
        """
        初始化敏感内容检查器
        
        Args:
            config_path: YAML配置文件路径
            model_path: BERT模型文件路径
            use_model: 是否优先使用模型，True表示使用BERT模型，False表示使用YAML配置
            cache_size: 敏感检查结果缓存大小
            threshold: BERT模型判断敏感的阈值（0-1之间）
        """
        self.config_path = config_path
        self.model_path = model_path
        self.use_model = use_model
        self.mode = None
        self.bert_model = None
        self.bert_tokenizer = None
        self.config = None
        self.all_keywords = []
        self.keyword_pattern = None
        self.max_length = 512  # 增加最大长度，提高检测效果
        self.threshold = threshold
        self.device = 'cuda' if torch.cuda.is_available() else 'cpu'
        self.result_cache = {}  # 缓存检查结果
        self.cache_size = cache_size
        self.check_count = 0  # 检查次数统计
        self.hit_count = 0  # 敏感内容命中次数统计
        
        # 初始化检测方法
        self._init_detector()
        
    def _init_detector(self):
        """初始化检测器，加载模型或配置"""
        if self.use_model:
            # 使用BERT模型
            if os.path.exists(self.model_path):
                try:
                    logger.info(f"尝试加载BERT模型: {self.model_path}")
                    # 加载分词器和模型
                    model_name = 'bert-base-multilingual-cased'
                    self.bert_tokenizer = BertTokenizer.from_pretrained(model_name)
                    self.bert_model = BertForSequenceClassification.from_pretrained(model_name, num_labels=2)
                    self.bert_model.load_state_dict(torch.load(self.model_path, map_location=self.device))
                    self.bert_model.to(self.device)
                    self.bert_model.eval()
                    
                    self.mode = 'bert'
                    logger.info(f"成功加载BERT模型，使用设备: {self.device}")
                    
                    # 如果有配置文件，同时加载作为补充
                    if os.path.exists(self.config_path):
                        try:
                            self._load_yaml_config()
                            logger.info("同时加载YAML配置作为补充检测")
                        except Exception as yaml_err:
                            logger.warning(f"加载YAML配置失败: {str(yaml_err)}")
                except Exception as e:
                    logger.error(f"加载BERT模型失败: {str(e)}")
                    # 尝试回退到YAML配置
                    if os.path.exists(self.config_path):
                        try:
                            self._load_yaml_config()
                            self.mode = 'yaml'
                            logger.info(f"BERT模型加载失败，回退到YAML配置: {self.config_path}")
                        except Exception as yaml_err:
                            logger.error(f"加载YAML配置也失败: {str(yaml_err)}")
                            self.mode = 'empty'
                    else:
                        self.mode = 'empty'
            else:
                logger.error(f"指定的模型文件不存在: {self.model_path}")
                # 尝试回退到YAML配置
                if os.path.exists(self.config_path):
                    try:
                        self._load_yaml_config()
                        self.mode = 'yaml'
                        logger.info(f"模型文件不存在，使用YAML配置: {self.config_path}")
                    except Exception as yaml_err:
                        logger.error(f"加载YAML配置失败: {str(yaml_err)}")
                        self.mode = 'empty'
                else:
                    self.mode = 'empty'
        else:
            # 使用YAML配置
            if os.path.exists(self.config_path):
                try:
                    self._load_yaml_config()
                    self.mode = 'yaml'
                    logger.info(f"使用YAML配置进行敏感内容检测: {self.config_path}")
                except Exception as e:
                    logger.error(f"加载YAML配置失败: {str(e)}")
                    self.mode = 'empty'
            else:
                logger.error(f"指定的配置文件不存在: {self.config_path}")
                self.mode = 'empty'
        
        if self.mode == 'empty':
            logger.warning("检测器初始化失败，将使用空检测器（不会检测任何敏感内容）")
    
    def _load_yaml_config(self):
        """加载YAML配置文件，优化正则表达式编译"""
        with open(self.config_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)
        
        # 收集所有关键词    
        self.all_keywords = []
        
        # 收集安全标记关键词
        security_marks = self.config.get("security_marks", [])
        if security_marks:
            self.all_keywords.extend(security_marks)
        
        # 收集各种敏感模式的关键词
        for cat in self.config.get("sensitive_patterns", {}).values():
            cat_keywords = cat.get("keywords", [])
            if cat_keywords:
                self.all_keywords.extend(cat_keywords)
        
        # 预编译正则表达式
        if self.all_keywords:
            # 对关键词按长度排序，保证优先匹配最长的关键词
            self.all_keywords.sort(key=len, reverse=True)
            # 转义特殊字符并编译
            escaped_keywords = [re.escape(kw) for kw in self.all_keywords]
            self.keyword_pattern = re.compile("|".join(escaped_keywords))
            
        # 预编译其他正则表达式
        self.structured_patterns = {}
        for pattern, weight in self.config.get("structured_patterns", {}).items():
            try:
                compiled = re.compile(pattern)
                self.structured_patterns[compiled] = weight
            except re.error as e:
                logger.error(f"正则表达式编译失败: {pattern}, 错误: {str(e)}")
                
        # 预编译数字模式
        self.number_patterns = []
        for pattern in self.config.get("number_patterns", []):
            try:
                compiled = re.compile(pattern)
                self.number_patterns.append(compiled)
            except re.error as e:
                logger.error(f"数字模式正则表达式编译失败: {pattern}, 错误: {str(e)}")

    def preprocess_text(self, text):
        """预处理输入文本用于BERT模型"""
        if not text or not isinstance(text, str):
            return None
        
        # 处理过长的文本，采用滑动窗口策略
        if len(text) > self.max_length * 3:
            # 文本过长，分段处理并选取最可能敏感的段落
            chunks = self._split_text_into_chunks(text, self.max_length)
            encodings = []
            
            for chunk in chunks:
                encoding = self.bert_tokenizer.encode_plus(
                    chunk,
                    add_special_tokens=True,
                    max_length=self.max_length,
                    truncation=True,
                    padding='max_length',
                    return_attention_mask=True,
                    return_tensors='pt'
                )
                encodings.append(encoding)
                
            return encodings
            
        # 正常长度文本的处理
        encoding = self.bert_tokenizer.encode_plus(
            text,
            add_special_tokens=True,
            max_length=self.max_length,
            truncation=True,
            padding='max_length',
            return_attention_mask=True,
            return_tensors='pt'
        )
        return encoding

    def _split_text_into_chunks(self, text, chunk_size):
        """将长文本分割成重叠的块"""
        words = text.split()
        chunks = []
        
        # 使用滑动窗口，步长为chunk_size的一半，确保重叠
        stride = chunk_size // 2
        
        for i in range(0, len(words), stride):
            chunk = ' '.join(words[i:i + chunk_size])
            chunks.append(chunk)
            
            # 限制最大块数，防止处理过多
            if len(chunks) >= 10:
                break
                
        return chunks

    def predict(self, text):
        """对输入文本进行分类预测"""
        if self.bert_model is None or self.bert_tokenizer is None:
            logger.error("BERT模型未成功加载，无法进行预测")
            return 0  # 默认为非敏感
        
        # 文本过长的预处理    
        encoding = self.preprocess_text(text)
        if encoding is None:
            return 0
            
        # 检查是否为分块编码
        if isinstance(encoding, list):
            # 多块处理，获取最高敏感度得分
            max_score = 0
            for enc in encoding:
                input_ids = enc['input_ids'].to(self.device)
                attention_mask = enc['attention_mask'].to(self.device)
                
                with torch.no_grad():
                    outputs = self.bert_model(input_ids=input_ids, attention_mask=attention_mask)
                    logits = outputs.logits
                    probabilities = torch.nn.functional.softmax(logits, dim=1)
                    sensitive_score = probabilities[0][1].item()  # 第二个类别为敏感的概率
                    
                    max_score = max(max_score, sensitive_score)
            
            # 根据最高得分判断敏感性
            return 1 if max_score >= self.threshold else 0
        else:
            # 单块处理
            input_ids = encoding['input_ids'].to(self.device)
            attention_mask = encoding['attention_mask'].to(self.device)

            # 进行预测
            with torch.no_grad():
                outputs = self.bert_model(input_ids=input_ids, attention_mask=attention_mask)
                logits = outputs.logits
                probabilities = torch.nn.functional.softmax(logits, dim=1)
                sensitive_score = probabilities[0][1].item()  # 第二个类别为敏感的概率

            return 1 if sensitive_score >= self.threshold else 0

    def _check_content_yaml(self, text):
        """使用YAML配置的正则表达式检查敏感内容，优化处理性能"""
        if not text or not isinstance(text, str) or len(text.strip()) == 0:
            return []
            
        results = []
        
        # 1. 检查关键词
        if self.keyword_pattern:
            keyword_matches = {}
            # 使用正则表达式一次性匹配所有关键词
            for match in self.keyword_pattern.finditer(text):
                keyword = match.group()
                if keyword not in keyword_matches:
                    keyword_matches[keyword] = []
                keyword_matches[keyword].append(match.start())
                
            # 将结果转换为列表格式
            for keyword, positions in keyword_matches.items():
                # 只保留前10个位置，防止位置列表过长
                if len(positions) > 10:
                    positions = positions[:10]
                results.append((keyword, positions))
                
        # 2. 检查结构化模式
        for pattern, weight in self.structured_patterns.items():
            matches = list(pattern.finditer(text))
            if matches:
                positions = [m.start() for m in matches[:10]]  # 限制位置数量
                pattern_str = pattern.pattern
                # 如果模式太长，截断显示
                if len(pattern_str) > 30:
                    pattern_str = pattern_str[:30] + "..."
                results.append((f"结构化模式[{pattern_str}]", positions))

        # 3. 检查数字模式
        for pattern in self.number_patterns:
            matches = list(pattern.finditer(text))
            if matches:
                positions = [m.start() for m in matches[:10]]  # 限制位置数量
                pattern_str = pattern.pattern
                # 如果模式太长，截断显示
                if len(pattern_str) > 30:
                    pattern_str = pattern_str[:30] + "..."
                results.append((f"数字模式[{pattern_str}]", positions))

        return results

    def _check_content_combined(self, text):
        """组合使用BERT模型和YAML配置进行检查"""
        # 先使用BERT模型检测是否敏感
        is_sensitive = self.predict(text)
        
        if is_sensitive == 1:
            # 如果BERT判断为敏感，再使用YAML配置进行详细检测
            yaml_results = self._check_content_yaml(text)
            
            # 如果YAML没有发现具体敏感内容，添加一个BERT检测结果
            if not yaml_results:
                return [("BERT模型检测为敏感内容", [0])]
            else:
                # 返回YAML检测的具体内容
                return yaml_results
        else:
            # BERT判断为非敏感，返回空结果
            return []

    def check_content(self, text):
        """
        检查文本中的敏感内容，根据初始化时的模式选择检测方法
        
        Args:
            text: 待检查的文本内容
            
        Returns:
            list: 包含敏感内容及其位置的列表，格式为[(敏感词1, [位置1, 位置2...]), ...]
        """
        # 增加检查计数
        self.check_count += 1
        
        # 基本检查
        if not text or not isinstance(text, str) or len(text.strip()) == 0:
            return []
            
        # 检查缓存
        cache_key = hashlib.md5(text[:1000].encode('utf-8')).hexdigest()
        if cache_key in self.result_cache:
            return self.result_cache[cache_key]
            
        try:
            result = []
            
            # 根据检测模式选择检测方法
            if self.mode == 'bert':
                # 单独使用BERT模型
                is_sensitive = self.predict(text)
                if is_sensitive == 1:
                    result = [("BERT模型检测为敏感内容", [0])]
            elif self.mode == 'yaml':
                # 使用YAML配置
                result = self._check_content_yaml(text)
            elif self.mode == 'combined':
                # 组合模式
                result = self._check_content_combined(text)
            else:
                # 空检测器
                result = []
                
            # 更新命中计数
            if result:
                self.hit_count += 1
                
            # 缓存结果，避免重复检查相同内容
            if len(self.result_cache) >= self.cache_size:
                # 如果缓存已满，清空一半
                keys = list(self.result_cache.keys())
                for k in keys[:self.cache_size//2]:
                    del self.result_cache[k]
                    
            self.result_cache[cache_key] = result
            
            return result
        except Exception as e:
            logger.error(f"敏感内容检查失败: {str(e)}")
            return []
            
    def get_stats(self):
        """获取检测器统计信息"""
        return {
            "mode": self.mode,
            "check_count": self.check_count,
            "hit_count": self.hit_count,
            "hit_ratio": f"{(self.hit_count / self.check_count * 100):.2f}%" if self.check_count > 0 else "0%",
            "cache_size": len(self.result_cache),
            "keywords_count": len(self.all_keywords) if hasattr(self, 'all_keywords') else 0,
            "device": self.device if self.mode == 'bert' else "N/A"
        }
        
    def check_file_content(self, file_content):
        """
        检查文件内容中的敏感信息
        
        Args:
            file_content: 文件内容字典，应包含'content'键
            
        Returns:
            list: 敏感内容列表
        """
        if not file_content or not isinstance(file_content, dict):
            return []
            
        content = file_content.get('content', '')
        if not content or not isinstance(content, str):
            return []
            
        # 对于超大文本，分块处理
        if len(content) > 100000:  # 10万字符以上的大文本
            chunks = []
            # 分块，每块5万字符，重叠1万字符
            for i in range(0, len(content), 40000):
                chunks.append(content[i:i+50000])
                
            # 如果块太多，只保留前10块
            if len(chunks) > 10:
                chunks = chunks[:10]
                
            # 检查每个块
            all_results = []
            for chunk in chunks:
                chunk_results = self.check_content(chunk)
                
                # 检查结果是否有重复，避免在重叠区域重复报告
                for result in chunk_results:
                    keyword, positions = result
                    if not any(kw == keyword for kw, _ in all_results):
                        all_results.append(result)
                    else:
                        # 更新已有关键词的位置
                        for i, (kw, pos) in enumerate(all_results):
                            if kw == keyword:
                                # 合并位置，保持唯一性和顺序
                                all_pos = list(set(pos + positions))
                                all_pos.sort()
                                # 限制位置数量
                                if len(all_pos) > 10:
                                    all_pos = all_pos[:10]
                                all_results[i] = (kw, all_pos)
                                break
                                
            return all_results
        else:
            # 常规大小文本的处理
            return self.check_content(content)


class ResultExporter:
    """增强版处理结果导出器，添加文本内容到Excel"""

    def export_to_json(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 JSON 文件"""
        export_data = [
            {
                "file_path": r.file_path,
                "mime_type": r.mime_type,
                "content_preview": self._get_content_preview(r, 200),
                "sensitive_words": [
                    {"word": w, "positions": p} for w, p in r.sensitive_words
                ],
                "error": r.error,
                "processing_time": r.processing_time,
            }
            for r in results
        ]
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)

    def export_to_excel(self, results: List[ProcessingResult], output_path: str):
        """导出结果到 Excel 文件，包含文本内容"""
        try:
            # 创建主工作表数据
            summary_data = [
                {
                    "文件路径": r.file_path,
                    "文件类型": r.mime_type,
                    "敏感词统计": "; ".join(
                        [f"{w}({len(p)}次)" for w, p in r.sensitive_words]
                    ),
                    "处理时间(秒)": round(r.processing_time, 3),
                    "错误信息": r.error or "",
                    "内容预览": self._get_content_preview(r, 500),
                }
                for r in results
            ]

            # 创建内容工作表数据
            content_data = [
                {
                    "文件路径": r.file_path,
                    "文件类型": r.mime_type,
                    "内容": self._get_full_content(r),
                }
                for r in results
            ]

            # 创建Excel工作簿并写入两个工作表
            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                # 摘要工作表
                pd.DataFrame(summary_data).to_excel(
                    writer, sheet_name="处理结果", index=False
                )
                # 内容工作表
                pd.DataFrame(content_data).to_excel(
                    writer, sheet_name="文本内容", index=False
                )

            logger.info(f"结果已导出到Excel（含文本内容）: {output_path}")
        except Exception as e:
            logger.error(f"导出到Excel失败: {e}")
            # 尝试导出简化版本，不包含内容
            try:
                simplified_path = output_path.replace(".xlsx", "_简化版.xlsx")
                # 创建简化版本
                simple_data = [
                    {
                        "文件路径": r.file_path,
                        "文件类型": r.mime_type,
                        "敏感词统计": "; ".join(
                            [f"{w}({len(p)}次)" for w, p in r.sensitive_words]
                        ),
                        "处理时间(秒)": round(r.processing_time, 3),
                        "错误信息": r.error or "",
                    }
                    for r in results
                ]
                pd.DataFrame(simple_data).to_excel(
                    simplified_path, index=False, engine="openpyxl"
                )
                logger.info(f"已导出简化版Excel: {simplified_path}")
            except Exception as e2:
                logger.error(f"导出简化版Excel也失败: {e2}")

    def _get_content_preview(
        self, result: ProcessingResult, max_length: int = 200
    ) -> str:
        """从处理结果中获取内容预览"""
        try:
            if isinstance(result.content, dict):
                content = result.content.get("content", "")
                if isinstance(content, str):
                    return content[:max_length] + (
                        "..." if len(content) > max_length else ""
                    )
            return ""
        except Exception:
            return ""

    def _get_full_content(self, result: ProcessingResult) -> str:
        """从处理结果中获取完整内容"""
        try:
            if isinstance(result.content, dict):
                content = result.content.get("content", "")
                if isinstance(content, str):
                    return content
            return ""
        except Exception:
            return ""


class ResultMonitor:
    """结果监控器，用于实时输出处理进度和结果"""

    def __init__(self, output_csv: str = "processing_results.csv"):
        self.output_csv = output_csv
        self._init_csv()

    def _init_csv(self):
        """初始化 CSV 文件"""
        with open(self.output_csv, "w", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow(["文件名", "识别类型"])

    def record_result(self, result: ProcessingResult):
        """记录单个处理结果"""
        file_name = os.path.basename(result.file_path)
        sensitive_words_count = len(result.sensitive_words)
        recognition_type = "敏感" if sensitive_words_count > 0 else "常规"

        with open(self.output_csv, "a", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow([file_name, recognition_type])

        print(f"[{datetime.now().strftime('%H:%M:%S')}] 处理文件: {result.file_path}")
        print(f"  类型: {result.mime_type}")
        print(f"  识别类型: {recognition_type}")
        if sensitive_words_count > 0:
            print(f"  发现敏感词: {'; '.join([f'{w}({len(p)}次)' for w, p in result.sensitive_words])}")
        if result.error:
            print(f"  错误: {result.error}")
        print("-" * 80)


class FileProcessor:
    """优化后的文件处理器主类"""

    def __init__(
        self,
        config_path: str = "sensitive_config.yaml",
        model_path: str = "best_model.pth",
        use_model: bool = False,
        monitor_output: str = "processing_results.csv",
        chunk_size: int = 1000,
        max_workers: Optional[int] = None,
        is_windows: bool = True,
    ):
        self.detector = FileTypeDetector()
        self.extractor = ContentExtractor(detector=self.detector, is_windows=is_windows)
        self.checker = SensitiveChecker(config_path=config_path, model_path=model_path, use_model=use_model)
        self.exporter = ResultExporter()
        self.monitor = ResultMonitor(monitor_output)
        self.chunk_size = chunk_size
        self.max_workers = max_workers or (os.cpu_count() or 1) * 2
        self._mime_cache = {}
        self._file_size_cache = {}

    def _scan_directory(self, directory: str) -> Iterator[List[str]]:
        """使用生成器分批扫描目录文件，确保正确关闭scandir迭代器"""
        current_chunk = []

        try:
            with os.scandir(directory) as dir_iter:
                for entry in dir_iter:
                    try:
                        if (
                            entry.is_file(follow_symlinks=False)
                            and not entry.name.startswith("~$")
                            and not self.detector._is_internal_stream(entry.name)
                        ):
                            ext = Path(entry.path).suffix.lower()
                            if ext not in self.detector.SKIP_EXTENSIONS:
                                current_chunk.append(entry.path)
                                if len(current_chunk) >= self.chunk_size:
                                    yield current_chunk
                                    current_chunk = []
                        elif entry.is_dir(follow_symlinks=False):
                            for sub_chunk in self._scan_directory(entry.path):
                                yield sub_chunk
                    except (PermissionError, OSError) as e:
                        logger.warning(f"访问文件/目录出错 {entry.path}: {e}")
                        continue
        except (PermissionError, OSError) as e:
            logger.warning(f"访问目录出错 {directory}: {e}")

        if current_chunk:
            yield current_chunk

    def _preload_file_info(self, file_paths: List[str]):
        """预加载文件信息以减少 I/O 操作"""
        for file_path in file_paths:
            try:
                stat = os.stat(file_path)
                self._file_size_cache[file_path] = stat.st_size
                mime_type = self.detector.detect_file_type(file_path)
                self._mime_cache[file_path] = mime_type
            except Exception as e:
                logger.warning(f"预加载文件信息失败 {file_path}: {e}")

    def _process_file_batch(self, file_paths: List[str]) -> List[ProcessingResult]:
        """处理一批文件，Office文件串行处理，其他文件并行处理"""
        results = []
        self._preload_file_info(file_paths)

        # 按类型对文件进行分类
        office_files = []
        other_files = []

        for fp in file_paths:
            mime_type = self._mime_cache.get(fp)
            ext = Path(fp).suffix.lower()

            # 检查是否为需要串行处理的Office文件
            if mime_type in (
                self.extractor.MIME_TYPE["DOC"],
                self.extractor.MIME_TYPE["PPT"],
            ) or ext in (".doc", ".ppt", ".xls"):
                office_files.append(fp)
            else:
                other_files.append(fp)


        return results

    def process_directory(self, directory: str) -> List[ProcessingResult]:
        """处理目录下的所有文件"""
        results = []
        total_files = len(self.detector.get_all_files(directory))
        completed = 0

        print(f"\n开始处理目录: {directory}")
        print(
            f"共发现 {total_files} 个文件"
        )
        print("=" * 80)

        try:
            for file_chunk in self._scan_directory(directory):
                chunk_results = self._process_file_batch(file_chunk)
                results.extend(chunk_results)
                completed += len(file_chunk)
                print(
                    f"\n已完成: {completed}/{total_files} ({completed / total_files * 100:.1f}%)"
                )
        except KeyboardInterrupt:
            print("\n用户中断处理...")
        except Exception as e:
            logger.error(f"处理目录时发生错误: {e}")
        finally:
            if results:
                try:
                    self.exporter.export_to_json(results, "interrupted_results.json")
                    self.exporter.export_to_excel(results, "interrupted_results.xlsx")
                except Exception as e:
                    logger.error(f"导出中断结果失败: {e}")

        print(f"\n处理完成！共处理 {len(results)}/{total_files} 个文件")
        print(f"详细结果已保存至: {self.monitor.output_csv}")
        print("=" * 80)
        return results

    def process_file(self, file_path: str) -> ProcessingResult:
        """处理单个文件，添加超时功能和容错性增强"""
        start_time = time.time()
        try:
            mime_type = self._mime_cache.get(
                file_path
            ) or self.detector.detect_file_type(file_path)
            ext = Path(file_path).suffix.lower()

            # 检查文件头部以辅助判断文件类型
            file_header = ""
            try:
                with open(file_path, "rb") as f:
                    file_header = f.read(16).hex().upper()
            except Exception:
                pass

            # 扩展不支持处理的MIME类型列表
            skip_mime_types = [
                "image/vnd.dwg",  # DWG图纸文件
                "application/octet-stream",  # 二进制文件
                "application/x-msdownload",  # 可执行文件
                "application/font-sfnt",  # 字体文件
                "font/ttf",  # TTF字体
                "font/otf",  # OTF字体
                "font/woff",  # WOFF字体
                "font/woff2",  # WOFF2字体
                "text/css",  # CSS文件
                "application/encrypted",
                "text/javascript",  # JS文件
                "image/jpeg",
                "image/png",
                "image/gif",
                "image/bmp",
                "image/webp",
                "image/svg+xml",  # 图片文件
                "audio/mpeg",
                "audio/wav",  # 音频文件
                "video/mp4",
                "video/x-msvideo",
                "video/quicktime",  # 视频文件
            ]

            # 扩展不支持的MIME类型前缀
            skip_mime_prefixes = ["font/", "image/", "audio/", "video/"]

            # 定义常规文件的条件 - 不再主要依赖文件扩展名
            regular_file_conditions = (
                mime_type in skip_mime_types
                or any(mime_type.startswith(prefix) for prefix in skip_mime_prefixes)
                or (ext and ext in self.detector.SKIP_EXTENSIONS)
            )

            if regular_file_conditions:
                logger.info(
                    f"识别为无需处理的文件类型，跳过内容提取: {file_path} (type: {mime_type}, header: {file_header})"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={
                        "content": "",
                        "metadata": {
                            "file_type": "regular",
                            "file_header": file_header,
                        },
                        "skipped": True,
                    },
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time,
                )

            # 检查文件是否为空
            file_size = self._file_size_cache.get(file_path, 0) or os.path.getsize(
                file_path
            )
            if file_size == 0:
                logger.info(f"空文件，跳过: {file_path}")
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={"content": "", "is_empty": True},
                    sensitive_words=[],
                    error=None,
                    processing_time=time.time() - start_time,
                )

            # 检查文件大小，过大的文件可能会导致处理缓慢
            size_mb = file_size / (1024 * 1024)
            if size_mb > 50:  # 大于50MB的文件
                logger.warning(
                    f"文件过大 ({size_mb:.2f} MB)，可能影响处理性能: {file_path}"
                )

            # =====================================================================
            # 特殊处理Office文档，避免线程锁序列化问题 - 使用超时功能
            # =====================================================================
            doc_timeout = 10  # 设置DOC文件处理超时时间为10秒
            ppt_timeout = 10  # 设置PPT文件处理超时时间为10秒

            # 特殊处理Office文档，避免线程锁序列化问题 - 使用超时功能
            # 因为没有扩展名，改为基于MIME类型判断
            if mime_type == "application/msword":
                # 使用带超时功能的函数处理DOC文件
                content = self.extractor._extract_doc_content(file_path)
            elif mime_type == "application/vnd.ms-powerpoint":
                # 使用带超时功能的函数处理PPT文件
                content = self.extractor._extract_ppt_content(file_path)
            elif mime_type == "application/vnd.ms-excel":
                # 针对Excel文件使用专用处理函数
                try:
                    content = self.extractor._extract_xls_content(file_path)
                except Exception as excel_error:
                    logger.error(f"处理Excel文件失败: {str(excel_error)}")
                    content = {
                        "content": "",
                        "metadata": {"file_type": "excel"},
                        "error": f"Excel处理失败: {str(excel_error)}"
                    }
            else:
                # 使用容错增强版的提取内容方法处理其他文件
                try:
                    content = self.extractor.extract_content(file_path, mime_type)
                except Exception as extract_error:
                    # 即使extract_content方法出现未捕获的异常，也不会导致程序崩溃
                    logger.error(
                        f"提取内容时发生严重错误: {str(extract_error)} - 文件: {file_path}"
                    )
                    content = {
                        "content": "",
                        "metadata": {"file_type": mime_type},
                        "error": f"提取内容时发生严重错误: {str(extract_error)}",
                    }
            # =====================================================================

            # 检查提取是否失败
            if content.get("error") and "UnsupportedFormatException" in content.get(
                "error"
            ):
                # 如果是不支持的格式，标记为跳过
                logger.info(
                    f"不支持的文件格式，标记为跳过: {file_path} (type: {mime_type})"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content={
                        "content": "",
                        "metadata": {"file_type": "unsupported"},
                        "skipped": True,
                    },
                    sensitive_words=[],
                    error=content.get("error"),
                    processing_time=time.time() - start_time,
                )
            elif content.get("error"):
                logger.warning(
                    f"提取内容失败: {file_path} - MIME: {mime_type} - 错误: {content.get('error')}"
                )
                return ProcessingResult(
                    file_path=file_path,
                    mime_type=mime_type,
                    content=content,
                    sensitive_words=[],
                    error=content.get("error"),
                    processing_time=time.time() - start_time,
                )

            # 敏感内容检查 - 包装在try-except中，确保不会因为敏感词检查失败而中断
            try:
                sensitive_result = self.checker.check_content(
                    content.get("content", "")
                )
                sensitive_words = sensitive_result
            except Exception as check_error:
                logger.error(f"敏感内容检查失败: {file_path} - {str(check_error)}")
                sensitive_words = []

            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type,
                content=content,
                sensitive_words=sensitive_words,
                error=None,
                processing_time=time.time() - start_time,
            )
        except Exception as e:
            error_msg = f"处理文件失败: {str(e)}"
            logger.error(f"{error_msg} - {file_path}")
            return ProcessingResult(
                file_path=file_path,
                mime_type=mime_type if "mime_type" in locals() else "unknown",
                content={
                    "content": "",
                    "metadata": {"file_type": "regular"},
                    "skipped": True,
                },
                sensitive_words=[],
                error=error_msg,
                processing_time=time.time() - start_time,
            )


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="文件敏感内容检测工具")
    parser.add_argument("path", help="要处理的文件或目录路径")
    parser.add_argument(
        "--config", default="sensitive_config.yaml", help="敏感词配置文件路径"
    )
    parser.add_argument(
        "--model", default="best_model.pth", help="BERT分类模型路径"
    )
    parser.add_argument(
        "--use-model", action="store_true", 
        help="使用BERT模型进行敏感内容检测（默认使用YAML配置）"
    )
    parser.add_argument(
        "--output", default="results", help="输出结果文件名(不含扩展名)"
    )
    parser.add_argument(
        "--chunk-size", type=int, default=1000, help="每批处理的文件数量"
    )
    parser.add_argument("--workers", type=int, default=None, help="最大工作线程数")
    parser.add_argument(
        "--no-windows",
        action="store_true",
        help="指定非Windows平台，禁用win32com相关功能",
    )

    args = parser.parse_args()

    try:
        processor = FileProcessor(
            config_path=args.config,
            model_path=args.model,
            use_model=args.use_model,
            monitor_output=f"{args.output}_processing.csv",
            chunk_size=args.chunk_size,
            max_workers=args.workers,
            is_windows=not args.no_windows,
        )

        results = []
        if Path(args.path).is_file():
            results = [processor.process_file(args.path)]
        else:
            results = processor.process_directory(args.path)

        processor.exporter.export_to_json(results, f"{args.output}.json")
        processor.exporter.export_to_excel(results, f"{args.output}.xlsx")

        logger.info(f"处理完成，共处理 {len(results)} 个文件")
    except KeyboardInterrupt:
        print("\n用户中断程序执行")
        sys.exit(1)
    except Exception as e:
        logger.error(f"程序执行出错: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
